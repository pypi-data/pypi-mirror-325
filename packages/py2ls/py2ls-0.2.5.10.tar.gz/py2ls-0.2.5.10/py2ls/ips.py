import numpy as np
import pandas as pd
import sys
import os
from IPython.display import display
from typing import Dict,List, Optional, Union,Any

from regex import X

try:
    get_ipython().run_line_magic("load_ext", "autoreload")
    get_ipython().run_line_magic("autoreload", "2")
except NameError:
    pass

import warnings

warnings.simplefilter("ignore", category=pd.errors.SettingWithCopyWarning)
warnings.filterwarnings("ignore", category=pd.errors.PerformanceWarning)
warnings.filterwarnings("ignore")
import shutil
import logging
from pathlib import Path
from datetime import datetime
import re
import stat
import platform


# only for backup these scripts
def backup(
    src="/Users/macjianfeng/Dropbox/github/python/py2ls/.venv/lib/python3.12/site-packages/py2ls/",
    tar="/Users/macjianfeng/Dropbox/github/python/py2ls/py2ls/",
    kind="py",
    overwrite=True,
    reverse=False,
    verbose=False,
):
    if reverse:
        src, tar = tar, src
        print(f"reversed")
    f = listdir(src, kind=kind, verbose=verbose)
    [copy(i, tar, overwrite=overwrite, verbose=verbose) for i in f.path]

    print(f"backup '*.{kind}'...\nfrom {src} \nto {tar}")


def run_once_within(duration=60, reverse=False):  # default 60s
    import time

    """
    如果reverse is True, 则在第一次运行时并不运行.但是在第二次运行时则运行
    usage:
    if run_once_within():
        print("This code runs once per minute.")
    else:
        print("The code has already been run in the last minute.")
    
    """
    if not hasattr(run_once_within, "time_last"):
        run_once_within.time_last = None
    time_curr = time.time()

    if (run_once_within.time_last is None) or (
        time_curr - run_once_within.time_last >= duration
    ):
        run_once_within.time_last = time_curr  # Update the last execution time
        return False if reverse else True
    else:
        return True if reverse else False


def plt_font(dir_font: str = "/System/Library/Fonts/Hiragino Sans GB.ttc"):
    """
    Add the Chinese (default) font to the font manager
    show chinese
    Args:
        dir_font (str, optional): _description_. Defaults to "/System/Library/Fonts/Hiragino Sans GB.ttc".
    """
    import matplotlib.pyplot as plt
    from matplotlib import font_manager

    slashtype = "/" if "mac" in get_os() else "\\"
    if slashtype in dir_font:
        font_manager.fontManager.addfont(dir_font)
        fontname = os.path.basename(dir_font).split(".")[0]
    else:
        if "cn" in dir_font.lower() or "ch" in dir_font.lower():
            fontname = "Hiragino Sans GB"  # default Chinese font
        else:
            fontname = dir_font

    plt.rcParams["font.sans-serif"] = [fontname]
    # plt.rcParams["font.family"] = "sans-serif"
    plt.rcParams["axes.unicode_minus"] = False
    fonts_in_system = font_manager.findSystemFonts(fontpaths=None, fontext="ttf")
    fontname_in_system = [os.path.basename(i).split(".")[0] for i in fonts_in_system]
    if fontname not in fontname_in_system:
        print(f"Font '{fontname}' not found. Falling back to default.")
        plt.rcParams["font.sans-serif"] = ["Arial"]
    return fontname


# set 'dir_save'
if "dar" in sys.platform:
    dir_save = "/Users/macjianfeng/Dropbox/Downloads/"
else:
    if "win" in sys.platform:
        dir_save = "Z:\\Jianfeng\\temp\\"
    elif "lin" in sys.platform:
        dir_data = "/Users/macjianfeng/Dropbox/github/python/py2ls/confidential_data/gmail_login.json"


def unique(lst, ascending=None):
    """
    移除列表中的重复元素，同时可以选择按升序或降序排序。

    参数:
    lst (list): 输入的列表，其中可能包含重复的元素。
    ascending (bool, 可选): 如果为 True，则按升序排序；如果为 False，则按降序排序；如果为 None，则不排序，只移除重复元素。

    返回:
    list: 一个列表，其中的元素是唯一的，顺序根据参数 `ascending` 进行排序。
    """
    if not lst:
        return []
    if ascending is not None:
        # 移除重复项
        unique_items = list(set(lst))
        # 统一元素类型（例如，转换为字符串）
        try:
            unique_items = sorted(
                unique_items, key=lambda x: str(x), reverse=not ascending
            )
        except TypeError:
            # 如果排序失败（例如，元素类型不一致），返回原始去重列表
            return unique_items
        return unique_items
    else:
        # 移除重复项同时保持原始顺序
        seen = set()  # 用于记录已见的元素
        result = []  # 用于存储结果
        for item in lst:
            if item not in seen:
                seen.add(item)  # 记录该元素
                result.append(item)  # 添加到结果列表中
        return result


# ************* below section: run_when *************
def run_when(when: str = "every 2 min", job=None, wait: int = 60):
    if "every" in when.lower():
        when = when.replace("every", "")
        run_every(when=when, job=job, wait=wait)
    elif any([i in when.lower() for i in ["at", "@", ":", "am", "pm"]]):
        time_words = ["at", "@", ":", "am", "pm"]
        # 判断'时间词'是否存在
        time_words_bool = [i in when.lower() for i in time_words]
        # 找到'时间词'的位置
        true_indices = [index for index, value in enumerate(time_words_bool) if value]
        time_word = time_words[true_indices[0]]  # 找到第一个'时间词'
        when = when.replace(time_word, "")  # 去除 时间词
        run_at(when=when, job=job, wait=wait)


def run_every(when: str = None, job=None, wait: int = 60):
    """
    Schedules a job to run at the given interval.

    :param when: String specifying the interval, e.g. '2 minutes', '4 hours', '1 day'.
    :param job: The function to be scheduled.

    # usage:
        def job():
            print("1 sec")
        run_every(when="1 sec", job=job)
    """
    import schedule
    import time

    if job is None:
        print("No job provided!")
        return

    interval, unit = (
        str2num(when),
        strcmp(when.replace("every", ""), ["seconds", "minutes", "hours", "days"])[0],
    )
    print(interval, unit)
    # Mapping the scheduling based on the unit1
    if unit == "seconds":
        schedule.every(interval).seconds.do(job)
    elif unit == "minutes":
        schedule.every(interval).minutes.do(job)
    elif unit == "hours":
        schedule.every(interval).hours.do(job)
    elif unit == "days":
        schedule.every(interval).days.do(job)
    else:
        print(f"Invalid time unit: {unit}")
        return

    print(f"Scheduled job when {interval} {unit}.")

    # Keep the script running to execute the schedule
    while True:
        schedule.run_pending()
        time.sleep(wait)  # in seconds
    time.sleep(wait)  # in seconds 
def run_at(when: str, job=None, wait: int = 60):
    """
    Schedules a job to run at an exact time of the day.

    # Example usage:
    def my_job():
        print("Job executed at the exact time!")
    # Schedule the job at 14:30 when day
    run_at(when="1.30 pm", job=my_job)

    :param when: String specifying the time, e.g. '1:30 pm','1.30 am','14:30', '1:30 pm', '8:45 am'.
    :param job: The function to be scheduled.
    :param wait: The sleep interval between checks in seconds.
    """
    from datetime import datetime
    import time

    if job is None:
        print("No job provided!")
        return
    when = when.replace("A.M.", "AM").replace("P.M.", "PM")
    when = when.replace(".", ":")
    when = when.strip()

    try:
        # Attempt to parse the time in both 24-hour and 12-hour format
        if "am" in when.lower() or "pm" in when.lower():
            scheduled_time = datetime.strptime(
                when, "%I:%M %p"
            ).time()  # 12-hour format with AM/PM
        else:
            scheduled_time = datetime.strptime(when, "%H:%M").time()  # 24-hour format
    except ValueError:
        print(
            f"Invalid time format: {when}. Use 'HH:MM' (24-hour) or 'H:MM AM/PM' format."
        )
        return
    print(f"Job scheduled to run at {scheduled_time}.")
    # Keep checking the current time
    while True:
        now = datetime.now()
        # Check if current time matches the scheduled time
        if (
            now.time().hour == scheduled_time.hour
            and now.time().minute == scheduled_time.minute
        ):
            job()  # Run the job
            time.sleep(
                wait
            )  # Sleep for a minute to avoid running the job multiple times in the same minute

        time.sleep(wait)  # wait to avoid excessive CPU usage
 
# ************* above section: run_when *************


def get_timezone(timezone: str | list = None):
    if timezone is None:
        usage = """
        usage:
        datetime.now().astimezone(get_timezone("shanghai")).strftime("%H:%M")
        """
        print(usage)
        return None
    from pytz import all_timezones
    import pytz

    if isinstance(timezone, str):
        timezone = [timezone]

    # Extract the part after the "/" in time zones (if exists)
    timezones = [ssplit(i, "/")[1] if "/" in i else i for i in all_timezones]

    # Print all available time zones for debugging purposes
    # print(timezones)

    # Find and return matched time zones using strcmp
    matched_timezones = [all_timezones[strcmp(i, timezones)[1]] for i in timezone]
    if len(matched_timezones) == 1:
        return pytz.timezone(matched_timezones[0])
    else:
        return matched_timezones


def is_package_installed(package_name):
    """Check if a package is installed."""
    import importlib.util

    package_spec = importlib.util.find_spec(package_name)
    return package_spec is not None


def upgrade(module="py2ls", uninstall=False):
    """
    Installs or upgrades a specified Python module.

    Parameters:
    module (str): The name of the module to install/upgrade.
    uninstall (bool): If True, uninstalls the webdriver-manager before upgrading.
    """
    import subprocess

    if not is_package_installed(module):
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", module])
        except subprocess.CalledProcessError as e:
            print(f"An error occurred while installing {module}: {e}")
    if uninstall:
        subprocess.check_call(["pip", "uninstall", "-y", "webdriver-manager"])
    try:
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install", "--upgrade", module]
        )
    except subprocess.CalledProcessError as e:
        print(f"An error occurred while upgrading py2ls: {e}")


def get_version(pkg):
    import importlib.metadata

    def get_v(pkg_name):
        try:
            version = importlib.metadata.version(pkg_name)
            print(f"version {pkg_name} == {version}")
        except importlib.metadata.PackageNotFoundError:
            print(f"Package '{pkg_name}' not found")

    if isinstance(pkg, str):
        get_v(pkg)
    elif isinstance(pkg, list):
        [get_v(pkg_) for pkg_ in pkg] 

def rm_folder(folder_path, verbose=True):
    import shutil

    try:
        shutil.rmtree(folder_path)
        if verbose:
            print(f"Successfully deleted {folder_path}")
    except Exception as e:
        if verbose:
            print(f"Failed to delete {folder_path}. Reason: {e}")


def fremove(path, verbose=True):
    """
    Remove a folder and all its contents or a single file.
    Parameters:
    path (str): The path to the folder or file to remove.
    verbose (bool): If True, print success or failure messages. Default is True.
    """
    try:
        if os.path.isdir(path):
            import shutil

            shutil.rmtree(path)
            if verbose:
                print(f"Successfully deleted folder {path}")
        elif os.path.isfile(path):
            os.remove(path)
            if verbose:
                print(f"Successfully deleted file {path}")
        else:
            if verbose:
                print(f"Path {path} does not exist")
    except Exception as e:
        if verbose:
            print(f"Failed to delete {path}. Reason: {e}")
 
def get_cwd():
    from pathlib import Path
    # Get the current script's directory as a Path object
    current_directory = Path(__file__).resolve().parent
    return current_directory


def search(
    query,
    limit=5,
    kind="text",
    output="df",
    verbose=False,
    download=False,
    dir_save=None,
    **kwargs,
):
    from duckduckgo_search import DDGS

    if "te" in kind.lower():
        results = DDGS().text(query, max_results=limit)
        res = pd.DataFrame(results)
        res.rename(columns={"href": "links"}, inplace=True)
    if verbose:
        print(f'searching "{query}": got the results below\n{res}')
    if download:
        try:
            downloader(
                url=res.links.tolist(), dir_save=dir_save, verbose=verbose, **kwargs
            )
        except:
            if verbose:
                print(f"failed link")
    return res


def echo(*args, **kwargs):
    """
    query, model="gpt", verbose=True, log=True, dir_save=dir_save
    a ai chat tool
    Args:
        query (str): _description_
        model (str, optional): _description_. Defaults to "gpt".
        verbose (bool, optional): _description_. Defaults to True.
        log (bool, optional): _description_. Defaults to True.
        dir_save (str, path, optional): _description_. Defaults to dir_save.

    Returns:
        str: the answer from ai
    """
    global dir_save
    from duckduckgo_search import DDGS

    query = None
    model = kwargs.get("model", "gpt")
    verbose = kwargs.get("verbose", True)
    log = kwargs.get("log", True)
    dir_save = kwargs.get("dir_save", dir_save)
    for arg in args:
        if isinstance(arg, str):
            if os.path.isdir(arg):
                dir_save = arg
            # elif os.path.isfile(arg):
            #     dir_save = dirname(arg)
            elif len(arg) <= 5:
                model = arg
            else:
                query = arg
        elif isinstance(arg, dict):
            verbose = arg.get("verbose", verbose)
            log = arg.get("log", log)

    def is_in_any(str_candi_short, str_full, ignore_case=True):
        if isinstance(str_candi_short, str):
            str_candi_short = [str_candi_short]
        res_bool = []
        if ignore_case:
            [res_bool.append(i in str_full.lower()) for i in str_candi_short]
        else:
            [res_bool.append(i in str_full) for i in str_candi_short]
        return any(res_bool)

    def valid_mod_name(str_fly):
        if is_in_any(str_fly, "claude-3-haiku"):
            return "claude-3-haiku"
        elif is_in_any(str_fly, "gpt-3.5"):
            return "gpt-3.5"
        elif is_in_any(str_fly, "llama-3-70b"):
            return "llama-3-70b"
        elif is_in_any(str_fly, "mixtral-8x7b"):
            return "mixtral-8x7b"
        else:
            print(
                f"not support your model{model}, supported models: 'claude','gpt(default)', 'llama','mixtral'"
            )
            return "gpt-3.5"  # default model

    model_valid = valid_mod_name(model)
    res = DDGS().chat(query, model=model_valid)
    if verbose:
        from pprint import pp

        pp(res)
    if log:
        from datetime import datetime
        import time

        dt_str = datetime.fromtimestamp(time.time()).strftime("%Y-%m-%d_%H:%M:%S")
        res_ = f"\n\n####Q:{query}\n\n#####Ans:{dt_str}\n\n>{res}\n"
        if bool(os.path.basename(dir_save)):
            fpath = dir_save
        else:
            os.makedirs(dir_save, exist_ok=True)
            fpath = os.path.join(dir_save, f"log_ai.md")
        fupdate(fpath=fpath, content=res_)
        print(f"log file:{fpath}")
    return res


def chat(*args, **kwargs):
    return echo(*args, **kwargs)
def ai(*args, **kwargs):
    return echo(*args, **kwargs)

def detect_lang(text, output="lang", verbose=True):
    from langdetect import detect

    dir_curr_script = os.path.dirname(os.path.abspath(__file__))
    dir_lang_code = dir_curr_script + "/data/lang_code_iso639.json"
    print(dir_curr_script, os.getcwd(), dir_lang_code)
    lang_code_iso639 = fload(dir_lang_code)
    l_lang, l_code = [], []
    [[l_lang.append(v), l_code.append(k)] for v, k in lang_code_iso639.items()]
    try:
        if is_text(text):
            code_detect = detect(text)
            if "c" in output.lower():  # return code
                return l_code[strcmp(code_detect, l_code, verbose=verbose)[1]]
            else:
                return l_lang[strcmp(code_detect, l_code, verbose=verbose)[1]]
        else:
            print(f"{text} is not supported")
            return "no"
    except:
        return "no"


def is_text(s):
    has_alpha = any(char.isalpha() for char in s)
    has_non_alpha = any(not char.isalpha() for char in s)
    # no_special = not re.search(r'[^A-Za-z0-9\s]', s)
    return has_alpha and has_non_alpha

def share(*args, strict=True, n_shared=2, verbose=True):
    """
    check the shared elelements in two list.
    usage:
        list1 = [1, 2, 3, 4, 5]
        list2 = [4, 5, 6, 7, 8]
        list3 = [5, 6, 9, 10]
        a = shared(list1, list2,list3)
    """
    if verbose:
        print("\n********* checking shared elements *********")

    if len(args) == 1 and isinstance(args[0], list):
        lists = args[0]  # Unpack the single list
    else:
        lists = args  # Use the provided arguments as lists
    flattened_lists = [flatten(lst, verbose=verbose) for lst in lists]
    # Ensure all arguments are lists
    if any(not isinstance(lst, list) for lst in flattened_lists):
        print(f"{' ' * 2}All inputs must be lists.")
        return []
    first_list = flattened_lists[0]
    shared_elements = [
        item for item in first_list if all(item in lst for lst in flattened_lists)
    ]
    if strict:
        # Strict mode: require elements to be in all lists
        shared_elements = set(flattened_lists[0])
        for lst in flattened_lists[1:]:
            shared_elements.intersection_update(lst)
    else:
        from collections import Counter

        all_elements = [item for sublist in flattened_lists for item in sublist]
        element_count = Counter(all_elements)
        # Get elements that appear in at least n_shared lists
        shared_elements = [
            item for item, count in element_count.items() if count >= n_shared
        ]

    shared_elements = flatten(shared_elements, verbose=verbose)
    if verbose:
        elements2show = (
            shared_elements if len(shared_elements) < 10 else shared_elements[:5]
        )
        tail = "" if len(shared_elements) < 10 else "......"
        elements2show.append(tail)
        print(f"{' '*2}{len(shared_elements)} elements shared: {' '*2}{elements2show}")
        print("********* checking shared elements *********")
    return shared_elements


def shared(*args, n_shared=None, verbose=True, **kwargs):
    """
    check the shared elelements in two list.
    usage:
        list1 = [1, 2, 3, 4, 5]
        list2 = [4, 5, 6, 7, 8]
        list3 = [5, 6, 9, 10]
        a = shared(list1, list2,list3)
    """
    if verbose:
        print("\n********* checking shared elements *********")

    if len(args) == 1 and isinstance(args[0], list):
        lists = args[0]  # Unpack the single list
    else:
        lists = args  # Use the provided arguments as lists
    flattened_lists = [flatten(lst, verbose=verbose) for lst in lists]

    if n_shared is None:
        n_shared = len(flattened_lists)
        strict = True
    else:
        strict = False
    # Ensure all arguments are lists
    if any(not isinstance(lst, list) for lst in flattened_lists):
        print(f"{' ' * 2}All inputs must be lists.")
        return []
    first_list = flattened_lists[0]
    shared_elements = [
        item for item in first_list if all(item in lst for lst in flattened_lists)
    ]
    if strict:
        # Strict mode: require elements to be in all lists
        shared_elements = set(flattened_lists[0])
        for lst in flattened_lists[1:]:
            shared_elements.intersection_update(lst)
    else:
        from collections import Counter

        all_elements = [item for sublist in flattened_lists for item in sublist]
        element_count = Counter(all_elements)
        # Get elements that appear in at least n_shared lists
        shared_elements = [
            item for item, count in element_count.items() if count >= n_shared
        ]

    shared_elements = flatten(shared_elements, verbose=verbose)
    if verbose:
        elements2show = (
            shared_elements if len(shared_elements) < 10 else shared_elements[:5]
        )
        print(f"{' '*2}{len(shared_elements)} elements shared: {' '*2}{elements2show}")
        print("********* checking shared elements *********")
    return shared_elements


def share_not(*args, n_shared=None, verbose=False):
    """
    To find the elements in list1 that are not shared with list2 while maintaining the original order of list1
    usage:
        list1 = [1, 8, 3, 3, 4, 5]
        list2 = [4, 5, 6, 7, 8]
        not_shared(list1,list2)# output [1,3]
    """
    _common = shared(*args, n_shared=n_shared, verbose=verbose)
    list1 = flatten(args[0], verbose=verbose)
    _not_shared = [item for item in list1 if item not in _common]
    return _not_shared


def not_shared(*args, n_shared=None, verbose=False):
    """
    To find the elements in list1 that are not shared with list2 while maintaining the original order of list1
    usage:
        list1 = [1, 8, 3, 3, 4, 5]
        list2 = [4, 5, 6, 7, 8]
        not_shared(list1,list2)# output [1,3]
    """
    _common = shared(*args, n_shared=n_shared, verbose=verbose)
    list1 = flatten(args[0], verbose=verbose)
    _not_shared = [item for item in list1 if item not in _common]
    return _not_shared


def flatten(nested: Any, unique_list=True, verbose=False):
    """
    Recursively flattens a nested structure (lists, tuples, dictionaries, sets) into a single list.
    Parameters:
        nested : Any, Can be a list, tuple, dictionary, or set.
    Returns: list, A flattened list.
    """
    flattened_list = []
    stack = [nested]
    while stack:
        current = stack.pop()
        if isinstance(current, dict):
            stack.extend(current.values())
        elif isinstance(current, (list, tuple, set)):
            stack.extend(current)
        elif isinstance(current, pd.Series):
            stack.extend(current)
        elif isinstance(
            current, (pd.Index, np.ndarray)
        ):  # df.columns df.index are object of type pd.Index
            stack.extend(current.tolist())
        else:
            flattened_list.append(current)
    if verbose:
        print(
            f"{' '*2}<in info: {len(unique(flattened_list))} elements after flattened>"
        )
    if unique_list:
        return unique(flattened_list)[::-1]
    else:
        return flattened_list


def strcmp(
    search_term,
    candidates,
    ignore_case=True,
    get_rank=False,
    verbose=False,
    scorer="WR",
):
    """
    Compares a search term with a list of candidate strings and finds the best match based on similarity score.

    Parameters:
    search_term (str): The term to be searched for.
    candidates (list of str): A list of candidate strings to compare against the search term.
    ignore_case (bool): If True, the comparison ignores case differences.
    verbose (bool): If True, prints the similarity score and the best match.

    Returns:
    tuple: A tuple containing the best match and its index in the candidates list.
    """
    from fuzzywuzzy import fuzz, process

    def to_lower(s, ignore_case=True):
        # Converts a string or list of strings to lowercase if ignore_case is True.
        if ignore_case:
            if isinstance(s, str):
                return s.lower()
            elif isinstance(s, list):
                s = [str(i) for i in s]  # convert all to str
                return [elem.lower() for elem in s]
        return s

    str1_, str2_ = to_lower(search_term, ignore_case), to_lower(candidates, ignore_case)
    if isinstance(str2_, list):
        if "part" in scorer.lower():
            similarity_scores = [fuzz.partial_ratio(str1_, word) for word in str2_]
        elif "W" in scorer.lower():
            similarity_scores = [fuzz.WRatio(str1_, word) for word in str2_]
        elif "ratio" in scorer.lower() or "stri" in scorer.lower():  # Ratio (Strictest)
            similarity_scores = [fuzz.ratio(str1_, word) for word in str2_]
        else:
            similarity_scores = [fuzz.WRatio(str1_, word) for word in str2_]
        if get_rank:
            idx = [
                similarity_scores.index(i)
                for i in sorted(similarity_scores, reverse=True)
            ]
            if verbose:
                display([candidates[ii] for ii in idx])
            return [candidates[ii] for ii in idx]
        best_match_index = similarity_scores.index(max(similarity_scores))
        best_match_score = similarity_scores[best_match_index]
    else:
        best_match_index = 0
        if "part" in scorer.lower():
            best_match_score = fuzz.partial_ratio(str1_, str2_)
        elif "W" in scorer.lower():
            best_match_score = fuzz.WRatio(str1_, str2_)
        elif "Ratio" in scorer.lower():
            best_match_score = fuzz.ratio(str1_, str2_)
        else:
            best_match_score = fuzz.WRatio(str1_, str2_)
    if verbose:
        print(f"\nbest_match is: {candidates[best_match_index],best_match_score}")
        best_match = process.extract(search_term, candidates)
        print(f"建议: {best_match}")
    return candidates[best_match_index], best_match_index


def imgcmp(
    img: list,
    method: str = "knn",
    thr: float = 0.75,
    detector: str = "sift",
    plot_: bool = True,
    figsize=[12, 6],
    grid_size=10,  # only for grid detector
    **kwargs,
):
    """
    Compare two images using SSIM, Feature Matching (SIFT), or KNN Matching.

    Parameters:
    - img (list): List containing two image file paths [img1, img2] or two numpy arrays.
    - method (str): Comparison method ('ssim', 'match', or 'knn').
    - detector (str): Feature detector ('sift', 'grid', 'pixel').
    - thr (float): Threshold for filtering matches.
    - plot_ (bool): Whether to display the results visually.
    - figsize (list): Size of the figure for plots.

    Returns:
    - For 'ssim': (diff, score): SSIM difference map and similarity score.
    - For 'match' or 'knn': (good_matches, len(good_matches), similarity_score): Matches and similarity score.
    """
    import cv2
    import matplotlib.pyplot as plt
    from skimage.metrics import structural_similarity as ssim

    # Load images
    if isinstance(img, list) and isinstance(img[0], str):
        image1 = cv2.imread(img[0])
        image2 = cv2.imread(img[1])
        bool_cvt = True
    else:
        image1, image2 = np.array(img[0]), np.array(img[1])
        bool_cvt = False

    if image1 is None or image2 is None:
        raise ValueError("Could not load one or both images. Check file paths.")
    methods = ["ssim", "match", "knn"]
    method = strcmp(method, methods)[0]
    if method == "ssim":
        # Convert images to grayscale
        gray1 = cv2.cvtColor(image1, cv2.COLOR_BGR2GRAY)
        gray2 = cv2.cvtColor(image2, cv2.COLOR_BGR2GRAY)

        # Compute SSIM
        score, diff = ssim(gray1, gray2, full=True)
        print(f"SSIM Score: {score:.4f}")

        # Convert diff to 8-bit for visualization
        diff = (diff * 255).astype("uint8")

        # Plot if needed
        if plot_:
            fig, ax = plt.subplots(1, 3, figsize=figsize)
            ax[0].imshow(gray1, cmap="gray")
            ax[0].set_title("Image 1")
            ax[1].imshow(gray2, cmap="gray")
            ax[1].set_title("Image 2")
            ax[2].imshow(diff, cmap="gray")
            ax[2].set_title("Difference (SSIM)")
            plt.tight_layout()
            plt.show()

        return diff, score

    elif method in ["match", "knn"]:
        # Convert images to grayscale
        gray1 = cv2.cvtColor(image1, cv2.COLOR_BGR2GRAY)
        gray2 = cv2.cvtColor(image2, cv2.COLOR_BGR2GRAY)

        if detector == "sift":
            # SIFT detector
            sift = cv2.SIFT_create()
            keypoints1, descriptors1 = sift.detectAndCompute(gray1, None)
            keypoints2, descriptors2 = sift.detectAndCompute(gray2, None)

        elif detector == "grid":
            # Grid-based detection
            keypoints1, descriptors1 = [], []
            keypoints2, descriptors2 = [], []

            for i in range(0, gray1.shape[0], grid_size):
                for j in range(0, gray1.shape[1], grid_size):
                    patch1 = gray1[i : i + grid_size, j : j + grid_size]
                    patch2 = gray2[i : i + grid_size, j : j + grid_size]
                    if patch1.size > 0 and patch2.size > 0:
                        keypoints1.append(
                            cv2.KeyPoint(
                                j + grid_size // 2, i + grid_size // 2, grid_size
                            )
                        )
                        keypoints2.append(
                            cv2.KeyPoint(
                                j + grid_size // 2, i + grid_size // 2, grid_size
                            )
                        )
                        descriptors1.append(np.mean(patch1))
                        descriptors2.append(np.mean(patch2))

            descriptors1 = np.array(descriptors1).reshape(-1, 1)
            descriptors2 = np.array(descriptors2).reshape(-1, 1)

        elif detector == "pixel":
            # Pixel-based direct comparison
            descriptors1 = gray1.flatten()
            descriptors2 = gray2.flatten()
            keypoints1 = [
                cv2.KeyPoint(x, y, 1)
                for y in range(gray1.shape[0])
                for x in range(gray1.shape[1])
            ]
            keypoints2 = [
                cv2.KeyPoint(x, y, 1)
                for y in range(gray2.shape[0])
                for x in range(gray2.shape[1])
            ]

        else:
            raise ValueError("Invalid detector. Use 'sift', 'grid', or 'pixel'.")

        # Handle missing descriptors
        if descriptors1 is None or descriptors2 is None:
            raise ValueError("Failed to compute descriptors for one or both images.")
        # Ensure descriptors are in the correct data type
        if descriptors1.dtype != np.float32:
            descriptors1 = descriptors1.astype(np.float32)
        if descriptors2.dtype != np.float32:
            descriptors2 = descriptors2.astype(np.float32)

        # BFMatcher initialization
        bf = cv2.BFMatcher()
        if method == "match":  # Cross-check matching
            bf = cv2.BFMatcher(cv2.NORM_L2, crossCheck=True)
            matches = bf.match(descriptors1, descriptors2)
            matches = sorted(matches, key=lambda x: x.distance)

            # Filter good matches
            good_matches = [
                m for m in matches if m.distance < thr * matches[-1].distance
            ]

        elif method == "knn":  # KNN matching with ratio test
            bf = cv2.BFMatcher()
            matches = bf.knnMatch(descriptors1, descriptors2, k=2)
            # Apply Lowe's ratio test
            good_matches = [m for m, n in matches if m.distance < thr * n.distance]

        # Calculate similarity score
        similarity_score = len(good_matches) / min(len(keypoints1), len(keypoints2))
        print(f"Number of good matches: {len(good_matches)}")
        print(f"Similarity Score: {similarity_score:.4f}")
        # Handle case where no good matches are found
        if len(good_matches) == 0:
            print("No good matches found.")
            return good_matches, 0.0, None

        # Identify matched keypoints
        src_pts = np.float32([keypoints1[m.queryIdx].pt for m in good_matches]).reshape(
            -1, 1, 2
        )
        dst_pts = np.float32([keypoints2[m.trainIdx].pt for m in good_matches]).reshape(
            -1, 1, 2
        )
        # Apply the homography to image2
        try:
            # Calculate Homography using RANSAC
            homography_matrix, mask = cv2.findHomography(
                src_pts, dst_pts, cv2.RANSAC, 5.0
            )
            h, w = image1.shape[:2]
            warped_image2 = cv2.warpPerspective(image2, homography_matrix, (w, h))

            # Plot result if needed
            if plot_:
                fig, ax = plt.subplots(1, 2, figsize=figsize)
                (
                    ax[0].imshow(cv2.cvtColor(image1, cv2.COLOR_BGR2RGB))
                    if bool_cvt
                    else ax[0].imshow(image1)
                )
                ax[0].set_title("Image 1")
                (
                    ax[1].imshow(cv2.cvtColor(warped_image2, cv2.COLOR_BGR2RGB))
                    if bool_cvt
                    else ax[1].imshow(warped_image2)
                )
                ax[1].set_title("Warped Image 2")
                plt.tight_layout()
                plt.show()
        except Exception as e:
            print(e)

        # Plot matches if needed
        if plot_:
            result = cv2.drawMatches(
                image1, keypoints1, image2, keypoints2, good_matches, None, flags=2
            )
            plt.figure(figsize=figsize)
            (
                plt.imshow(cv2.cvtColor(result, cv2.COLOR_BGR2RGB))
                if bool_cvt
                else plt.imshow(result)
            )
            plt.title(
                f"Feature Matches ({len(good_matches)} matches, Score: {similarity_score:.4f})"
            )
            plt.axis("off")
            plt.show()
        # Identify unmatched keypoints
        matched_idx1 = [m.queryIdx for m in good_matches]
        matched_idx2 = [m.trainIdx for m in good_matches]
        matched_kp1 = [kp for i, kp in enumerate(keypoints1) if i in matched_idx1]
        matched_kp2 = [kp for i, kp in enumerate(keypoints2) if i in matched_idx2]
        unmatched_kp1 = [kp for i, kp in enumerate(keypoints1) if i not in matched_idx1]
        unmatched_kp2 = [kp for i, kp in enumerate(keypoints2) if i not in matched_idx2]

        # Mark keypoints on the images
        img1_match = cv2.drawKeypoints(
            image1,
            matched_kp1,
            None,
            color=(0, 0, 255),
            flags=cv2.DRAW_MATCHES_FLAGS_DRAW_RICH_KEYPOINTS,
        )
        img2_match = cv2.drawKeypoints(
            image2,
            matched_kp2,
            None,
            color=(0, 0, 255),
            flags=cv2.DRAW_MATCHES_FLAGS_DRAW_RICH_KEYPOINTS,
        )
        img1_unmatch = cv2.drawKeypoints(
            image1,
            unmatched_kp1,
            None,
            color=(0, 0, 255),
            flags=cv2.DRAW_MATCHES_FLAGS_DRAW_RICH_KEYPOINTS,
        )
        img2_unmatch = cv2.drawKeypoints(
            image2,
            unmatched_kp2,
            None,
            color=(0, 0, 255),
            flags=cv2.DRAW_MATCHES_FLAGS_DRAW_RICH_KEYPOINTS,
        )

        if plot_:
            fig, ax = plt.subplots(1, 2, figsize=figsize)
            (
                ax[0].imshow(cv2.cvtColor(img1_unmatch, cv2.COLOR_BGR2RGB))
                if bool_cvt
                else ax[0].imshow(img1_unmatch)
            )
            ax[0].set_title("Unmatched Keypoints (Image 1)")
            (
                ax[1].imshow(cv2.cvtColor(img2_unmatch, cv2.COLOR_BGR2RGB))
                if bool_cvt
                else ax[1].imshow(img2_unmatch)
            )
            ax[1].set_title("Unmatched Keypoints (Image 2)")
            ax[0].axis("off")
            ax[1].axis("off")
            plt.tight_layout()
            plt.show()
        if plot_:
            fig, ax = plt.subplots(1, 2, figsize=figsize)
            (
                ax[0].imshow(cv2.cvtColor(img1_match, cv2.COLOR_BGR2RGB))
                if bool_cvt
                else ax[0].imshow(img1_match)
            )
            ax[0].set_title("Matched Keypoints (Image 1)")
            (
                ax[1].imshow(cv2.cvtColor(img2_match, cv2.COLOR_BGR2RGB))
                if bool_cvt
                else ax[1].imshow(img2_match)
            )
            ax[1].set_title("Matched Keypoints (Image 2)")
            ax[0].axis("off")
            ax[1].axis("off")
            plt.tight_layout()
            plt.show()
        return good_matches, similarity_score  # , homography_matrix

    else:
        raise ValueError("Invalid method. Use 'ssim', 'match', or 'knn'.")

def fcmp(file1, file2, kind= None, verbose=True, **kwargs):
    import pandas as pd
    import os
    from concurrent.futures import ThreadPoolExecutor
    from datetime import datetime
    import json

    # --- Compare excel files ---
    def cmp_excel(
        file1,# base
        file2,  # new
        sheet_name=None,  # list or strings; default:"common" sheet
        key_columns=None,
        ignore_columns=None,
        numeric_tolerance=0,
        ignore_case=False,
        detect_reordered_rows=False,
        verbose=True,
        **kwargs,
    ):
        """
        Compare two Excel files and identify differences across specified sheets.

        Parameters:
        - file1 (Base/Reference): str, path to the first Excel file.
        - file2: str, path to the second Excel file.
        - sheet_name: list of str, specific sheets to compare (default: all common sheets).
        - key_columns: list of str, columns to use as unique identifiers (default: None, compares all columns).
        - ignore_columns: list of str, columns to exclude from comparison (default: None).
        - numeric_tolerance: float, tolerance for numeric column differences (default: 0, exact match).
        - ignore_case: bool, whether to ignore case differences (default: False).  # Changed here
        - detect_reordered_rows: bool, whether to detect reordered rows (default: False).
        - verbose: bool, whether to print progress messages (default: True).

        Returns:
        - dict, summary of differences for each sheet.
        """
        # Define output directory based on file1 basename
        file1_basename = os.path.splitext(os.path.basename(file1))[0]
        output_dir = f"CMP_{file1_basename}"
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        # Load both files into a dictionary of DataFrames
        xl1 = pd.ExcelFile(file1)
        xl2 = pd.ExcelFile(file2)

        # Get the sheets to compare
        sheets1 = set(xl1.sheet_names)
        sheets2 = set(xl2.sheet_names)
        if sheet_name is None:
            sheet_name = list(sheets1 & sheets2)  # Compare only common sheets
        else:
            sheet_name = [sheet for sheet in sheet_name if sheet in sheets1 and sheets2]

        summary = {}
        print(f"Reference file: '{os.path.basename(file1)}'")
        def compare_sheet(sheet):
            
            if verbose:
                print(f"Comparing sheet: {sheet}...")

            # Read sheets as DataFrames
            df1 = xl1.parse(sheet).fillna("NA")
            df2 = xl2.parse(sheet).fillna("NA")

            # Handle case insensitivity
            if ignore_case:
                df1.columns = [col.lower() for col in df1.columns]
                df2.columns = [col.lower() for col in df2.columns]
                df1 = df1.applymap(lambda x: x.lower() if isinstance(x, str) else x)
                df2 = df2.applymap(lambda x: x.lower() if isinstance(x, str) else x)

            # Drop ignored columns
            if ignore_columns:
                df1 = df1.drop(
                    columns=[col for col in ignore_columns if col in df1.columns],
                    errors="ignore",
                )
                df2 = df2.drop(
                    columns=[col for col in ignore_columns if col in df2.columns],
                    errors="ignore",
                )

            # Normalize column order for comparison
            common_cols = df1.columns.intersection(df2.columns)
            df1 = df1[common_cols]
            df2 = df2[common_cols]

            # Specify key columns for comparison
            if key_columns:
                df1 = df1.set_index(key_columns)
                df2 = df2.set_index(key_columns)
            # Identify added and deleted rows based on entire row comparison, not just index
            added_rows = df2[~df2.apply(tuple, 1).isin(df1.apply(tuple, 1))]
            deleted_rows = df1[~df1.apply(tuple, 1).isin(df2.apply(tuple, 1))]

            # Detect reordered rows
            reordered_rows = pd.DataFrame()
            if detect_reordered_rows:
                # Find rows that exist in both DataFrames but are in different positions
                for idx in df1.index:
                    if idx in df2.index:
                        if not df1.loc[idx].equals(df2.loc[idx]):
                            reordered_rows = reordered_rows.append(df1.loc[idx])

            # Detect modified rows (in case of exact matches between the two files)
            aligned_df1 = df1[df1.index.isin(df2.index)]
            aligned_df2 = df2[df2.index.isin(df1.index)]

            if numeric_tolerance > 0:
                modified_rows = aligned_df1.compare(
                    aligned_df2,
                    keep_shape=False,
                    keep_equal=False,
                    result_names=["left", "right"],
                ).pipe(
                    lambda df: df[
                        ~df.apply(
                            lambda row: (
                                abs(row["left"] - row["right"]) <= numeric_tolerance
                                if pd.api.types.is_numeric_dtype(row["left"])
                                else False
                            ),
                            axis=1,
                        )
                    ]
                )
            else:
                modified_rows = aligned_df1.compare(
                    aligned_df2, keep_shape=False, keep_equal=False
                )

            # Save differences to Excel files
            sheet_dir = os.path.join(output_dir, sheet)
            os.makedirs(sheet_dir, exist_ok=True)
            added_path = os.path.join(sheet_dir, f"{sheet}_added.xlsx")
            deleted_path = os.path.join(sheet_dir, f"{sheet}_deleted.xlsx")
            modified_path = os.path.join(sheet_dir, f"{sheet}_modified.xlsx")
            reordered_path = os.path.join(sheet_dir, f"{sheet}_reordered.xlsx")

            if not added_rows.empty:
                added_rows.to_excel(added_path)
            if not deleted_rows.empty:
                deleted_rows.to_excel(deleted_path)
            if not modified_rows.empty:
                modified_rows.to_excel(modified_path)
            if not reordered_rows.empty:
                reordered_rows.to_excel(reordered_path)

            # Return the summary
            return {
                "added_rows": len(added_rows),
                "deleted_rows": len(deleted_rows),
                "modified_rows": len(modified_rows),
                "reordered_rows": len(reordered_rows),
                "added_file": added_path if not added_rows.empty else None,
                "deleted_file": deleted_path if not deleted_rows.empty else None,
                "modified_file": modified_path if not modified_rows.empty else None,
                "reordered_file": reordered_path if not reordered_rows.empty else None,
            }

        # Use ThreadPoolExecutor for parallel processing
        with ThreadPoolExecutor() as executor:
            results = executor.map(compare_sheet, sheet_name)

        # Collect results
        summary = {sheet: result for sheet, result in zip(sheet_name, results)}

        # Save JSON log
        json_path = os.path.join(output_dir, "comparison_summary.json")
        if os.path.exists(json_path):
            with open(json_path, "r") as f:
                existing_data = json.load(f)
        else:
            existing_data = {}

        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        existing_data[timestamp] = summary
        # Sort the existing data by the timestamp in descending order (latest first)
        existing_data = dict(sorted(existing_data.items(), reverse=True))

        with open(json_path, "w") as f:
            json.dump(existing_data, f, indent=4)
        if verbose:
            print(f"Comparison complete. Results saved in '{output_dir}'")

        return summary

    # --- Compare CSV files ---
    def cmp_csv(
        file1,
        file2,
        ignore_case=False,
        numeric_tolerance=0,
        ignore_columns=None,
        verbose=True,
        **kwargs,
    ):
        import pandas as pd

        # Load data and fill NaNs
        df1 = pd.read_csv(file1).fillna("NA")
        df2 = pd.read_csv(file2).fillna("NA")

        # Standardize case if needed
        if ignore_case:
            df1.columns = df1.columns.str.lower()
            df2.columns = df2.columns.str.lower()
            df1 = df1.applymap(lambda x: x.lower() if isinstance(x, str) else x)
            df2 = df2.applymap(lambda x: x.lower() if isinstance(x, str) else x)

        # Drop ignored columns
        if ignore_columns:
            ignore_columns = [col.lower() if ignore_case else col for col in ignore_columns]
            df1.drop(columns=[col for col in ignore_columns if col in df1.columns], errors="ignore", inplace=True)
            df2.drop(columns=[col for col in ignore_columns if col in df2.columns], errors="ignore", inplace=True)

        # Reset index to ensure alignment
        df1.reset_index(drop=True, inplace=True)
        df2.reset_index(drop=True, inplace=True)

        # Align DataFrames by columns
        df1, df2 = df1.align(df2, join="inner", axis=1)

        # Compare rows
        added_rows = df2[~df2.apply(tuple, axis=1).isin(df1.apply(tuple, axis=1))]
        deleted_rows = df1[~df1.apply(tuple, axis=1).isin(df2.apply(tuple, axis=1))]

        # Compare modified rows
        if numeric_tolerance > 0:
            def numeric_diff(row):
                if pd.api.types.is_numeric_dtype(row["left"]):
                    return abs(row["left"] - row["right"]) > numeric_tolerance
                return row["left"] != row["right"]

            modified_rows = df1.compare(df2, keep_shape=True, keep_equal=False)
            modified_rows = modified_rows[modified_rows.apply(numeric_diff, axis=1)]
        else:
            modified_rows = df1.compare(df2, keep_shape=True, keep_equal=False)

        # Return results
        return {
            "added_rows": len(added_rows),
            "deleted_rows": len(deleted_rows),
            "modified_rows": len(modified_rows),
            "added_file": added_rows if not added_rows.empty else pd.DataFrame(),
            "deleted_file": deleted_rows if not deleted_rows.empty else pd.DataFrame(),
            "modified_file": modified_rows if not modified_rows.empty else pd.DataFrame(),
        }

    # --- Compare JSON files ---
    def cmp_json(
        file1, file2, ignore_case=False, numeric_tolerance=0, verbose=True, **kwargs
    ):
        import json

        with open(file1, "r") as f1:
            json1 = json.load(f1)
        with open(file2, "r") as f2:
            json2 = json.load(f2)

        # Normalize case and compare JSONs
        if ignore_case:

            def normalize(obj):
                if isinstance(obj, dict):
                    return {k.lower(): normalize(v) for k, v in obj.items()}
                elif isinstance(obj, list):
                    return [normalize(item) for item in obj]
                elif isinstance(obj, str):
                    return obj.lower()
                else:
                    return obj

            json1 = normalize(json1)
            json2 = normalize(json2)

        # Compare JSONs
        def compare_json(obj1, obj2):
            if isinstance(obj1, dict) and isinstance(obj2, dict):
                added_keys = {k: obj2[k] for k in obj2 if k not in obj1}
                deleted_keys = {k: obj1[k] for k in obj1 if k not in obj2}
                modified_keys = {
                    k: (obj1[k], obj2[k])
                    for k in obj1
                    if k in obj2 and obj1[k] != obj2[k]
                }
                return added_keys, deleted_keys, modified_keys

            elif isinstance(obj1, list) and isinstance(obj2, list):
                added_items = [item for item in obj2 if item not in obj1]
                deleted_items = [item for item in obj1 if item not in obj2]
                modified_items = [
                    (item1, item2) for item1, item2 in zip(obj1, obj2) if item1 != item2
                ]
                return added_items, deleted_items, modified_items

            else:
                if obj1 != obj2:
                    return obj1, obj2, None
                else:
                    return None, None, None

        added, deleted, modified = compare_json(json1, json2)

        return {"added_keys": added, "deleted_keys": deleted, "modified_keys": modified}

    # --- Compare Text files ---
    def cmp_txt(
        file1, file2, ignore_case=False, numeric_tolerance=0, verbose=True, **kwargs
    ):
        def read_lines(file):
            with open(file, "r") as f:
                return f.readlines()

        lines1 = read_lines(file1)
        lines2 = read_lines(file2)

        if ignore_case:
            lines1 = [line.lower() for line in lines1]
            lines2 = [line.lower() for line in lines2]

        added_lines = [line for line in lines2 if line not in lines1]
        deleted_lines = [line for line in lines1 if line not in lines2]

        modified_lines = []
        if numeric_tolerance > 0:
            for line1, line2 in zip(lines1, lines2):
                if abs(float(line1) - float(line2)) > numeric_tolerance:
                    modified_lines.append((line1, line2))
        else:
            for line1, line2 in zip(lines1, lines2):
                if line1 != line2:
                    modified_lines.append((line1, line2))

        return {
            "added_lines": added_lines,
            "deleted_lines": deleted_lines,
            "modified_lines": modified_lines,
        }

    if kind is None:
        kind = os.path.splitext(file1)[1].lower()[1:]
    # Compare based on the file type
    if kind == "xlsx":
        return cmp_excel(file1=file1, file2=file2, verbose=verbose, **kwargs)

    elif kind == "csv":
        return cmp_csv(file1=file1, file2=file2, verbose=verbose, **kwargs)

    elif kind == "json":
        return cmp_json(file1=file1, file2=file2, verbose=verbose, **kwargs)

    elif kind == "txt":
        return cmp_txt(file1=file1, file2=file2, verbose=verbose, **kwargs)

    else:
        raise ValueError(f"Unsupported file type: {kind}")
def cn2pinyin(
    cn_str: Union[str, list] = None,
    sep: str = " ",
    fmt: str = "normal",  # which style you want to set
):
    from pypinyin import pinyin, Style

    """
    Converts Chinese characters to Pinyin.
    usage: 
        cn2pinyin(cn_str, sep="_", fmt="tone")
    Args:
        cn_str (str): Chinese string to convert.
        sep (str): Separator for the output Pinyin string.
        fmt (Style): "normal","tone", "tone2","tone3","finals","finals_tone","finals_tone2","finals_tone3","initials","bopomofo","bopomofo_first","cyrillic","pl",
    Returns:
        cn_str: The Pinyin representation of the Chinese string.
    """
    fmts = [
        "normal",
        "tone",
        "tone2",
        "tone3",
        "finals",
        "finals_tone",
        "finals_tone2",
        "finals_tone3",
        "initials",
        "bopomofo",
        "bopomofo_first",
        "cyrillic",
        "pl",
    ]
    fmt = strcmp(fmt, fmts)[0]
    if fmt == "normal":
        style = Style.NORMAL
    elif fmt == "tone":
        style = Style.TONE
    elif fmt == "tone2":
        style = Style.TONE2
    elif fmt == "tone3":
        style = Style.TONE3
    elif fmt == "finals":
        style = Style.FINALS
    elif fmt == "finals_tone":
        style = Style.FINALS_TONE
    elif fmt == "finals_tone2":
        style = Style.FINALS_TONE2
    elif fmt == "finals_tone3":
        style = Style.FINALS_TONE3
    elif fmt == "initials":
        style = Style.INITIALS
    elif fmt == "bopomofo":
        style = Style.BOPOMOFO
    elif fmt == "bopomofo_first":
        style = Style.BOPOMOFO_FIRST
    elif fmt == "cyrillic":
        style = Style.CYRILLIC
    elif fmt == "pl":
        style = Style.PL
    else:
        style = Style.NORMAL
    if not isinstance(cn_str, list):
        cn_str = [cn_str]
    pinyin_flat = []
    for cn_str_ in cn_str:
        pinyin_string = pinyin(cn_str_, style=style)
        pinyin_flat.append(sep.join([item[0] for item in pinyin_string]))
    if len(pinyin_flat) == 1:
        return pinyin_flat[0]
    else:
        return pinyin_flat


def counter(list_, verbose=True):
    from collections import Counter

    c = Counter(list_)
    # Print the name counts
    for item, count in c.items():
        if verbose:
            print(f"{item}: {count}")
    return c


# usage:
# print(f"Return an iterator over elements repeating each as many times as its count:\n{sorted(c.elements())}")
# print(f"Return a list of the n most common elements:\n{c.most_common()}")
# print(f"Compute the sum of the counts:\n{c.total()}")

def dict2df(dict_, fill=None, axis=0):
    """
    Convert a dictionary to a DataFrame with flexible axis and padding options.

    Parameters:
    - dict_: The dictionary to convert (keys are columns or index).
    - fill: Value to fill in case of shorter lists.
    - axis: Axis for DataFrame construction (0 for columns, 1 for rows).

    Returns:
    - DataFrame created from the dictionary.
    """
    for key, value in dict_.items():
        if not isinstance(value, list):
            dict_[key] = [value]
            print(f"'{key}' is not a list. trying to convert it to 'list'")

    # Get the maximum length of values
    len_max = max(len(value) for value in dict_.values())

    # Extend lists to match the length of the longest list
    for key, value in dict_.items():
        if isinstance(value, list):
            value.extend([fill] * (len_max - len(value)))  # Fill shorter lists
        dict_[key] = value

    # If axis=0, the dictionary keys will be treated as column names
    if axis == 0:
        return pd.DataFrame(dict_)
    # If axis=1, the dictionary keys will be treated as index names (rows)
    else:
        return pd.DataFrame(dict_).transpose()

def text2audio(
    text,
    method=None,  # "pyttsx3","gTTS"
    rate=200,
    slow=False,  # "gTTS"
    volume=1.0,
    voice=None,
    lang=None,
    gender=None,
    age=None,
    dir_save=None,
):
    """
    # sample_text = "Hello! This is a test of the pyttsx3 text-to-speech system."
    # sample_text = "这个是中文, 测试"
    # sample_text = "Hallo, ich bin echo, Wie Heissen Sie"

    # text2audio(
    #     text=sample_text,
    #     rate=150,
    #     volume=0.9,
    #     # voice=None,  # Replace with a voice name or ID available on your system
    # )
    """
    if method is not None:
        methods = ["gTTS", "pyttsx3", "google"]
        method = strcmp(method, methods)[0]
    else:
        try:
            text2audio(
                text,
                method="google",
                rate=rate,
                slow=slow,
                volume=volume,
                voice=voice,
                lang=lang,
                gender=gender,
                age=age,
                dir_save=dir_save,
            )
        except Exception as e:
            print(e)
            text2audio(
                text,
                method="pyttsx3",
                rate=rate,
                slow=slow,
                volume=volume,
                voice=voice,
                lang=lang,
                gender=gender,
                age=age,
                dir_save=dir_save,
            )

    if method == "pyttsx3":
        import pyttsx3

        try:
            engine = pyttsx3.init()
            engine.setProperty("rate", rate)
            if 0.0 <= volume <= 1.0:
                engine.setProperty("volume", volume)
            else:
                raise ValueError("Volume must be between 0.0 and 1.0")

            if gender is not None:
                gender = strcmp(gender, ["male", "female"])[0]
            if age is not None:
                if isinstance(age, (float, int)):
                    if age <= 10:
                        age = "child"
                    elif 10 < age < 18:
                        age = "senior"
                    else:
                        age = "adult"
                elif isinstance(age, str):
                    age = strcmp(age, ["child", "adult", "senior"])[0]
                else:
                    raise ValueError("age: should be in ['child', 'adult', 'senior']")
            voices = engine.getProperty("voices")
            if voice is None:
                if lang is None:
                    voice = strcmp(detect_lang(text), [v.name for v in voices])[0]
                else:
                    if run_once_within():
                        print([v.name for v in voices])
                    print(f"lang:{lang}")
                    voice = strcmp(lang, [v.name for v in voices])[0]
            selected_voice = None

            for v in voices:
                # Check if the voice matches the specified gender or age
                if voice and (voice.lower() in v.name.lower() or voice in v.id):
                    selected_voice = v
                    break
                if gender and gender.lower() in v.name.lower():
                    selected_voice = v
                if age and age.lower() in v.name.lower():
                    selected_voice = v

            if selected_voice:
                engine.setProperty("voice", selected_voice.id)
            else:
                if voice or gender or age:
                    raise ValueError(
                        f"No matching voice found for specified criteria. Available voices: {[v.name for v in voices]}"
                    )
            # Generate audio
            if dir_save:
                engine.save_to_file(text, dir_save)
                print(f"Audio saved to {dir_save}")
            else:
                engine.say(text)

            engine.runAndWait()
        except Exception as e:
            print(f"An error occurred: {e}")
            #     # Explicitly terminate the pyttsx3 engine to release resources
            try:
                engine.stop()
            except RuntimeError:
                pass
                # Safely exit the script if running interactively to avoid kernel restarts
            try:
                import sys

                sys.exit()
            except SystemExit:
                pass
    elif method.lower() in ["google", "gtts"]:
        from gtts import gTTS

        try:
            if lang is None:
                from langdetect import detect

                lang = detect(text)
            # Initialize gTTS with the provided parameters
            tts = gTTS(text=text, lang=lang, slow=slow)
        except Exception as e:
            print(f"An error occurred: {e}")

        print("not realtime reading...")
        if dir_save:
            if "." not in dir_save:
                dir_save = dir_save + ".mp3"
            tts.save(dir_save)
            print(f"Audio saved to {dir_save}")
        else:
            dir_save = "temp_audio.mp3"
            if "." not in dir_save:
                dir_save = dir_save + ".mp3"
            tts.save(dir_save)
        try:
            fopen(dir_save)
        except Exception as e:
            print(f"Error opening file: {e}")
    print("done")


def str2time(time_str, fmt="24"):
    """
    Convert a time string into the specified format.
    Parameters:
    - time_str (str): The time string to be converted.
    - fmt (str): The format to convert the time to. Defaults to '%H:%M:%S'.
    Returns:
        %I represents the hour in 12-hour format.
        %H represents the hour in 24-hour format (00 through 23).
        %M represents the minute.
        %S represents the second.
        %p represents AM or PM.
    - str: The converted time string.
    """
    from datetime import datetime

    def time_len_corr(time_str):
        time_str_ = (
            ssplit(time_str, by=[":", " ", "digital_num"]) if ":" in time_str else None
        )
        time_str_split = []
        [time_str_split.append(i) for i in time_str_ if is_num(i)]
        if time_str_split:
            if len(time_str_split) == 2:
                H, M = time_str_split
                time_str_full = H + ":" + M + ":00"
            elif len(time_str_split) == 3:
                H, M, S = time_str_split
                time_str_full = H + ":" + M + ":" + S
        else:
            time_str_full = time_str_
        if "am" in time_str.lower():
            time_str_full += " AM"
        elif "pm" in time_str.lower():
            time_str_full += " PM"
        return time_str_full

    if "12" in fmt:
        fmt = "%I:%M:%S %p"
    elif "24" in fmt:
        fmt = "%H:%M:%S"

    try:
        # Try to parse the time string assuming it could be in 24-hour or 12-hour format
        time_obj = datetime.strptime(time_len_corr(time_str), "%H:%M:%S")
    except ValueError:
        try:
            time_obj = datetime.strptime(time_len_corr(time_str), "%I:%M:%S %p")
        except ValueError as e:
            raise ValueError(f"Unable to parse time string: {time_str}. Error: {e}")

    # Format the time object to the desired output format
    formatted_time = time_obj.strftime(fmt)
    return formatted_time


# # Example usage:
# time_str1 = "14:30:45"
# time_str2 = "02:30:45 PM"

# formatted_time1 = str2time(time_str1, fmt='12')  # Convert to 12-hour format
# formatted_time2 = str2time(time_str2, fmt='24')    # Convert to 24-hour format

# print(formatted_time1)  # Output: 02:30:45 PM
# print(formatted_time2)  # Output: 14:30:45


def str2date(date_str, original_fmt=None, fmt="%Y-%m-%d"):
    """
    Convert a date string to the desired format and extract components if needed.
    Usage:
        str2date(x, fmt="%d.%m.%y",original_fmt="%d.%m.%y")
    Parameters:
    - date_str (str): The input date string.
    - original_fmt (str, optional): The original format of the date string. If not provided, it will be auto-detected.
    - fmt (str): The desired format for the output date string. Defaults to '%Y-%m-%d'.

    Returns:
    - dict: A dictionary containing the converted date string and its components (year, month, day).
    
    Raises:
    - ValueError: If the date cannot be parsed.
    """ 
    from dateutil import parser
    try:
        if not isinstance(date_str,str):
            date_str=str(date_str) 
        # Parse the date using the provided original format or auto-detect
        if original_fmt:
            try:
                date_obj = datetime.strptime(date_str, original_fmt)
            except Exception as e:
                print(e)
                date_obj=None
        else:
            try:
                date_obj = parser.parse(date_str)
            except Exception as e:
                print(e)
                date_obj=None
        # Return formatted string if `fmt` is specified, otherwise return the datetime object
        if date_obj is not None:
            if fmt: 
                date_obj=date_obj.strftime(fmt)
        else:
            date_obj=date_str
        return date_obj

    except (ValueError, TypeError) as e:
        raise ValueError(f"Unable to process date string: '{date_str}'. Error: {e}")


# str1=str2date(num2str(20240625),fmt="%a %d-%B-%Y")
# print(str1)
# str2=str2num(str2date(str1,fmt='%a %Y%m%d'))
# print(str2)
 
def str2num(
    s: str,
    *args,
    sep: Optional[Union[str, List[str]]] = None,
    round_digits: Optional[int] = None,
    return_list: bool = True,
    handle_text: bool = True
) -> Union[float, int, List[Union[float, int]], None]:
    """
    # Examples
        print(str2num("123"))                # Output: 123
        print(str2num("123.456", 2))         # Output: 123.46
        print(str2num("one hundred and twenty three"))  # Output: 123
        print(str2num("seven million"))      # Output: 7000000
        print(str2num('one thousand thirty one',','))  # Output: 1,031
        print(str2num("12345.6789", ","))    # Output: 12,345.6789
        print(str2num("12345.6789", " ", 2)) # Output: 12 345.68
        print(str2num('111113.34555',3,',')) # Output: 111,113.346
        print(str2num("123.55555 sec miniuets",3)) # Output: 1.3
        print(str2num("every 3,300.55 hours and 5.045555 min", sep=",", round=1))
        print(str2num("five hundred fourty one"), str2num(
            "this is 5.9435 euros for 10.04499 killograme", round=3
        )[0])
    Convert a string containing numeric or textual data into an integer, float, or list of numbers.

    Parameters:
    - s (str): Input string containing a number or textual representation of a number.
    - *args: Additional arguments for delimiter or rounding digits.
    - sep (str or list): Delimiter(s) to remove from the string (e.g., ',' or ['.', ',']).
    - round_digits (int): Number of decimal places to round the result to.
    - return_list (bool): Whether to return a list of numbers if multiple are found.
    - handle_text (bool): Whether to process textual numbers using the numerizer library.

    Returns:
    - Union[float, int, List[Union[float, int]], None]: Converted number(s) or None if conversion fails.
    """
    import re
    from numerizer import numerize

    if not isinstance(s, str):
        return None

    # Merge args with explicit parameters
    if sep is None:
        sep = []
    elif isinstance(sep, str):
        sep = [sep]
    for arg in args:
        if isinstance(arg, str):
            sep.append(arg)
        elif isinstance(arg, int) and round_digits is None:
            round_digits = arg

    # Remove all specified delimiters
    for delimiter in sep:
        s = s.replace(delimiter, "")

    # Attempt conversion
    def try_convert(segment: str) -> Union[float, int, None]:
        try:
            return int(segment)
        except ValueError:
            try:
                return float(segment)
            except ValueError:
                return None

    # Handle textual numbers
    if handle_text:
        try:
            s = numerize(s)
        except Exception:
            pass

    # Extract numeric segments
    number_segments = re.findall(r"[-+]?\d*\.\d+|\d+", s)
    numbers = [try_convert(seg) for seg in number_segments if seg]
    numbers = [num for num in numbers if num is not None]

    if not numbers:
        return None  # No valid numbers found

    # Single or multiple numbers
    if len(numbers) == 1 and not return_list:
        result = numbers[0]
    else:
        result = (
            numbers[0] if len(numbers) == 1 else numbers if return_list else numbers[0]
        )

    # Apply rounding if necessary
    if round_digits is not None:
        if isinstance(result, list):
            result = [round(num + 1e-10, round_digits) for num in result]
        else:
            result = round(result + 1e-10, round_digits)

        # Convert to int if rounding to 0 digits
        if round_digits == 0:
            if isinstance(result, list):
                result = [int(num) for num in result]
            else:
                result = int(result)

    return result 
def num2str(num, *args, **kwargs):
    delimiter = kwargs.get("sep", None)
    round_digits = kwargs.get("round", None)

    # Parse additional arguments
    for arg in args:
        if isinstance(arg, str):
            delimiter = arg
        elif isinstance(arg, int):
            round_digits = arg

    # Apply rounding if specified
    if round_digits is not None:
        num = round(num, round_digits)

    # Convert number to string
    num_str = f"{num}"

    # Apply delimiter if specified
    if delimiter is not None:
        num_str = num_str.replace(".", ",")  # Replace decimal point with comma
        num_str_parts = num_str.split(",")
        if len(num_str_parts) > 1:
            integer_part = num_str_parts[0]
            decimal_part = num_str_parts[1]
            integer_part = "{:,}".format(int(integer_part))
            num_str = integer_part + "." + decimal_part
        else:
            num_str = "{:,}".format(int(num_str_parts[0]))

    return num_str


# Examples
# print(num2str(123), type(num2str(123)))  # Output: "123"
# print(num2str(123.456, round=2), type(num2str(123.456, 2)))  # Output: "123.46"
# print(num2str(7000.125, round=1), type(num2str(7000.125, 2)))  # Output: "7000.13"
# print(
#     num2str(12345333.6789, sep=","), type(num2str(12345.6789, ","))
# )  # Output: "12,345.6789"
# print(num2str(7000.00, sep=","), type(num2str(7000.00, ",")))  # Output: "7,000.00"


# Helper to convert text or list of text to HTML
def str2html(text_list, strict=False):
    if not isinstance(text_list, list):
        text_list = [text_list]

    # Define a mapping for escape sequences to their HTML representations
    escape_mapping = {
        "\\": "&bsol;",  # Backslash
        "'": "&apos;",  # Single quote
        '"': "&quot;",  # Double quote
        "\n": "<br>",  # Newline
        "\r": "",  # Carriage return (not represented in HTML)
        "\t": "&nbsp;&nbsp;&nbsp;&nbsp;",  # Tab (4 spaces)
        "\b": "",  # Backspace (not typically represented)
        "\f": "",  # Form feed (not typically represented)
        "\v": "",  # Vertical tab (not typically represented)
        "\a": "",  # Bell/alert sound (not typically represented)
        "\0": "",  # Null character (not typically represented)
    }

    # Convert text by replacing escape sequences with their HTML equivalents
    html_content = ""
    for text in text_list:
        for escape_seq, html_rep in escape_mapping.items():
            text = text.replace(escape_seq, html_rep)
        html_content += text.replace("\n", "<br>")  # Add line breaks for newlines

    if strict:
        html_content = "<html><body>\n" + html_content + "\n</body></html>"

    return html_content


def cm2px(*cm, dpi=300) -> list:
    # Case 1: When the user passes a single argument that is a list or tuple, such as cm2px([8, 5]) or inch2cm((8, 5))
    if len(cm) == 1 and isinstance(cm[0], (list, tuple)):
        # If the input is a single list or tuple, we unpack its elements and convert each to cm
        return [i / 2.54 * dpi for i in cm[0]]
    # Case 2: When the user passes multiple arguments directly, such as cm2px(8, 5)
    else:
        return [i / 2.54 * dpi for i in cm]


def px2cm(*px, dpi=300) -> list:
    # Case 1: When the user passes a single argument that is a list or tuple, such as px2cm([8, 5]) or inch2cm((8, 5))
    if len(px) == 1 and isinstance(px[0], (list, tuple)):
        # If the input is a single list or tuple, we unpack its elements and convert each to cm
        return [i * 2.54 / dpi for i in px[0]]
    # Case 2: When the user passes multiple arguments directly, such as px2cm(8, 5)
    else:
        # Here, we convert each individual argument directly to cm
        return [i * 2.54 / dpi for i in px]


def px2inch(*px, dpi=300) -> list:
    """
    px2inch: converts pixel measurements to inches based on the given dpi.
    Usage:
    px2inch(300, 600, dpi=300); px2inch([300, 600], dpi=300)
    Returns:
        list: in inches
    """
    # Case 1: When the user passes a single argument that is a list or tuple, such as px2inch([300, 600]) or px2inch((300, 600))
    if len(px) == 1 and isinstance(px[0], (list, tuple)):
        # If the input is a single list or tuple, we unpack its elements and convert each to inches
        return [i / dpi for i in px[0]]
    # Case 2: When the user passes multiple arguments directly, such as px2inch(300, 600)
    else:
        # Here, we convert each individual argument directly to inches
        return [i / dpi for i in px]


def inch2cm(*cm) -> list:
    """
    cm2inch: converts centimeter measurements to inches.
    Usage:
    cm2inch(10, 12.7); cm2inch((10, 12.7)); cm2inch([10, 12.7])
    Returns:
        list: in inches
    """
    # Case 1: When the user passes a single argument that is a list or tuple, such as cm2inch([10, 12.7]) or cm2inch((10, 12.7))
    if len(cm) == 1 and isinstance(cm[0], (list, tuple)):
        # If the input is a single list or tuple, we unpack its elements and convert each to inches
        return [i * 2.54 for i in cm[0]]
    # Case 2: When the user passes multiple arguments directly, such as cm2inch(10, 12.7)
    else:
        # Here, we convert each individual argument directly to inches
        return [i * 2.54 for i in cm]


def inch2px(*inch, dpi=300) -> list:
    """
    inch2px: converts inch measurements to pixels based on the given dpi.

    Usage:
    inch2px(1, 2, dpi=300); inch2px([1, 2], dpi=300)

    Parameters:
    inch : float, list, or tuple
        Single or multiple measurements in inches to convert to pixels.
    dpi : int, optional (default=300)
        Dots per inch (DPI), representing the pixel density.

    Returns:
        list: Converted measurements in pixels.
    """
    # Case 1: When the user passes a single argument that is a list or tuple, e.g., inch2px([1, 2]) or inch2px((1, 2))
    if len(inch) == 1 and isinstance(inch[0], (list, tuple)):
        return [i * dpi for i in inch[0]]

    # Case 2: When the user passes multiple arguments directly, e.g., inch2px(1, 2)
    else:
        return [i * dpi for i in inch]


def cm2inch(*inch) -> list:
    """
    Usage:
    inch2cm(8,5); inch2cm((8,5)); inch2cm([8,5])
    Returns:
        list: in centimeters
    """
    # Case 1: When the user passes a single argument that is a list or tuple, such as inch2cm([8, 5]) or inch2cm((8, 5))
    if len(inch) == 1 and isinstance(inch[0], (list, tuple)):
        # If the input is a single list or tuple, we unpack its elements and convert each to cm
        return [i / 2.54 for i in inch[0]]
    # Case 2: When the user passes multiple arguments directly, such as inch2cm(8, 5)
    else:
        # Here, we convert each individual argument directly to cm
        return [i / 2.54 for i in inch]



def sqlite2sql(db_path, sql_path):
    """
    Export an SQLite database to an SQL file, including schema and data for all tables.

    :param db_path: Path to the SQLite .db file
    :param output_file: Path to the output .sql file

    # Usage
        db_path = "your_database.db"  # Replace with the path to your SQLite database
        sql_path = "output.sql"  # Replace with your desired output file name
        export_sqlite_to_sql(db_path, sql_path)

    """
    import sqlite3
    try:
        # Connect to the SQLite database
        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()

        with open(sql_path, 'w') as f:
            # Write a header for the SQL dump
            f.write("-- SQLite Database Dump\n")
            f.write(f"-- Source: {db_path}\n\n")
            
            # Retrieve all table names
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table' AND name NOT LIKE 'sqlite_%';")
            tables = [row[0] for row in cursor.fetchall()]

            for table in tables:
                # Write the schema for the table
                cursor.execute(f"SELECT sql FROM sqlite_master WHERE type='table' AND name='{table}';")
                schema = cursor.fetchone()
                if schema:
                    f.write(f"{schema[0]};\n\n")
                
                # Write data for the table
                cursor.execute(f"SELECT * FROM {table};")
                rows = cursor.fetchall()
                if rows:
                    cursor.execute(f"PRAGMA table_info({table});")
                    column_names = [info[1] for info in cursor.fetchall()]
                    column_list = ', '.join(f'"{col}"' for col in column_names)

                    for row in rows:
                        values = ', '.join(f"'{str(val).replace('\'', '\'\'')}'" if val is not None else 'NULL' for val in row)
                        f.write(f"INSERT INTO {table} ({column_list}) VALUES ({values});\n")
                
                f.write("\n")
            
            print(f"Database exported successfully to {sql_path}")
    
    except sqlite3.Error as e:
        print(f"SQLite error: {e}")
    except Exception as e:
        print(f"Unexpected error: {e}")
    finally:
        # Ensure the connection is closed
        if conn:
            conn.close()


def sreplace(*args, **kwargs):
    """
    sreplace(text, by=None, robust=True)
    Replace specified substrings in the input text with provided replacements.
    Args:
        text (str): The input text where replacements will be made.
        by (dict, optional): A dictionary containing substrings to be replaced as keys
            and their corresponding replacements as values. Defaults to {".com": "..come", "\n": " ", "\t": " ", "  ": " "}.
        robust (bool, optional): If True, additional default replacements for newline and tab characters will be applied.
                                Default is False.
    Returns:
        str: The text after replacements have been made.
    """
    text = None
    by = kwargs.get("by", None)
    robust = kwargs.get("robust", True)

    for arg in args:
        if isinstance(arg, str):
            text = arg
        elif isinstance(arg, dict):
            by = arg
        elif isinstance(arg, bool):
            robust = arg
        else:
            Error(f"{type(arg)} is not supported")

    # Default replacements for newline and tab characters
    default_replacements = {
        "\a": "",
        "\b": "",
        "\f": "",
        "\n": "",
        "\r": "",
        "\t": "",
        "\v": "",
        "\\": "",  # Corrected here
        # "\?": "",
        "�": "",
        "\\x": "",  # Corrected here
        "\\x hhhh": "",
        "\\ ooo": "",  # Corrected here
        "\xa0": "",
        "  ": " ",
    }

    # If dict_replace is None, use the default dictionary
    if by is None:
        by = {}
    # If robust is True, update the dictionary with default replacements
    if robust:
        by.update(default_replacements)

    # Iterate over each key-value pair in the dictionary and replace substrings accordingly
    for k, v in by.items():
        text = text.replace(k, v)
    return text


# usage:
# sreplace(text, by=dict(old_str='new_str'), robust=True)


def paper_size(paper_type_str="a4"):
    df = pd.DataFrame(
        {
            "a0": [841, 1189],
            "a1": [594, 841],
            "a2": [420, 594],
            "a3": [297, 420],
            "a4": [210, 297],
            "a5": [148, 210],
            "a6": [105, 148],
            "a7": [74, 105],
            "b0": [1028, 1456],
            "b1": [707, 1000],
            "b2": [514, 728],
            "b3": [364, 514],
            "b4": [257, 364],
            "b5": [182, 257],
            "b6": [128, 182],
            "letter": [215.9, 279.4],
            "legal": [215.9, 355.6],
            "business card": [85.6, 53.98],
            "photo china passport": [33, 48],
            "passport single": [125, 88],
            "visa": [105, 74],
            "sim": [25, 15],
        }
    )
    for name in df.columns:
        if paper_type_str in name.lower():
            paper_type = name
    if not paper_type:
        paper_type = "a4"  # default
    return df[paper_type].tolist()


def docx2pdf(dir_docx, dir_pdf=None):
    from docx2pdf import convert

    if dir_pdf:
        convert(dir_docx, dir_pdf)
    else:
        convert(dir_docx)


def img2pdf(dir_img, kind=None, page=None, dir_save=None, page_size="a4", dpi=300):
    import img2pdf as image2pdf

    def mm_to_point(size):
        return (image2pdf.mm_to_pt(size[0]), image2pdf.mm_to_pt(size[1]))

    def set_dpi(x):
        dpix = dpiy = x
        return image2pdf.get_fixed_dpi_layout_fun((dpix, dpiy))

    if kind is None:
        _, kind = os.path.splitext(dir_img)
    if not kind.startswith("."):
        kind = "." + kind
    if dir_save is None:
        dir_save = dir_img.replace(kind, ".pdf")
    imgs = []
    if os.path.isdir(dir_img):
        if not dir_save.endswith(".pdf"):
            dir_save += "#merged_img2pdf.pdf"
        if page is None:
            select_range = listdir(dir_img, kind=kind).fpath
        else:
            if not isinstance(page, (np.ndarray, list, range)):
                page = [page]
            select_range = listdir(dir_img, kind=kind)["fpath"][page]
        for fname in select_range:
            if not fname.endswith(kind):
                continue
            path = os.path.join(dir_img, fname)
            if os.path.isdir(path):
                continue
            imgs.append(path)
    else:
        imgs = [
            # os.path.isdir(dir_img),
            dir_img
        ]
    print(imgs)
    if page_size:
        if isinstance(page_size, str):
            pdf_in_mm = mm_to_point(paper_size(page_size))
        else:
            print("default: page_size = (210,297)")
            pdf_in_mm = mm_to_point(page_size)
            print(f"page size was set to {page_size}")
        p_size = image2pdf.get_layout_fun(pdf_in_mm)
    else:
        p_size = set_dpi(dpi)
    with open(dir_save, "wb") as f:
        f.write(image2pdf.convert(imgs, layout_fun=p_size))


# usage:
# dir_img="/Users/macjianfeng/Dropbox/00-Personal/2015-History/2012-2015_兰州大学/120901-大学课件/生物统计学 陆卫/复习题/"
# img2pdf(dir_img,kind='tif', page=range(3,7,2))


def pdf2ppt(dir_pdf, dir_ppt):
    from PyPDF2 import PdfReader
    from pptx.util import Inches
    from pptx import Presentation

    prs = Presentation()

    # Open the PDF file
    with open(dir_pdf, "rb") as f:
        reader = PdfReader(f)
        num_pages = len(reader.pages)

        # Iterate through each page in the PDF
        for page_num in range(num_pages):
            page = reader.pages[page_num]
            text = page.extract_text()

            # Add a slide for each page's content
            slide_layout = prs.slide_layouts[
                5
            ]  # Use slide layout that suits your needs
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = f"Page {page_num + 1}"
            slide.shapes.add_textbox(
                Inches(1), Inches(1.5), Inches(8), Inches(5)
            ).text = text

    # Save the PowerPoint presentation
    prs.save(dir_ppt)
    print(f"Conversion from {dir_pdf} to {dir_ppt} complete.")


def ssplit(text, by="space", verbose=False, strict=False, **kws):
    import re

    if isinstance(text, list):
        nested_list = [ssplit(i, by=by, verbose=verbose, **kws) for i in text]
        flat_list = [item for sublist in nested_list for item in sublist]
        return flat_list

    def split_by_word_length(text, length):
        return [word for word in text.split() if len(word) == length]

    def split_by_multiple_delimiters(text, delimiters):
        regex_pattern = "|".join(map(re.escape, delimiters))
        return re.split(regex_pattern, text)

    def split_by_camel_case(text):
        return re.findall(r"[A-Z](?:[a-z]+|[A-Z]*(?=[A-Z]|$))", text)

    def split_at_upper_fl_lower(text):
        return re.findall(r"[A-Z](?:[a-z]+|[A-Z]+(?=[A-Z]|$))", text)

    def split_at_lower_fl_upper(text):
        split_text = re.split(r"(?<=[a-z])(?=[A-Z])", text)
        return split_text

    def split_at_upper(text):
        split_text = re.split(r"(?=[A-Z])", text)
        split_text = [part for part in split_text if part]
        return split_text

    def split_by_regex_lookahead(text, pattern):
        return re.split(f"(?<={pattern})", text)

    def split_by_regex_end(text, pattern):
        return re.split(f"(?={pattern})", text)

    # def split_by_sentence_endings(text):
    #     return re.split(r"(?<=[.!?])", text)
    def split_non_ascii(text):
        # return re.split(r"([^\x00-\x7F\w\s,.!?:\"'()\-]+)", text)
        # return re.split(r"[^\x00-\x7F]+", text)
        return re.split(r"([^\x00-\x7F]+)", text)

    def split_by_consecutive_non_alphanumeric(text):
        return re.split(r"\W+", text)

    def split_by_fixed_length_chunks(text, length):
        return [text[i : i + length] for i in range(0, len(text), length)]

    def split_by_sent_num(text, n=10):
        from nltk.tokenize import sent_tokenize
        from itertools import pairwise

        # split text into sentences
        text_split_by_sent = sent_tokenize(text)
        cut_loc_array = np.arange(0, len(text_split_by_sent), n)
        if cut_loc_array[-1] != len(text_split_by_sent):
            cut_loc = np.append(cut_loc_array, len(text_split_by_sent))
        else:
            cut_loc = cut_loc_array
        # get text in section (e.g., every 10 sentences)
        text_section = []
        for i, j in pairwise(cut_loc):
            text_section.append(text_split_by_sent[i:j])
        return text_section

    def split_general(text, by, verbose=False, ignore_case=False):
        if ignore_case:
            if verbose:
                print(f"used {by} to split, ignore_case=True")
            pattern = re.compile(re.escape(by), re.IGNORECASE)
            split_text = pattern.split(text)
            return split_text
        else:
            if verbose:
                print(f"used {by} to split, ignore_case=False")
            return text.split(by)

    def reg_split(text, pattern):
        return re.split(pattern, text)

    if ("sp" in by or "white" in by) and not strict:
        if verbose:
            print(f"splited by space")
        return text.split()
    elif ("word" in by and "len" in by) and not strict:
        if verbose:
            print(f"split_by_word_length(text, length)")
        return split_by_word_length(text, **kws)  # split_by_word_length(text, length)
    # elif "," in by:
    #     if verbose:
    #         print(f"splited by ','")
    #     return text.split(",")
    elif isinstance(by, list):
        if verbose:
            print(f"split_by_multiple_delimiters: ['|','&']")
        return split_by_multiple_delimiters(text, by)
    elif (
        all([("digi" in by or "num" in by), not "sent" in by, not "str" in by])
        and not strict
    ):
        if verbose:
            print(f"splited by digital (numbers)")
        return re.split(r"(\d+)", text)
    elif all([("digi" in by or "num" in by), "str" in by]) and not strict:
        if verbose:
            print(f"Splitting by (number strings)")
        pattern = re.compile(
            r"\b((?:one|two|three|four|five|six|seven|eight|nine|ten|eleven|twelve|thirteen|fourteen|fifteen|sixteen|seventeen|eighteen|nineteen|twenty|thirty|forty|fifty|sixty|seventy|eighty|ninety|hundred|thousand|million|billion|trillion|and|[\d,]+(?:\.\d+)?)(?:[-\s]?(?:one|two|three|four|five|six|seven|eight|nine|ten|eleven|twelve|thirteen|fourteen|fifteen|sixteen|seventeen|eighteen|nineteen|twenty|thirty|forty|fifty|sixty|seventy|eighty|ninety|hundred|thousand|million|billion|trillion|and|[\d,]+(?:\.\d+)?))*)\b",
            re.IGNORECASE,
        )
        return re.split(pattern, text)
    elif ("pun" in by) and not strict:
        if verbose:
            print(f"splited by 标点('.!?;')")
        return re.split(r"[.!?;]", text)
    elif ("\n" in by or "li" in by) and not strict:
        if verbose:
            print(f"splited by lines('\n')")
        return text.splitlines()
    elif ("cam" in by) and not strict:
        if verbose:
            print(f"splited by camel_case")
        return split_by_camel_case(text)
    elif ("word" in by) and not strict:
        from nltk.tokenize import word_tokenize

        if verbose:
            print(f"splited by word")
        return word_tokenize(text)
    elif ("sen" in by and not "num" in by) and not strict:
        from nltk.tokenize import sent_tokenize

        if verbose:
            print(f"splited by sentence")
        return sent_tokenize(text)
    elif ("sen" in by and "num" in by) and not strict:
        return split_by_sent_num(text, **kws)
    elif ("cha" in by) and not strict:
        if verbose:
            print(f"splited by chracters")
        return list(text)
    elif ("up" in by or "cap" in by) and ("l" not in by) and not strict:
        if verbose:
            print(f"splited by upper case")
        return split_at_upper(text)
    elif ("u" in by and "l" in by) and not strict:
        if by.find("u") < by.find("l"):
            if verbose:
                print(f"splited by upper followed by lower case")
            return split_at_upper_fl_lower(text)
        else:
            if verbose:
                print(f"splited by lower followed by upper case")
            return split_at_lower_fl_upper(text)
    elif ("start" in by or "head" in by) and not strict:
        if verbose:
            print(f"splited by lookahead")
        return split_by_regex_lookahead(text, **kws)
    elif ("end" in by or "tail" in by) and not strict:
        if verbose:
            print(f"splited by endings")
        return split_by_regex_end(text, **kws)
    elif ("other" in by or "non_alp" in by) and not strict:
        if verbose:
            print(f"splited by non_alphanumeric")
        return split_by_consecutive_non_alphanumeric(text)
    elif ("len" in by) and not strict:
        if verbose:
            print(f"splited by fixed length")
        return split_by_fixed_length_chunks(text, **kws)
    elif ("re" in by or "cus" in by or "cos" in by) and not strict:
        if verbose:
            print(f"splited by customed, re; => {by}")
        return reg_split(text, **kws)
    elif ("lang" in by or "eng" in by) and not strict:
        return split_non_ascii(text)
    else:
        return split_general(text, by, verbose=verbose, **kws)


def pdf2img(dir_pdf, dir_save=None, page=None, kind="png", verbose=True, **kws):
    from pdf2image import convert_from_path, pdfinfo_from_path

    df_dir_img_single_page = pd.DataFrame()
    dir_single_page = []
    if verbose:
        from pprint import pp

        pp(pdfinfo_from_path(dir_pdf))
    if isinstance(page, tuple) and page:
        page = list(page)
    if isinstance(page, int):
        page = [page]
    if page is None:
        page = [pdfinfo_from_path(dir_pdf)["Pages"]]
    if len(page) == 1 and page != [pdfinfo_from_path(dir_pdf)["Pages"]]:
        page = [page[0], page[0]]
    else:
        page = [1, page[0]]
    print(page)
    pages = convert_from_path(dir_pdf, first_page=page[0], last_page=page[-1], **kws)
    if dir_save is None:
        dir_save = mkdir(dirname(dir_pdf), basename(dir_pdf).split(".")[0] + "_img")
    for i, page in enumerate(pages):
        if verbose:
            print(f"processing page: {i+1}")
        if i < 9:
            dir_img_each_page = dir_save + f"page_0{i+1}.png"
        else:
            dir_img_each_page = dir_save + f"page_{i+1}.png"
        dir_single_page.append(dir_img_each_page)
        page.save(dir_img_each_page, kind.upper())
    df_dir_img_single_page["fpath"] = dir_single_page
    return df_dir_img_single_page


# dir_pdf = "/Users/macjianfeng/Dropbox/github/python/240308_Python Data Science Handbook.pdf"
# df_page = pdf2img(dir_pdf, page=[1, 5],dpi=300)
def get_encoding(fpath, alternative_encodings=None, verbose=False):
    """
    Attempt to determine the encoding of a file by trying multiple encodings.

    Parameters:
    fpath (str): The path to the file.
    alternative_encodings (list): List of encodings to try. If None, uses a default list.
    verbose (bool): If True, print detailed information about each attempted encoding.

    Returns:
    str: The encoding that successfully read the file, or None if no encoding worked.
    """
    if alternative_encodings is None:
        alternative_encodings = [
            "utf-8",
            "latin1",
            "windows-1252",
            "iso-8859-1",
            "iso-8859-2",
            "iso-8859-3",
            "iso-8859-4",
            "iso-8859-5",
            "iso-8859-6",
            "iso-8859-7",
            "iso-8859-8",
            "iso-8859-9",
            "windows-1250",
            "windows-1251",
            "windows-1253",
            "windows-1254",
            "windows-1255",
            "windows-1256",
            "windows-1257",
            "windows-1258",
            "big5",
            "gb18030",
            "shift_jis",
            "euc_jp",
            "koi8_r",
            "mac_roman",
            "mac_central_europe",
            "mac_greek",
            "mac_cyrillic",
            "mac_arabic",
            "mac_hebrew",
        ]

    if not os.path.isfile(fpath):
        raise FileNotFoundError(f"The file {fpath} does not exist.")

    for enc in alternative_encodings:
        try:
            with open(fpath, mode="r", encoding=enc) as file:
                file.read()  # Try to read the file
            if verbose:
                print(f"Successfully detected encoding: {enc}")
            return enc
        except UnicodeDecodeError:
            if verbose:
                print(f"Failed to decode with encoding: {enc}")
            continue

    # If no encoding worked
    print("No suitable encoding found.")
    return None


def unzip(dir_path, output_dir=None):
    """
    Unzips or extracts various compressed file formats (.gz, .zip, .7z, .tar, .bz2, .xz, .rar).
    If the output directory already exists, it will be replaced.

    # Example usage:
    output_dir = unzip('data.tar.gz')
    output_file = unzip('file.csv.gz')
    output_dir_zip = unzip('archive.zip')
    output_dir_7z = unzip('archive.7z')

    Parameters:
    dir_path (str): Path to the compressed file.
    output_dir (str): Directory where the extracted files will be saved.
                      If None, it extracts to the same directory as the file, with the same name.

    Returns:
    str: The path to the output directory where files are extracted.
    """

    # Set default output directory to the same as the input file
    if output_dir is None:
        output_dir = os.path.splitext(dir_path)[0]

    # If the output directory already exists, remove it and replace it
    if os.path.exists(output_dir):
        if os.path.isdir(output_dir):  # check if it is a folder
            import shutil

            shutil.rmtree(output_dir)  # remove folder
        else:
            os.remove(output_dir)  # remove file

    # Handle .tar.gz files
    if dir_path.endswith(".tar.gz") or dir_path.endswith(".tgz"):
        import tarfile

        with tarfile.open(dir_path, "r:gz") as tar_ref:
            tar_ref.extractall(output_dir)
        return output_dir
    # Handle .gz files
    if dir_path.endswith(".gz") or dir_path.endswith(".gzip"):
        import gzip

        output_file = os.path.splitext(dir_path)[0]  # remove the .gz extension
        try:
            import shutil

            with gzip.open(dir_path, "rb") as gz_file:
                with open(output_file, "wb") as out_file:
                    shutil.copyfileobj(gz_file, out_file)
            print(f"unzipped '{dir_path}' to '{output_file}'")
        except FileNotFoundError:
            print(f"Error: The file '{dir_path}' was not found.")
        except PermissionError:
            print(
                f"Error: Permission denied when accessing '{dir_path}' or writing to '{output_file}'."
            )
        except Exception as e:
            try:
                import tarfile

                with tarfile.open(dir_path, "r:gz") as tar:
                    tar.extractall(path=output_file)
            except Exception as final_e:
                print(f"An final unexpected error occurred: {final_e}")
        return output_file

    # Handle .zip files
    elif dir_path.endswith(".zip"):
        import zipfile

        with zipfile.ZipFile(dir_path, "r") as zip_ref:
            zip_ref.extractall(output_dir)
        return output_dir

    # Handle .7z files (requires py7zr)
    elif dir_path.endswith(".7z"):
        import py7zr

        with py7zr.SevenZipFile(dir_path, mode="r") as z:
            z.extractall(path=output_dir)
        return output_dir

    # Handle .tar files
    elif dir_path.endswith(".tar"):
        import tarfile

        with tarfile.open(dir_path, "r") as tar_ref:
            tar_ref.extractall(output_dir)
        return output_dir

    # Handle .tar.bz2 files
    elif dir_path.endswith(".tar.bz2"):
        import tarfile

        with tarfile.open(dir_path, "r:bz2") as tar_ref:
            tar_ref.extractall(output_dir)
        return output_dir

    # Handle .bz2 files
    elif dir_path.endswith(".bz2"):
        import bz2

        output_file = os.path.splitext(dir_path)[0]  # remove the .bz2 extension
        with bz2.open(dir_path, "rb") as bz_file:
            with open(output_file, "wb") as out_file:
                shutil.copyfileobj(bz_file, out_file)
        return output_file

    # Handle .xz files
    elif dir_path.endswith(".xz"):
        import lzma

        output_file = os.path.splitext(dir_path)[0]  # remove the .xz extension
        with lzma.open(dir_path, "rb") as xz_file:
            with open(output_file, "wb") as out_file:
                shutil.copyfileobj(xz_file, out_file)
        return output_file

    # Handle .rar files (requires rarfile)
    elif dir_path.endswith(".rar"):
        import rarfile

        with rarfile.RarFile(dir_path) as rar_ref:
            rar_ref.extractall(output_dir)
        return output_dir

    else:
        raise ValueError(f"Unsupported file format: {os.path.splitext(dir_path)[1]}")

def is_df_abnormal(df: pd.DataFrame, verbose=False) -> bool:
    """
    Usage
    is_abnormal = is_df_abnormal(df, verbose=1)
    True: abnormal
    False: normal
    """
    if not isinstance(df, pd.DataFrame):
        if verbose:
            print("not pd.DataFrame")
        return False
    df.columns = df.columns.astype(str)  # 把它变成str, 这样就可以进行counts运算了
    # Initialize a list to hold messages about abnormalities
    messages = []
    is_abnormal = False
    # Check the shape of the DataFrame
    actual_shape = df.shape
    messages.append(f"Shape of DataFrame: {actual_shape}")

    # Check column names
    column_names = df.columns.tolist()

    # Count of delimiters and their occurrences
    delimiter_counts = {"\t": 0, ",": 0, "\n": 0, "": 0}  # Count of empty strings

    for name in column_names:
        # Count occurrences of each delimiter
        delimiter_counts["\t"] += name.count("\t")
        delimiter_counts[","] += name.count(",")
        delimiter_counts["\n"] += name.count("\n")
        if name.strip() == "":
            delimiter_counts[""] += 1

    # Check for abnormalities based on delimiter counts
    if len(column_names) == 1 and delimiter_counts["\t"] > 1:
        messages.append("Abnormal: Column names are not split correctly.")
        is_abnormal = True
        if verbose:
            print(f'len(column_names) == 1 and delimiter_counts["\t"] > 1')
    if verbose:
        print("1", is_abnormal)
    if any(delimiter_counts[d] > 3 for d in delimiter_counts if d != ""):
        messages.append("Abnormal: Too many delimiters in column names.")
        is_abnormal = True
        if verbose:
            print(f'any(delimiter_counts[d] > 3 for d in delimiter_counts if d != "")')
    if verbose:
        print("2", is_abnormal)
    if delimiter_counts[""] > 3:
        messages.append("Abnormal: There are empty column names.")
        is_abnormal = True
        if verbose:
            print(f'delimiter_counts[""] > 3')
    if verbose:
        print("3", is_abnormal)
    if any(delimiter_counts[d] > 3 for d in ["\t", ",", "\n"]):
        messages.append("Abnormal: Some column names contain unexpected characters.")
        is_abnormal = True
        if verbose:
            print(f'any(delimiter_counts[d] > 3 for d in ["\t", ",", "\n"])')
    if verbose:
        print("4", is_abnormal)
    # # Check for missing values
    # missing_values = df.isnull().sum()
    # if missing_values.any():
    #     messages.append("Missing values in columns:")
    #     messages.append(missing_values[missing_values > 0].to_string())
    #     is_abnormal = True
    #     print(f'missing_values.any()')

    # Check data types
    data_types = df.dtypes
    # messages.append(f"Data types of columns:\n{data_types}")

    # Check for an unreasonable number of rows or columns
    if actual_shape[0] < 2 or actual_shape[1] < 2:
        messages.append(
            "Abnormal: DataFrame is too small (less than 2 rows or columns)."
        )
        is_abnormal = True
        if verbose:
            print(f"actual_shape[0] < 2 or actual_shape[1] < 2")
    if verbose:
        print("6", is_abnormal)
    # Compile results
    if verbose:
        print("\n".join(messages))
    return is_abnormal  # Data is abnormal
def decrypt_excel(fpath, password):
    # * needs a password?
    import msoffcrypto  # pip install msoffcrypto-tool
    from io import BytesIO

    # Open the encrypted Excel file
    with open(fpath, "rb") as f:
        try:
            office_file = msoffcrypto.OfficeFile(f)
            office_file.load_key(password=password)  # Provide the password
            decrypted = BytesIO()
            office_file.decrypt(decrypted)
        except:
            office_file = msoffcrypto.OfficeFile(f)
            office_file.load_key(password=depass(password))  # Provide the password
            decrypted = BytesIO()
            office_file.decrypt(decrypted)
    decrypted.seek(0) # reset pointer to start
    return decrypted

def fload(fpath, kind=None, **kwargs):
    """
    Load content from a file with specified file type.
    Parameters:
        fpath (str): The file path from which content will be loaded.
        kind (str): The file type to load. Supported options: 'docx', 'txt', 'md', 'html', 'json', 'yaml', 'xml', 'csv', 'xlsx', 'pdf'.
        **kwargs: Additional parameters for 'csv' and 'xlsx' file types.
    Returns:
        content: The content loaded from the file.
    """

    def read_mplstyle(style_file):
        import matplotlib.pyplot as plt

        # Load the style file
        plt.style.use(style_file)

        # Get the current style properties
        style_dict = plt.rcParams

        # Convert to dictionary
        style_dict = dict(style_dict)
        # Print the style dictionary
        for i, j in style_dict.items():
            print(f"\n{i}::::{j}")
        return style_dict

    # #example usage:
    # style_file = "/ std-colors.mplstyle"
    # style_dict = read_mplstyle(style_file)

    def load_txt_md(fpath):
        with open(fpath, "r") as file:
            content = file.read()
        return content

    # def load_html(fpath):
    #     with open(fpath, "r") as file:
    #         content = file.read()
    #     return content
    def load_html(fpath, **kwargs):
        return pd.read_html(fpath, **kwargs)

    def load_json(fpath, **kwargs):
        output = kwargs.pop("output", "json")
        if output == "json":
            import json

            with open(fpath, "r") as file:
                content = json.load(file)
            return content
        else:
            return pd.read_json(fpath, **kwargs)

    def load_yaml(fpath):
        import yaml

        with open(fpath, "r") as file:
            content = yaml.safe_load(file)
        return content

    def load_xml(fpath, fsize_thr: int = 100):
        from lxml import etree

        def load_small_xml(fpath):
            tree = etree.parse(fpath)
            root = tree.getroot()
            return etree.tostring(root, pretty_print=True).decode()

        def load_large_xml(fpath):
            xml_parts = []
            context = etree.iterparse(
                fpath, events=("start", "end"), recover=True, huge_tree=True
            )

            for event, elem in context:
                if event == "end":
                    xml_parts.append(etree.tostring(elem, pretty_print=True).decode())
                    elem.clear()
                    while elem.getprevious() is not None:
                        del elem.getparent()[0]
            del context
            return "".join(xml_parts)

        file_size = os.path.getsize(fpath) / 1024 / 1024  # in MB

        if file_size > fsize_thr:
            print(f"reading a small file:{file_size} Mb")
            return load_large_xml(fpath)
        else:
            print(f"reading a big file:{file_size} Mb")
            return load_small_xml(fpath)

    def get_comment(fpath, comment=None, encoding="utf-8", lines_to_check=5):
        """
        Detect comment characters in a file.

        Parameters:
        - fpath: str, the file path of the CSV file.
        - encoding: str, the encoding of the file (default is 'utf-8').
        - lines_to_check: int, number of lines to check for comment characters (default is 5).

        Returns:
        - str or None: the detected comment character, or None if no comment character is found.
        """
        comment_chars = [
            "#",
            "!",
            "//",
            ";",
        ]  # can use any character or string as a comment
        try:
            with open(fpath, "r", encoding=encoding) as f:
                lines = [next(f) for _ in range(lines_to_check)]
        except (UnicodeDecodeError, ValueError):
            with open(fpath, "r", encoding=get_encoding(fpath)) as f:
                lines = [next(f) for _ in range(lines_to_check)]
        for line in lines:
            for char in comment_chars:
                if line.startswith(char):
                    return char
        return None

    def _get_chunks(df_fake):
        """
        helper func for 'load_csv'
        """
        chunks = []
        for chunk in df_fake:
            chunks.append(chunk)
        return pd.concat(chunks, ignore_index=True)

    def load_csv(fpath, **kwargs):
        from pandas.errors import EmptyDataError

        engine = kwargs.pop("engine", "pyarrow")  # default: None
        sep = kwargs.pop("sep", None)  # default: ','
        index_col = kwargs.pop("index_col", None)  # default: None
        memory_map = kwargs.pop("memory_map", False)  # default: False
        skipinitialspace = kwargs.pop("skipinitialspace", False)  # default: False
        encoding = kwargs.pop("encoding", "utf-8")  # default: "utf-8"
        on_bad_lines = kwargs.pop("on_bad_lines", "skip")  # default: 'error'
        comment = kwargs.pop("comment", None)  # default: None
        fmt = kwargs.pop("fmt", False)  # default:
        chunksize = kwargs.pop("chunksize", None)  # default: None

        # check filesize
        f_size = round(os.path.getsize(fpath) / 1024 / 1024, 3)
        if f_size >= 50:  # 50 MB
            if chunksize is None:
                chunksize = 5000
                print(
                    f"file size is {f_size}MB, then set the chunksize with {chunksize}"
                )
        engine = "c" if chunksize else engine  # when chunksize, recommend 'c'
        low_memory = kwargs.pop("low_memory", True)  # default: True
        low_memory = (
            False if chunksize else True
        )  # when chunksize, recommend low_memory=False # default:
        verbose = kwargs.pop("verbose", False)
        if run_once_within(reverse=True) and verbose:
            use_pd("read_csv", verbose=verbose)

        if comment is None:  # default: None
            comment = get_comment(
                fpath, comment=None, encoding="utf-8", lines_to_check=5
            )
        try:
            df = pd.read_csv(
                fpath,
                engine=engine,
                index_col=index_col,
                memory_map=memory_map,
                encoding=encoding,
                comment=comment,
                skipinitialspace=skipinitialspace,
                sep=sep,
                on_bad_lines=on_bad_lines,
                chunksize=chunksize,
                low_memory=low_memory,
                **kwargs,
            )
            if chunksize:
                df = _get_chunks(df)
                print(df.shape)
            if is_df_abnormal(df, verbose=0):  # raise error
                raise ValueError("the df is abnormal")
        except:
            try:
                try:
                    if engine == "pyarrow" and not chunksize:
                        df = pd.read_csv(
                            fpath,
                            engine=engine,
                            index_col=index_col,
                            encoding=encoding,
                            sep=sep,
                            on_bad_lines=on_bad_lines,
                            comment=comment,
                            low_memory=low_memory,
                            **kwargs,
                        )
                    else:
                        df = pd.read_csv(
                            fpath,
                            engine=engine,
                            index_col=index_col,
                            memory_map=memory_map,
                            encoding=encoding,
                            sep=sep,
                            skipinitialspace=skipinitialspace,
                            on_bad_lines=on_bad_lines,
                            comment=comment,
                            chunksize=chunksize,
                            low_memory=low_memory,
                            **kwargs,
                        )
                    if chunksize:
                        df = _get_chunks(df)
                        print(df.shape)
                    if is_df_abnormal(df, verbose=0):
                        raise ValueError("the df is abnormal")
                except (UnicodeDecodeError, ValueError):
                    encoding = get_encoding(fpath)
                    # print(f"utf-8 failed. Retrying with detected encoding: {encoding}")
                    if engine == "pyarrow" and not chunksize:
                        df = pd.read_csv(
                            fpath,
                            engine=engine,
                            index_col=index_col,
                            encoding=encoding,
                            sep=sep,
                            on_bad_lines=on_bad_lines,
                            comment=comment,
                            low_memory=low_memory,
                            **kwargs,
                        )
                    else:
                        df = pd.read_csv(
                            fpath,
                            engine=engine,
                            index_col=index_col,
                            memory_map=memory_map,
                            encoding=encoding,
                            sep=sep,
                            skipinitialspace=skipinitialspace,
                            on_bad_lines=on_bad_lines,
                            comment=comment,
                            chunksize=chunksize,
                            low_memory=low_memory,
                            **kwargs,
                        )
                    if chunksize:
                        df = _get_chunks(df)
                        print(df.shape)
                    if is_df_abnormal(df, verbose=0):
                        raise ValueError("the df is abnormal")
            except Exception as e:
                separators = [",", "\t", ";", "|", " "]
                for sep in separators:
                    sep2show = sep if sep != "\t" else "\\t"
                    if verbose:
                        print(f'trying with: engine=pyarrow, sep="{sep2show}"')
                    try:
                        df = pd.read_csv(
                            fpath,
                            engine=engine,
                            skipinitialspace=skipinitialspace,
                            sep=sep,
                            on_bad_lines=on_bad_lines,
                            comment=comment,
                            chunksize=chunksize,
                            low_memory=low_memory,
                            **kwargs,
                        )
                        if chunksize:
                            df = _get_chunks(df)
                            print(df.shape)
                        if not is_df_abnormal(df, verbose=0) and verbose:  # normal
                            display(df.head(2))
                            print(f"shape: {df.shape}")
                            return df
                    except:
                        pass
                else:
                    if not chunksize:
                        engines = [None, "c", "python"]
                        for engine in engines:
                            separators = [",", "\t", ";", "|", " "]
                            for sep in separators:
                                try:
                                    sep2show = sep if sep != "\t" else "\\t"
                                    if verbose:
                                        print(
                                            f"trying with: engine={engine}, sep='{sep2show}'"
                                        )
                                    # print(".")
                                    df = pd.read_csv(
                                        fpath,
                                        engine=engine,
                                        sep=sep,
                                        on_bad_lines=on_bad_lines,
                                        comment=comment,
                                        chunksize=chunksize,
                                        low_memory=low_memory,
                                        **kwargs,
                                    )
                                    # display(df.head(2))
                                    # print(f"is_df_abnormal:{is_df_abnormal(df, verbose=0)}")
                                    if chunksize:
                                        df = _get_chunks(df)
                                        print(df.shape)
                                    if not is_df_abnormal(df, verbose=0):
                                        if verbose:
                                            (
                                                display(df.head(2))
                                                if isinstance(df, pd.DataFrame)
                                                else display("it is not a DataFrame")
                                            )
                                            (
                                                print(f"shape: {df.shape}")
                                                if isinstance(df, pd.DataFrame)
                                                else display("it is not a DataFrame")
                                            )
                                        return df
                                except EmptyDataError as e:
                                    continue
                            else:
                                pass
        # print(kwargs)
        # if is_df_abnormal(df,verbose=verbose):
        #     df=pd.read_csv(fpath,**kwargs)
        if verbose:
            display(df.head(2))
            print(f"shape: {df.shape}")
        return df
 
    def load_excel(fpath, **kwargs):
        engine = kwargs.get("engine", "openpyxl")
        verbose = kwargs.pop("verbose", False)
        password = kwargs.pop("password", None)
        output = kwargs.pop("output", "DataFrame").lower()
        sheet_name = kwargs.pop("sheet_name", None)

        def print_sheet_info(fpath):
            try:
                meta = pd.ExcelFile(fpath)
                print(f"n_sheet={len(meta.sheet_names)},\t'sheetname = 0 (default)':")
                [print(f"{i}:\t{name}") for i, name in enumerate(meta.sheet_names)]
            except Exception as e:
                if verbose:
                    print(f"Error retrieving sheet info: {e}")
        if output in ["dataframe", "df"]:
            if not password:
                if verbose:
                    print("Reading Excel without password protection...")
                df = pd.read_excel(fpath, engine=engine, sheet_name=sheet_name, **kwargs)
                if verbose:
                    print_sheet_info(fpath)
                return df
            # Handle password-protected DataFrame case
            else:
                if verbose:
                    print("Decrypting and loading DataFrame...")
                decrypted = decrypt_excel(fpath, password=password)
                df = pd.read_excel(decrypted, engine=engine,sheet_name=sheet_name, **kwargs)
                if verbose:
                    print_sheet_info(fpath)
                return df
        # Handle cases for non-dataframe output
        else:
            from openpyxl import load_workbook
            if verbose:
                print("Returning worksheet (non-DataFrame output)...")
            if password:
                decrypted = decrypt_excel(fpath, password=password)
                workbook = load_workbook(decrypted)
            else:
                workbook = load_workbook(fpath)

            # If sheet_name is specified, keep only that sheet
            if sheet_name:
                if sheet_name in workbook.sheetnames:
                    worksheet = workbook[sheet_name]
                    # Remove all other sheets
                    for sheet in workbook.sheetnames:
                        if sheet != sheet_name:
                            del workbook[sheet]
                else:
                    raise ValueError(f"Sheet '{sheet_name}' not found in the workbook.")
            else:
                worksheet = workbook.active  # Default to active sheet

            return workbook

    def load_parquet(fpath, **kwargs):
        """
        Load a Parquet file into a Pandas DataFrame with advanced options.

        Parameters:
        - fpath (str): The file path to the Parquet file.
        - engine (str): The engine to use for reading the Parquet file (default is 'pyarrow').
        - columns (list): List of columns to load. If None, loads all columns.
        - verbose (bool): If True, prints additional information about the loading process.
        - filters (list): List of filter conditions for predicate pushdown.
        - **kwargs: Additional keyword arguments for `pd.read_parquet`.

        Returns:
        - df (DataFrame): The loaded DataFrame.
        """

        engine = kwargs.get("engine", "pyarrow")
        verbose = kwargs.pop("verbose", False)

        if run_once_within(reverse=True) and verbose:
            use_pd("read_parquet", verbose=verbose)
        try:
            df = pd.read_parquet(fpath, engine=engine, **kwargs)
            if verbose:
                if "columns" in kwargs:
                    print(f"Loaded columns: {kwargs['columns']}")
                else:
                    print("Loaded all columns.")
            print(f"shape: {df.shape}")
        except Exception as e:
            print(f"An error occurred while loading the Parquet file: {e}")
            df = None

        return df

    def load_ipynb(fpath, **kwargs):
        import nbformat
        from nbconvert import MarkdownExporter

        as_version = kwargs.get("as_version", 4)
        with open(fpath, "r") as file:
            nb = nbformat.read(file, as_version=as_version)
            md_exporter = MarkdownExporter()
            md_body, _ = md_exporter.from_notebook_node(nb)
        return md_body

    def load_pdf(fpath, page="all", verbose=False, **kwargs):
        """
        Parameters:
        fpath: The path to the PDF file to be loaded.
        page (optional):
            Specifies which page or pages to extract text from. By default, it's set to "all", which means text from all
            pages will be returned. It can also be an integer to specify a single page number or a list of integers to
            specify multiple pages.
        verbose (optional):
            If True, prints the total number of pages processed.
        Functionality:
        It initializes an empty dictionary text_dict to store page numbers as keys and their corresponding text as values.
        It iterates through each page of the PDF file using a for loop.
        For each page, it extracts the text using PyPDF2's extract_text() method and stores it in text_dict with the page number incremented by 1 as the key.
        If the page parameter is an integer, it converts it into a list containing that single page number to ensure consistency in handling.
        If the page parameter is a NumPy array, it converts it to a list using the tolist() method to ensure compatibility with list operations.
        If verbose is True, it prints the total number of pages processed.
        If page is a list, it combines the text of the specified pages into a single string combined_text and returns it.
        If page is set to "all", it returns the entire text_dict containing text of all pages.
        If page is an integer, it returns the text of the specified page number.
        If the specified page is not found, it returns the string "Page is not found".
        """
        from PyPDF2 import PdfReader

        text_dict = {}
        with open(fpath, "rb") as file:
            pdf_reader = PdfReader(file)
            num_pages = len(pdf_reader.pages)
            for page_num in range(num_pages):
                if verbose:
                    print(f"processing page {page_num}")
                page_ = pdf_reader.pages[page_num]
                text_dict[page_num + 1] = page_.extract_text()
        if isinstance(page, int):
            page = [page]
        elif isinstance(page, np.ndarray):
            page = page.tolist()
        if verbose:
            print(f"total pages: {page_num}")
        if isinstance(page, list):
            combined_text = ""
            for page_num in page:
                combined_text += text_dict.get(page_num, "")
            return combined_text
        elif "all" in page.lower():
            combined_text = ""
            for i in text_dict.values():
                combined_text += i
            return combined_text
        else:
            return text_dict.get(int(page), "Page is not found")

    def load_docx(fpath):
        from docx import Document

        doc = Document(fpath)
        content = [para.text for para in doc.paragraphs]
        return content

    def load_rtf(file_path):
        from striprtf.striprtf import rtf_to_text

        try:
            with open(file_path, "r") as file:
                rtf_content = file.read()
                text = rtf_to_text(rtf_content)
                return text
        except Exception as e:
            print(f"Error loading RTF file: {e}")

    if kind is None:
        _, kind = os.path.splitext(fpath)
        kind = kind.lower()
    kind = kind.lstrip(".").lower()
    img_types = [
        "bmp",
        "eps",
        "gif",
        "png",
        "jpg",
        "jpeg",
        "jpeg2000",
        "tiff",
        "tif",
        "icns",
        "ico",
        "im",
        "msp",
        "pcx",
        "ppm",
        "sgi",
        "spider",
        "tga",
        "webp",
    ]
    doc_types = [
        "docx",
        "pdf",
        "txt",
        "csv",
        "xlsx",
        "tsv",
        "parquet",
        "snappy",
        "md",
        "html",
        "json",
        "yaml",
        "xml",
        "ipynb",
        "mtx",
        "rtf",
    ]
    zip_types = [
        "gz",
        "zip",
        "7z",
        "rar",
        "tgz",
        "tar",
        "tar.gz",
        "tar.bz2",
        "bz2",
        "xz",
        "gzip",
    ]
    other_types = ["fcs"]
    supported_types = [*doc_types, *img_types, *zip_types, *other_types]
    if kind not in supported_types:
        print(
            f'Warning:\n"{kind}" is not in the supported list '
        )  # {supported_types}')

    if kind == "docx":
        return load_docx(fpath)
    elif kind == "txt" or kind == "md":
        return load_txt_md(fpath)
    elif kind == "html":
        return load_html(fpath, **kwargs)
    elif kind == "json":
        return load_json(fpath, **kwargs)
    elif kind == "yaml":
        return load_yaml(fpath)
    elif kind == "xml":
        return load_xml(fpath)
    elif kind in ["csv", "tsv"]:
        # verbose = kwargs.pop("verbose", False)
        # if run_once_within(reverse=True) and verbose:
        #     use_pd("read_csv")
        content = load_csv(fpath, **kwargs)
        return content
    elif kind == "pkl":
        verbose = kwargs.pop("verbose", False)
        if run_once_within(reverse=True) and verbose:
            use_pd("read_pickle")
        try:
            res_ = pd.read_pickle(fpath, **kwargs)
        except Exception as e:
            import pickle

            with open("sgd_classifier.pkl", "rb") as f:
                res_ = pickle.load(f)
        return res_
    elif kind in ["ods", "ods", "odt"]:
        engine = kwargs.get("engine", "odf")
        kwargs.pop("engine", None)
        return load_excel(fpath, engine=engine, **kwargs)
    elif kind == "xls":
        verbose = kwargs.pop("verbose", False)
        engine = kwargs.get("engine", "xlrd")
        kwargs.pop("engine", None)
        content = load_excel(fpath, engine=engine, **kwargs)
        (
            print(f"shape: {content.shape}")
            if isinstance(content, pd.DataFrame) and verbose
            else None
        )
        display(content.head(3)) if isinstance(content, pd.DataFrame) else None
        return content
    elif kind == "xlsx":
        verbose = kwargs.pop("verbose", False)
        content = load_excel(fpath, verbose=verbose,**kwargs)
        # (
        #     display(content.head(3))
        #     if isinstance(content, pd.DataFrame) and verbose
        #     else None
        # )
        print(f"shape: {content.shape}") if isinstance(content, pd.DataFrame) else None
        return content
    elif kind == "mtx":
        from scipy.io import mmread

        verbose = kwargs.pop("verbose", False)
        dat_mtx = mmread(fpath)
        content = pd.DataFrame.sparse.from_spmatrix(dat_mtx, **kwargs)
        (
            display(content.head(3))
            if isinstance(content, pd.DataFrame) and verbose
            else None
        )
        print(f"shape: {content.shape}")
        return content
    elif kind == "ipynb":
        return load_ipynb(fpath, **kwargs)
    elif kind in ["parquet", "snappy"]:
        verbose = kwargs.pop("verbose", False)
        if run_once_within(reverse=True):
            use_pd("read_parquet")
        return load_parquet(fpath, **kwargs)
    elif kind == "feather":
        verbose = kwargs.pop("verbose", False)
        if run_once_within(reverse=True):
            use_pd("read_feather")
        content = pd.read_feather(fpath, **kwargs)
        return content
    elif kind == "h5":
        content = pd.read_hdf(fpath, **kwargs)
        return content
    elif kind == "pkl":
        content = pd.read_pickle(fpath, **kwargs)
        return content
    elif kind == "pdf":
        # print('usage:load_pdf(fpath, page="all", verbose=False)')
        return load_pdf(fpath, **kwargs)
    elif kind.lower() in img_types:
        print(f'Image ".{kind}" is loaded.')
        return load_img(fpath)
    elif kind == "gz" and fpath.endswith(".soft.gz"):
        import GEOparse

        return GEOparse.get_GEO(filepath=fpath)
    elif kind.lower() in zip_types:
        from pprint import pp

        keep = kwargs.get("keep", False)
        fpath_unzip = unzip(fpath)
        if os.path.isdir(fpath_unzip):
            print(f"{fpath_unzip} is a folder. fload stoped.")
            fpath_list = os.listdir("./datasets/GSE10927_family.xml")
            print(f"{len(fpath_list)} files within the folder")
            if len(fpath_list) > 5:
                pp(fpath_list[:5])
                print("there are more ...")
            else:
                pp(fpath_list)
            return fpath_list
        elif os.path.isfile(fpath_unzip):
            print(f"{fpath_unzip} is a file.")
            content_unzip = fload(fpath_unzip, **kwargs)
            if not keep:
                os.remove(fpath_unzip)
            return content_unzip
        else:
            print(f"{fpath_unzip} does not exist or is a different type.")

    elif kind.lower() == "gmt":
        import gseapy as gp

        gene_sets = gp.read_gmt(fpath)
        return gene_sets

    elif kind.lower() == "fcs":
        import fcsparser

        # https://github.com/eyurtsev/fcsparser
        meta, data = fcsparser.parse(fpath, reformat_meta=True)
        print("meta, data = fload(*.fcs)")
        return meta, data

    elif kind == "mplstyle":
        return read_mplstyle(fpath)
    elif kind == "rtf":
        return load_rtf(fpath)

    else:
        print("direct reading...")
        try:
            try:
                with open(fpath, "r", encoding="utf-8") as f:
                    content = f.readlines()
            except UnicodeDecodeError:
                print("Failed to read as utf-8, trying different encoding...")
                with open(
                    fpath, "r", encoding=get_encoding(fpath)
                ) as f:  # Trying with a different encoding
                    content = f.readlines()
        except:
            try:
                with open(fpath, "r", encoding="utf-8") as f:
                    content = f.read()
            except UnicodeDecodeError:
                print("Failed to read as utf-8, trying different encoding...")
                with open(
                    fpath, "r", encoding=get_encoding(fpath)
                ) as f:  # Trying with a different encoding
                    content = f.read()
        return content


# Example usage
# txt_content = fload('sample.txt')
# md_content = fload('sample.md')
# html_content = fload('sample.html')
# json_content = fload('sample.json')
# yaml_content = fload('sample.yaml')
# xml_content = fload('sample.xml')
# csv_content = fload('sample.csv')
# xlsx_content = fload('sample.xlsx')
# docx_content = fload('sample.docx')


def fopen(fpath):
    import os
    import platform
    import sys

    try:
        # Check if the file exists
        if not os.path.isfile(fpath):
            print(f"Error: The file does not exist - {fpath}")
            return

        # Get the system platform
        system = platform.system()

        # Platform-specific file opening commands
        if system == "Darwin":  # macOS
            os.system(f'open "{fpath}"')
        elif system == "Windows":  # Windows
            # Ensure the path is handled correctly in Windows, escape spaces
            os.system(f'start "" "{fpath}"')
        elif system == "Linux":  # Linux
            os.system(f'xdg-open "{fpath}"')
        elif system == "Java":  # Java (or other unhandled systems)
            print(f"Opening {fpath} on unsupported system.")
        else:
            print(f"Unsupported OS: {system}")

        print(f"Successfully opened {fpath} with the default application.")
    except Exception as e:
        print(f"Error opening file {fpath}: {e}")


def fupdate(fpath, content=None, how="head"):
    """
    Update a file by adding new content at the top and moving the old content to the bottom.
    If the file is a JSON file, merge the new content with the old content.

    Parameters
    ----------
    fpath : str
        The file path where the content should be updated.
    content : str or dict, optional
        The new content to add at the top of the file (for text) or merge (for JSON).
        If not provided, the function will not add any new content.

    Notes
    -----
    - If the file at `fpath` does not exist, it will be created.
    - For text files, the new content will be added at the top, followed by the old content.
    - For JSON files, the new content will be merged with the existing JSON content.
    """
    content = content or ""
    file_ext = os.path.splitext(fpath)[1]
    how_s = ["head", "tail", "start", "end", "beginning", "stop", "last", "before"]
    how = strcmp(how, how_s)[0]
    print(how)
    add_where = "head" if how in ["head", "start", "beginning", "before"] else "tail"
    if "json" in file_ext.lower():
        old_content = fload(fpath, kind="json") if os.path.exists(fpath) else {}
        updated_content = (
            {**content, **old_content}
            if add_where == "head"
            else (
                {**old_content, **content} if isinstance(content, dict) else old_content
            )
        )
        fsave(fpath, updated_content)
    else:
        # Handle text file
        if os.path.exists(fpath):
            with open(fpath, "r") as file:
                old_content = file.read()
        else:
            old_content = ""

        # Write new content at the top followed by old content
        with open(fpath, "w") as file:
            if add_where == "head":
                file.write(content + "\n")
                file.write(old_content)
            else:
                file.write(old_content)
                file.write(content + "\n")


def fappend(fpath, content=None):
    """
    append new content at the end.
    """
    content = content or ""
    if os.path.exists(fpath):
        with open(fpath, "r") as file:
            old_content = file.read()
    else:
        old_content = ""

    with open(fpath, "w") as file:
        file.write(old_content)
        file.write(content)


def filter_kwargs(kws, valid_kwargs):
    if isinstance(valid_kwargs, dict):
        kwargs_filtered = {
            key: value for key, value in kws.items() if key in valid_kwargs.keys()
        }
    elif isinstance(valid_kwargs, list):
        kwargs_filtered = {
            key: value for key, value in kws.items() if key in valid_kwargs
        }
    return kwargs_filtered


str_space_speed = 'sapce cmp:parquet(0.56GB)<feather(1.14GB)<csv(6.55GB)<pkl=h5("26.09GB")\nsaving time: pkl=feather("13s")<parquet("35s")<h5("2m31s")<csv("58m")\nloading time: pkl("6.9s")<parquet("16.1s")=feather("15s")<h5("2m 53s")<csv(">>>30m")'


def fsave(
    fpath,
    content,
    how="overwrite",
    kind=None,
    font_name="Times",
    font_size=10,
    spacing=6,
    **kwargs,
):
    """
    Save content into a file with specified file type and formatting.
    Parameters:
        fpath (str): The file path where content will be saved.
        content (list of str or dict): The content to be saved, where each string represents a paragraph or a dictionary for tabular data.
        kind (str): The file type to save. Supported options: 'docx', 'txt', 'md', 'html', 'pdf', 'csv', 'xlsx', 'json', 'xml', 'yaml'.
        font_name (str): The font name for text formatting (only applicable for 'docx', 'html', and 'pdf').
        font_size (int): The font size for text formatting (only applicable for 'docx', 'html', and 'pdf').
        spacing (int): The space after each paragraph (only applicable for 'docx').
        **kwargs: Additional parameters for 'csv', 'xlsx', 'json', 'yaml' file types.
    Returns:
        None
    """

    def save_content(fpath, content, mode="w", how="overwrite"):
        if "wri" in how.lower():
            with open(fpath, mode, encoding="utf-8") as file:
                file.write(content)
        elif "upd" in how.lower():
            fupdate(fpath, content=content)
        elif "app" in how.lower():
            fappend(fpath, content=content)

    def save_docx(fpath, content, font_name, font_size, spacing):
        import docx

        if isinstance(content, str):
            content = content.split(". ")
        doc = docx.Document()
        for i, paragraph_text in enumerate(content):
            paragraph = doc.add_paragraph()
            run = paragraph.add_run(paragraph_text)
            font = run.font
            font.name = font_name
            font.size = docx.shared.Pt(font_size)
            if i != len(content) - 1:  # Add spacing for all but the last paragraph
                paragraph.space_after = docx.shared.Pt(spacing)
        doc.save(fpath)

    def save_txt_md(fpath, content, sep="\n", mode="w"):
        # Ensure content is a single string
        if isinstance(content, list):
            content = sep.join(content)
        save_content(fpath, sep.join(content), mode)

    def save_html(fpath, content, font_name, font_size, mode="w"):
        html_content = "<html><body>"
        for paragraph_text in content:
            html_content += f'<p style="font-family:{font_name}; font-size:{font_size}px;">{paragraph_text}</p>'
        html_content += "</body></html>"
        save_content(fpath, html_content, mode)

    def save_pdf(fpath, content, font_name, font_size):
        from fpdf import FPDF

        pdf = FPDF()
        pdf.add_page()
        # pdf.add_font('Arial','',r'/System/Library/Fonts/Supplemental/Arial.ttf',uni=True)
        pdf.set_font(font_name, "", font_size)
        for paragraph_text in content:
            pdf.multi_cell(0, 10, paragraph_text)
            pdf.ln(h="")
        pdf.output(fpath, "F")

    def save_csv(fpath, data, **kwargs):
        # https://pandas.pydata.org/docs/reference/api/pandas.DataFrame.to_csv.html

        verbose = kwargs.pop("verbose", False)
        if run_once_within(reverse=True):
            use_pd("to_csv", verbose=verbose)
        kwargs_csv = dict(
            path_or_buf=None,
            sep=",",
            na_rep="",
            float_format=None,
            columns=None,
            header=True,
            index=True,
            index_label=None,
            mode="w",
            encoding="UTF-8",
            compression="infer",
            quoting=None,
            quotechar='"',
            lineterminator=None,
            chunksize=None,
            date_format=None,
            doublequote=True,
            escapechar=None,
            decimal=".",
            errors="strict",
            storage_options=None,
        )
        kwargs_valid = filter_kwargs(kwargs, kwargs_csv)
        df = pd.DataFrame(data)
        df.to_csv(fpath, **kwargs_valid)

    def save_xlsx(fpath, data, password=None,apply_format=None, **kwargs):
        import msoffcrypto
        from io import BytesIO
        import openpyxl
        import pandas.io.formats.style

        verbose = kwargs.pop("verbose", False)
        sheet_name = kwargs.pop("sheet_name", "Sheet1")
        engine = kwargs.pop("engine", "xlsxwriter")
        mode = kwargs.pop("mode","a")
        if_sheet_exists = strcmp(kwargs.get("if_sheet_exists","new"),['error', 'new', 'replace', 'overlay'])[0]
        kwargs.pop("if_sheet_exists",None)
        if run_once_within(reverse=True):
            use_pd("to_excel", verbose=verbose)
 
        if apply_format is None:
            kwargs_format=list(extract_kwargs(format_excel).keys())[4:]
            apply_format=True if any([i in kwargs_format for i in kwargs]) else False
            print(f"apply format: {apply_format}")
        if apply_format or any([
                isinstance(data, openpyxl.worksheet.worksheet.Worksheet), 
                isinstance(data, openpyxl.workbook.workbook.Workbook),
                isinstance(data, pd.io.formats.style.Styler)
             ]):
            format_excel(df=data, 
                         filename=fpath,
                         sheet_name=sheet_name,
                         password=password,
                         if_sheet_exists=if_sheet_exists,
                         mode=mode, 
                         engine=engine,
                         verbose=verbose,
                         **kwargs)
        else:
            # Remove non-relevant kwargs
            irrelevant_keys=list(extract_kwargs(format_excel).keys())[4:]
            [kwargs.pop(key, None) for key in irrelevant_keys]
            df = pd.DataFrame(data)
            # Write to Excel without password first
            temp_file = BytesIO()
            
            df.to_excel(
                temp_file,
                sheet_name=sheet_name,
                index=False,
                engine="xlsxwriter",
                **kwargs,
            )
            # If a password is provided, encrypt the file
            if password:
                temp_file.seek(0)
                office_file = msoffcrypto.OfficeFile(temp_file) 
                with open(fpath, "wb") as encrypted_file:
                    office_file.encrypt(outfile=encrypted_file,password=password)
            else: # Save the file without encryption if no password is provided
                try:
                    # Use ExcelWriter with append mode if the file exists
                    engine="openpyxl" if mode=="a" else "xlsxwriter"
                    if mode=="a":
                        with pd.ExcelWriter(fpath, engine=engine, mode=mode,if_sheet_exists=if_sheet_exists) as writer:
                            df.to_excel(writer, sheet_name=sheet_name, index=False, **kwargs)
                    else:
                        with pd.ExcelWriter(fpath, engine=engine, mode=mode) as writer:
                            df.to_excel(writer, sheet_name=sheet_name, index=False, **kwargs)
                except FileNotFoundError:
                    # If file doesn't exist, create a new one
                    df.to_excel(fpath, sheet_name=sheet_name, index=False, **kwargs)

    def save_ipynb(fpath, data, **kwargs):
        # Split the content by code fences to distinguish between code and markdown
        import nbformat

        parts = data.split("```")
        cells = []

        for i, part in enumerate(parts):
            if i % 2 == 0:
                # Even index: markdown content
                cells.append(nbformat.v4.new_markdown_cell(part.strip()))
            else:
                # Odd index: code content
                cells.append(nbformat.v4.new_code_cell(part.strip()))
        # Create a new notebook
        nb = nbformat.v4.new_notebook()
        nb["cells"] = cells
        # Write the notebook to a file
        with open(fpath, "w", encoding="utf-8") as ipynb_file:
            nbformat.write(nb, ipynb_file) 
    def save_json(fpath_fname, var_dict_or_df):
        import json
        def _convert_js(data):
            if isinstance(data, pd.DataFrame):
                return data.to_dict(orient="list")
            elif isinstance(data, np.ndarray):
                return data.tolist()
            elif isinstance(data, dict):
                return {key: _convert_js(value) for key, value in data.items()}
            return data

        serializable_data = _convert_js(var_dict_or_df)
        # Save the serializable data to the JSON file
        with open(fpath_fname, "w") as f_json:
            json.dump(serializable_data, f_json, indent=4) 

    def save_yaml(fpath, data, **kwargs):
        import yaml

        with open(fpath, "w") as file:
            yaml.dump(data, file, **kwargs)

    def save_xml(fpath, data):
        from lxml import etree

        root = etree.Element("root")
        if isinstance(data, dict):
            for key, val in data.items():
                child = etree.SubElement(root, key)
                child.text = str(val)
        else:
            raise ValueError("XML saving only supports dictionary data")
        tree = etree.ElementTree(root)
        tree.write(fpath, pretty_print=True, xml_declaration=True, encoding="UTF-8")

    def save_parquet(fpath: str, data: pd.DataFrame, **kwargs):
        engine = kwargs.pop(
            "engine", "auto"
        )  # auto先试pyarrow, 不行就转为fastparquet, {‘auto’, ‘pyarrow’, ‘fastparquet’}
        compression = kwargs.pop(
            "compression", None
        )  # Use None for no compression. Supported options: ‘snappy’, ‘gzip’, ‘brotli’, ‘lz4’, ‘zstd’
        try:
            # Attempt to save with "pyarrow" if engine is set to "auto"
            data.to_parquet(fpath, engine=engine, compression=compression, **kwargs)
            print(
                f"DataFrame successfully saved to {fpath} with engine '{engine}' and {compression} compression."
            )
        except Exception as e:
            print(
                f"Error using with engine '{engine}' and {compression} compression: {e}"
            )
            if "Sparse" in str(e):
                try:
                    # Handle sparse data by converting columns to dense
                    print("Attempting to convert sparse columns to dense format...")
                    data = data.apply(
                        lambda x: (
                            x.sparse.to_dense() if pd.api.types.is_sparse(x) else x
                        )
                    )
                    save_parquet(fpath, data=data, **kwargs)
                except Exception as last_e:
                    print(
                        f"After converted sparse columns to dense format, Error using with engine '{engine}' and {compression} compression: {last_e}"
                    )

    if kind is None:
        _, kind = os.path.splitext(fpath)
        kind = kind.lower()

    kind = kind.lstrip(".").lower()

    if kind not in [
        "docx",
        "txt",
        "md",
        "html",
        "pdf",
        "csv",
        "xlsx",
        "json",
        "xml",
        "yaml",
        "ipynb",
    ]:
        print(
            f"Warning:\n{kind} is not in the supported list ['docx', 'txt', 'md', 'html', 'pdf', 'csv', 'xlsx', 'json', 'xml', 'yaml']"
        )
    mode = kwargs.get("mode", "w")
    if kind == "docx" or kind == "doc":
        save_docx(fpath, content, font_name, font_size, spacing)
    elif kind == "txt":
        save_txt_md(fpath, content, sep="", mode=mode)
    elif kind == "md":
        save_txt_md(fpath, content, sep="", mode=mode)
    elif kind == "html":
        save_html(fpath, content, font_name, font_size)
    elif kind == "pdf":
        save_pdf(fpath, content, font_name, font_size)
    elif kind == "csv":
        save_csv(fpath, content, **kwargs)
    elif kind == "xlsx":
        save_xlsx(fpath, content, **kwargs)
    elif kind == "json":
        save_json(fpath, content)
    elif kind == "xml":
        save_xml(fpath, content)
    elif kind == "yaml":
        save_yaml(fpath, content, **kwargs)
    elif kind == "ipynb":
        save_ipynb(fpath, content, **kwargs)
    elif kind.lower() in ["parquet", "pq", "big", "par"]:
        verbose = kwargs.pop("verbose", False)
        if verbose:
            print(str_space_speed)
            use_pd("to_parquet")
            return None
        compression = kwargs.pop(
            "compression", None
        )  # Use None for no compression. Supported options: ‘snappy’, ‘gzip’, ‘brotli’, ‘lz4’, ‘zstd’
        # fix the fpath ends
        _fpath, _ext = os.path.splitext(fpath)
        fpath = _fpath + _ext.replace(kind, "parquet")
        if compression is not None:
            if not fpath.endswith(compression):
                fpath = fpath + f".{compression}"
        save_parquet(fpath=fpath, data=content, compression=compression, **kwargs)
    elif kind.lower() in ["pkl", "pk", "pickle", "pick"]:
        # Pickle: Although not as efficient in terms of I/O speed and storage as Parquet or Feather,
        # Pickle is convenient if you want to preserve exact Python object types.
        verbose = kwargs.pop("verbose", False)
        if verbose:
            print(str_space_speed)
            use_pd("to_pickle")
            return None
        _fpath, _ext = os.path.splitext(fpath)
        fpath = _fpath + _ext.replace(kind, "pkl")
        compression = kwargs.pop("compression", None)
        if compression is not None:
            if not fpath.endswith(compression["method"]):
                fpath = fpath + f".{compression['method']}"
        if isinstance(content, pd.DataFrame):
            content.to_pickle(fpath, **kwargs)
        else:
            try:
                content = pd.DataFrame(content)
                content.to_pickle(fpath, **kwargs)
            except Exception as e:
                try:
                    import pickle

                    with open(fpath, "wb") as f:
                        pickle.dump(content, f)
                    print("done!", fpath)
                except Exception as e:
                    raise ValueError(
                        f"content is not a DataFrame, cannot be saved as a 'pkl' format: {e}"
                    )
    elif kind.lower() in ["fea", "feather", "ft", "fe", "feat", "fether"]:
        # Feather: The Feather format, based on Apache Arrow, is designed for fast I/O operations. It's
        # optimized for data analytics tasks and is especially fast when working with Pandas.

        verbose = kwargs.pop("verbose", False)
        if verbose:
            print(str_space_speed)
            use_pd("to_feather")
            return None
        _fpath, _ext = os.path.splitext(fpath)
        fpath = _fpath + _ext.replace(kind, "feather")
        if isinstance(content, pd.DataFrame):
            content.to_feather(fpath, **kwargs)
        else:
            try:
                print("trying to convert it as a DataFrame...")
                content = pd.DataFrame(content)
                content.to_feather(fpath, **kwargs)
            except Exception as e:
                raise ValueError(
                    f"content is not a DataFrame, cannot be saved as a 'pkl' format: {e}"
                )
    elif kind.lower() in ["hd", "hdf", "h", "h5"]:
        # particularly useful for large datasets and can handle complex data structures
        verbose = kwargs.pop("verbose", False)
        if verbose:
            print(str_space_speed)
            use_pd("to_hdf")
        _fpath, _ext = os.path.splitext(fpath)
        fpath = _fpath + _ext.replace(kind, "h5")
        compression = kwargs.pop("compression", None)
        if compression is not None:
            if not fpath.endswith(compression):
                fpath = fpath + f".{compression}"
        if isinstance(content, pd.DataFrame):
            content.to_hdf(fpath, key="content", **kwargs)
        else:
            try:
                print("trying to convert it as a DataFrame...")
                content = pd.DataFrame(content)
                content.to_hdf(fpath, **kwargs)
            except Exception as e:
                raise ValueError(
                    f"content is not a DataFrame, cannot be saved as a 'pkl' format: {e}"
                )
    else:
        from . import netfinder

        try:
            netfinder.downloader(url=content, dir_save=dirname(fpath), kind=kind)
        except:
            print(
                f"Error:\n{kind} is not in the supported list ['docx', 'txt', 'md', 'html', 'pdf', 'csv', 'xlsx', 'json', 'xml', 'yaml']"
            )


# # Example usage
# text_content = ["Hello, this is a sample text file.", "This is the second paragraph."]
# tabular_content = {"Name": ["Alice", "Bob"], "Age": [24, 30]}
# json_content = {"name": "Alice", "age": 24}
# yaml_content = {"Name": "Alice", "Age": 24}
# xml_content = {"Name": "Alice", "Age": 24}
# dir_save = "/Users/macjianfeng/Dropbox/Downloads/"
# fsave(dir_save + "sample.txt", text_content)
# fsave(dir_save + "sample.md", text_content)
# fsave(dir_save + "sample.html", text_content)
# fsave(dir_save + "sample.pdf", text_content)
# fsave(dir_save + "sample.docx", text_content)
# fsave(dir_save + "sample.csv", tabular_content, index=False)
# fsave(dir_save + "sample.xlsx", tabular_content, sheet_name="Sheet1", index=False)
# fsave(dir_save + "sample.json", json_content, indent=4)
# fsave(dir_save + "sample.yaml", yaml_content)
# fsave(dir_save + "sample.xml", xml_content)


def addpath(fpath):
    sys.path.insert(0, dir)


def dirname(fpath):
    """
    dirname: Extracting Directory Name from a File Path
    Args:
        fpath (str): the file or directory path
    Returns:
        str: directory, without filename
    """
    dirname_ = os.path.dirname(fpath)
    if not dirname_.endswith("/"):
        dirname_ = dirname_ + "/"
    return dirname_


def dir_name(fpath):  # same as "dirname"
    return dirname(fpath)


def basename(fpath):
    """
    basename: # Output: file.txt
    Args:
        fpath (str): the file or directory path
    Returns:
        str: # Output: file.txt
    """
    return os.path.basename(fpath)


def flist(fpath, contains="all"):
    all_files = [
        os.path.join(fpath, f)
        for f in os.listdir(fpath)
        if os.path.isfile(os.path.join(fpath, f))
    ]
    if isinstance(contains, list):
        filt_files = []
        for filter_ in contains:
            filt_files.extend(flist(fpath, filter_))
        return filt_files
    else:
        if "all" in contains.lower():
            return all_files
        else:
            filt_files = [f for f in all_files if isa(f, contains)]
            return filt_files


def sort_kind(df, by="name", ascending=True):
    if df[by].dtype == "object":  # Check if the column contains string values
        if ascending:
            sorted_index = df[by].str.lower().argsort()
        else:
            sorted_index = df[by].str.lower().argsort()[::-1]
    else:
        if ascending:
            sorted_index = df[by].argsort()
        else:
            sorted_index = df[by].argsort()[::-1]
    sorted_df = df.iloc[sorted_index].reset_index(drop=True)
    return sorted_df


def isa(content, kind):
    """
    content, kind='img'
    kinds file paths based on the specified kind.
    Args:
        content (str): Path to the file.
        kind (str): kind of file to kind. Default is 'img' for images. Other options include 'doc' for documents,
                    'zip' for ZIP archives, and 'other' for other types of files.
    Returns:
        bool: True if the file matches the kind, False otherwise.
    """
    if "img" in kind.lower() or "image" in kind.lower():
        return is_image(content)
    elif "vid" in kind.lower():
        return is_video(content)
    elif "aud" in kind.lower():
        return is_audio(content)
    elif "doc" in kind.lower():
        return is_document(content)
    elif "zip" in kind.lower():
        return is_zip(content)
    elif "dir" in kind.lower() or ("f" in kind.lower() and "d" in kind.lower()):
        return os.path.isdir(content)
    elif "code" in kind.lower():  # file
        return is_code(content)
    elif "fi" in kind.lower():  # file
        return os.path.isfile(content)
    elif "num" in kind.lower():  # file
        return isnum(content)
    elif "text" in kind.lower() or "txt" in kind.lower():  # file
        return is_text(content)
    elif "color" in kind.lower():  # file
        return is_str_color(content)
    elif "html" in kind.lower():
        import re

        if content is None or not isinstance(content, str):
            return False
        # Remove leading and trailing whitespace
        content = content.strip()
        # Check for common HTML tags using regex
        # This pattern matches anything that looks like an HTML tag
        tag_pattern = r"<[a-zA-Z][^>]*>(.*?)</[a-zA-Z][^>]*>"
        # Check for any opening and closing tags
        if re.search(tag_pattern, content):
            return True
        # Additional checks for self-closing tags
        self_closing_tags = ["img", "br", "hr", "input", "meta", "link"]
        for tag in self_closing_tags:
            if f"<{tag}" in content:
                return True
        return False
    else:
        print(f"{kind} was not set up correctly")
        return False


def get_os(full=False, verbose=False):
    """Collects comprehensive system information.
    full(bool): True, get more detailed info
    verbose(bool): True, print it
    usage:
        info = get_os(full=True, verbose=False)
    """
    import sys
    import platform
    import psutil
    # import GPUtil
    import socket
    import uuid
    import cpuinfo
    import os
    import subprocess
    from datetime import datetime, timedelta

    def get_os_type():
        os_name = sys.platform
        if "dar" in os_name:
            return "macOS"
        else:
            if "win" in os_name:
                return "Windows"
            elif "linux" in os_name:
                return "Linux"
            else:
                print(f"{os_name}, returned 'None'")
                return None

    if not full:
        return get_os_type()

    def get_os_info():
        """Get the detailed OS name, version, and other platform-specific details."""

        def get_mac_os_info():
            """Get detailed macOS version and product name."""
            try:
                sw_vers = subprocess.check_output(["sw_vers"]).decode("utf-8")
                product_name = (
                    [
                        line
                        for line in sw_vers.split("\n")
                        if line.startswith("ProductName")
                    ][0]
                    .split(":")[1]
                    .strip()
                )
                product_version = (
                    [
                        line
                        for line in sw_vers.split("\n")
                        if line.startswith("ProductVersion")
                    ][0]
                    .split(":")[1]
                    .strip()
                )
                build_version = (
                    [
                        line
                        for line in sw_vers.split("\n")
                        if line.startswith("BuildVersion")
                    ][0]
                    .split(":")[1]
                    .strip()
                )

                # Return the formatted macOS name, version, and build
                return f"{product_name} {product_version} (Build {build_version})"
            except Exception as e:
                return f"Error retrieving macOS name: {str(e)}"

        def get_windows_info():
            """Get detailed Windows version and edition."""
            try:
                # Get basic Windows version using platform
                windows_version = platform.version()
                release = platform.release()
                version = platform.win32_ver()[0]

                # Additional information using Windows-specific system commands
                edition_command = "wmic os get caption"
                edition = (
                    subprocess.check_output(edition_command, shell=True)
                    .decode("utf-8")
                    .strip()
                    .split("\n")[1]
                )

                # Return Windows information
                return f"Windows {version} {release} ({edition})"
            except Exception as e:
                return f"Error retrieving Windows information: {str(e)}"

        def get_linux_info():
            """Get detailed Linux version and distribution info."""
            try:
                # Check /etc/os-release for modern Linux distros
                with open("/etc/os-release") as f:
                    os_info = f.readlines()

                os_name = (
                    next(line for line in os_info if line.startswith("NAME"))
                    .split("=")[1]
                    .strip()
                    .replace('"', "")
                )
                os_version = (
                    next(line for line in os_info if line.startswith("VERSION"))
                    .split("=")[1]
                    .strip()
                    .replace('"', "")
                )

                # For additional info, check for the package manager (e.g., apt, dnf)
                package_manager = "Unknown"
                if os.path.exists("/usr/bin/apt"):
                    package_manager = "APT (Debian/Ubuntu)"
                elif os.path.exists("/usr/bin/dnf"):
                    package_manager = "DNF (Fedora/RHEL)"

                # Return Linux distribution, version, and package manager
                return f"{os_name} {os_version} (Package Manager: {package_manager})"
            except Exception as e:
                return f"Error retrieving Linux information: {str(e)}"

        os_name = platform.system()

        if os_name == "Darwin":
            return get_mac_os_info()
        elif os_name == "Windows":
            return get_windows_info()
        elif os_name == "Linux":
            return get_linux_info()
        else:
            return f"Unknown OS: {os_name} {platform.release()}"

    def get_os_name_and_version():
        os_name = platform.system()
        if os_name == "Darwin":
            try:
                # Run 'sw_vers' command to get macOS details like "macOS Sequoia"
                sw_vers = subprocess.check_output(["sw_vers"]).decode("utf-8")
                product_name = (
                    [
                        line
                        for line in sw_vers.split("\n")
                        if line.startswith("ProductName")
                    ][0]
                    .split(":")[1]
                    .strip()
                )
                product_version = (
                    [
                        line
                        for line in sw_vers.split("\n")
                        if line.startswith("ProductVersion")
                    ][0]
                    .split(":")[1]
                    .strip()
                )

                # Return the formatted macOS name and version
                return f"{product_name} {product_version}"

            except Exception as e:
                return f"Error retrieving macOS name: {str(e)}"

        # For Windows, we use platform to get the OS name and version
        elif os_name == "Windows":
            os_version = platform.version()
            return f"Windows {os_version}"

        # For Linux, check for distribution info using platform and os-release file
        elif os_name == "Linux":
            try:
                # Try to read Linux distribution info from '/etc/os-release'
                with open("/etc/os-release") as f:
                    os_info = f.readlines()

                # Find fields like NAME and VERSION
                os_name = (
                    next(line for line in os_info if line.startswith("NAME"))
                    .split("=")[1]
                    .strip()
                    .replace('"', "")
                )
                os_version = (
                    next(line for line in os_info if line.startswith("VERSION"))
                    .split("=")[1]
                    .strip()
                    .replace('"', "")
                )
                return f"{os_name} {os_version}"

            except Exception as e:
                return f"Error retrieving Linux name: {str(e)}"

        # Default fallback (for unknown OS or edge cases)
        return f"{os_name} {platform.release()}"

    def get_system_uptime():
        """Returns system uptime as a human-readable string."""
        try:
            boot_time = datetime.fromtimestamp(psutil.boot_time())
            uptime = datetime.now() - boot_time
            return str(uptime).split(".")[0]  # Remove microseconds
        except:
            return None

    def get_active_processes(limit=10):
        try:
            processes = []
            for proc in psutil.process_iter(
                ["pid", "name", "cpu_percent", "memory_percent"]
            ):
                try:
                    processes.append(proc.info)
                except psutil.NoSuchProcess:
                    pass
            # Handle NoneType values by treating them as 0
            processes.sort(key=lambda x: x["cpu_percent"] or 0, reverse=True)
            return processes[:limit]
        except:
            return None

    def get_virtual_environment_info():
        """Checks if the script is running in a virtual environment and returns details."""
        try:
            # Check if running in a virtual environment
            if hasattr(sys, "real_prefix") or (
                hasattr(sys, "base_prefix") and sys.base_prefix != sys.prefix
            ):
                return {
                    "Virtual Environment": sys.prefix,
                    "Site-Packages Path": os.path.join(
                        sys.prefix,
                        "lib",
                        "python{}/site-packages".format(sys.version_info.major),
                    ),
                }
            else:
                return {"Virtual Environment": "Not in a virtual environment"}
        except Exception as e:
            return {"Error": str(e)}

    def get_temperatures():
        """Returns temperature sensor readings."""
        try:
            return psutil.sensors_temperatures(fahrenheit=False)
        except AttributeError:
            return {"Error": "Temperature sensors not available"}

    def get_battery_status():
        """Returns battery status."""
        try:
            battery = psutil.sensors_battery()
            if battery:
                time_left = (
                    str(timedelta(seconds=battery.secsleft))
                    if battery.secsleft != psutil.POWER_TIME_UNLIMITED
                    else "Charging/Unlimited"
                )
                return {
                    "Percentage": battery.percent,
                    "Plugged In": battery.power_plugged,
                    "Time Left": time_left,
                }
            return {"Status": "No battery detected"}
        except:
            return {"Status": "No battery detected"}

    def get_disk_io():
        """Returns disk I/O statistics."""
        disk_io = psutil.disk_io_counters()
        return {
            "Read (GB)": disk_io.read_bytes / (1024**3),
            "Write (GB)": disk_io.write_bytes / (1024**3),
            "Read Count": disk_io.read_count,
            "Write Count": disk_io.write_count,
        }

    def get_network_io():
        """Returns network I/O statistics."""
        net_io = psutil.net_io_counters()
        return {
            "Bytes Sent (GB)": net_io.bytes_sent / (1024**3),
            "Bytes Received (GB)": net_io.bytes_recv / (1024**3),
            "Packets Sent": net_io.packets_sent,
            "Packets Received": net_io.packets_recv,
        }

    def run_shell_command(command):
        """Runs a shell command and returns its output."""
        try:
            result = subprocess.run(
                command,
                shell=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
            )
            return (
                result.stdout.strip()
                if result.returncode == 0
                else result.stderr.strip()
            )
        except Exception as e:
            return f"Error running command: {e}"

    system_info = {
        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "os": get_os_type(),
        "system": {
            "os": get_os_info(),
            "platform": f"{platform.system()} {platform.release()}",
            "version": platform.version(),
            "machine": platform.machine(),
            "processor": platform.processor(),
            "architecture": platform.architecture()[0],
            "hostname": socket.gethostname(),
            "ip address": socket.gethostbyname(socket.gethostname()),
            "mac address": ":".join(
                ["{:02x}".format((uuid.getnode() >> i) & 0xFF) for i in range(0, 48, 8)]
            ),
            "cpu brand": cpuinfo.get_cpu_info().get("brand_raw", "Unknown"),
            "python version": platform.python_version(),
            "uptime": get_system_uptime(),
        },
        "cpu": {
            "physical cores": psutil.cpu_count(logical=False),
            "logical cores": psutil.cpu_count(logical=True),
            "max frequency (MHz)": psutil.cpu_freq().max,
            "min frequency (MHz)": psutil.cpu_freq().min,
            "current frequency (MHz)": psutil.cpu_freq().current,
            "usage per core (%)": psutil.cpu_percent(percpu=True),
            "total cpu Usage (%)": psutil.cpu_percent(),
            "load average (1m, 5m, 15m)": (
                os.getloadavg() if hasattr(os, "getloadavg") else "N/A"
            ),
        },
        "memory": {
            "total memory (GB)": psutil.virtual_memory().total / (1024**3),
            "available memory (GB)": psutil.virtual_memory().available / (1024**3),
            "used memory (GB)": psutil.virtual_memory().used / (1024**3),
            "memory usage (%)": psutil.virtual_memory().percent,
            "swap total (GB)": psutil.swap_memory().total / (1024**3),
            "swap free (GB)": psutil.swap_memory().free / (1024**3),
            "swap used (GB)": psutil.swap_memory().used / (1024**3),
            "swap usage (%)": psutil.swap_memory().percent,
        },
        "disk": {},
        "disk io": get_disk_io(),
        "network": {},
        "network io": get_network_io(),
        "gpu": [],
        # "temperatures": get_temperatures(),
        # "battery": get_battery_status(),
        "active processes": get_active_processes(),
        "environment": {
            "user": os.getenv("USER", "Unknown"),
            "environment variables": dict(os.environ),
            "virtual environment info": get_virtual_environment_info(),  # Virtual env details
            "docker running": os.path.exists("/.dockerenv"),  # Check for Docker
            "shell": os.environ.get("SHELL", "Unknown"),
            "default terminal": run_shell_command("echo $TERM"),
            "kernel version": platform.uname().release,
            "virtualization type": run_shell_command("systemd-detect-virt"),
        },
        "additional info": {
            "Shell": os.environ.get("SHELL", "Unknown"),
            "default terminal": run_shell_command("echo $TERM"),
            "kernel version": platform.uname().release,
            "virtualization type": run_shell_command("systemd-detect-virt"),
            "running in docker": os.path.exists("/.dockerenv"),
        },
    }

    # Disk Information
    for partition in psutil.disk_partitions():
        try:
            usage = psutil.disk_usage(partition.mountpoint)
            system_info["disk"][partition.device] = {
                "mountpoint": partition.mountpoint,
                "file system type": partition.fstype,
                "total size (GB)": usage.total / (1024**3),
                "used (GB)": usage.used / (1024**3),
                "free (GB)": usage.free / (1024**3),
                "usage (%)": usage.percent,
            }
        except PermissionError:
            system_info["disk"][partition.device] = "Permission Denied"

    # Network Information
    if_addrs = psutil.net_if_addrs()
    for interface_name, interface_addresses in if_addrs.items():
        system_info["network"][interface_name] = []
        for address in interface_addresses:
            if str(address.family) == "AddressFamily.AF_INET":
                system_info["network"][interface_name].append(
                    {
                        "ip address": address.address,
                        "netmask": address.netmask,
                        "broadcast ip": address.broadcast,
                    }
                )
            elif str(address.family) == "AddressFamily.AF_PACKET":
                system_info["network"][interface_name].append(
                    {
                        "mac address": address.address,
                        "netmask": address.netmask,
                        "broadcast mac": address.broadcast,
                    }
                )

    # # GPU Information
    # gpus = GPUtil.getGPUs()
    # for gpu in gpus:
    #     gpu_info = {
    #         "name": gpu.name,
    #         "load (%)": gpu.load * 100,
    #         "free memory (MB)": gpu.memoryFree,
    #         "used memory (MB)": gpu.memoryUsed,
    #         "total memory (MB)": gpu.memoryTotal,
    #         "driver version": gpu.driver,
    #         "temperature (°C)": gpu.temperature,
    #     }
    #     if hasattr(gpu, "powerDraw"):
    #         gpu_info["Power Draw (W)"] = gpu.powerDraw
    #     if hasattr(gpu, "powerLimit"):
    #         gpu_info["Power Limit (W)"] = gpu.powerLimit
    #     system_info["gpu"].append(gpu_info)

    res = system_info if full else get_os_type()
    if verbose:
        try:
            preview(res)
        except Exception as e:
            pnrint(e)
    return res


def listdir(
    rootdir,
    kind=None,
    sort_by="name",
    ascending=True,
    contains=None,  # filter filenames using re
    booster=False,  # walk in subfolders
    depth=0,  # 0: no subfolders; None: all subfolders; [int 1,2,3]: levels of subfolders
    hidden=False,  # Include hidden files/folders
    orient="list",
    output="df",  # "df", 'list','dict','records','index','series'
    verbose=True,
):
    def is_hidden(filepath):
        """Check if a file or folder is hidden."""
        system = platform.system()
        if system == "Windows":
            import ctypes

            attribute = ctypes.windll.kernel32.GetFileAttributesW(filepath)
            if attribute == -1:
                raise FileNotFoundError(f"File {filepath} not found.")
            return bool(attribute & 2)  # FILE_ATTRIBUTE_HIDDEN
        else:  # macOS/Linux: Hidden if the name starts with a dot
            return os.path.basename(filepath).startswith(".")

    def get_user():
        """Retrieve the username of the current user."""
        system = platform.system()
        if system == "Windows":
            return os.environ.get("USERNAME", "Unknown")
        else:
            import pwd

            return pwd.getpwuid(os.getuid()).pw_name

    if isinstance(kind, list):
        f_ = []
        for kind_ in kind:
            f_tmp = listdir(
                rootdir=rootdir,
                kind=kind_,
                sort_by=sort_by,
                ascending=ascending,
                contains=contains,
                depth=depth,  # walk in subfolders
                hidden=hidden,
                orient=orient,
                output=output,
                verbose=verbose,
            )
            f_.append(f_tmp)
        if f_:
            return pd.concat(f_, ignore_index=True)
    if kind is not None:
        if not kind.startswith("."):
            kind = "." + kind
    fd = [".fd", ".fld", ".fol", ".fd", ".folder"]
    i = 0
    f = {
        "name": [],
        "kind": [],
        "length": [],
        "basename": [],
        "path": [],
        "created_time": [],
        "modified_time": [],
        "last_open_time": [],
        "size": [],
        "permission": [],
        "owner": [],
        "rootdir": [],
        "fname": [],
        "fpath": [],
        "num": [],
        "os": [],
    }
    root_depth = rootdir.rstrip(os.sep).count(os.sep)
    for dirpath, dirnames, ls in os.walk(rootdir):
        current_depth = dirpath.rstrip(os.sep).count(os.sep) - root_depth
        # Check depth limit
        if depth is not None and current_depth > depth:
            dirnames[:] = []  # Prevent further traversal into subfolders
            continue

        if not hidden:
            dirnames[:] = [
                d for d in dirnames if not is_hidden(os.path.join(dirpath, d))
            ]
            ls = [i for i in ls if not is_hidden(os.path.join(dirpath, i))]

        for dirname in dirnames:
            if kind is not None and kind not in fd:  # do not check folders
                continue
            if contains and not re.search(contains, dirname):
                continue
            dirname_path = os.path.join(dirpath, dirname)
            fpath = os.path.join(os.path.dirname(dirname_path), dirname)
            try:
                stats_file = os.stat(fpath)
            except Exception as e:
                print(e)
                continue
            filename, file_extension = os.path.splitext(dirname)
            file_extension = file_extension if file_extension != "" else None
            f["name"].append(filename)
            f["kind"].append(file_extension)
            f["length"].append(len(filename))
            f["size"].append(round(os.path.getsize(fpath) / 1024 / 1024, 3))
            f["basename"].append(os.path.basename(dirname_path))
            f["path"].append(os.path.join(os.path.dirname(dirname_path), dirname))
            f["created_time"].append(
                pd.to_datetime(int(os.path.getctime(dirname_path)), unit="s")
            )
            f["modified_time"].append(
                pd.to_datetime(int(os.path.getmtime(dirname_path)), unit="s")
            )
            f["last_open_time"].append(
                pd.to_datetime(int(os.path.getatime(dirname_path)), unit="s")
            )
            f["permission"].append(stat.filemode(stats_file.st_mode)),
            f["owner"].append(get_user()),
            f["rootdir"].append(dirpath)
            f["fname"].append(filename)  # will be removed
            f["fpath"].append(fpath)  # will be removed
            i += 1
        for item in ls:
            if kind in fd:  # only check folders
                continue
            if contains and not re.search(contains, item):
                continue
            item_path = os.path.join(dirpath, item)
            fpath = os.path.join(os.path.dirname(item_path), item)
            try:
                stats_file = os.stat(fpath)
            except Exception as e:
                print(e)
                continue
            filename, file_extension = os.path.splitext(item)
            if kind is not None:
                is_folder = kind.lower() in fd and os.path.isdir(item_path)
                is_file = kind.lower() in file_extension.lower() and (
                    os.path.isfile(item_path)
                )
                if kind in [
                    ".doc",
                    ".img",
                    ".zip",
                    ".code",
                    ".file",
                    ".image",
                    ".video",
                    ".audio",
                ]:  # 选择大的类别
                    if kind != ".folder" and not isa(item_path, kind):
                        continue
                elif kind in [".all"]:
                    return flist(fpath, contains=contains)
                else:  # 精确到文件的后缀
                    if not is_folder and not is_file:
                        continue
            file_extension = file_extension if file_extension != "" else None
            f["name"].append(filename)
            f["kind"].append(file_extension)
            f["length"].append(len(filename))
            f["size"].append(round(os.path.getsize(fpath) / 1024 / 1024, 3))
            f["basename"].append(os.path.basename(item_path))
            f["path"].append(os.path.join(os.path.dirname(item_path), item))
            f["created_time"].append(
                pd.to_datetime(int(os.path.getctime(item_path)), unit="s")
            )
            f["modified_time"].append(
                pd.to_datetime(int(os.path.getmtime(item_path)), unit="s")
            )
            f["last_open_time"].append(
                pd.to_datetime(int(os.path.getatime(item_path)), unit="s")
            )
            f["permission"].append(stat.filemode(stats_file.st_mode)),
            f["owner"].append(
                os.getlogin() if platform.system() != "Windows" else "N/A"
            ),
            f["fname"].append(filename)  # will be removed
            f["fpath"].append(fpath)  # will be removed
            f["rootdir"].append(dirpath)
            i += 1

        f["num"] = i
        f["os"] = get_os()  # os.uname().machine
        # if not booster: # go deeper subfolders
        #     break
    # * convert to pd.DataFrame
    f = pd.DataFrame(f)
    f = f[
        [
            "basename",
            "name",
            "kind",
            "length",
            "size",
            "num",
            "path",
            "created_time",
            "modified_time",
            "last_open_time",
            "rootdir",
            "permission",
            "owner",
            "os",
            "fname",
            "fpath",
        ]
    ]
    if "nam" in sort_by.lower():
        f = sort_kind(f, by="name", ascending=ascending)
    elif "crea" in sort_by.lower():
        f = sort_kind(f, by="created_time", ascending=ascending)
    elif "modi" in sort_by.lower():
        f = sort_kind(f, by="modified_time", ascending=ascending)
    elif "s" in sort_by.lower() and "z" in sort_by.lower():
        f = sort_kind(f, by="size", ascending=ascending)

    if "df" in output:
        if verbose:
            display(f.head())
            print(f"shape: {f.shape}")
        return f
    else:
        from box import Box

        if "l" in orient.lower():  # list # default
            res_output = Box(f.to_dict(orient="list"))
            return res_output
        if "d" in orient.lower():  # dict
            return Box(f.to_dict(orient="dict"))
        if "r" in orient.lower():  # records
            return Box(f.to_dict(orient="records"))
        if "in" in orient.lower():  # records
            return Box(f.to_dict(orient="index"))
        if "se" in orient.lower():  # records
            return Box(f.to_dict(orient="series"))


def listpkg(where="env", verbose=False):
    """list all pacakages"""

    def listfunc_(lib_name, opt="call"):
        """list functions in specific lib"""
        if opt == "call":
            funcs = [
                func
                for func in dir(lib_name)
                if callable(getattr(lib_name, func))
                if not func.startswith("__")
            ]
        else:
            funcs = dir(lib_name)
        return funcs

    if any([i in where.lower() for i in ["all", "env"]]):
        import pkg_resources

        lst_pkg = [pkg.key for pkg in pkg_resources.working_set]
    else:
        lst_pkg = listfunc_(where)
    print(lst_pkg) if verbose else None
    return lst_pkg


def copy(src, dst, overwrite=False, verbose=True):
    """Copy a file from src to dst."""
    try:
        dir_par_dst = os.path.dirname(dst)
        if not os.path.isdir(dir_par_dst):
            mkdir(dir_par_dst)
        print(dir_par_dst) if verbose else None
        src = Path(src)
        dst = Path(dst)
        if not src.is_dir():
            if dst.is_dir():
                dst = dst / src.name

            if dst.exists():
                if overwrite:
                    dst.unlink()
                else:
                    dst = dst.with_name(
                        f"{dst.stem}_{datetime.now().strftime('%y%m%d_%H%M%S')}{dst.suffix}"
                    )
            shutil.copy(src, dst)
            print(f"\n Done! copy to {dst}\n") if verbose else None
            return dst
        else:
            dst = dst / src.name
            if dst.exists():
                if overwrite:
                    shutil.rmtree(dst)  # Remove existing directory
                else:
                    dst = dst.with_name(
                        f"{dst.stem}_{datetime.now().strftime('%y%m%d%H%M%S')}"
                    )
            shutil.copytree(src, dst)
            print(f"\n Done! copy to {dst}\n") if verbose else None
            return dst

    except Exception as e:
        logging.error(f"Failed {e}")
def local_path(fpath,station=r"Q:\\IM\\AGLengerke\\Jeff\\# testing\\temp\\"):
    """copy file to a specific folder first, to aviod file conflict""" 
    try:
        f=listdir(station) 
        if listdir(station ,verbose=False).loc[0,"num"]>=10:
            for fpath_ in f['path']:
                if os.path.basename(fpath)[:5] in fpath_:
                    if fpath== fpath_:
                        pass
                    else:
                        delete(fpath_)
    except:
        pass
    try:
        new_path=copy(fpath, station)
    except Exception as e:
        print(f"Path did not update because: Error:{e}")
        new_path=fpath
    return new_path

def cut(src, dst, overwrite=False):
    return move(src=src, dst=dst, overwrite=overwrite)


def move(src, dst, overwrite=False, verbose=True):
    try:
        dir_par_dst = os.path.dirname(dst)
        if not os.path.isdir(dir_par_dst):
            mkdir(dir_par_dst)
        src = Path(src)
        dst = Path(dst)
        if dst.is_dir():
            dst = dst / src.name
        if dst.exists():
            if overwrite:
                # dst.unlink()  # Delete the existing file
                pass
            else:
                dst = dst.with_name(
                    f"{dst.stem}_{datetime.now().strftime('_%H%M%S')}{dst.suffix}"
                )
        shutil.move(src, dst)
        print(f"\n Done! moved to {dst}\n") if verbose else None
    except Exception as e:
        logging.error(f"Failed to move file from {src} to {dst}: {e}")


def delete(fpath):
    """Delete a file/folder."""
    try:
        fpath = Path(fpath)
        if not fpath.is_dir():  # file
            if fpath.exists():
                fpath.unlink()
                print(f"\n Done! delete {fpath}\n")
            else:
                print(f"File '{fpath}' does not exist.")
        else:  # folder
            if fpath.exists():
                shutil.rmtree(fpath)  # Remove existing directory
                print(f"\n Done! delete {fpath}\n")
            else:
                print(f"Folder '{fpath}' does not exist.")
    except Exception as e:
        logging.error(f"Failed to delete {fpath}: {e}")


def rename(fpath, dst, smart=True):
    """Rename a file or folder."""
    try:
        src_kind, dst_kind = None, None
        if smart:
            dir_name_src = os.path.dirname(fpath)
            dir_name_dst = os.path.dirname(dst)
            src_kind = os.path.splitext(fpath)[1]
            dst_kind = os.path.splitext(dst)[1]
            if dir_name_dst != dir_name_src:
                dst = os.path.join(dir_name_src, dst)
            if dst_kind is not None and src_kind is not None:
                if dst_kind != src_kind:
                    dst = dst + src_kind
        if os.path.exists(fpath):
            os.rename(fpath, dst)
            print(f"Done! rename to {dst}")
        else:
            print(f"Failed: {fpath} does not exist.")
    except Exception as e:
        logging.error(f"Failed to rename {fpath} to {dst}: {e}")


def mkdir_nest(fpath: str) -> str:
    """
    Create nested directories based on the provided file path.

    Parameters:
    - fpath (str): The full file path for which the directories should be created.

    Returns:
    - str: The path of the created directory.
    """
    # Split the full path into directories
    f_slash = "/" if "mac" in get_os().lower() else "\\"
    if os.path.isdir(fpath):
        fpath = fpath + f_slash if not fpath.endswith(f_slash) else fpath
        return fpath
    dir_parts = fpath.split(f_slash)  # Split the path by the OS-specific separator

    # Start creating directories from the root to the desired path
    root_dir = os.path.splitdrive(fpath)[
        0
    ]  # Get the root drive on Windows (e.g., 'C:')
    current_path = (
        root_dir if root_dir else f_slash
    )  # Start from the root directory or POSIX '/'

    for part in dir_parts:
        if part:
            current_path = os.path.join(current_path, part)
            if not os.path.isdir(current_path):
                os.makedirs(current_path)
    if not current_path.endswith(f_slash):
        current_path += f_slash
    return current_path


def mkdir(pardir: str = None, chdir: str | list = None, overwrite=False):
    """
    Create a directory.

    Parameters:
    - pardir (str): Parent directory where the new directory will be created. If None, uses the current working directory.
    - chdir (str | list): Name of the new directory or a list of directories to create.
                          If None, a default name 'new_directory' will be used.
    - overwrite (bool): If True, overwrite the directory if it already exists. Defaults to False.

    Returns:
    - str: The path of the created directory or an error message.
    """
    rootdir = []
    pardir = mkdir_nest(pardir)
    if chdir is None:
        return pardir
    else:
        pass
    print(pardir)
    if isinstance(chdir, str):
        chdir = [chdir]
    chdir = list(set(chdir))
    if isinstance(pardir, str):  # Dir_parents should be 'str' type
        pardir = os.path.normpath(pardir)
    if "mac" in get_os().lower() or "lin" in get_os().lower():
        stype = "/"
    elif "win" in get_os().lower():
        stype = "\\"
    else:
        stype = "/"

    if os.path.isdir(pardir):
        os.chdir(pardir)  # Set current path
        # Check if subdirectories are not empty
        if chdir:
            chdir.sort()
            for folder in chdir:
                child_tmp = os.path.join(pardir, folder)
                if not os.path.isdir(child_tmp):
                    os.mkdir("./" + folder)
                    print(f"\n {folder} was created successfully!\n")
                else:
                    if overwrite:
                        shutil.rmtree(child_tmp)
                        os.mkdir("./" + folder)
                        print(f"\n {folder} overwrite! \n")
                    else:
                        print(f"\n {folder} already exists! \n")
                rootdir.append(child_tmp + stype)  # Note down
        else:
            print("\nWarning: Dir_child doesn't exist\n")
    else:
        print("\nWarning: Dir_parent is not a directory path\n")
    # Dir is the main output, if only one dir, then str type is inconvenient
    if len(rootdir) == 1:
        rootdir = rootdir[0]
        rootdir = rootdir + stype if not rootdir.endswith(stype) else rootdir

    return rootdir


def split_path(fpath):
    f_slash = "/" if "mac" in get_os().lower() else "\\"
    dir_par = f_slash.join(fpath.split(f_slash)[:-1])
    dir_ch = "".join(fpath.split(f_slash)[-1:])
    return dir_par, dir_ch

from pathlib import Path
import matplotlib.pyplot as plt
from PIL import Image
import numpy as np
import cv2

def figsave(*args, dpi=300, **kwargs):
    bbox_inches = kwargs.pop("bbox_inches", "tight")
    pad_inches = kwargs.pop("pad_inches", 0)
    facecolor = kwargs.pop("facecolor", "white")
    edgecolor = kwargs.pop("edgecolor", "auto")

    dir_save = None
    fname = None
    img = None

    for arg in args:
        if isinstance(arg, str):
            path = Path(arg)
            if path.suffix:  # Has file extension
                fname = path.name
                dir_save = path.parent
            else:
                dir_save = path
        elif isinstance(arg, (Image.Image, np.ndarray)):
            img = arg  # Store PIL image or numpy array

    # Set default save directory
    dir_save = Path(dir_save) if dir_save else Path(".")
    dir_save.mkdir(parents=True, exist_ok=True)

    # Handle filename and extension
    if fname is None:
        fname = "figure"
    fname = dir_save / fname
    if fname.suffix == "":
        fname = fname.with_suffix(".pdf")  # Default format

    ftype = fname.suffix.lstrip(".").lower()

    # Save figure based on file type
    if ftype == "eps":
        plt.savefig(fname, format="eps", bbox_inches=bbox_inches)
        plt.savefig(fname.with_suffix(".pdf"), format="pdf", dpi=dpi,
                    pad_inches=pad_inches, bbox_inches=bbox_inches,
                    facecolor=facecolor, edgecolor=edgecolor)
    elif ftype == "pdf":
        plt.savefig(fname, format="pdf", dpi=dpi, pad_inches=pad_inches,
                    bbox_inches=bbox_inches, facecolor=facecolor, edgecolor=edgecolor)
    elif ftype in ["jpg", "jpeg", "png", "tiff", "tif"]:
        if img is not None:  # If an image is provided
            if isinstance(img, Image.Image):
                img = img.convert("RGB") if img.mode == "RGBA" else img
                img.save(fname, format=ftype.upper(), dpi=(dpi, dpi))
            elif isinstance(img, np.ndarray):
                if img.ndim == 2:
                    Image.fromarray(img).save(fname, format=ftype.upper(), dpi=(dpi, dpi))
                elif img.ndim == 3:
                    if img.shape[2] == 3:
                        img = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
                    elif img.shape[2] == 4:
                        img = cv2.cvtColor(img, cv2.COLOR_BGRA2RGBA)
                    Image.fromarray(img).save(fname, format=ftype.upper(), dpi=(dpi, dpi))
                else:
                    raise ValueError("Unexpected image dimensions.")
        else:
            plt.savefig(fname, format=ftype, dpi=dpi, pad_inches=pad_inches,
                        bbox_inches=bbox_inches, facecolor=facecolor, edgecolor=edgecolor)
    elif ftype == "ico":
        if img is None:
            plt.savefig(fname, dpi=dpi, pad_inches=pad_inches,
                        bbox_inches=bbox_inches, facecolor=facecolor, edgecolor=edgecolor)
            img = Image.open(fname)
        img = img.convert("RGBA")
        icon_sizes = [(32, 32), (64, 64), (128, 128), (256, 256)]
        img.save(fname, format="ICO", sizes=icon_sizes)
        print(f"Icon saved @: {fname} with sizes: {icon_sizes}")
    else:
        raise ValueError(f"Unsupported file format: {ftype}")

    print(f"\nSaved @ {fname} (dpi={dpi})")

def figsave(*args, dpi=300, **kwargs):
    """
    Save a Matplotlib figure or image file in various formats.

    This function automatically determines whether to save a Matplotlib figure 
    or an image (PIL or NumPy array) and handles different file formats, including:
    - PDF, EPS, PNG, JPG, TIFF, ICO, EMF

    Parameters:
    -----------
    *args : str, PIL.Image, np.ndarray
        - File path (directory and/or filename) to save the figure.
        - If an image is provided (PIL or NumPy), it will be saved accordingly.
    
    dpi : int, optional (default=300)
        - Resolution (dots per inch) for saved figures.
    
    **kwargs : dict
        - Additional keyword arguments for `matplotlib.pyplot.savefig()`, including:
            - bbox_inches (str, default="tight"): Bounding box for figure.
            - pad_inches (float, default=0): Padding around figure.
            - facecolor (str, default="white"): Background color.
            - edgecolor (str, default="auto"): Edge color.

    Supported Formats:
    ------------------
    - Vector: `pdf`, `eps`, `emf`
    - Raster: `png`, `jpg`, `jpeg`, `tiff`, `tif`, `ico`
    
    Example Usage:
    --------------
    >>> figsave("output_plot.pdf")
    >>> figsave("figs/plot.png", dpi=600)
    >>> figsave("./results/figure", format="pdf")
    >>> figsave("icons/logo.ico", image)  # Save an image file as an icon

    Returns:
    --------
    None
    """
    import matplotlib.pyplot as plt
    from PIL import Image
    from pathlib import Path
    
    bbox_inches=kwargs.pop("bbox_inches","tight")
    pad_inches=kwargs.pop("pad_inches",0)
    facecolor=kwargs.pop("facecolor",'white')
    edgecolor=kwargs.pop("edgecolor",'auto')
    
    dir_save = None
    fname = None
    img = None
    for arg in args:
        if isinstance(arg, str):
            path = Path(arg)
            if path.suffix:  # Has file extension
                fname = path.name
                dir_save = path.parent
            else:
                dir_save = path
        elif isinstance(arg, (Image.Image, np.ndarray)):
            img = arg  # Store PIL image or numpy array

    dir_save = Path(dir_save) if dir_save else Path(".")
    dir_save.mkdir(parents=True, exist_ok=True)
    # Handle filename and extension
    if fname is None:
        fname = dir_save
    else:
        fname = dir_save / fname
    if fname.suffix == "":
        fname = fname.with_suffix(".pdf")  # Default format

    ftype = fname.suffix.lstrip(".").lower()

    # Save figure based on file type
    if ftype == "eps":
        plt.savefig(fname, format="eps", bbox_inches=bbox_inches)
        plt.savefig(fname.with_suffix(".pdf"), format="pdf", dpi=dpi,
                    pad_inches=pad_inches, bbox_inches=bbox_inches,
                    facecolor=facecolor, edgecolor=edgecolor) 
    elif ftype.lower() in ["jpg", "jpeg", "png", "tiff", "tif"]:
        if img is not None:  # If a PIL image is provided
            if isinstance(img, Image.Image):
                if img.mode == "RGBA":
                    img = img.convert("RGB")
                img.save(fname, format=ftype.upper(), dpi=(dpi, dpi))
            elif isinstance(img, np.ndarray):
                import cv2
                if img.ndim == 2:
                    # Grayscale image
                    Image.fromarray(img).save(fname, format=ftype.upper(), dpi=(dpi, dpi))
                elif img.ndim == 3:
                    if img.shape[2] == 3:
                        # RGB image
                        img = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
                        Image.fromarray(img).save(fname, format=ftype.upper(), dpi=(dpi, dpi))
                    elif img.shape[2] == 4:
                        # RGBA image
                        img = cv2.cvtColor(img, cv2.COLOR_BGRA2RGBA)  # Convert BGRA to RGBA
                        Image.fromarray(img).save(fname, format=ftype.upper(), dpi=(dpi, dpi))
                    else:
                        raise ValueError("Unexpected number of channels in the image array.")
                else:
                    raise ValueError("Image array has an unexpected number of dimensions.")
        else:
            plt.savefig(fname,format=ftype.lower(),dpi=dpi,pad_inches=pad_inches,bbox_inches=bbox_inches,facecolor=facecolor,edgecolor=edgecolor)
    elif ftype.lower() in ["emf","pdf","fig"]:
        plt.savefig(fname,format=ftype.lower(),dpi=dpi,pad_inches=pad_inches,bbox_inches=bbox_inches,facecolor=facecolor,edgecolor=edgecolor)
    elif ftype.lower() == "ico":
        # Ensure the image is in a format that can be saved as an icon (e.g., 32x32, 64x64, etc.)
        if img is None:  # If no image is provided, use the matplotlib figure
            img = plt.figure() 
            print(fname)
            img.savefig(fname, 
                        format="png",
                        dpi=dpi,
                        pad_inches=pad_inches,
                        bbox_inches=bbox_inches,
                        facecolor=facecolor,
                        edgecolor=edgecolor )
            img = Image.open(fname)  # Load the saved figure image

        # Resize the image to typical icon sizes and save it as .ico
        icon_sizes = [(32, 32), (64, 64), (128, 128), (256, 256)]
        img = img.convert("RGBA") 
        img.save(fname, format="ICO", sizes=icon_sizes)
        print(f"Icon saved @: {fname} with sizes: {icon_sizes}")
    print(f"\n✅ Saved @: dpi={dpi}\n{fname}")


def is_str_color(s):
    # Regular expression pattern for hexadecimal color codes
    if isinstance(s, str):
        import re

        color_code_pattern = r"^#([A-Fa-f0-9]{6}|[A-Fa-f0-9]{8})$"
        return re.match(color_code_pattern, s) is not None
    else:
        return True


def is_num(s):
    """
    Check if a string can be converted to a number (int or float).
    Parameters:
    - s (str): The string to check.
    Returns:
    - bool: True if the string can be converted to a number, False otherwise.
    """
    try:
        float(s)  # Try converting the string to a float
        return True
    except ValueError:
        return False


def isnum(s):
    return is_num(s)


def is_image(fpath):
    """
    Determine if a given file is an image based on MIME type and file extension.

    Args:
        fpath (str): Path to the file.

    Returns:
        bool: True if the file is a recognized image, False otherwise.
    """
    from PIL import Image

    if isinstance(fpath, str):
        import mimetypes

        # Known image MIME types
        image_mime_types = {
            "image/jpeg",
            "image/png",
            "image/gif",
            "image/bmp",
            "image/webp",
            "image/tiff",
            "image/x-icon",
            "image/svg+xml",
            "image/heic",
            "image/heif",
        }

        # Known image file extensions
        image_extensions = {
            ".jpg",
            ".jpeg",
            ".png",
            ".gif",
            ".bmp",
            ".webp",
            ".tif",
            ".tiff",
            ".ico",
            ".svg",
            ".heic",
            ".heif",
            ".fig",
            ".jpg",
        }

        # Get MIME type using mimetypes
        mime_type, _ = mimetypes.guess_type(fpath)

        # Check MIME type
        if mime_type in image_mime_types:
            return True

        # Fallback: Check file extension
        ext = os.path.splitext(fpath)[
            -1
        ].lower()  # Get the file extension and ensure lowercase
        if ext in image_extensions:
            return True

        return False

    elif isinstance(fpath, Image.Image):
        # If the input is a PIL Image object
        return True

    return False


def is_video(fpath):
    """
    Determine if a given file is a video based on MIME type and file extension.

    Args:
        fpath (str): Path to the file.

    Returns:
        bool: True if the file is a recognized video, False otherwise.
    """
    import mimetypes

    # Known video MIME types
    video_mime_types = {
        "video/mp4",
        "video/quicktime",
        "video/x-msvideo",
        "video/x-matroska",
        "video/x-flv",
        "video/webm",
        "video/ogg",
        "video/x-ms-wmv",
        "video/x-mpeg",
        "video/3gpp",
        "video/avi",
        "video/mpeg",
        "video/x-mpeg2",
        "video/x-ms-asf",
    }

    # Known video file extensions
    video_extensions = {
        ".mp4",
        ".mov",
        ".avi",
        ".mkv",
        ".flv",
        ".webm",
        ".ogv",
        ".wmv",
        ".mpg",
        ".mpeg",
        ".3gp",
        ".mpeg2",
        ".asf",
        ".ts",
        ".m4v",
        ".divx",
    }

    # Get MIME type using mimetypes
    mime_type, _ = mimetypes.guess_type(fpath)

    # Check MIME type
    if mime_type in video_mime_types:
        return True

    # Fallback: Check file extension
    ext = os.path.splitext(fpath)[
        -1
    ].lower()  # Get the file extension and ensure lowercase
    if ext in video_extensions:
        return True

    return False


def is_document(fpath):
    """
    Determine if a given file is a document based on MIME type and file extension.

    Args:
        fpath (str): Path to the file.

    Returns:
        bool: True if the file is a recognized document, False otherwise.
    """
    import mimetypes

    # Define known MIME types for documents
    document_mime_types = {
        "text/",
        "application/pdf",
        "application/msword",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/rtf",
        "application/x-latex",
        "application/vnd.oasis.opendocument.text",
        "application/vnd.oasis.opendocument.spreadsheet",
        "application/vnd.oasis.opendocument.presentation",
    }

    # Define extensions for fallback
    document_extensions = {
        ".txt",
        ".log",
        ".csv",
        ".json",
        ".xml",
        ".pdf",
        ".doc",
        ".docx",
        ".xls",
        ".xlsx",
        ".ppt",
        ".pptx",
        ".odt",
        ".ods",
        ".odp",
        ".rtf",
        ".tex",
    }

    # Get MIME type
    mime_type, _ = mimetypes.guess_type(fpath)

    # Check MIME type
    if mime_type and any(
        mime_type.startswith(doc_type) for doc_type in document_mime_types
    ):
        return True

    # Fallback: Check file extension
    ext = os.path.splitext(fpath)[
        -1
    ].lower()  # Get the extension, ensure it's lowercase
    if ext in document_extensions:
        return True

    return False


def is_audio(fpath):
    """
    Determine if a given file is an audio file based on MIME type and file extension.

    Args:
        fpath (str): Path to the file.

    Returns:
        bool: True if the file is a recognized audio file, False otherwise.
    """
    import mimetypes

    # Known audio MIME types
    audio_mime_types = {
        "audio/mpeg",
        "audio/wav",
        "audio/ogg",
        "audio/aac",
        "audio/flac",
        "audio/midi",
        "audio/x-midi",
        "audio/x-wav",
        "audio/x-flac",
        "audio/pcm",
        "audio/x-aiff",
        "audio/x-m4a",
    }

    # Known audio file extensions
    audio_extensions = {
        ".mp3",
        ".wav",
        ".ogg",
        ".aac",
        ".flac",
        ".midi",
        ".m4a",
        ".aiff",
        ".pcm",
        ".wma",
        ".ape",
        ".alac",
        ".opus",
    }

    # Get MIME type using mimetypes
    mime_type, _ = mimetypes.guess_type(fpath)

    # Check MIME type
    if mime_type in audio_mime_types:
        return True

    # Fallback: Check file extension
    ext = os.path.splitext(fpath)[
        -1
    ].lower()  # Get the file extension and ensure lowercase
    if ext in audio_extensions:
        return True

    return False


def is_code(fpath):
    """
    Determine if a given file is a code file based on file extension and optionally MIME type.

    Args:
        fpath (str): Path to the file.
        check_mime (bool): Whether to perform a MIME type check in addition to file extension check.

    Returns:
        bool: True if the file is a recognized code file, False otherwise.
    """
    # Known programming and scripting file extensions
    code_extensions = {
        ".m",
        ".py",
        ".ipynb",
        ".js",
        ".html",
        ".css",
        ".java",
        ".cpp",
        ".h",
        ".cs",
        ".go",
        ".rs",
        ".sh",
        ".rb",
        ".swift",
        ".ts",
        ".json",
        ".xml",
        ".yaml",
        ".toml",
        ".bash",
        ".r",
    }

    # Check file extension
    ext = os.path.splitext(fpath)[-1].lower()
    if ext in code_extensions:
        return True
    return False


def is_zip(fpath):
    import mimetypes

    mime_type, _ = mimetypes.guess_type(fpath)
    if mime_type == "application/zip":
        return True
    else:
        return False


def adjust_spines(ax=None, spines=["left", "bottom"], distance=2):
    import matplotlib.pyplot as plt

    if ax is None:
        ax = plt.gca()
    for loc, spine in ax.spines.items():
        if loc in spines:
            spine.set_position(("outward", distance))  # outward by 2 points
            # spine.set_smart_bounds(True)
        else:
            spine.set_color("none")  # don't draw spine
    # turn off ticks where there is no spine
    if "left" in spines:
        ax.yaxis.set_ticks_position("left")
    else:
        ax.yaxis.set_ticks([])
    if "bottom" in spines:
        ax.xaxis.set_ticks_position("bottom")
    else:
        # no xaxis ticks
        ax.xaxis.set_ticks([])


# And then plot the data:


def add_colorbar(im, width=None, pad=None, **kwargs):
    # usage: add_colorbar(im, width=0.01, pad=0.005, label="PSD (dB)", shrink=0.8)
    l, b, w, h = im.axes.get_position().bounds  # get boundaries
    width = width or 0.1 * w  # get width of the colorbar
    pad = pad or width  # get pad between im and cbar
    fig = im.axes.figure  # get figure of image
    cax = fig.add_axes([l + w + pad, b, width, h])  # define cbar Axes
    return fig.colorbar(im, cax=cax, **kwargs)  # draw cbar


# =============================================================================
# # for plot figures: setting rcParams
# usage: set_pub()
# or by setting sns.set_theme...see below:
# sns.set_theme(style="ticks", rc=params)      # 白色无刻度线，有坐标轴标度
# # sns.set_theme(style="whitegrid", rc=params)# 白色＋刻度线，无坐标轴标度
# # sns.set_theme(style="white", rc=params)    # 白色无刻度线，无坐标轴标度
# # sns.set_theme(style="dark", rc=params)     # 深色无刻度线，无坐标轴标度
# =============================================================================


def list2str(x_str):
    s = "".join(str(x) for x in x_str)
    return s


def str2list(str_):
    l = []
    [l.append(x) for x in str_]
    return l


def str2words(
    content,
    method="combined",
    custom_dict=None,
    sym_spell_params=None,
    use_threading=True,
):
    """
    Ultimate text correction function supporting multiple methods,
    lists or strings, and domain-specific corrections.

    Parameters:
        content (str or list): Input text or list of strings to correct.
        method (str): Correction method ('textblob', 'sym', 'combined').
        custom_dict (dict): Custom dictionary for domain-specific corrections.
        sym_spell_params (dict): Parameters for initializing SymSpell.

    Returns:
        str or list: Corrected text or list of corrected strings.
    """
    from textblob import TextBlob
    from symspellpy import SymSpell, Verbosity
    from functools import lru_cache
    import pkg_resources
    from concurrent.futures import ThreadPoolExecutor

    def initialize_symspell(max_edit_distance=2, prefix_length=7):
        """Initialize SymSpell for advanced spelling correction."""
        sym_spell = SymSpell(max_edit_distance, prefix_length)
        dictionary_path = pkg_resources.resource_filename(
            "symspellpy",
            # "frequency_bigramdictionary_en_243_342.txt",
            "frequency_dictionary_en_82_765.txt",
        )

        sym_spell.load_dictionary(dictionary_path, term_index=0, count_index=1)
        return sym_spell

    def segment_words(text, sym_spell):
        """Segment concatenated words into separate words."""
        segmented = sym_spell.word_segmentation(text)
        return segmented.corrected_string

    @lru_cache(maxsize=1000)  # Cache results for repeated corrections
    def advanced_correction(word, sym_spell):
        """Correct a single word using SymSpell."""
        suggestions = sym_spell.lookup(word, Verbosity.CLOSEST, max_edit_distance=2)
        return suggestions[0].term if suggestions else word

    def apply_custom_corrections(word, custom_dict):
        """Apply domain-specific corrections using a custom dictionary."""
        return custom_dict.get(word.lower(), word)

    def preserve_case(original_word, corrected_word):
        """
        Preserve the case of the original word in the corrected word.
        """
        if original_word.isupper():
            return corrected_word.upper()
        elif original_word[0].isupper():
            return corrected_word.capitalize()
        else:
            return corrected_word.lower()

    def process_string(text, method, sym_spell=None, custom_dict=None):
        """
        Process a single string for spelling corrections.
        Handles TextBlob, SymSpell, and custom corrections.
        """
        if method in ("sym", "combined") and sym_spell:
            text = segment_words(text, sym_spell)

        if method in ("textblob", "combined"):
            text = str(TextBlob(text).correct())

        corrected_words = []
        for word in text.split():
            original_word = word
            if method in ("sym", "combined") and sym_spell:
                word = advanced_correction(word, sym_spell)

            # Step 3: Apply custom corrections
            if custom_dict:
                word = apply_custom_corrections(word, custom_dict)
            # Preserve original case
            word = preserve_case(original_word, word)
            corrected_words.append(word)

        return " ".join(corrected_words)

    # Initialize SymSpell if needed
    sym_spell = None
    if method in ("sym", "combined"):
        if not sym_spell_params:
            sym_spell_params = {"max_edit_distance": 2, "prefix_length": 7}
        sym_spell = initialize_symspell(**sym_spell_params)

    # Process lists or strings
    if isinstance(content, list):
        if use_threading:
            with ThreadPoolExecutor() as executor:
                corrected_content = list(
                    executor.map(
                        lambda x: process_string(x, method, sym_spell, custom_dict),
                        content,
                    )
                )
            return corrected_content
        else:
            return [
                process_string(item, method, sym_spell, custom_dict) for item in content
            ]
    else:
        return process_string(content, method, sym_spell, custom_dict)


def load_img(fpath):
    """
    Load an image from the specified file path.
    Args:
        fpath (str): The file path to the image.
    Returns:
        PIL.Image: The loaded image.
    Raises:
        FileNotFoundError: If the specified file is not found.
        OSError: If the specified file cannot be opened or is not a valid image file.
    """
    from PIL import Image

    try:
        img = Image.open(fpath)
        return img
    except FileNotFoundError:
        raise FileNotFoundError(f"The file '{fpath}' was not found.")
    except OSError:
        raise OSError(f"Unable to open file '{fpath}' or it is not a valid image file.")


def apply_filter(img, *args, verbose=True):
    # def apply_filter(img, filter_name, filter_value=None):
    """
    Apply the specified filter to the image.
    Args:
        img (PIL.Image): The input image.
        filter_name (str): The name of the filter to apply.
        **kwargs: Additional parameters specific to the filter.
    Returns:
        PIL.Image: The filtered image.
    """
    from PIL import ImageFilter

    def correct_filter_name(filter_name):
        if all(
            [
                "b" in filter_name.lower(),
                "ur" in filter_name.lower(),
                "box" not in filter_name.lower(),
            ]
        ):
            return "BLUR"
        elif "cont" in filter_name.lower():
            return "Contour"
        elif "det" in filter_name.lower():
            return "Detail"
        elif (
            "edg" in filter_name.lower()
            and "mo" not in filter_name.lower()
            and "f" not in filter_name.lower()
        ):
            return "EDGE_ENHANCE"
        elif "edg" in filter_name.lower() and "mo" in filter_name.lower():
            return "EDGE_ENHANCE_MORE"
        elif "emb" in filter_name.lower():
            return "EMBOSS"
        elif "edg" in filter_name.lower() and "f" in filter_name.lower():
            return "FIND_EDGES"
        elif "sh" in filter_name.lower() and "mo" not in filter_name.lower():
            return "SHARPEN"
        elif "sm" in filter_name.lower() and "mo" not in filter_name.lower():
            return "SMOOTH"
        elif "sm" in filter_name.lower() and "mo" in filter_name.lower():
            return "SMOOTH_MORE"
        elif "min" in filter_name.lower():
            return "MIN_FILTER"
        elif "max" in filter_name.lower():
            return "MAX_FILTER"
        elif "mod" in filter_name.lower():
            return "MODE_FILTER"
        elif "mul" in filter_name.lower():
            return "MULTIBAND_FILTER"
        elif "gau" in filter_name.lower():
            return "GAUSSIAN_BLUR"
        elif "box" in filter_name.lower():
            return "BOX_BLUR"
        elif "med" in filter_name.lower():
            return "MEDIAN_FILTER"
        else:
            supported_filters = [
                "BLUR",
                "CONTOUR",
                "DETAIL",
                "EDGE_ENHANCE",
                "EDGE_ENHANCE_MORE",
                "EMBOSS",
                "FIND_EDGES",
                "SHARPEN",
                "SMOOTH",
                "SMOOTH_MORE",
                "MIN_FILTER",
                "MAX_FILTER",
                "MODE_FILTER",
                "MULTIBAND_FILTER",
                "GAUSSIAN_BLUR",
                "BOX_BLUR",
                "MEDIAN_FILTER",
            ]
            raise ValueError(
                f"Unsupported filter: {filter_name}, should be one of: {supported_filters}"
            )

    for arg in args:
        if isinstance(arg, str):
            filter_name = correct_filter_name(arg)
        else:
            filter_value = arg
    if verbose:
        print(f"processing {filter_name}")
    filter_name = filter_name.upper()  # Ensure filter name is uppercase

    # Supported filters
    supported_filters = {
        "BLUR": ImageFilter.BLUR,
        "CONTOUR": ImageFilter.CONTOUR,
        "DETAIL": ImageFilter.DETAIL,
        "EDGE_ENHANCE": ImageFilter.EDGE_ENHANCE,
        "EDGE_ENHANCE_MORE": ImageFilter.EDGE_ENHANCE_MORE,
        "EMBOSS": ImageFilter.EMBOSS,
        "FIND_EDGES": ImageFilter.FIND_EDGES,
        "SHARPEN": ImageFilter.SHARPEN,
        "SMOOTH": ImageFilter.SMOOTH,
        "SMOOTH_MORE": ImageFilter.SMOOTH_MORE,
        "MIN_FILTER": ImageFilter.MinFilter,
        "MAX_FILTER": ImageFilter.MaxFilter,
        "MODE_FILTER": ImageFilter.ModeFilter,
        "MULTIBAND_FILTER": ImageFilter.MultibandFilter,
        "GAUSSIAN_BLUR": ImageFilter.GaussianBlur,
        "BOX_BLUR": ImageFilter.BoxBlur,
        "MEDIAN_FILTER": ImageFilter.MedianFilter,
    }
    # Check if the filter name is supported
    if filter_name not in supported_filters:
        raise ValueError(
            f"Unsupported filter: {filter_name}, should be one of: {[i.lower() for i in supported_filters.keys()]}"
        )

    # Apply the filter
    if filter_name.upper() in [
        "BOX_BLUR",
        "GAUSSIAN_BLUR",
        "MEDIAN_FILTER",
        "MIN_FILTER",
        "MAX_FILTER",
        "MODE_FILTER",
    ]:
        radius = filter_value if filter_value is not None else 2
        return img.filter(supported_filters[filter_name](radius))
    elif filter_name in ["MULTIBAND_FILTER"]:
        bands = filter_value if filter_value is not None else None
        return img.filter(supported_filters[filter_name](bands))
    else:
        if filter_value is not None and verbose:
            print(
                f"{filter_name} doesn't require a value for {filter_value}, but it remains unaffected"
            )
        return img.filter(supported_filters[filter_name])


def detect_angle(image, by="median", template=None):
    """Detect the angle of rotation using various methods."""
    from sklearn.decomposition import PCA
    from skimage import transform, feature, filters, measure
    from skimage.color import rgb2gray
    from scipy.fftpack import fftshift, fft2
    import numpy as np
    import cv2

    # Convert to grayscale
    if np.array(image).shape[-1] > 3:
        image = np.array(image)[:, :, :3]
    gray_image = rgb2gray(image)

    # Detect edges using Canny edge detector
    edges = feature.canny(gray_image, sigma=2)

    # Use Hough transform to detect lines
    lines = transform.probabilistic_hough_line(edges)

    if not lines and any(["me" in by, "pca" in by]):
        print("No lines detected. Adjust the edge detection parameters.")
        return 0
    methods = [
        "mean",
        "median",
        "pca",
        "gradient orientation",
        "template matching",
        "moments",
        "fft",
    ]
    by = strcmp(by, methods)[0]
    # Hough Transform-based angle detection (Median/Mean)
    if "me" in by.lower():
        angles = []
        for line in lines:
            (x0, y0), (x1, y1) = line
            angle = np.arctan2(y1 - y0, x1 - x0) * 180 / np.pi
            if 80 < abs(angle) < 100:
                angles.append(angle)
        if not angles:
            return 0
        if "di" in by:
            median_angle = np.median(angles)
            rotation_angle = (
                90 - median_angle if median_angle > 0 else -90 - median_angle
            )

            return rotation_angle
        else:
            mean_angle = np.mean(angles)
            rotation_angle = 90 - mean_angle if mean_angle > 0 else -90 - mean_angle

            return rotation_angle

    # PCA-based angle detection
    elif "pca" in by.lower():
        y, x = np.nonzero(edges)
        if len(x) == 0:
            return 0
        pca = PCA(n_components=2)
        pca.fit(np.vstack((x, y)).T)
        angle = np.arctan2(pca.components_[0, 1], pca.components_[0, 0]) * 180 / np.pi
        return angle

    # Gradient Orientation-based angle detection
    elif "gra" in by.lower():
        gx, gy = np.gradient(gray_image)
        angles = np.arctan2(gy, gx) * 180 / np.pi
        hist, bin_edges = np.histogram(angles, bins=360, range=(-180, 180))
        return bin_edges[np.argmax(hist)]

    # Template Matching-based angle detection
    elif "temp" in by.lower():
        if template is None:
            # Automatically extract a template from the center of the image
            height, width = gray_image.shape
            center_x, center_y = width // 2, height // 2
            size = (
                min(height, width) // 4
            )  # Size of the template as a fraction of image size
            template = gray_image[
                center_y - size : center_y + size, center_x - size : center_x + size
            ]
        best_angle = None
        best_corr = -1
        for angle in range(0, 180, 1):  # Checking every degree
            rotated_template = transform.rotate(template, angle)
            res = cv2.matchTemplate(gray_image, rotated_template, cv2.TM_CCOEFF)
            _, max_val, _, _ = cv2.minMaxLoc(res)
            if max_val > best_corr:
                best_corr = max_val
                best_angle = angle
        return best_angle

    # Image Moments-based angle detection
    elif "mo" in by.lower():
        moments = measure.moments_central(gray_image)
        angle = (
            0.5
            * np.arctan2(2 * moments[1, 1], moments[0, 2] - moments[2, 0])
            * 180
            / np.pi
        )
        return angle

    # Fourier Transform-based angle detection
    elif "fft" in by.lower():
        f = fft2(gray_image)
        fshift = fftshift(f)
        magnitude_spectrum = np.log(np.abs(fshift) + 1)
        rows, cols = magnitude_spectrum.shape
        r, c = np.unravel_index(np.argmax(magnitude_spectrum), (rows, cols))
        angle = np.arctan2(r - rows // 2, c - cols // 2) * 180 / np.pi
        return angle

    else:
        print(f"Unknown method {by}: supported methods: {methods}")
        return 0


def imgsets(
    img,
    auto: bool = True,
    size=None,
    figsize=None,
    dpi: int = 200,
    show_axis: bool = False,
    plot_: bool = True,
    verbose: bool = False,
    model: str = "isnet-general-use",
    **kwargs,
):
    """
    Apply various enhancements and filters to an image using PIL's ImageEnhance and ImageFilter modules.

    Args:
        img (PIL.Image): The input image.
        sets (dict): A dictionary specifying the enhancements, filters, and their parameters.
        show (bool): Whether to display the enhanced image.
        show_axis (bool): Whether to display axes on the image plot.
        size (tuple): The size of the thumbnail, cover, contain, or fit operation.
        dpi (int): Dots per inch for the displayed image.
        figsize (tuple): The size of the figure for displaying the image.
        auto (bool): Whether to automatically enhance the image based on its characteristics.

    Returns:
        PIL.Image: The enhanced image.

    Supported enhancements and filters:
        - "sharpness": Adjusts the sharpness of the image. Values > 1 increase sharpness, while values < 1 decrease sharpness.
        - "contrast": Adjusts the contrast of the image. Values > 1 increase contrast, while values < 1 decrease contrast.
        - "brightness": Adjusts the brightness of the image. Values > 1 increase brightness, while values < 1 decrease brightness.
        - "color": Adjusts the color saturation of the image. Values > 1 increase saturation, while values < 1 decrease saturation.
        - "rotate": Rotates the image by the specified angle.
        - "crop" or "cut": Crops the image. The value should be a tuple specifying the crop box as (left, upper, right, lower).
        - "size": Resizes the image to the specified dimensions.
        - "thumbnail": Resizes the image to fit within the given size while preserving aspect ratio.
        - "cover": Resizes and crops the image to fill the specified size.
        - "contain": Resizes the image to fit within the specified size, adding borders if necessary.
        - "fit": Resizes and pads the image to fit within the specified size.
        - "filter": Applies various filters to the image (e.g., BLUR, CONTOUR, EDGE_ENHANCE).

    Note:
        The "color" and "enhance" enhancements are not implemented in this function.
    """

    import matplotlib.pyplot as plt
    from PIL import ImageEnhance, ImageOps, Image

    supported_filters = [
        "BLUR",
        "CONTOUR",
        "DETAIL",
        "EDGE_ENHANCE",
        "EDGE_ENHANCE_MORE",
        "EMBOSS",
        "FIND_EDGES",
        "SHARPEN",
        "SMOOTH",
        "SMOOTH_MORE",
        "MIN_FILTER",
        "MAX_FILTER",
        "MODE_FILTER",
        "MULTIBAND_FILTER",
        "GAUSSIAN_BLUR",
        "BOX_BLUR",
        "MEDIAN_FILTER",
    ]
    # *Rembg is a tool to remove images background.
    # https://github.com/danielgatis/rembg
    rem_models = {
        "u2net": "general use cases.",
        "u2netp": "A lightweight version of u2net model.",
        "u2net_human_seg": "human segmentation.",
        "u2net_cloth_seg": "Cloths Parsing from human portrait. Here clothes are parsed into 3 category: Upper body, Lower body and Full body.",
        "silueta": "Same as u2net but the size is reduced to 43Mb.",
        "isnet-general-use": "A new pre-trained model for general use cases.",
        "isnet-anime": "A high-accuracy segmentation for anime character.",
        "sam": "any use cases.",
        "birefnet-general": "general use cases.",
        "birefnet-general-lite": "A light pre-trained model for general use cases.",
        "birefnet-portrait": "human portraits.",
        "birefnet-dis": "dichotomous image segmentation (DIS).",
        "birefnet-hrsod": "high-resolution salient object detection (HRSOD).",
        "birefnet-cod": "concealed object detection (COD).",
        "birefnet-massive": "A pre-trained model with massive dataset.",
    }
    models_support_rem = list(rem_models.keys())
    str_usage = """
    imgsets(dir_img, auto=1, color=1.5, plot_=0)
    imgsets(dir_img, color=2)
    imgsets(dir_img, pad=(300, 300), bgcolor=(73, 162, 127), plot_=0)
    imgsets(dir_img, contrast=0, color=1.2, plot_=0)
    imgsets(get_clip(), flip="tb")# flip top and bottom
    imgsets(get_clip(), contrast=1, rm=[100, 5, 2]) #'foreground_threshold', 'background_threshold' and 'erode_structure_size'
    imgsets(dir_img, rm="birefnet-portrait") # with using custom model
    """
    if run_once_within():
        print(str_usage)

    def gamma_correction(image, gamma=1.0, v_max=255):
        # adjust gama value
        inv_gamma = 1.0 / gamma
        lut = [
            int((i / float(v_max)) ** inv_gamma * int(v_max)) for i in range(int(v_max))
        ]
        return lut  # image.point(lut)

    def auto_enhance(img):
        """
        Automatically enhances the image based on its characteristics, including brightness,
        contrast, color range, sharpness, and gamma correction.

        Args:
            img (PIL.Image): The input image.

        Returns:
            dict: A dictionary containing the optimal enhancement values applied.
            PIL.Image: The enhanced image.
        """
        from PIL import Image, ImageEnhance, ImageOps, ImageFilter
        import numpy as np

        # Determine the bit depth based on the image mode
        try:
            if img.mode in ["1", "L", "P", "RGB", "YCbCr", "LAB", "HSV"]:
                bit_depth = 8
            elif img.mode in ["RGBA", "CMYK"]:
                bit_depth = 8
            elif img.mode in ["I", "F"]:
                bit_depth = 16
            else:
                raise ValueError("Unsupported image mode")
        except:
            bit_depth = 8

        # Initialize enhancement factors
        enhancements = {
            "brightness": 1.0,
            "contrast": 0,  # autocontrasted
            "color": 1.35,
            "sharpness": 1.0,
            "gamma": 1.0,
        }

        # Calculate brightness and contrast for each channel
        num_channels = len(img.getbands())
        brightness_factors = []
        contrast_factors = []
        for channel in range(num_channels):
            channel_histogram = img.split()[channel].histogram()
            total_pixels = sum(channel_histogram)
            brightness = (
                sum(i * w for i, w in enumerate(channel_histogram)) / total_pixels
            )
            channel_min, channel_max = img.split()[channel].getextrema()
            contrast = channel_max - channel_min
            # Adjust calculations based on bit depth
            normalization_factor = 2**bit_depth - 1
            brightness_factor = (
                1.0 + (brightness - normalization_factor / 2) / normalization_factor
            )
            contrast_factor = (
                1.0 + (contrast - normalization_factor / 2) / normalization_factor
            )
            brightness_factors.append(brightness_factor)
            contrast_factors.append(contrast_factor)

        # Calculate average brightness and contrast factors across channels
        enhancements["brightness"] = sum(brightness_factors) / num_channels
        # Adjust brightness and contrast
        img = ImageEnhance.Brightness(img).enhance(enhancements["brightness"])

        # # Automatic color enhancement (saturation)
        # if img.mode == "RGB":
        #     color_enhancer = ImageEnhance.Color(img)
        #     color_histogram = np.array(img.histogram()).reshape(3, -1)
        #     avg_saturation = np.mean([np.std(channel) for channel in color_histogram]) / normalization_factor
        #     print(avg_saturation)
        #     enhancements["color"] = min(0, max(0.5, 1.0 + avg_saturation))  # Clamp to a reasonable range
        #     # img = color_enhancer.enhance(enhancements["color"])

        # Adjust sharpness
        sharpness_enhancer = ImageEnhance.Sharpness(img)
        # Use edge detection to estimate sharpness need
        edges = img.filter(ImageFilter.FIND_EDGES).convert("L")
        avg_edge_intensity = np.mean(np.array(edges))
        enhancements["sharpness"] = min(
            2.0, max(0.5, 1.0 + avg_edge_intensity / normalization_factor)
        )
        # img = sharpness_enhancer.enhance(enhancements["sharpness"])

        # # Apply gamma correction
        # def gamma_correction(image, gamma):
        #     inv_gamma = 1.0 / gamma
        #     lut = [min(255, max(0, int((i / 255.0) ** inv_gamma * 255))) for i in range(256)]
        #     return image.point(lut)

        # avg_brightness = np.mean(np.array(img.convert("L"))) / 255
        # enhancements["gamma"] = min(2.0, max(0.5, 1.0 if avg_brightness > 0.5 else 1.2 - avg_brightness))
        # img = gamma_correction(img, enhancements["gamma"])

        # Return the enhancements and the enhanced image
        return enhancements

    # Load image if input is a file path
    if isinstance(img, str):
        img = load_img(img)
    img_update = img.copy()

    if auto:
        kwargs = {**auto_enhance(img_update), **kwargs}
    params = [
        "sharp",
        "color",
        "contrast",
        "bright",
        "crop",
        "rotate",
        "size",
        "resize",
        "thumbnail",
        "cover",
        "contain",
        "filter",
        "fit",
        "pad",
        "rem",
        "rm",
        "back",
        "bg_color",
        "cut",
        "gamma",
        "flip",
        "booster",
    ]
    for k, value in kwargs.items():
        k = strcmp(k, params)[0]  # correct the param name
        if "shar" in k.lower():
            enhancer = ImageEnhance.Sharpness(img_update)
            img_update = enhancer.enhance(value)
        elif all(
            ["col" in k.lower(), "bg" not in k.lower(), "background" not in k.lower()]
        ):
            # *color
            enhancer = ImageEnhance.Color(img_update)
            img_update = enhancer.enhance(value)
        elif "contr" in k.lower():
            if value and isinstance(value, (float, int)):
                enhancer = ImageEnhance.Contrast(img_update)
                img_update = enhancer.enhance(value)
            else:
                try:
                    img_update = ImageOps.autocontrast(img_update)
                    print("autocontrasted")
                except Exception as e:
                    print(f"Failed 'auto-contrasted':{e}")
        elif "bri" in k.lower():
            enhancer = ImageEnhance.Brightness(img_update)
            img_update = enhancer.enhance(value)
        elif "cro" in k.lower() or "cut" in k.lower():
            img_update = img_update.crop(value)
        elif "rota" in k.lower():
            if isinstance(value, str):
                value = detect_angle(img_update, by=value)
                print(f"rotated by {value}°")
            img_update = img_update.rotate(value)
        elif "flip" in k.lower():
            if "l" in value and "r" in value:
                # left/right
                img_update = img_update.transpose(Image.FLIP_LEFT_RIGHT)
            elif any(["u" in value and "d" in value, "t" in value and "b" in value]):
                # up/down or top/bottom
                img_update = img_update.transpose(Image.FLIP_TOP_BOTTOM)
        elif "si" in k.lower():
            if isinstance(value, tuple):
                value = list(value)
            value = [int(i) for i in value]
            img_update = img_update.resize(value)
        elif "thum" in k.lower():
            img_update.thumbnail(value)
        elif "cover" in k.lower():
            img_update = ImageOps.cover(img_update, size=value)
        elif "contain" in k.lower():
            img_update = ImageOps.contain(img_update, size=value)
        elif "fi" in k.lower() and "t" in k.lower():  # filter
            if isinstance(value, dict):
                if verbose:
                    print(f"supported filter: {supported_filters}")
                for filter_name, filter_value in value.items():
                    img_update = apply_filter(
                        img_update, filter_name, filter_value, verbose=verbose
                    )
            else:
                img_update = ImageOps.fit(img_update, size=value)
        elif "pad" in k.lower():
            # *ImageOps.pad ensures that the resized image has the exact size specified by the size parameter while maintaining the aspect ratio.
            # size: A tuple specifying the target size (width, height).
            img_update = ImageOps.pad(img_update, size=value)
        elif "rem" in k.lower() or "rm" in k.lower() or "back" in k.lower():
            from rembg import remove, new_session

            if verbose:
                preview(rem_models)

            print(f"supported modles: {models_support_rem}")
            model = strcmp(model, models_support_rem)[0]
            session = new_session(model)
            if isinstance(value, bool):
                print(f"using model:{model}")
                img_update = remove(img_update, session=session)
            elif value and isinstance(value, (int, float, list)):
                if verbose:
                    print("https://github.com/danielgatis/rembg/blob/main/USAGE.md")
                    print(
                        f"rm=True # using default setting;\nrm=(240,10,10)\n'foreground_threshold'(240) and 'background_threshold' (10) values used to determine foreground and background pixels. \nThe 'erode_structure_size'(10) parameter specifies the size of the erosion structure to be applied to the mask."
                    )
                if isinstance(value, int):
                    value = [value]
                if len(value) < 2:
                    img_update = remove(
                        img_update,
                        alpha_matting=True,
                        alpha_matting_background_threshold=value,
                        session=session,
                    )
                elif 2 <= len(value) < 3:
                    img_update = remove(
                        img_update,
                        alpha_matting=True,
                        alpha_matting_background_threshold=value[0],
                        alpha_matting_foreground_threshold=value[1],
                        session=session,
                    )
                elif 3 <= len(value) < 4:
                    img_update = remove(
                        img_update,
                        alpha_matting=True,
                        alpha_matting_background_threshold=value[0],
                        alpha_matting_foreground_threshold=value[1],
                        alpha_matting_erode_size=value[2],
                        session=session,
                    )
            elif isinstance(value, tuple):  # replace the background color
                if len(value) == 3:
                    value += (255,)
                img_update = remove(img_update, bgcolor=value, session=session)
            elif isinstance(value, str):
                # use custom model
                print(f"using model:{strcmp(value, models_support_rem)[0]}")
                img_update = remove(
                    img_update,
                    session=new_session(strcmp(value, models_support_rem)[0]),
                )
        elif "bg" in k.lower() and "color" in k.lower():
            from rembg import remove

            if isinstance(value, list):
                value = tuple(value)
            if isinstance(value, tuple):  # replace the background color
                if len(value) == 3:
                    value += (255,)
                img_update = remove(img_update, bgcolor=value)
        elif "boost" in k.lower():
            import torch
            from realesrgan import RealESRGANer

            if verbose:
                print("Applying Real-ESRGAN for image reconstruction...")
            if isinstance(value, bool):
                scale = 4
            elif isinstance(value, (float, int)):
                scale = value
            else:
                scale = 4

            # try:
            device = "cuda" if torch.cuda.is_available() else "cpu"
            dir_curr_script = os.path.dirname(os.path.abspath(__file__))
            model_path = dir_curr_script + "/data/RealESRGAN_x4plus.pth"
            model_RealESRGAN = RealESRGANer(
                device=device,
                scale=scale,
                model_path=model_path,
                model="RealESRGAN_x4plus",
            )
            # https://github.com/xinntao/Real-ESRGAN?tab=readme-ov-file#python-script

            img_update = model_RealESRGAN.enhance(np.array(img_update))[0]
            # except Exception as e:
            #     print(f"Failed to apply Real-ESRGAN: {e}")

        # elif "ga" in k.lower() and "m" in k.lower():
        #     img_update = gamma_correction(img_update, gamma=value)
    # Display the image if requested
    if plot_:
        if figsize is None:
            plt.figure(dpi=dpi)
        else:
            plt.figure(figsize=figsize, dpi=dpi)
        plt.imshow(img_update)
        if show_axis:
            plt.axis("on")  # Turn on axis
            plt.minorticks_on()
            plt.grid(
                which="both", linestyle="--", linewidth=0.5, color="gray", alpha=0.7
            )

        else:
            plt.axis("off")  # Turn off axis
    return img_update


def thumbnail(dir_img_list, figsize=(10, 10), dpi=100, dir_save=None, kind=".png"):
    """
    Display a thumbnail figure of all images in the specified directory.
    Args:
        dir_img_list (list): List of the Directory containing the images.
    """
    import matplotlib.pyplot as plt
    from PIL import Image

    num_images = len(dir_img_list)
    if not kind.startswith("."):
        kind = "." + kind

    if num_images == 0:
        print("No images found to display.")
        return
    grid_size = int(num_images**0.5) + 1  # Determine grid size
    fig, axs = plt.subplots(grid_size, grid_size, figsize=figsize, dpi=dpi)
    for ax, image_file in zip(axs.flatten(), dir_img_list):
        try:
            img = Image.open(image_file)
            ax.imshow(img)
            ax.axis("off")
        except:
            continue
    # for ax in axs.flatten():
    #     ax.axis('off')
    [ax.axis("off") for ax in axs.flatten()]
    plt.tight_layout()
    if dir_save is None:
        plt.show()
    else:
        if basename(dir_save):
            fname = basename(dir_save) + kind
        else:
            fname = "_thumbnail_" + basename(dirname(dir_save)[:-1]) + ".png"
        if dirname(dir_img_list[0]) == dirname(dir_save):
            figsave(dirname(dir_save[:-1]), fname)
        else:
            figsave(dirname(dir_save), fname)


# usage:
# fpath = "/Users/macjianfeng/Dropbox/github/python/py2ls/tests/xample_netfinder/images/"
# thumbnail(listdir(fpath,'png').fpath.to_list(),dir_save=dirname(fpath))


# search and fine the director of the libary, which installed at local
def dir_lib(lib_oi):
    """
    # example usage:
    # dir_lib("seaborn")
    """
    import site

    # Get the site-packages directory
    f = listdir(site.getsitepackages()[0], "folder")

    # Find Seaborn directory within site-packages
    dir_list = []
    for directory in f.fpath:
        if lib_oi in directory.lower():
            dir_list.append(directory)

    if dir_list != []:
        print(f"{lib_oi} directory:", dir_list)
    else:
        print(f"Cannot find the {lib_oi} in site-packages directory.")
    return dir_list


class FileInfo:
    def __init__(
        self,
        size,
        creation_time,
        ctime,
        mod_time,
        mtime,
        parent_dir,
        fname,
        kind,
        extra_info=None,
    ):
        self.size = size
        self.creation_time = creation_time
        self.ctime = ctime
        self.mod_time = mod_time
        self.mtime = mtime
        self.parent_dir = parent_dir
        self.fname = fname
        self.kind = kind
        if extra_info:
            for key, value in extra_info.items():
                setattr(self, key, value)
        print("to show the res: 'finfo(fpath).show()'")

    def __repr__(self):
        return (
            f"FileInfo(size={self.size} MB, creation_time='{self.creation_time}', "
            f"ctime='{self.ctime}', mod_time='{self.mod_time}', mtime='{self.mtime}', "
            f"parent_dir='{self.parent_dir}', fname='{self.fname}', kind='{self.kind}')"
        )

    def __str__(self):
        return (
            f"FileInfo:\n"
            f"  Size: {self.size} MB\n"
            f"  Creation Time: {self.creation_time}\n"
            f"  CTime: {self.ctime}\n"
            f"  Modification Time: {self.mod_time}\n"
            f"  MTime: {self.mtime}\n"
            f"  Parent Directory: {self.parent_dir}\n"
            f"  File Name: {self.fname}\n"
            f"  Kind: {self.kind}"
        )

    def show(self):
        # Convert the object to a dictionary
        return {
            "size": self.size,
            "creation_time": self.creation_time,
            "ctime": self.ctime,
            "mod_time": self.mod_time,
            "mtime": self.mtime,
            "parent_dir": self.parent_dir,
            "fname": self.fname,
            "kind": self.kind,
            **{
                key: getattr(self, key)
                for key in vars(self)
                if key
                not in [
                    "size",
                    "creation_time",
                    "ctime",
                    "mod_time",
                    "mtime",
                    "parent_dir",
                    "fname",
                    "kind",
                ]
            },
        }


def finfo(fpath):
    import time

    fname, fmt = os.path.splitext(fpath)
    dir_par = os.path.dirname(fpath) + "/"
    data = {
        "size": round(os.path.getsize(fpath) / 1024 / 1024, 3),
        "creation_time": time.ctime(os.path.getctime(fpath)),
        "ctime": time.ctime(os.path.getctime(fpath)),
        "mod_time": time.ctime(os.path.getmtime(fpath)),
        "mtime": time.ctime(os.path.getmtime(fpath)),
        "parent_dir": dir_par,
        "fname": fname.replace(dir_par, ""),
        "kind": fmt,
    }
    extra_info = {}
    if data["kind"] == ".pdf":
        from pdf2image import pdfinfo_from_path

        extra_info = pdfinfo_from_path(fpath)

    return FileInfo(
        size=data["size"],
        creation_time=data["creation_time"],
        ctime=data["ctime"],
        mod_time=data["mod_time"],
        mtime=data["mtime"],
        parent_dir=data["parent_dir"],
        fname=data["fname"],
        kind=data["kind"],
        extra_info=extra_info,
    )


# ! format excel file


def hex2argb(color):
    """
    Convert a color name or hex code to aARGB format required by openpyxl.

    :param color: A color in the format: 'blue', '#RRGGBB', 'RRGGBB', 'aARRGGBB'
    :return: A hex color code in the format aARRGGBB.

    Example:
        print(hex2argb("blue"))      # Output: FF0000FF
        print(hex2argb("FFFF00"))    # Output: FFFFFF00
        print(hex2argb("#DF4245"))   # Output: FFDf4245
        print(hex2argb("FF00FF00"))  # Output: FF00FF00 (already in aARGB format)
    """
    import matplotlib.colors as mcolors
    import re
    color = color.lower().replace(" ", "") # 'light blue'
    # Convert color name (e.g., "blue") to hex
    if color.lower() in mcolors.CSS4_COLORS:
        color = mcolors.CSS4_COLORS[color.lower()].lstrip("#")
    color = color.lstrip("#").upper()# Remove '#' if present
    
    # Validate hex format
    if not re.fullmatch(r"[A-F0-9]{6,8}", color):
        raise ValueError(f"格式错误❌: {color}, 应该使用 RRGGBB, #RRGGBB, or aARRGGBB format.")

    # If already in aARRGGBB format (8 chars), return as is
    if len(color) == 8:
        return color
    
    # If in RRGGBB format, add FF (full opacity) as alpha
    return f"FF{color}"

def extract_kwargs(func):
    import inspect

    # Get the signature of the function
    signature = inspect.signature(func)
    # Extract parameters that are kwargs (parameters with default values or **kwargs)
    kwargs = {
        param.name: param.default
        for param in signature.parameters.values()
        if param.default is not inspect.Parameter.empty
    }

    return kwargs
def format_excel(
    df: pd.DataFrame=None,
    filename:str=None,
    sheet_name:Union[str, int]=0,
    insert_img:dict=None,# {"A1":img_path}
    usage:bool=False,
    text_color:Union[dict,bool]=False, # dict: set the text color
    bg_color:Union[dict,bool]=False, # dict: set the back_ground color
    cell:Union[dict, list]=None,  # dict: or list for multiple locs setting:
    width:Union[bool, dict]=None,  # dict
    width_factor:int=2,# calculated with plus this factor
    height:Union[bool, dict]=None,  # dict e.g., {2: 50, 3: 25},  keys are columns
    height_max:int=25,
    merge:tuple=None,  # tuple e.g.,  (slice(0, 1), slice(1, 3)),
    shade:Union[dict, list]=None,  # dict
    comment:Union[dict, list]=None,  # dict e.g., {(2, 4): "This is a comment"},
    comment_always_visible:bool=True,# always display comment
    link:Union[dict, list]=None,  # dict e.g., {(2, 2): "https://example.com"},
    protect:dict=None,  # dict
    number_format:dict=None,  # dict: e.g., {1:"0.00", 2:"#,##0",3:"0%",4:"$#,##0.00"}
    data_validation=None,  # dict
    apply_filter:bool=True, # add filter 
    freeze :str= False,#"A2",
    conditional_format:dict=None,  # dict
    verbose:bool=False,
    **kwargs,
):
    """
    Parameters:
        df : pandas.DataFrame, optional
            DataFrame to be written to the Excel file.
        filename : str, optional
            Path to the output Excel file.
        sheet_name : str or int, default 0
            Name or index of the sheet where data will be written.
        insert_img : dict, optional
            Dictionary specifying image insert locations, e.g., {"A1": "path/to/image.png"}.
        usage : bool, default False
            If True, display usage examples.
        cell : dict or list, optional
            Specifies cell formatting options.
        width : dict, optional
            Dictionary specifying column widths, e.g., {1: 20, 2: 30}.
        width_factor : int, default 2
            Additional factor to adjust column width dynamically.
        height : dict, optional
            Dictionary specifying row heights, e.g., {2: 50, 3: 25}.
        height_max : int, default 25
            Maximum row height allowed.
        merge : tuple, optional
            Specifies cell merging, e.g., (slice(0, 1), slice(1, 3)).
        shade : dict, optional
            Dictionary defining cell shading/styling.
        comment : dict, optional
            Dictionary adding comments, e.g., {(2, 4): "This is a comment"}.
        comment_always_visible : bool, default True
            Whether comments should always be visible.
        link : dict, optional
            Dictionary specifying hyperlinks, e.g., {(2, 2): "https://example.com"}.
        protect : dict, optional
            Dictionary defining cell protection settings.
        number_format : dict, optional
            Dictionary specifying number formats, e.g., {1: "0.00", 2: "#,##0"}.
        data_validation : dict, optional
            Dictionary setting data validation rules.
        apply_filter : bool, default True
            Whether to apply filters to the header row.
        freeze : str, optional
            Cell reference (e.g., "A2") to freeze rows/columns.
        conditional_format : dict, optional
            Dictionary defining conditional formatting rules.
        verbose : bool, default False
            Whether to print detailed execution logs.
        **kwargs : dict
            Additional parameters for advanced customization. 
    """

    usage_str="""
        Formats an Excel file with various styling options.
        Usage:
        fsave(
                dir_save,
                fload(dir_save, output="bit", sheet_name=sheet_name),
                sheet_name=sheet_name,
                if_sheet_exists="overlay",
                mode="a",
                width_factor=0,
                height={1: 50},
                cell=[
                    {
                        (slice(0, 1), slice(0, df_exists.shape[1])): {
                            "fill": {
                                "start_color": "61AFEF",  # Starting color
                                "end_color": "61AFEF",  # Ending color (useful for gradients)
                                "fill_type": "solid",  # Fill type (solid, gradient, etc.)
                            },
                            "font": {
                                "name": "Arial",  # Font name
                                "size": 11,  # Font size
                                "bold": True,  # Bold text
                                "italic": False,  # Italic text
                                # "underline": "single",  # Underline (single, double)
                                "color": "#000000",  # Font color
                            },
                            "alignment": {
                                "horizontal": "center",  # Horizontal alignment (left, center, right)
                                "vertical": "center",  # Vertical alignment (top, center, bottom)
                                "wrap_text": True,  # Wrap text in the cell
                                "shrink_to_fit": True,  # Shrink text to fit within cell
                                "text_rotation": 0,  # Text rotation angle
                            },
                        }
                    },
                    {
                        (
                            slice(0, df_exists.shape[0]),
                            slice(0, df_exists.shape[1]),
                        ): {
                            "alignment": {
                                "horizontal": "center",  # Horizontal alignment (left, center, right)
                                "vertical": "center",  # Vertical alignment (top, center, bottom)
                                "wrap_text": True,  # Wrap text in the cell
                                "shrink_to_fit": True,  # Shrink text to fit within cell
                                "text_rotation": 0,  # Text rotation angle
                            },
                        }
                    },
                    {
                        (slice(0, df_exists.shape[0]), slice(2, 3)): {
                            "alignment": {
                                "horizontal": "left",  # Horizontal alignment (left, center, right)
                            },
                        }
                    },
                    {
                        (slice(0, df_exists.shape[0]), slice(7, 8)): {
                            "alignment": {
                                "horizontal": "left",  # Horizontal alignment (left, center, right)
                            },
                        }
                    },
                ],
                password=False,  # depass("ogB3B7y3xR9iuH4QIQbyy6VXG14I0A8DlsTxyiGqg1U="),
            )
            """
    if verbose:
        print(usage_str)
    import pandas as pd
    from datetime import datetime
    import openpyxl
    from openpyxl import load_workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    from openpyxl.worksheet.datavalidation import DataValidation
    from openpyxl.comments import Comment
    from openpyxl.formatting.rule import ColorScaleRule, DataBarRule, IconSetRule,IconSet
    from openpyxl.utils import get_column_letter
    def convert_indices_to_range(row_slice, col_slice):
        """Convert numerical row and column slices to Excel-style range strings."""
        start_row = row_slice.start + 1
        end_row = row_slice.stop if row_slice.stop is not None else None
        start_col = col_slice.start + 1
        end_col = col_slice.stop if col_slice.stop is not None else None

        start_col_letter = get_column_letter(start_col)
        end_col_letter = get_column_letter(end_col) if end_col else None
        return (
            f"{start_col_letter}{start_row}:{end_col_letter}{end_row}"
            if end_col_letter
            else f"{start_col_letter}{start_row}"
        )
    def apply_color_to_worksheet(ws=None, sheet_name=None, conditions=None, cell_idx=None,where="text"):
        """
        Apply text color formatting to a specific cell range in an openpyxl workbook based on conditions.

        Parameters:
        ws : worrksheet
            The openpyxl workbook object to style.
        sheet_name : str
            The name of the sheet to style.
        conditions : dict
            Dictionary defining conditions for text or background coloring.
                Example:
                {
                    ">10": "#FF0000",           # Red if value is greater than 10
                    "contains:Error": "#FFFF00", # Yellow if text contains 'Error'
                    "startswith:Warn": "#FFA500" # Orange if text starts with 'Warn'
                }
        cell_idx : tuple, optional
            A tuple of slices defining the selected row and column range (only for DataFrame).
        where : str, default="text"
            "text" -> Apply color to text, "bg" -> Apply color to background.

        Returns:
        openpyxl.workbook.workbook.Workbook
            The workbook with applied formatting.
        """ 
        def evaluate_condition(value, condition):
            """Evaluate the condition dynamically."""
            if not isinstance(conditions, dict):
                raise ValueError(f"condition必须是dict格式:e.g., {'x>=20':'#DD0531', 'startswith:Available':'#DD0531'}")
            try:
                if "x" in condition and re.search(r"[<>=!]=*", condition):
                    expr = condition.replace("x", str(value))
                    return eval(expr)
                elif condition.startswith("startswith:") or condition.startswith("startwith:"):
                    return value.startswith(condition.split(":", 1)[1])
                elif condition.startswith("endswith:") or condition.startswith("endwith:"):
                    return value.endswith(condition.split(":", 1)[1])
                elif condition.startswith("contains:") or condition.startswith("contain:") or condition.startswith("include:"):
                    return condition.split(":", 1)[1] in value
                elif condition.startswith("matches:") or condition.startswith("match:"):
                    return re.search(condition.split(":", 1)[1], value) is not None 
                else:
                    expr = condition
                return False
            except Exception as e: 
                return False

        def apply_condition_to_cell_text_color(cell, value):
            """Apply color to a cell if it matches any condition."""
            for condition, color in conditions.items():
                if evaluate_condition(value, condition):
                    # Apply color to font
                    cell.font = openpyxl.styles.Font(
                        color=openpyxl.styles.Color(rgb=hex2argb(color))
                    )
                    return
        def apply_condition_to_cell_bg_color(cell, value):
            """Apply background color to a cell if it matches any condition."""
            for condition, color in conditions.items():
                if evaluate_condition(value, condition):
                    if not isinstance(color,list):
                        color=[color]
                    if len(color)==1:
                        cell.fill = PatternFill(
                            start_color=hex2argb(color[0]),
                            end_color=hex2argb(color[0]),
                            fill_type="solid"
                        )
                    elif len(color)==2:
                        cell.fill = PatternFill(
                            start_color=hex2argb(color[0]),
                            end_color=hex2argb(color[1]),
                            fill_type="solid"
                        )
                    return
        if isinstance(cell_idx, tuple):
            # If cell_idx is provided, select a range based on the slice
            row_slice, col_slice = cell_idx
            rows = list(
                ws.iter_rows(
                    min_row=row_slice.start + 1,
                    max_row=row_slice.stop,
                    min_col=col_slice.start + 1,
                    max_col=col_slice.stop,
                )
            )
            for row in rows:
                for cell in row:
                    if where=="text":
                        apply_condition_to_cell_text_color(cell, cell.value)
                    elif where=="bg":
                        apply_condition_to_cell_bg_color(cell, cell.value)
        else:
            # If no cell_idx is provided, apply to all cells
            for row in ws.iter_rows():
                for cell in row:
                    if where=="text":
                        apply_condition_to_cell_text_color(cell, cell.value)
                    elif where=="bg":
                        apply_condition_to_cell_bg_color(cell,cell.value)
        return ws
    
    def apply_format(ws, cell, cell_range):
        """Apply cell formatting to a specified range."""
        cell_font, cell_fill, cell_alignment, border = None, None, None, None
        kws_cell = ["font", "fill", "alignment", "border"]
        for K, _ in cell.items():
            if strcmp(K, kws_cell)[0] == "font":
                #! font
                font_color = "000000"
                font_name = "Arial"
                font_underline = "none"
                font_size = 11
                font_bold = False
                font_strike = False
                font_italic = False
                kws_font = [
                    "name",
                    "size",
                    "bold",
                    "underline",
                    "color",
                    "strike",
                    "italic",
                ]
                for k_, v_ in cell.get(K, {}).items():
                    if strcmp(k_, kws_font)[0] == "name":
                        font_name = v_
                    elif strcmp(k_, kws_font)[0] == "size":
                        font_size = v_
                    elif strcmp(k_, kws_font)[0] == "bold":
                        font_bold = v_
                    elif strcmp(k_, kws_font)[0] == "underline":
                        font_underline = strcmp(v_, ["none", "single", "double"])[0]
                    elif strcmp(k_, kws_font)[0] == "color":
                        font_color = hex2argb(v_)
                    elif strcmp(k_, kws_font)[0] == "strike":
                        font_strike = v_
                    elif strcmp(k_, kws_font)[0] == "italic":
                        font_italic = v_

                cell_font = Font(
                    name=font_name,
                    size=font_size,
                    bold=font_bold,
                    italic=font_italic,
                    underline=font_underline,
                    strike=font_strike,
                    color=font_color,
                )

            if strcmp(K, kws_cell)[0] == "fill":
                #! fill
                kws_fill = ["start_color", "end_color", "fill_type", "color"]
                kws_fill_type = [
                    "darkVertical",
                    "lightDown",
                    "lightGrid",
                    "solid",
                    "darkDown",
                    "lightGray",
                    "lightUp",
                    "gray0625",
                    "lightVertical",
                    "lightHorizontal",
                    "darkHorizontal",
                    "gray125",
                    "darkUp",
                    "mediumGray",
                    "darkTrellis",
                    "darkGray",
                    "lightTrellis",
                    "darkGrid",
                ]
                start_color, end_color, fill_type = (
                    "FFFFFF",
                    "FFFFFF",
                    "solid",
                )  # default
                for k, v in cell.get(K, {}).items():
                    if strcmp(k, kws_fill)[0] == "color":
                        start_color, end_color = hex2argb(v), hex2argb(v)
                        break
                for k, v in cell.get(K, {}).items():
                    if strcmp(k, kws_fill)[0] == "start_color":
                        start_color = hex2argb(v)
                    elif strcmp(k, kws_fill)[0] == "end_color":
                        end_color = hex2argb(v)
                    elif strcmp(k, kws_fill)[0] == "fill_type":
                        fill_type = strcmp(v, kws_fill_type)[0]
                cell_fill = PatternFill(
                    start_color=start_color,
                    end_color=end_color,
                    fill_type=fill_type,
                )

            if strcmp(K, kws_cell)[0] == "alignment":
                #! alignment
                # default
                align_horizontal = "general"
                align_vertical = "center"
                align_rot = 0
                align_wrap = False
                align_shrink = False
                align_indent = 0
                kws_align = [
                    "horizontal",
                    "ha",
                    "vertical",
                    "va",
                    "text_rotation",
                    "rotat",
                    "rot",
                    "wrap_text",
                    "wrap",
                    "shrink_to_fit",
                    "shrink",
                    "indent",
                ]
                for k, v in cell.get(K, {}).items():
                    if strcmp(k, kws_align)[0] in ["horizontal", "ha"]:
                        align_horizontal = strcmp(
                            v, ["general", "left", "right", "center"]
                        )[0]
                    elif strcmp(k, kws_align)[0] in ["vertical", "va"]:
                        align_vertical = strcmp(v, ["top", "center", "bottom"])[0]
                    elif strcmp(k, kws_align)[0] in ["text_rotation", "rotat", "rot"]:
                        align_rot = v
                    elif strcmp(k, kws_align)[0] in ["wrap_text", "wrap"]:
                        align_wrap = v
                    elif strcmp(k, kws_align)[0] in [
                        "shrink_to_fit",
                        "shrink",
                        "wrap_text",
                        "wrap",
                    ]:
                        align_shrink = v
                    elif strcmp(k, kws_align)[0] in ["indent"]:
                        align_indent = v
                cell_alignment = Alignment(
                    horizontal=align_horizontal,
                    vertical=align_vertical,
                    text_rotation=align_rot,
                    wrap_text=align_wrap,
                    shrink_to_fit=align_shrink,
                    indent=align_indent,
                )

            if strcmp(K, kws_cell)[0] == "border":
                #! border
                kws_border = [
                    "color_left",
                    "color_l",
                    "color_right",
                    "color_r",
                    "color_top",
                    "color_t",
                    "color_bottom",
                    "color_b",
                    "color_diagonal",
                    "color_d",
                    "color_outline",
                    "color_o",
                    "color_vertical",
                    "color_v",
                    "color_horizontal",
                    "color_h",
                    "color",
                    "style_left",
                    "style_l",
                    "style_right",
                    "style_r",
                    "style_top",
                    "style_t",
                    "style_bottom",
                    "style_b",
                    "style_diagonal",
                    "style_d",
                    "style_outline",
                    "style_o",
                    "style_vertical",
                    "style_v",
                    "style_horizontal",
                    "style_h",
                    "style",
                ]
                # * border color
                border_color_l, border_color_r, border_color_t, border_color_b = (
                    "FF000000",
                    "FF000000",
                    "FF000000",
                    "FF000000",
                )
                border_color_d, border_color_o, border_color_v, border_color_h = (
                    "FF000000",
                    "FF000000",
                    "FF000000",
                    "FF000000",
                )
                # get colors config
                for k, v in cell.get(K, {}).items():
                    if strcmp(k, kws_border)[0] in ["color"]:
                        border_color_all = hex2argb(v)
                        # 如果设置了color,表示其它的所有的都设置成为一样的
                        # 然后再才开始自己定义其它的color
                        (
                            border_color_l,
                            border_color_r,
                            border_color_t,
                            border_color_b,
                        ) = (
                            border_color_all,
                            border_color_all,
                            border_color_all,
                            border_color_all,
                        )
                        (
                            border_color_d,
                            border_color_o,
                            border_color_v,
                            border_color_h,
                        ) = (
                            border_color_all,
                            border_color_all,
                            border_color_all,
                            border_color_all,
                        )
                    elif strcmp(k, kws_border)[0] in ["color_left", "color_l"]:
                        border_color_l = hex2argb(v)
                    elif strcmp(k, kws_border)[0] in ["color_right", "color_r"]:
                        border_color_r = hex2argb(v)
                    elif strcmp(k, kws_border)[0] in ["color_top", "color_t"]:
                        border_color_t = hex2argb(v)
                    elif strcmp(k, kws_border)[0] in ["color_bottom", "color_b"]:
                        border_color_b = hex2argb(v)
                    elif strcmp(k, kws_border)[0] in ["color_diagonal", "color_d"]:
                        border_color_d = hex2argb(v)
                    elif strcmp(k, kws_border)[0] in ["color_outline", "color_o"]:
                        border_color_o = hex2argb(v)
                    elif strcmp(k, kws_border)[0] in ["color_vertical", "color_v"]:
                        border_color_v = hex2argb(v)
                    elif strcmp(k, kws_border)[0] in ["color_horizontal", "color_h"]:
                        border_color_h = hex2argb(v)
                # *border style
                border_styles = [
                    "thin",
                    "medium",
                    "thick",
                    "dotted",
                    "dashed",
                    "hair",
                    "mediumDashed",
                    "dashDot",
                    "dashDotDot",
                    "slantDashDot",
                    "none",
                ]
                border_style_l, border_style_r, border_style_t, border_style_b = (
                    None,
                    None,
                    None,
                    None,
                )
                border_style_d, border_style_o, border_style_v, border_style_h = (
                    None,
                    None,
                    None,
                    None,
                )
                # get styles config
                for k, v in cell.get(K, {}).items():
                    # if not "style" in k:
                    #     break
                    if strcmp(k, kws_border)[0] in ["style"]:
                        border_style_all = strcmp(v, border_styles)[0]
                        # 如果设置了style,表示其它的所有的都设置成为一样的
                        # 然后再才开始自己定义其它的style
                        (
                            border_style_l,
                            border_style_r,
                            border_style_t,
                            border_style_b,
                        ) = (
                            border_style_all,
                            border_style_all,
                            border_style_all,
                            border_style_all,
                        )
                        (
                            border_style_d,
                            border_style_o,
                            border_style_v,
                            border_style_h,
                        ) = (
                            border_style_all,
                            border_style_all,
                            border_style_all,
                            border_style_all,
                        )
                    elif strcmp(k, kws_border)[0] in ["style_left", "style_l"]:
                        border_style_l = strcmp(v, border_styles)[0]
                    elif strcmp(k, kws_border)[0] in ["style_right", "style_r"]:
                        border_style_r = strcmp(v, border_styles)[0]
                    elif strcmp(k, kws_border)[0] in ["style_top", "style_t"]:
                        border_style_t = strcmp(v, border_styles)[0]
                    elif strcmp(k, kws_border)[0] in ["style_bottom", "style_b"]:
                        border_style_b = strcmp(v, border_styles)[0]
                    elif strcmp(k, kws_border)[0] in ["style_diagonal", "style_d"]:
                        border_style_d = strcmp(v, border_styles)[0]
                    elif strcmp(k, kws_border)[0] in ["style_outline", "style_o"]:
                        border_style_o = strcmp(v, border_styles)[0]
                    elif strcmp(k, kws_border)[0] in ["style_vertical", "style_v"]:
                        border_style_v = strcmp(v, border_styles)[0]
                    elif strcmp(k, kws_border)[0] in ["style_horizontal", "style_h"]:
                        border_style_h = strcmp(v, border_styles)[0]
                # * apply border config
                border = Border(
                    left=Side(border_style=border_style_l, color=border_color_l),
                    right=Side(border_style=border_style_r, color=border_color_r),
                    top=Side(border_style=border_style_t, color=border_color_t),
                    bottom=Side(border_style=border_style_b, color=border_color_b),
                    diagonal=Side(border_style=border_style_d, color=border_color_d),
                    diagonal_direction=0,
                    outline=Side(border_style=border_style_o, color=border_color_o),
                    vertical=Side(border_style=border_style_v, color=border_color_v),
                    horizontal=Side(border_style=border_style_h, color=border_color_h),
                )

        #! final apply configs
        for row in ws[cell_range]:
            for cell_ in row:
                if cell_font:
                    cell_.font = cell_font
                if cell_fill:
                    cell_.fill = cell_fill
                if cell_alignment:
                    cell_.alignment = cell_alignment
                if border:
                    cell_.border = border
    def generate_unique_sheet_name(wb, sheet_name):
        """Generate a unique sheet name if the given name already exists in the workbook."""
        if sheet_name not in wb.sheetnames:
            return sheet_name
        counter = 1
        unique_name = f"{sheet_name}_{counter}"
        while unique_name in wb.sheetnames:
            counter += 1
            unique_name = f"{sheet_name}_{counter}"
        return unique_name


    # if it is already worksheet format
    if isinstance(df, pd.DataFrame):
        pass
    elif isinstance(df, openpyxl.worksheet.worksheet.Worksheet) or isinstance(df, openpyxl.workbook.workbook.Workbook):
        pass
    elif df is None:
        if any(filename):
            df = fload(filename, output="bit")
    else:
        try:
            print(f"is loading file {os.path.basename(df)}")
            df = fload(df)
        except Exception as e:
            print(e)
            print(f"check the fpath:{df}")
    if filename is None:
        filename = str(datetime.now().strftime("%y%m%d_%H.xlsx"))
    
    kwargs.pop("format", None)  # 更好地跟fsave结合使用
    kwargs.pop("sheet_name", 0)  # 更好地跟df.to_excel结合使用
    # 只有openpyxl才支持 append
    mode = strcmp(kwargs.get("mode", "a"), ["a", "w","auto"])[0]
    # print(f'mode="{mode}"')
    kwargs.pop("mode", None)
    engine = strcmp(kwargs.get("engine", "openpyxl"), ["xlsxwriter", "openpyxl"])[0]
    # corr engine
    engine="openpyxl" if mode=="a" else "xlsxwriter" 
    # print(f'engine="{engine}"')
    if_sheet_exists=kwargs.get("if_sheet_exists","replace")
    # 通常是不需要保存index的
    index = kwargs.get("index", False)
    # header
    header=kwargs.pop("header",False)
    password = kwargs.pop("password", None)  # Use kwargs if provided
    
    kwargs.pop("password", None)
    kwargs.pop("header", None)
    kwargs.pop("index", None)
    kwargs.pop("if_sheet_exists", None)
    if isinstance(df, openpyxl.workbook.workbook.Workbook):
        """打开Sheet_name指定的表格，如果该表不存在，则创建一个新的或从现有文件中加载数据"""
        wb=df
        try:
            ws = wb.worksheets[sheet_name]
        except Exception as e:
            print(f'mode="{mode}"')
            if not os.path.exists(filename) or mode=="w":
                ws=wb.active
                ws.title = sheet_name
            else:# file exists
                wb = load_workbook(filename)
                with pd.ExcelWriter(filename, mode="a", engine=engine, if_sheet_exists=if_sheet_exists) as writer:
                    for ws in df.worksheets:  # Iterate through worksheets in the input workbook
                        ws_df = pd.DataFrame(ws.values)
                        ws_df.to_excel(writer,sheet_name=sheet_name,index=index,header=header,**kwargs) 
                # 重新打开刚更新过的数据
                wb = load_workbook(filename) 
                if sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    if not sheet_name==sheet_name:
                        wb.remove(wb[sheet_name])
                else:
                    raise KeyError(f"Worksheet {sheet_name} does not exist.") 
    else:
        if not os.path.exists(filename) or mode=="w": # or overwrite 
            with pd.ExcelWriter(filename, mode="w", engine=engine) as writer:
                df.to_excel(writer, sheet_name=sheet_name, index=index, header=header,**kwargs)
            wb = load_workbook(filename)
            if isinstance(sheet_name, str):
                ws = wb[sheet_name]
            elif isinstance(sheet_name, int):
                ws = wb.worksheets[sheet_name]
            else:
                ws = wb.worksheets[sheet_name]  # the index of worksheets
        else:# file exists
            wb = load_workbook(filename)
            with pd.ExcelWriter(filename, mode="a", engine=engine, if_sheet_exists=if_sheet_exists) as writer:
                df.to_excel(writer, sheet_name=sheet_name, index=index, header=header,**kwargs)
            wb = load_workbook(filename)
            if sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
            else:
                raise KeyError(f"Worksheet {sheet_name} does not exist.") 
    # ! Apply Text color
    if text_color:
        if verbose:
            text_color_str="""
            text_color=[
                {
                    (slice(1, 2), slice(0, 3)): {
                        "x>20": "#DD0531",  # Numbers > 20 → red
                        "x<=8": "#35B20C",  # Numbers ≤ 10 → blue
                        "'x'!='available'": "#0510DD",  # 'available' → green
                        "10<x<=30": "#EAB107",  # 10 < value ≤ 30 → orange
                        "10<=x<30": "#C615BE",  # 10 ≤ value < 30 → purple
                    }
                },
                {
                    (slice(3, df.shape[0] + 1), slice(0, 3)): {
                        "x>20": "#DD0531",  # Numbers > 20 → red
                        "x<=10": "#35B20C",  # Numbers ≤ 10 → blue
                        "'x'!='available'": "#0510DD",  # 'available' → green
                        "10<x<=30": "#EAB107",  # 10 < value ≤ 30 → orange
                        "10<=x<30": "#C615BE",  # 10 ≤ value < 30 → purple
                    }
                },
            ],
            """
            print(text_color_str)
        if not isinstance(text_color, list):
            text_color=[text_color]
        for text_color_ in text_color:
            for indices, dict_text_conditions in text_color_.items():
                ws = apply_color_to_worksheet(ws, sheet_name=sheet_name, conditions=dict_text_conditions, cell_idx=indices,where="text")
    # ! Apply Text color
    if bg_color:
        if verbose:
            bg_color_str="""
            bg_color=[
                {
                    (slice(1, 2), slice(0, 3)): {
                        "x>20": ["#DD0531","#35B20C"],  # Numbers > 20 → red
                        "x<=8": "#35B20C",  # Numbers ≤ 10 → blue
                        "'x'!='available'": "#0510DD",  # 'available' → green
                        "10<x<=30": "#EAB107",  # 10 < value ≤ 30 → orange
                        "10<=x<30": "#C615BE",  # 10 ≤ value < 30 → purple
                    }
                },
                {
                    (slice(3, df.shape[0] + 1), slice(0, 3)): {
                        "x>20": "#DD0531",  # Numbers > 20 → red
                        "x<=10": "#35B20C",  # Numbers ≤ 10 → blue
                        "'x'!='available'": "#0510DD",  # 'available' → green
                        "10<x<=30": "#EAB107",  # 10 < value ≤ 30 → orange
                        "10<=x<30": "#C615BE",  # 10 ≤ value < 30 → purple
                    }
                },
            ],
            """
            print(bg_color_str)
        if not isinstance(bg_color, list):
            bg_color=[bg_color]
        for bg_color_ in bg_color:
            for indices, dict_text_conditions in bg_color_.items():
                ws = apply_color_to_worksheet(ws, sheet_name=sheet_name, conditions=dict_text_conditions, cell_idx=indices,where="bg")
    # !Apply cell formatting
    if cell:
        if not isinstance(cell, list):
            cell = [cell]
        for cell_ in cell:
            for indices, format_options in cell_.items():
                cell_range = convert_indices_to_range(*indices)
                apply_format(ws, format_options, cell_range)

        if verbose:
            cell_tmp="""cell=[
                    {
                        (slice(0, 1), slice(0, len(df.columns))): {
                            "font": {
                                "name": "Calibri",  # Font name
                                "size": 14,  # Font size
                                "bold": True,  # Bold text
                                "italic": False,  # Italic text
                                "underline": "single",  # Underline (single, double)
                                "color": "#FFFFFF",  # Font color
                            },
                            "fill": {
                                "start_color": "a1cd1e",  # Starting color
                                "end_color": "4F81BD",  # Ending color (useful for gradients)
                                "fill_type": "solid",  # Fill type (solid, gradient, etc.)
                            },
                            "alignment": {
                                "horizontal": "center",  # Horizontal alignment (left, center, right)
                                "vertical": "center",  # Vertical alignment (top, center, bottom)
                                "wrap_text": True,  # Wrap text in the cell
                                "shrink_to_fit": True,  # Shrink text to fit within cell
                                "text_rotation": 0,  # Text rotation angle
                            },
                            "border": {
                                "left": "thin",
                                "right": "thin",
                                "top": "thin",
                                "bottom": "thin",
                                "color": "000000",  # Border color
                            },
                        }
                    },{}]"""
            print(cell_tmp)
    # !Apply cell shading
    if shade:
        if not isinstance(shade, list):
            shade = [shade]
        for shade_ in shade:
            for indices, shading in shade_.items():
                cell_range = convert_indices_to_range(*indices)
                fill = PatternFill(
                    start_color=hex2argb(shading.get("bg_color", "FFFFFF")),
                    end_color=hex2argb(shading.get("end_color", "FFFFFF")),
                    fill_type=shading.get("fill_type", "solid"),
                    patternType=shading.get("pattern_type", "solid"),
                    fgColor=hex2argb(shading.get("fg_color", "0000FF")),
                )
                for row in ws[cell_range]:
                    for cell in row:
                        cell.fill = fill
        if verbose:
            shade_temp="""shade={        
                        (slice(1, 4), slice(1, 3)): {
                        "bg_color": "#63C187",  # Background color
                        "pattern_type": "solid",  # Fill pattern (e.g., solid, darkGrid, lightGrid)
                        "fg_color": "#0000FF",  # Foreground color, used in patterns
                        "end_color": "0000FF",  # End color, useful for gradients
                        "fill_type": "solid",  # Type of fill (solid, gradient, etc.)
                    }}"""
            print(shade_temp)
    # !number formatting
    if number_format:
        if not isinstance(number_format, list):
            number_format = [number_format]
        for number_format_ in number_format:
            for col_idx, fmt in number_format_.items():
                col_letter = get_column_letter(col_idx)
                for cell in ws[col_letter][1:]:  # Skip the header
                    cell.number_format = fmt
        if verbose:
            number_format_temp="""number_format={
                                1: "0.00",  # Two decimal places for column index 1
                                2: "#,##0",  # Thousands separator
                                3: "0%",  # Percentage format
                                4: "$#,##0.00",  # Currency format
                            }""" 
            print(number_format_temp)
    
    if freeze:
        if isinstance(freeze,bool):
            freeze='A2'
        ws.freeze_panes = freeze  # Freeze everything above and to the left of A2
    if apply_filter:
        if isinstance(apply_filter, bool):
            # Default: Apply filter to the entire first row (header)
            filter_range = f"A1:{get_column_letter(ws.max_column)}1"
            ws.auto_filter.ref = filter_range
            if not freeze:
                ws.freeze_panes = "A2"  # Freeze everything above and to the left of A2
        elif isinstance(apply_filter, tuple):
            row_slice, col_slice = apply_filter
            # Extract the start and end indices for rows and columns
            start_row, end_row = row_slice.start, row_slice.stop
            start_col_idx, end_col_idx = col_slice.start, col_slice.stop

            # Ensure valid row and column indices
            if start_row < 1: start_row = 1  # Row should not be less than 1
            if end_row > ws.max_row: end_row = ws.max_row  # Ensure within sheet's row limits
            if start_col_idx < 1: start_col_idx = 1  # Column should not be less than 1
            if end_col_idx > ws.max_column: end_col_idx = ws.max_column  # Ensure within sheet's column limits
            
            # Get column letters based on indices
            start_col = get_column_letter(start_col_idx)
            end_col = get_column_letter(end_col_idx)

            # Define the filter range based on specific rows and columns
            filter_range = f"{start_col}{start_row}:{end_col}{end_row}"

            # Apply the filter
            ws.auto_filter.ref = filter_range
            if freeze:
                ws.freeze_panes = freeze  # Freeze everything above and to the left of A2
    # !widths 
    if isinstance(width,bool):
        width=None if width else False
    if isinstance(height,bool):
        height=None if height else False
    if width is None:  # automatic adust width
        for col in ws.columns:
            max_length = 0
            """column = col[0].column_letter  # Get the column letter"""
                    # Check the first cell in the column to get the column letter
            cell_first = col[0]

            # Check if the cell is part of a merged range
            if not any(cell_first.coordinate in range_ for range_ in ws.merged_cells.ranges):
                column = get_column_letter(cell_first.column)  # Get the column letter from the first cell
            else:
                # Skip the column if the first cell is merged
                continue
            for cell_ in col:
                try:
                    if cell_.value:
                        max_length = max(max_length, len(str(cell_.value)))
                except Exception:
                    pass
            adjusted_width = max_length + width_factor  # You can adjust the padding value as needed
            ws.column_dimensions[column].width = adjusted_width
    elif isinstance(width,bool):
        pass
    else:
        for col_idx, width_ in width.items():
            col_letter = get_column_letter(col_idx)
            ws.column_dimensions[col_letter].width = width_

    # !heights
    if height is None:  # automatic adust height
        for row in ws.iter_rows(min_row=1, max_row=ws.max_row):
            max_height = height_max
            for cell in row:
                if cell.value:
                    lines = str(cell.value).split("\n")
                    max_line_length = max(len(line) for line in lines)
                    estimated_height = 15 * len(lines)
                    if max_line_length > 20:
                        estimated_height += 5 * (max_line_length // 20)
                    max_height = max(max_height, estimated_height)
            ws.row_dimensions[row[0].row].height = max_height
    elif isinstance(height,bool) and not height:
        pass
    else:
        for row, height_ in height.items():
            ws.row_dimensions[row].height = height_

    # !Merge cells using slice indices
    if merge:
        if isinstance(merge, tuple):
            merge = [merge]
        for indices in merge:
            # Ensure indices are slice objects
            if len(indices) == 2:
                row_slice, col_slice = indices
                merge_range = convert_indices_to_range(row_slice, col_slice)
                ws.merge_cells(merge_range)
            elif len(indices) == 4:
                start_row, start_col, end_row, end_col = indices
                # Convert column numbers to letters (e.g., 1 -> 'A')
                start_cell = f"{get_column_letter(start_col)}{start_row}"
                end_cell = f"{get_column_letter(end_col)}{end_row}"
                merge_range = f"{start_cell}:{end_cell}"
                ws.merge_cells(merge_range)
            else:
                raise ValueError(
                    f"两种方式: 1. format: (start_row, start_col, end_row, end_col), 2. format: (slice(0, 3), slice(1, 2))"
                )

    # !Add comment
    if comment:
        if not isinstance(comment, list):
            comment = [comment]
        for comment_ in comment:
            if not isinstance(comment_, dict):
                raise TypeError("Each item in the `comments` list must be a dictionary.")
        
            for (row, col), comment_str in comment_.items():
                if not isinstance(row, int) or not isinstance(col, int):
                    raise ValueError("Row and column indices must be integers.")
                if not isinstance(comment_str, str):
                    raise ValueError("Comment text must be a string.")

                comment_curr = Comment(comment_str, "Author")
                comment_curr.visible = comment_always_visible
                # if comment_always_visible:
                #     comment_curr.width = 200  # Adjust width
                #     comment_curr.height = 100  # Adjust height
                ws.cell(row=row + 1, column=col + 1).comment = comment_curr
        if verbose:
            comment_tmp="""comment=[
                        {(0, 0): "This is a comment for A1"},
                        {(1, 1): "This is a comment for B2"},
                        {(2, 2): "This is a comment for C3"},
                    ]"""
            print(comment_tmp)
    # !Add link
    if link:
        if not isinstance(link, list):
            link = [link]
        for link_ in link:
            for (row, col), link_str in link_.items():
                ws.cell(row=row + 1, column=col + 1).hyperlink = link_str
        if verbose:
            print('link={(2, 2): "https://example.com"}')
    # !Apply data validation
    if data_validation:
        for indices, validation in data_validation.items():
            cell_range = convert_indices_to_range(*indices)
            dv = DataValidation(**validation)
            ws.add_data_validation(dv)
            dv.add(cell_range)
        if verbose:
            print("""data_validation={
                        (slice(1, 2), slice(2, 10)): {
                            "type": "list",
                            "formula1": '"Option1,Option2,Option3"',  # List of options
                            "allow_blank": True,
                            "showDropDown": True,
                            "showErrorMessage": True,
                            "errorTitle": "Invalid input",
                            "error": "Please select a valid option.",
                        }
                    }"""
                  )
    # !Protect sheet with a password
    # Fetch the password

    if all([password is not None, any([protect, isinstance(password, (str, list, tuple)) and any(password)])]):  # Check if protection options are provided
        if protect is None:
            protect={}
        password = password or protect.get("password")  # Default to 'protect' password if not set
        if password:  # Apply password protection if provided
            ws.protection.password = password
        ws.protection.sheet = protect.get("sheet", bool(password))
        ws.protection.objects = protect.get("objects", True)
        # Formatting options
        ws.protection.scenarios = protect.get("scenarios", False)
        ws.protection.formatCells = protect.get("formatCells", False)
        ws.protection.formatColumns = protect.get("formatColumns", False)
        ws.protection.formatRows = protect.get("formatRows", False)
        # Insert and delete options
        ws.protection.insertColumns = protect.get("insertColumns", True)
        ws.protection.insertRows = protect.get("insertRows", True)
        ws.protection.deleteColumns = protect.get("deleteColumns", True)
        ws.protection.deleteRows = protect.get("deleteRows", True)
        # Select locked or unlocked cells
        ws.protection.selectLockedCells = protect.get("selectLockedCells", False)
        ws.protection.selectUnlockedCells = protect.get("selectUnlockedCells", False)


    # !conditional formatting
    if conditional_format:
        if not isinstance(conditional_format, list):
            conditional_format = [conditional_format]
        print(f"conditional_format dict setting: 'color_scale', 'data_bar' and 'icon_set[not-ready]'")
        for conditional_format_ in conditional_format:
            for indices, rules in conditional_format_.items():
                cell_range = convert_indices_to_range(*indices)
                if not isinstance(rules, list):
                    rules=[rules]
                for rule in rules:
                    # Handle color scale
                    if "color_scale" in rule:
                        if verbose:
                            color_scale_tmp="""
                                    conditional_format={
                                        (slice(1, df.shape[0] + 1), slice(1, 2)):
                                            {
                                                "color_scale": {
                                                    "start_type": "min",
                                                    "start_value": 0,
                                                    "start_color": "#74ADE9",
                                                    "mid_type": "percentile",
                                                    "mid_value": 50,
                                                    "mid_color": "74ADE9",
                                                    "end_type": "max",
                                                    "end_value": 100,
                                                    "end_color": "#B62833",
                                                }
                                            }}
                            """
                            print(color_scale_tmp)
                        color_scale = rule["color_scale"] 

                        color_scale_rule = ColorScaleRule(
                            start_type=color_scale.get("start_type", "min"),
                            start_value=color_scale.get("start_value",None),
                            start_color=hex2argb(color_scale.get("start_color", "#74ADE9")),
                            mid_type=color_scale.get("mid_type","percentile"),
                            mid_value=color_scale.get("mid_value",None),
                            mid_color=hex2argb(color_scale.get("mid_color", "FFFFFF")),
                            end_type=color_scale.get("end_type", "max"),
                            end_value=color_scale.get("end_value",None),
                            end_color=hex2argb(color_scale.get("end_color", "#B62833")),
                        )
                        ws.conditional_formatting.add(cell_range, color_scale_rule)
                    # Handle data bar
                    if "data_bar" in rule:
                        if verbose:
                            data_bar_tmp="""
                                    conditional_format={
                                        (slice(1, df.shape[0] + 1), slice(1, 2)):
                                            {
                                                "data_bar": {
                                                    "start_type": "min",
                                                    "start_value": None,
                                                    "end_type": "max",
                                                    "end_value": None,
                                                    "color": "F6C9CE",
                                                    "show_value": True,
                                                }
                                            }}
                            """
                            print(data_bar_tmp)
                        data_bar = rule["data_bar"]
                        bar_color = hex2argb(data_bar.get("color", "638EC6"))

                        data_bar_rule = DataBarRule(
                            start_type=data_bar.get("start_type", "min"),
                            start_value=data_bar.get("start_value",None),
                            end_type=data_bar.get("end_type", "max"),
                            end_value=data_bar.get("end_value",None),
                            color=bar_color,
                            showValue=data_bar.get("show_value", True),
                        )
                        ws.conditional_formatting.add(cell_range, data_bar_rule)

                    # Handle icon setse
                    if "icon_set" in rule:
                        icon_set = rule["icon_set"]
                        icon_set_rule = IconSet(
                            iconSet=icon_set.get("iconSet", "3TrafficLights1"),  # Corrected
                            showValue=icon_set.get("show_value", True),        # Corrected
                            reverse=icon_set.get("reverse", False)            # Corrected
                        )
                        ws.conditional_formatting.add(cell_range, icon_set_rule)
                    # Handle text-based conditions
                    if "text_color" in rule: # not work
                        from openpyxl.styles.differential import DifferentialStyle
                        from openpyxl.formatting.rule import Rule
                        from openpyxl.styles import PatternFill

                        # Extract the fill properties from the rule
                        fill = rule.get("fill", {})
                        start_color = fill.get("start_color", "FFFFFF")  # Default to white if not specified
                        end_color = fill.get("end_color", "FFFFFF")  # Default to white if not specified
                        fill_type = fill.get("fill_type", "solid")  # Default to solid fill if not specified

                        # Extract the text condition or default to a space if 'text' is not provided
                        text = rule.get("text", " ")  

                        # Create the DifferentialStyle using the extracted fill settings
                        dxf = DifferentialStyle(
                            fill=PatternFill(start_color=start_color, end_color=end_color, fill_type=fill_type)
                        )
                        
                        # Create the text rule based on the text condition
                        text_rule = Rule(
                            type="containsText",  # The type of condition
                            operator=rule.get("operator", "equal"),  # Default operator is "equal"
                            text=text,
                            dxf=dxf,  # Apply the fill color from DifferentialStyle
                        )
                        ws.conditional_formatting.add(cell_range, text_rule)
        if verbose:
            conditional_format_temp="""
                    conditional_format={
                            (slice(1, 3), slice(1, 4)): [
                                {
                                    "data_bar": {
                                        "start_type": "min",
                                        "start_value": 100,
                                        "end_type": "max",
                                        "end_value": None,
                                        "color": "F6C9CE",
                                        "show_value": True,
                                    }
                                },
                                {
                                    "color_scale": {
                                        "start_type": "min",
                                        "start_value": 0,
                                        "start_color": "#74ADE9",
                                        "mid_type": "percentile",
                                        "mid_value": 50,
                                        "mid_color": "74ADE9",
                                        "end_type": "max",
                                        "end_value": 100,
                                        "end_color": "#B62833",
                                    }
                                },
                            ]
                        }
            """
            print(conditional_format_temp)
    if insert_img:
        if not isinstance(insert_img, dict):
            raise ValueError(f'insert_img 需要dict格式: e.g., insert_img={"A1":"example.png"}')
        try:
            from openpyxl import drawing
            from PIL import Image
            import PIL
            for img_cell, img_data in insert_img.items():
                img_width = img_height = None
                pil_img=img_path = None
                if isinstance(img_data, dict):
                    if "path" in img_data:
                        img_ = drawing.image.Image(img_data["path"])# File path
                    elif "image" in img_data:
                        img_ = drawing.image.Image(img_data["image"])# PIL Image object
                    elif "array" in img_data:
                        img_ = drawing.image.Image(Image.fromarray(img_data["array"]))# Convert NumPy array to PIL Image

                    img_width = img_data.get("width", None)
                    img_height = img_data.get("height", None)
                elif isinstance(img_data, str):
                    img_ = drawing.image.Image(img_data)# Direct file path
                elif isinstance(img_data, (PIL.Image.Image,PIL.PngImagePlugin.PngImageFile)):
                    img_ = drawing.image.Image(img_data)# Direct PIL Image object
                elif isinstance(img_data, np.ndarray):
                    img_ = drawing.image.Image(Image.fromarray(img_data))# Convert NumPy array to PIL Image
                elif pil_img:
                    img_ = drawing.image.Image(pil_img)

                # Set width and height if provided
                if img_width is not None:
                    img_.width = img_width
                if img_height is not None:
                    img_.height = img_height
                ws.add_image(img_, img_cell) 
                print(f"✅ at {img_cell} inserted image:  {os.path.basename(img_path)}")

        except Exception as e:
            print(e)
            
    # Save the workbook
    wb.save(filename)


def preview(var):
    """Master function to preview formatted variables in Jupyter."""
    from bs4 import BeautifulSoup
    from IPython.display import display, HTML, Markdown

    if isinstance(var, str):
        if isa(var, "html"):
            display(HTML(var))  # Render as HTML
        # Check if it's a valid markdown
        elif var.startswith("#"):
            display(Markdown(var))
        else:
            # Otherwise, display as plain text
            print(var)
    elif isinstance(var, BeautifulSoup):
        preview(str(var))
    elif isinstance(var, pd.DataFrame):
        # Display pandas DataFrame
        display(var)

    elif isinstance(var, list) or isinstance(var, dict):
        import json

        # Display JSON
        json_str = json.dumps(var, indent=4)
        display(Markdown(f"```json\n{json_str}\n```"))

    elif isinstance(var, bytes):
        # Display image if it's in bytes format
        display(Image(data=var))

    elif isinstance(var, str) and (var.endswith(".png") or var.endswith(".jpg")):
        # Display image from file path
        display(Image(filename=var))

    elif isinstance(var, dict):
        import json

        # Handle dictionary formatting
        json_str = json.dumps(var, indent=4)
        display(Markdown(f"```json\n{json_str}\n```"))

    else:
        # If the format is not recognized, print a message
        print("Format not recognized or unsupported.")


# # Example usages:
# preview("This is a plain text message.")
# preview("# This is a Markdown header")
# preview(pd.DataFrame({"Name": ["Alice", "Bob"], "Age": [25, 30]}))
# preview({"key": "value", "numbers": [1, 2, 3]})

def df2db(df: pd.DataFrame, db_url: str, table_name: str, if_exists: str = "replace"):
    """
    Save a pandas DataFrame to a SQL database (SQLite, MySQL, PostgreSQL).

    Usage:
    case 1. SQLite:
        df2db(df, 'sqlite:///sample_manager.db', 'users', 'replace')
    case 2. MySQL:
        df2db(df, 'mysql+mysqlconnector://user:password@localhost/mydatabase', 'users', 'replace')
    case 3. PostgreSQL:
        df2db(df, 'postgresql://user:password@localhost/mydatabase', 'users', 'replace')

    Parameters:
    - df (pd.DataFrame): The DataFrame to save.
    - db_url (str): The SQLAlchemy connection string to the database.
    - table_name (str): The name of the table to save the DataFrame to.
    - if_exists (str): What to do if the table already exists. Options:
        - 'replace': Drop the table before inserting new values.
        - 'append': Append new values to the table.
        - 'fail': Do nothing if the table exists. Default is 'replace'.
    """
    from sqlalchemy import create_engine
    from sqlalchemy.exc import SQLAlchemyError

    try:
        # Create the SQLAlchemy engine based on the provided db_url
        engine = create_engine(db_url)

        # Save the DataFrame to the SQL database (table)
        df.to_sql(table_name, con=engine, index=False, if_exists=if_exists)

        print(f"Data successfully saved to {db_url} in the {table_name} table.")

    except SQLAlchemyError as e:
        # Handle SQLAlchemy-related errors (e.g., connection issues, query issues)
        print(f"Error saving DataFrame to database: {str(e)}")

    except Exception as e:
        # Handle other unexpected errors
        print(f"An unexpected error occurred: {str(e)}")

def _df_outlier(
    data,
    columns=None,
    method=["zscore", "iqr", "percentile", "iforest"],
    min_outlier_method=3,  # 至少两种方法检查出outlier
    zscore_threshold=3,
    iqr_threshold=1.5,
    lower_percentile=5,
    upper_percentile=95,
):
    from scipy.stats import zscore
    from sklearn.ensemble import IsolationForest
    from sklearn.preprocessing import StandardScaler

    # Fill completely NaN columns with a default value (e.g., 0)
    data = data.copy()
    data.loc[:, data.isna().all()] = 0
    if columns is not None:
        if isinstance(columns, (list, pd.core.indexes.base.Index)):
            data = data[columns]
    col_names_org = data.columns.tolist()
    index_names_org = data.index.tolist()
    # Separate numeric and non-numeric columns
    numeric_data = data.select_dtypes(include=[np.number])
    non_numeric_data = data.select_dtypes(exclude=[np.number])

    # if columns is not None:
    #     numeric_data = numeric_data[columns]
    if numeric_data.empty:
        raise ValueError("Input data must contain numeric columns.")

    outliers_df = pd.DataFrame(index=numeric_data.index)
    if isinstance(method, str):
        method = [method]

    # Z-score method
    if "zscore" in method:
        z_scores = np.abs(zscore(numeric_data))
        outliers_df["zscore"] = np.any(z_scores > zscore_threshold, axis=1)

    # IQR method
    if "iqr" in method:
        Q1 = numeric_data.quantile(0.25)
        Q3 = numeric_data.quantile(0.75)
        IQR = Q3 - Q1
        lower_bound = Q1 - iqr_threshold * IQR
        upper_bound = Q3 + iqr_threshold * IQR
        outliers_df["iqr"] = (
            (numeric_data < lower_bound) | (numeric_data > upper_bound)
        ).any(axis=1)

    # Percentile method
    if "percentile" in method:
        lower_bound = numeric_data.quantile(lower_percentile / 100)
        upper_bound = numeric_data.quantile(upper_percentile / 100)
        outliers_df["percentile"] = (
            (numeric_data < lower_bound) | (numeric_data > upper_bound)
        ).any(axis=1)

    # Isolation Forest method
    if "iforest" in method:
        # iforest method cannot handle NaNs, then fillna with mean
        numeric_data_ = numeric_data.fillna(numeric_data.mean())
        scaler = StandardScaler()
        scaled_data = scaler.fit_transform(numeric_data_)
        iso_forest = IsolationForest(contamination=0.05)
        outliers_df["iforest"] = iso_forest.fit_predict(scaled_data) == -1

    # Combine all outlier detections
    if len(method) == 4:  # all method are used:
        outliers_df["outlier"] = outliers_df.sum(axis=1) >= min_outlier_method
    else:
        outliers_df["outlier"] = outliers_df.any(axis=1)

    # Handling Outliers: Remove or Winsorize or Replace with NaN
    processed_data = numeric_data.copy()

    processed_data.loc[outliers_df["outlier"]] = np.nan

    return processed_data

def df_group(
    data: pd.DataFrame,
    columns: Union[str, list, None] = None,
    by: str = None,
    param: Dict[str, Any] = None,
    sep: Union[str, list] = [", ",","],
    dropna: bool = True,
    unique: bool = False,
    astype: type = str,
    merge: List[str] = None,
    merge_symbo_column:str=' & ',
    merge_symbo_cell:str='[]',# ["{}","()","[]"]
) -> pd.DataFrame:
    """
    Groups a dataframe based on a specified column and applies aggregation functions dynamically.
    
    Parameters:
    data (pd.DataFrame): The dataframe to be grouped.
    columns (Union[str, list, None]): Columns to select; if None, all columns are selected.
    by (str): The column name to group by.
    param (dict): A dictionary specifying aggregation rules.
    sep (Union[str, list]): Separator for concatenated values. when sep is a list, then sep[0] used for general, sep[1] used in the merging
    dropna (bool): Whether to drop NaN values before aggregation.
    unique (bool): Whether to apply uniqueness before concatenation.
    astype (type): Data type to cast values before aggregation.
    merge (List[str]): List of columns to merge into a single paired column.
    merge_symbo_column:str: indicate in the columns, default ("&")
    merge_symbo_cell:str=default: '{}' or can select from ["{}","()","[]"]

    Usage:
    data = pd.DataFrame({
                        "Cage Tag": [1, 1, 2, 2, 3],
                        "Physical Tag": ["A1", "A2", "B1", "B2", "C1"],
                        "Sex": ["M", "F", "M", "M", "F"],
                        "Date of Birth": ["2021-06-01", "2021-06-02", "2021-07-01", "2021-07-02", "2021-08-01"],
                        "Age": [34, 35, 30, 31, 29],
                        "State": ["Mating", "Resting", "Mating", "Resting", "Mating"],
                        "Owner": ["Dr. Smith", "Dr. Smith", "Dr. Brown", "Dr. Brown", "Dr. Lee"],
                        })
    display(data)
    result = df_group(data, 
                    #   columns=["Sex", "Date of Birth", "Age", "State"], 
                    by="Cage Tag", 
                    merge=["Age", "State"], 
                    merge_symbo_column="|",
                    #   astype=str
                    # sep=[',',    '_'],
                    merge_symbo_cell=None
                    )  
    result

    
    """
    if param is None:
        param = {}
    if columns is None:
        columns = data.columns.tolist()
    elif isinstance(columns, str):
        columns = [columns]
    if not isinstance(sep, list):
        sep = [sep]
    sep.extend(sep) if len(sep)==1 else None

    # Merge specified columns into a single column
    if merge:
        merge_col_name = merge_symbo_column.join(merge)
        # data[merge_col_name] = data[merge].apply(lambda row: tuple(map(astype, row.dropna() if dropna else row)), axis=1)
        if merge_symbo_cell is None: 
            data[merge_col_name] = data[merge].apply(lambda row: f"{sep[1].join(map(astype, row.dropna()))}" if dropna else f"{sep[1].join(map(astype, row))}", axis=1)
        elif len(merge_symbo_cell)==2:
            data[merge_col_name] = data[merge].apply(lambda row: f"{merge_symbo_cell[0]}{sep[1].join(map(astype, row.dropna()))}{merge_symbo_cell[1]}" if dropna else f"{merge_symbo_cell[0]}{sep[1].join(map(astype, row))}{merge_symbo_cell[1]}", axis=1)
        else:
            data[merge_col_name] = data[merge].apply(lambda row: f"[{sep[1].join(map(astype, row.dropna()))}]" if dropna else f"[{sep[1].join(map(astype, row))}]", axis=1)
        columns.append(merge_col_name)
        
    default_aggregations = {
        col: (lambda x: sep[0].join(map(astype, x.dropna().unique() if unique else x.dropna())) if dropna else sep[0].join(map(astype, x.unique() if unique else x)))
        for col in columns if col != by and (merge is None or col not in merge)
    } 
    aggregation_rules = {**default_aggregations, **param}
    
    grouped_df = data.groupby(by).agg(aggregation_rules).reset_index() 
    return grouped_df
def df_outlier(
    data,
    columns=None,
    method=["zscore", "iqr", "percentile", "iforest"],
    min_outlier_method=2,  # 至少两种方法检查出outlier
    zscore_threshold=3,
    iqr_threshold=1.5,
    lower_percentile=5,
    upper_percentile=95,
):
    """
    Usage:
    data_out = df_outlier(
        data,
        columns=["income"],
        method="iforest",
        min_outlier_method=1)

    Advanced outlier detection and handling function.

    Parameters:
    - data: DataFrame, the input data (numerical).
    - method: List, the outlier detection method to use. Options: 'zscore', 'iqr', 'percentile', 'iforest'.
    - zscore_threshold: float, threshold for Z-score outlier detection (default 3).
    - iqr_threshold: float, threshold for IQR method (default 1.5).
    - lower_percentile: float, lower percentile for percentile-based outliers (default 5).
    - upper_percentile: float, upper percentile for percentile-based outliers (default 95).
    - keep_nan: bool, whether to replace outliers with NaN (default True).
    - plot: bool, whether to visualize the outliers (default False).
    - min_outlier_method: int, minimum number of method that need to flag a row as an outlier (default 2).
    - inplace: bool, whether to modify the original `data` DataFrame (default False).

    Returns:
    - processed_data: DataFrame with outliers handled based on method (if winsorize/remove is True).
    """
    col_names_org = data.columns.tolist()
    index_names_org = data.index.tolist()

    numeric_data = data.select_dtypes(include=[np.number])
    non_numeric_data = data.select_dtypes(exclude=[np.number])

    _outlier_df_tmp = pd.DataFrame()
    for col in numeric_data.columns:
        _outlier_df_tmp = pd.concat(
            [
                _outlier_df_tmp,
                _df_outlier(
                    data=data,
                    columns=[col],
                    method=method,
                    min_outlier_method=min_outlier_method,  # 至少两种方法检查出outlier
                    zscore_threshold=zscore_threshold,
                    iqr_threshold=iqr_threshold,
                    lower_percentile=lower_percentile,
                    upper_percentile=upper_percentile,
                ),
            ],
            axis=1,
            # join="inner",
        )
    processed_data = pd.concat([_outlier_df_tmp, non_numeric_data], axis=1)
    processed_data = processed_data[col_names_org]
    return processed_data


def df_extend(data: pd.DataFrame, column, axis=0, sep=None, prefix="col"):
    """
    Extend a DataFrame by the list elecments in the column.

    Parameters:
    ----------
    data : pd.DataFrame
        The input DataFrame to be extended.

    column : str
        The name of the column to be split.

    axis : int, optional
        The axis along which to expand the DataFrame.
        - 0 (default): Expand the specified column into multiple rows.
        - 1: Expand the specified column into multiple columns.

    sep : str, optional
        The separator used to split the values in the specified column.
        Must be provided for the function to work correctly.
    """

    data = data.copy()
    mask = data[column].str.contains(sep, na=False)
    data = data.copy()
    if mask.any():
        data[column] = data[column].apply(
            lambda x: x.split(sep) if isinstance(x, str) else x
        )  # Only split if x is a string

        # Strip spaces from each item in the lists
        data[column] = data[column].apply(
            lambda x: [item.strip() for item in x] if isinstance(x, list) else x
        )

    data = data.explode(column, ignore_index=True)
    return data


def df_cycle(data: pd.DataFrame, columns=None, max_val=None, inplace=False):
    """
    Purpose: transforms a datetime feature (like month or day) into a cyclic encoding for use in machine learning models, particularly neural networks.
    Usage:
        data = pd.DataFrame({'month': [1, 4, 7, 10, 12]})  # Just months as an example
        # df_cycle month cyclically
        data = df_cycle(data, 'month', 12)
    """
    if columns is None:
        columns = list(
            data.select_dtypes(include=np.number).columns
        )  # If no columns specified, use all columns
    if max_val is None:
        max_val = np.max(
            data[columns]
        )  # If no max_val specified, use the maximum value across all columns
    if isinstance(columns, str):
        columns = [
            columns
        ]  # If a single column name is provided as a string, convert it to a list

    # Check if inplace is True, so we modify the original dataframe
    if inplace:
        # Modify the data in place, no return statement needed
        for col in columns:
            data[col + "_sin"] = np.sin(2 * np.pi * data[col] / max_val)
            data[col + "_cos"] = np.cos(2 * np.pi * data[col] / max_val)
    else:
        # If inplace is False, return the modified dataframe
        new_data = data.copy()
        for col in columns:
            new_data[col + "_sin"] = np.sin(2 * np.pi * new_data[col] / max_val)
            new_data[col + "_cos"] = np.cos(2 * np.pi * new_data[col] / max_val)
        return new_data


# ! DataFrame
def df_astype(
    data: pd.DataFrame,
    columns: Optional[Union[str, List[str]]] = None,
    astype: str = None,  # "datetime",
    skip_row: Union[str, list] = None,
    original_fmt:str=None,
    fmt: Optional[str] = None,
    inplace: bool = False,
    errors: str = "coerce",  # Can be "ignore", "raise", or "coerce"
    verbose:bool=True,
    **kwargs,
) -> Optional[pd.DataFrame]:
    """
    Convert specified columns of a DataFrame to a specified type (e.g., datetime, float, int, numeric, timedelta).
    If columns is None, all columns in the DataFrame will be converted.

    Parameters:
    - df: DataFrame containing the columns to convert.
    - columns: Either a single column name, a list of column names, or None to convert all columns.
    - astype: The target type to convert the columns to ('datetime', 'float', 'int', 'numeric', 'timedelta', etc.).
    - fmt: Optional; format to specify the datetime format (only relevant for 'datetime' conversion).
    - inplace: Whether to modify the DataFrame in place or return a new one. Defaults to False.
    - errors: Can be "ignore", "raise", or "coerce"
    - **kwargs: Additional keyword arguments to pass to the conversion function (e.g., errors='ignore' for pd.to_datetime or pd.to_numeric).

    Returns:
    - If inplace=False: DataFrame with the specified columns (or all columns if columns=None) converted to the specified type.
    - If inplace=True: The original DataFrame is modified in place, and nothing is returned.
    """
    astypes = [
        "datetime",
        "timedelta",
        "numeric",
        "int",
        "int8",
        "int16",
        "int32",
        "int64",
        "uint8",
        "uint16",
        "uint32",
        "uint64",
        "float",
        "float16",
        "float32",
        "float64",
        "complex",
        "complex64",
        "complex128",
        "str",
        "string",
        "bool",
        "datetime64",
        "datetime64[ns]",
        "timedelta64",
        "timedelta64[ns]",
        "category",
        "object",
        "Sparse",
        "hour",
        "minute",
        "second",
        "time",
        "week",
        "date",
        "day",
        "month",
        "year",
        "circular",
    ]
    # If inplace is False, make a copy of the DataFrame
    if not inplace:
        data = data.copy()
    if skip_row is not None:
        data = data.drop(index=skip_row, errors="ignore")
    # If columns is None, apply to all columns
    if columns is None:
        columns = data.columns.tolist()
    # correct the astype input
    if isinstance(astype, str):
        astype = strcmp(astype, astypes)[0]
        print(f"converting as type: {astype}") 
    elif isinstance(astype, dict):
        for col, dtype in astype.items():
            dtype = "date" if dtype == "day" else dtype
            target_dtype = strcmp(dtype, astypes)[0]
            try:
                if target_dtype == "datetime":
                    data[col] = pd.to_datetime(data[col], format=original_fmt, errors=errors)
                elif target_dtype == "timedelta":
                    data[col] = pd.to_timedelta(data[col], errors=errors)
                else:
                    data[col] = data[col].astype(target_dtype)
            except Exception as e:
                print(f"Error converting column '{col}' to {target_dtype}: {e}")
        return data if not inplace else None
    # Ensure columns is a list
    if isinstance(columns, str):
        columns = [columns]

    # Convert specified columns
    for column in columns:
        try:
            if astype in [
                "datetime",
                "hour",
                "minute",
                "second",
                "time",
                "week",
                "date",
                "month",
                "year",
            ]:
                kwargs.pop("errors", None)
                # convert it as type: datetime
                if isinstance(column, int):
                    data.iloc[:, column] = pd.to_datetime(data.iloc[:, column], format=original_fmt, errors=errors, **kwargs) if original_fmt is not None else pd.to_datetime(data[column], errors=errors, **kwargs)

                    try:
                        if fmt is not None:
                            # data[column] = data[column].apply(lambda x: f"{x:{fmt}}")
                            data[column] = data[column].apply(
                                lambda x: x.strftime(fmt) if pd.notnull(x) else None
                            )
                    except Exception as e:
                        print(f"设置格式的时候有误: {e}")

                    # further convert:
                    if astype == "time":
                        data.iloc[:, column] = data.iloc[:, column].dt.time
                    elif astype == "month":
                        data.iloc[:, column] = data.iloc[:, column].dt.month
                    elif astype == "year":
                        data.iloc[:, column] = data.iloc[:, column].dt.year
                    elif astype == "date" or astype == "day":
                        data.iloc[:, column] = data.iloc[:, column].dt.date
                    elif astype == "hour":
                        data.iloc[:, column] = data.iloc[:, column].dt.hour
                    elif astype == "minute":
                        data.iloc[:, column] = data.iloc[:, column].dt.minute
                    elif astype == "second":
                        data.iloc[:, column] = data.iloc[:, column].dt.second
                    elif astype == "week":
                        data.iloc[:, column] = data.iloc[:, column].dt.day_name()
                else:
                    data[column] = (
                        pd.to_datetime(
                            data[column], format=original_fmt, errors=errors, **kwargs
                        )
                        if original_fmt is not None
                        else pd.to_datetime(data[column], errors=errors, **kwargs)
                    )

                    try:
                        if fmt is not None:
                            # data[column] = data[column].apply(lambda x: f"{x:{fmt}}")
                            data[column] = data[column].apply(
                                lambda x: x.strftime(fmt) if pd.notnull(x) else None
                            )
                    except Exception as e:
                        print(f"设置格式的时候有误: {e}")
                    # further convert:
                    if astype == "time":
                        data[column] = data[column].dt.time
                    elif astype == "month":
                        data[column] = data[column].dt.month
                    elif astype == "year":
                        data[column] = data[column].dt.year
                    elif astype == "date":
                        data[column] = data[column].dt.date
                    elif astype == "hour":
                        data[column] = data[column].dt.hour
                    elif astype == "minute":
                        data[column] = data[column].dt.minute
                    elif astype == "second":
                        data[column] = data[column].dt.second
                    elif astype == "week":
                        data[column] = data[column].dt.day_name()

            elif astype == "numeric":
                kwargs.pop("errors", None)
                data[column] = pd.to_numeric(data[column], errors=errors, **kwargs)
                # print(f"Successfully converted '{column}' to numeric.")
            elif astype == "timedelta":
                kwargs.pop("errors", None)
                data[column] = pd.to_timedelta(data[column], errors=errors, **kwargs)
                # print(f"Successfully converted '{column}' to timedelta.")
            elif astype == "circular":
                max_val = kwargs.get("max_val", None)
                data[column] = df_cycle(data=data, columns=column, max_val=max_val)
            else:
                # Convert to other types (e.g., float, int)
                if astype == "int":
                    data[column] = data[column].astype("float").astype("int")
                else:
                    data[column] = data[column].astype(astype)
                # print(f"Successfully converted '{column}' to {astype}.")

        except Exception as e:
            print(f"Error converting '{column}' to {astype}: {e}")
    try:
        if verbose:
            display(data.info()[:10])
    except:
        pass
    return data


# ! DataFrame
def df_sort_values(data, column, by=None, ascending=True, inplace=True, **kwargs):
    """
    Sort a DataFrame by a specified column based on a custom order or by count.

    Parameters:
    - data: DataFrame to be sorted.
    - column: The name of the column to sort by.
    - by: List specifying the custom order for sorting or 'count' to sort by frequency.
    - ascending: Boolean or list of booleans, default True.
                 Sort ascending vs. descending.
    - inplace: If True, perform operation in place and return None.
    - **kwargs: Additional arguments to pass to sort_values.

    Returns:
    - Sorted DataFrame if inplace is False, otherwise None.
    """
    if column not in data.columns:
        raise ValueError(f"Column '{column}' does not exist in the DataFrame.")

    if isinstance(by, str) and "count" in by.lower():
        # Count occurrences of each value in the specified column
        value_counts = data[column].value_counts()

        # Determine the order based on counts
        count_ascending = kwargs.pop("count_ascending", ascending)
        sorted_counts = value_counts.sort_values(
            ascending=count_ascending
        ).index.tolist()

        # Convert to a categorical type with the new order
        data[column] = pd.Categorical(
            data[column], categories=sorted_counts, ordered=True
        )
        # Set ascending to count_ascending for sorting
        ascending = count_ascending  # Adjust ascending for the final sort
    elif isinstance(by, list):
        # Convert the specified column to a categorical type with the custom order
        data[column] = pd.Categorical(data[column], categories=by, ordered=True)
    else:
        raise ValueError("Custom order must be a list or 'count'.")

    try:
        if inplace:  # replace the original
            data.sort_values(column, ascending=ascending, inplace=True, **kwargs)
            print(f"Successfully sorted DataFrame by '{column}'")
            return None
        else:
            sorted_df = data.sort_values(column, ascending=ascending, **kwargs)
            print(f"Successfully sorted DataFrame by '{column}' using custom order.")
            return sorted_df
    except Exception as e:
        print(f"Error sorting DataFrame by '{column}': {e}")
        return data


# # Example usage:
# # Sample DataFrame
# data = {
#     "month": ["March", "January", "February", "April", "December"],
#     "Amount": [200, 100, 150, 300, 250],
# }
# df_month = pd.DataFrame(data)

# # Define the month order
# month_order = [
#     "January",
#     "February",
#     "March",
#     "April",
#     "May",
#     "June",
#     "July",
#     "August",
#     "September",
#     "October",
#     "November",
#     "December",
# ]
# display(df_month)
# sorted_df_month = df_sort_values(df_month, "month", month_order, ascending=True)
# display(sorted_df_month)
# df_sort_values(df_month, "month", month_order, ascending=True, inplace=True)
# display(df_month)


def df_merge(
    df1: pd.DataFrame,
    df2: pd.DataFrame,
    use_index: bool = False,
    columns: list = ["col_left", "col_right"],
    how: str = "left",
) -> pd.DataFrame:
    """
    Merges two DataFrames based on either the index or shared columns with matching data types.
    usage:
        #(1) if the index are the same
            df_merged = df_merge(df1, df2, use_index=True(defalut), how='outer')
        #(2) if there are shaed columns, then based on shared columns
            df_merged = df_merge(df1, df2, how='outer')
        #(3) if columns: then based on the specific columns
            df_merged = df_merge(df1, df2, columns=["col_left", "col_right"],how='outer')
    Parameters:
    - df1 (pd.DataFrame): The first DataFrame.
    - df2 (pd.DataFrame): The second DataFrame.
    - use_index (bool): If True, first try to merge by index if they are comparable; otherwise, fall back to column-based merge.
    - how (str): Type of merge to perform: 'inner', 'outer', 'left', or 'right'. Default is 'inner'.
    'inner': only the rows that have matching values in both DataFrames (intersection)
    'outer': keeps all rows from both DataFrames and fills in missing values with NaN
    'left': keeps all rows from the left DataFrame and matches rows from the right DataFrame
    'right': keeps all rows from the right DataFrame and matches rows from the left DataFrame, filling with NaN if there is no match.

    Returns:
    - pd.DataFrame: The merged DataFrame.
    """

    # 1. Check if indices are comparable (same length and types)
    if use_index:
        print(f"Merging based on index using '{how}' join...")
        df_merged = pd.merge(df1, df2, left_index=True, right_index=True, how=how)
        return df_merged

    # 2. Find common columns with the same dtype
    common_columns = df1.columns.intersection(df2.columns)
    shared_columns = []
    for col in common_columns:
        if df1[col].dtype == df2[col].dtype:
            shared_columns.append(col)
    if not isinstance(columns, list):
        columns = [columns]
    if len(columns) != 2:
        raise ValueError(
            "'columns':list shoule be a list: columns=['col_left','col_right']"
        )
    if all(columns):
        print(f"Merging based on columns: {columns} using '{how}' join...")
        df_merged = pd.merge(df1, df2, left_on=columns[0], right_on=columns[1], how=how)
    elif shared_columns:
        print(
            f"Merging based on shared columns: {shared_columns} using '{how}' join..."
        )
        df_merged = pd.merge(df1, df2, on=shared_columns, how=how)
    else:
        raise ValueError(
            "No common columns with matching data types to merge on, and indices are not comparable."
        )
    return df_merged


def df_drop_duplicates(
    data: pd.DataFrame,
    by: Union[
        str, List[str]
    ] = "index",  # Options: 'index', or column name(s) for 'rows'
    keep="first",  # Options: 'first', 'last', or False (drop all duplicates)
    ignore_index=True,
    inplace: bool = False,
    verbose=True,
):
    """
    data (pd.DataFrame): DataFrame to drop duplicates from.
    by (str): Specify by to drop duplicates:
                 - 'index': Drop duplicates based on the DataFrame index.
                 - Column name(s) for row-wise duplicate checking.
    keep (str): Which duplicates to keep:
        'first',
        'last',
        False (drop all duplicates).
    inplace (bool): Whether to modify the original DataFrame in place.
    """
    original_shape = data.shape
    if by == "index":
        # Drop duplicates in the index
        result = data[~data.index.duplicated(keep=keep)]
    else:
        # Drop duplicates row-wise based on column(s)
        result = data.drop_duplicates(subset=by, keep=keep, ignore_index=ignore_index)
    if original_shape != result.shape or verbose:
        print(f"\nshape:{original_shape} (before drop_duplicates)")
        print(f"shape:{result.shape} (after drop_duplicates)")
    if inplace:
        # Modify the original DataFrame in place
        data.drop(data.index, inplace=True)  # Drop all rows first
        data[data.columns] = result  # Refill the DataFrame
        return None
    else:
        return result


#! fillna()
def df_fillna(
    data: pd.DataFrame,
    method: str = "knn",
    axis: int = 0,  # column-wise
    constant: float = None,
    n_neighbors: int = 5,  # KNN-specific
    max_iter: int = 10,  # Iterative methods specific
    inplace: bool = False,
    random_state: int = 1,
) -> pd.DataFrame:
    """
    Fill missing values in a DataFrame using specified imputation method.

    Parameters:
    data (pd.DataFrame): The DataFrame to fill missing values.
    method (str): The imputation method to use. Options are:
        - 'mean': Replace missing values with the mean of the column.
        - 'median': Replace missing values with the median of the column.
        - 'most_frequent': Replace missing values with the most frequent value in the column.
        - 'constant': Replace missing values with a constant value provided by the `constant` parameter.
        - 'knn': Use K-Nearest Neighbors imputation; replaces missing values based on the values of the nearest neighbors
        - 'iterative': Use Iterative imputation; each feature with missing values as a function of other features and estimates them iteratively
        - 'mice' (Multivariate Imputation by Chained Equations): A special case of iterative imputation.
        # - 'missforest': A random forest-based imputation method. Uses a random forest model to predict and fill missing values
        # - 'softimpute': Matrix factorization imputation.A matrix factorization technique where missing values are imputed by
        #       reconstructing the data matrix using low-rank approximation
        # - EM (Expectation-Maximization): Often used in advanced statistics to estimate missing values in a probabilistic framework.
        # - 'svd': Use IterativeSVD (matrix factorization via Singular Value Decomposition).

    axis (int): The axis along which to impute:
        - 0: Impute column-wise (default).
        - 1: Impute row-wise.
    constant (float, optional): Constant value to use for filling NaNs if method is 'constant'.
    inplace (bool): If True, modify the original DataFrame. If False, return a new DataFrame.

    """
    if isinstance(data, pd.Series):
        data = pd.DataFrame(data)
    # handle None
    for col in data.columns:
        data[col] = data[col].apply(lambda x: np.nan if x is None else x)

    # Fill completely NaN columns with a default value (e.g., 0)
    data = data.copy()
    data.loc[:, data.isna().all()] = 0

    col_names_org = data.columns.tolist()
    index_names_org = data.index.tolist()
    # Separate numeric and non-numeric columns
    numeric_data = data.select_dtypes(include=[np.number])
    non_numeric_data = data.select_dtypes(exclude=[np.number])

    if data.empty:
        raise ValueError("Input DataFrame is empty.")

    # Validate method
    methods = [
        "mean",
        "median",
        "most_frequent",
        "constant",
        "knn",
        "iterative",
    ]  # ,"missforest","softimpute","svd"]
    method = strcmp(method, methods)[0]

    # If using constant method, ask for a constant value
    if constant is not None:
        method = "constant"
        try:
            constant = float(constant)
        except ValueError:
            raise ValueError("Constant value must be a number.")

    # Initialize SimpleImputer with the chosen method
    if method == "constant":
        from sklearn.impute import SimpleImputer

        imputer = SimpleImputer(strategy=method, fill_value=constant)
    elif method == "knn":
        from sklearn.impute import KNNImputer

        imputer = KNNImputer(n_neighbors=n_neighbors)
    elif method == "iterative" or method == "mice":
        from sklearn.experimental import enable_iterative_imputer
        from sklearn.impute import IterativeImputer

        imputer = IterativeImputer(max_iter=max_iter, random_state=random_state)
    else:  # mean, median, most_frequent
        from sklearn.impute import SimpleImputer

        imputer = SimpleImputer(strategy=method)

    # Fit and transform the data
    if axis == 0:
        # Impute column-wise
        imputed_data = imputer.fit_transform(numeric_data)
    elif axis == 1:
        # Impute row-wise
        imputed_data = imputer.fit_transform(numeric_data.T)
    else:
        raise ValueError("Invalid axis. Use 0 for columns or 1 for rows.")

    imputed_data = pd.DataFrame(
        imputed_data if axis == 0 else imputed_data.T,
        index=numeric_data.index if axis == 0 else numeric_data.columns,
        columns=numeric_data.columns if axis == 0 else numeric_data.index,
    )
    for col in imputed_data.select_dtypes(include=[np.number]).columns:
        imputed_data[col] = imputed_data[col].astype(numeric_data[col].dtype)

    # Handle non-numeric data imputation
    if not non_numeric_data.empty:
        from sklearn.impute import SimpleImputer

        if method == "constant":
            non_numeric_imputer = SimpleImputer(
                strategy="constant", fill_value=constant
            )
        else:
            non_numeric_imputer = SimpleImputer(strategy="most_frequent")

        # Impute non-numeric columns column-wise (axis=0)
        imputed_non_numeric = non_numeric_imputer.fit_transform(non_numeric_data)

        # Convert imputed non-numeric array back to DataFrame with original index and column names
        imputed_non_numeric_df = pd.DataFrame(
            imputed_non_numeric,
            index=non_numeric_data.index,
            columns=non_numeric_data.columns,
        )
    else:
        imputed_non_numeric_df = pd.DataFrame(index=data.index)

    imputed_data = pd.concat([imputed_data, imputed_non_numeric_df], axis=1).reindex(
        columns=data.columns
    )

    if inplace:
        # Modify the original DataFrame
        data[:] = imputed_data[col_names_org]
        return None
    else:
        # Return the modified DataFrame
        return imputed_data[col_names_org]


# # example
# data = {
#     "A": [1, 2, np.nan, 4, 5],
#     "B": [np.nan, 2, 3, 4, np.nan],
#     "C": [1, np.nan, 3, 4, 5],
#     "D": [1, 2, 3, 4, np.nan],
# }

# # Define a function to test each imputation method
# methods = [
#     "mean",
#     "median",
#     "most_frequent",
#     "constant",
#     "knn",
#     "iterative",
#     # "missforest",
#     # "softimpute",
#     # "svd",
# ]

# # Create a dictionary to hold results
# results = {}

# for method_name in methods:
#     print(method_name)
#     display(df)
#     display(df_fillna(data=df, method=method_name, inplace=False, axis=0))


def df_encoder(
    data: pd.DataFrame,
    method: str = "dummy",  #'dummy', 'onehot', 'ordinal', 'label', 'target', 'binary'
    columns=None,
    target_column=None,  # Required for 'target' encoding method
    **kwargs,
) -> pd.DataFrame:
    """
    Methods explained:
    - 'dummy': pandas' `get_dummies` to create dummy variables for categorical columns, which is another form of one-hot encoding, but with a simpler interface.

    - 'onehot': One-hot encoding is used when there is no inherent order in categories. It creates a binary column for each category and is useful for nominal categorical variables. However, it increases dimensionality significantly if there are many unique categories.

    - 'ordinal': Ordinal encoding is used when there is an inherent order in the categories. It assigns integers to categories based on their order. Use this when the categories have a ranking (e.g., 'low', 'medium', 'high').

    - 'label': Label encoding is used for converting each unique category to a numeric label. It can be useful when working with algorithms that can handle categorical data natively (e.g., decision trees). However, it might introduce unintended ordinal relationships between the categories.

    - 'target': Target encoding is used when you encode a categorical feature based on the mean of the target variable. This is useful when there is a strong correlation between the categorical feature and the target variable. It is often used in predictive modeling to capture relationships that are not directly encoded in the feature.

    - 'binary': Binary encoding is a more efficient alternative to one-hot encoding when dealing with high-cardinality categorical variables. It converts categories into binary numbers and then splits them into multiple columns, reducing dimensionality compared to one-hot encoding.
    """

    # Select categorical columns
    categorical_cols = data.select_dtypes(exclude=np.number).columns.tolist()
    methods = ["dummy", "onehot", "ordinal", "label", "target", "binary"]
    method = strcmp(method, methods)[0]

    if columns is None:
        columns = categorical_cols

    # pd.get_dummies()
    if method == "dummy":
        dtype = kwargs.pop("dtype", int)
        drop_first = kwargs.pop("drop_first", True)
        try:
            encoded_df = pd.get_dummies(
                data[columns], drop_first=drop_first, dtype=dtype, **kwargs
            )
            return pd.concat([data.drop(columns, axis=1), encoded_df], axis=1)
        except Exception as e:
            # print(f"Warning, 没有进行转换, 因为: {e}")
            return data
    # One-hot encoding
    elif method == "onehot":
        from sklearn.preprocessing import OneHotEncoder

        encoder = OneHotEncoder(drop="first", sparse_output=False, **kwargs)
        encoded_data = encoder.fit_transform(data[columns])
        encoded_df = pd.DataFrame(
            encoded_data,
            columns=encoder.get_feature_names_out(columns),
            index=data.index,
        )
        return pd.concat([data.drop(columns, axis=1), encoded_df], axis=1)

    # Ordinal encoding
    elif method == "ordinal":
        from sklearn.preprocessing import OrdinalEncoder

        encoder = OrdinalEncoder(**kwargs)
        encoded_data = encoder.fit_transform(data[columns])
        encoded_df = pd.DataFrame(encoded_data, columns=columns, index=data.index)
        return pd.concat([data.drop(columns, axis=1), encoded_df], axis=1)

    # Label encoding
    elif method == "label":
        from sklearn.preprocessing import LabelEncoder

        encoder = LabelEncoder()
        # Apply LabelEncoder only to non-numeric columns
        non_numeric_columns = [
            col for col in columns if not pd.api.types.is_numeric_dtype(data[col])
        ]

        if not non_numeric_columns:
            return data
        encoded_data = data[non_numeric_columns].apply(
            lambda col: encoder.fit_transform(col)
        )
        return pd.concat([data.drop(non_numeric_columns, axis=1), encoded_data], axis=1)

    # Target encoding (Mean of the target for each category)
    elif method == "target":
        if target_column is None:
            raise ValueError("target_column must be provided for target encoding.")
        from category_encoders import TargetEncoder

        encoder = TargetEncoder(cols=columns, **kwargs)
        encoded_data = encoder.fit_transform(data[columns], data[target_column])
        return pd.concat([data.drop(columns, axis=1), encoded_data], axis=1)

    # Binary encoding (for high-cardinality categorical variables)
    elif method == "binary":
        from category_encoders import BinaryEncoder

        encoder = BinaryEncoder(cols=columns, **kwargs)
        encoded_data = encoder.fit_transform(data[columns])
        return pd.concat([data.drop(columns, axis=1), encoded_data], axis=1)


def df_scaler(
    data: pd.DataFrame,  # should be numeric dtype
    scaler=None,
    method="standard",
    columns=None,  # default, select all numeric col/row
    feature_range=None,  # specific for 'minmax'
    vmin=0,
    vmax=1,
    inplace=False,
    verbose=False,  # show usage
    axis=0,  # defalut column-wise
    return_scaler: bool = False,  # True: return both: return df, scaler
    **kwargs,
):
    """
    df_scaler(data, scaler="standard", inplace=False, axis=0, verbose=True)

    Parameters:
    - data: pandas DataFrame to be scaled.
    - method: Scaler type ('standard', 'minmax', 'robust'). Default is 'standard'.
    - columns: List of columns (for axis=0) or rows (for axis=1) to scale.
               If None, all numeric columns/rows will be scaled.
    - inplace: If True, modify the DataFrame in place. Otherwise, return a new DataFrame.
    - axis: Axis along which to scale. 0 for column-wise, 1 for row-wise. Default is 0.
    - verbose: If True, prints logs of the process.
    - kwargs: Additional arguments to be passed to the scaler.
    """
    if verbose:
        print('df_scaler(data, scaler="standard", inplace=False, axis=0, verbose=True)')
    if scaler is None:
        methods = ["standard", "minmax", "robust", "maxabs"]
        method = strcmp(method, methods)[0]
        if method == "standard":
            from sklearn.preprocessing import StandardScaler

            if verbose:
                print(
                    "performs z-score normalization: This will standardize each feature to have a mean of 0 and a standard deviation of 1."
                )
                print(
                    "Use when the data is approximately normally distributed (Gaussian).\nWorks well with algorithms sensitive to feature distribution, such as SVMs, linear regression, logistic regression, and neural networks."
                )
            scaler = StandardScaler(**kwargs)
        elif method == "minmax":
            from sklearn.preprocessing import MinMaxScaler

            if feature_range is None:
                feature_range = (vmin, vmax)
            if verbose:
                print(
                    "don't forget to define the range: e.g., 'feature_range=(0, 1)'. "
                )
                print(
                    "scales the features to the range [0, 1]. Adjust feature_range if you want a different range, like [-1, 1]."
                )
                print(
                    "Use when the data does not follow a normal distribution and you need all features in a specific range (e.g., [0, 1]).\nIdeal for algorithms that do not assume a particular distribution, such as k-nearest neighbors and neural networks."
                )
            scaler = MinMaxScaler(feature_range=feature_range, **kwargs)
        elif method == "robust":
            from sklearn.preprocessing import RobustScaler

            if verbose:
                print(
                    "scales the data based on the median and interquartile range, which is robust to outliers."
                )
                print(
                    "Use when the dataset contains outliers.\nThis method is useful because it scales based on the median and the interquartile range (IQR), which are more robust to outliers than the mean and standard deviation."
                )
            scaler = RobustScaler(**kwargs)
        elif method == "maxabs":
            from sklearn.preprocessing import MaxAbsScaler

            if verbose:
                print(
                    "This scales each feature by its maximum absolute value, resulting in values within the range [-1, 1] for each feature."
                )
                print(
                    "Use for data that is already sparse or when features have positive or negative values that need scaling without shifting the data.\nOften used with sparse data (data with many zeros), where preserving zero entries is essential, such as in text data or recommendation systems."
                )
            scaler = MaxAbsScaler(**kwargs)
        if axis not in [0, 1]:
            raise ValueError("Axis must be 0 (column-wise) or 1 (row-wise).")
    if verbose:
        print(scaler)
    if axis == 0:
        # Column-wise scaling (default)
        if columns is None:
            columns = data.select_dtypes(include=np.number).columns.tolist()
        non_numeric_columns = data.columns.difference(columns)

        # scaled_data = scaler.fit_transform(data[columns])
        if scaler is None or not hasattr(scaler, "mean_"):
            scaled_data = scaler.fit_transform(data[columns])
        else:
            scaled_data = scaler.transform(data[columns])

        if inplace:
            data[columns] = scaled_data
            print("Original DataFrame modified in place (column-wise).")
        else:
            scaled_df = pd.concat(
                [
                    pd.DataFrame(scaled_data, columns=columns, index=data.index),
                    data[non_numeric_columns],
                ],
                axis=1,
            )
            scaled_df = scaled_df[data.columns]  # Maintain column order
            if return_scaler:
                return scaled_df, scaler
            else:
                return scaled_df

    elif axis == 1:
        # Row-wise scaling
        if columns is None:
            columns = data.index.tolist()
        numeric_rows = data.loc[columns].select_dtypes(include=np.number)
        if numeric_rows.empty:
            raise ValueError("No numeric rows to scale.")

        print(f"Scaling rows")

        # scaled_data = scaler.fit_transform(
        #     numeric_rows.T
        # ).T  # Transpose for scaling and then back
        scaled_data = (
            scaler.fit_transform(numeric_rows.T).T
            if scaler is None or not hasattr(scaler, "mean_")
            else scaler.transform(numeric_rows.T).T
        )

        if inplace:
            data.loc[numeric_rows.index] = scaled_data
            print("Original DataFrame modified in place (row-wise).")
        else:
            scaled_df = data.copy()
            scaled_df.loc[numeric_rows.index] = scaled_data
            if return_scaler:
                return scaled_df, scaler
            else:
                return scaled_df


def df_special_characters_cleaner(
    data: pd.DataFrame, where=["column", "content", "index"]
) -> pd.DataFrame:
    """
    to clean special characters:
    usage:
        df_special_characters_cleaner(data=df, where='column')
    """
    if not isinstance(where, list):
        where = [where]
    where_to_clean = ["column", "content", "index"]
    where_ = [strcmp(i, where_to_clean)[0] for i in where]

    # 1. Clean column names by replacing special characters with underscores
    if "column" in where_:
        try:
            data.columns = data.columns.str.replace(r"[^\w\s]", "_", regex=True)
        except Exception as e:
            print(e)

    # 2. Clean only object-type columns (text columns)
    try:
        if "content" in where_:
            for col in data.select_dtypes(include=["object"]).columns:
                data[col] = data[col].str.replace(r"[^\w\s]", "", regex=True)
        if data.index.dtype == "object" and index in where_:
            data.index = data.index.str.replace(r"[^\w\s]", "_", regex=True)
    except:
        pass
    return data


def df_cluster(
    data: pd.DataFrame,
    columns: Optional[list] = None,
    n_clusters: Optional[int] = None,
    range_n_clusters: Union[range, np.ndarray] = range(2, 11),
    scale: bool = True,
    plot: Union[str, list] = "all",
    inplace: bool = True,
    ax=None,
):
    from sklearn.preprocessing import StandardScaler
    from sklearn.cluster import KMeans
    from sklearn.metrics import silhouette_score, silhouette_samples
    import seaborn as sns
    import numpy as np
    import pandas as pd
    import matplotlib.pyplot as plt

    """
    Performs clustering analysis on the provided feature matrix using K-Means.

    Parameters:
        X (np.ndarray):
            A 2D numpy array or DataFrame containing numerical feature data,
            where each row corresponds to an observation and each column to a feature.

        range_n_clusters (range):
            A range object specifying the number of clusters to evaluate for K-Means clustering.
            Default is range(2, 11), meaning it will evaluate from 2 to 10 clusters.

        scale (bool):
            A flag indicating whether to standardize the features before clustering.
            Default is True, which scales the data to have a mean of 0 and variance of 1.

        plot (bool):
            A flag indicating whether to generate visualizations of the clustering analysis.
            Default is True, which will plot silhouette scores, inertia, and other relevant plots.
    Returns:
        tuple: 
            A tuple containing the modified DataFrame with cluster labels, 
            the optimal number of clusters, and the Axes object (if any).
    """
    X = data[columns].values if columns is not None else data.values

    silhouette_avg_scores = []
    inertia_scores = []

    # Standardize the features
    if scale:
        scaler = StandardScaler()
        X = scaler.fit_transform(X)

    for n_cluster in range_n_clusters:
        kmeans = KMeans(n_clusters=n_cluster, random_state=1)
        cluster_labels = kmeans.fit_predict(X)

        silhouette_avg = silhouette_score(X, cluster_labels)
        silhouette_avg_scores.append(silhouette_avg)
        inertia_scores.append(kmeans.inertia_)
        print(
            f"For n_clusters = {n_cluster}, the average silhouette_score is : {silhouette_avg:.4f}"
        )

    # Determine the optimal number of clusters based on the maximum silhouette score
    if n_clusters is None:
        n_clusters = range_n_clusters[np.argmax(silhouette_avg_scores)]
    print(f"n_clusters = {n_clusters}")

    # Apply K-Means Clustering with Optimal Number of Clusters
    kmeans = KMeans(n_clusters=n_clusters, random_state=1)
    cluster_labels = kmeans.fit_predict(X)

    if plot:
        # ! Interpreting the plots from your K-Means clustering analysis
        # ! 1. Silhouette Score and Inertia vs Number of Clusters
        # Description:
        # This plot has two y-axes: the left y-axis represents the Silhouette Score, and the right y-axis
        # represents Inertia.
        # The x-axis represents the number of clusters (k).

        # Interpretation:

        # Silhouette Score:
        # Ranges from -1 to 1, where a score close to 1 indicates that points are well-clustered, while a
        # score close to -1 indicates that points might be incorrectly clustered.
        # A higher silhouette score generally suggests that the data points are appropriately clustered.
        # Look for the highest value to determine the optimal number of clusters.

        # Inertia:
        # Represents the sum of squared distances from each point to its assigned cluster center.
        # Lower inertia values indicate tighter clusters.
        # As the number of clusters increases, inertia typically decreases, but the rate of decrease
        # may slow down, indicating diminishing returns for additional clusters.

        # Optimal Number of Clusters:
        # You can identify an optimal number of clusters where the silhouette score is maximized and
        # inertia starts to plateau (the "elbow" point).
        # This typically suggests that increasing the number of clusters further yields less meaningful
        # separations.
        if ax is None:
            _, ax = plt.subplots(figsize=inch2cm(10, 6))
        color = "tab:blue"
        ax.plot(
            range_n_clusters,
            silhouette_avg_scores,
            marker="o",
            color=color,
            label="Silhouette Score",
        )
        ax.set_xlabel("Number of Clusters")
        ax.set_ylabel("Silhouette Score", color=color)
        ax.tick_params(axis="y", labelcolor=color)
        # add right axis: inertia
        ax2 = ax.twinx()
        color = "tab:red"
        ax2.set_ylabel("Inertia", color=color)
        ax2.plot(
            range_n_clusters,
            inertia_scores,
            marker="x",
            color=color,
            label="Inertia",
        )
        ax2.tick_params(axis="y", labelcolor=color)

        plt.title("Silhouette Score and Inertia vs Number of Clusters")
        plt.xticks(range_n_clusters)
        plt.grid()
        plt.axvline(x=n_clusters, linestyle="--", color="r", label="Optimal n_clusters")
        # ! 2. Elbow Method Plot
        # Description:
        # This plot shows the Inertia against the number of clusters.

        # Interpretation:
        # The elbow point is where the inertia begins to decrease at a slower rate. This point suggests that
        # adding more clusters beyond this point does not significantly improve the clustering performance.
        # Look for a noticeable bend in the curve to identify the optimal number of clusters, indicated by the
        # vertical dashed line.
        # Inertia plot
        plt.figure(figsize=inch2cm(10, 6))
        plt.plot(range_n_clusters, inertia_scores, marker="o")
        plt.title("Elbow Method for Optimal k")
        plt.xlabel("Number of clusters")
        plt.ylabel("Inertia")
        plt.grid()
        plt.axvline(
            x=np.argmax(silhouette_avg_scores) + 2,
            linestyle="--",
            color="r",
            label="Optimal n_clusters",
        )
        plt.legend()
        # ! Silhouette Plots
        # 3. Silhouette Plot for Various Clusters
        # Description:
        # This horizontal bar plot shows the silhouette coefficient values for each sample, organized by cluster.

        # Interpretation:
        # Each bar represents the silhouette score of a sample within a specific cluster. Longer bars indicate
        # that the samples are well-clustered.
        # The height of the bars shows how similar points within the same cluster are to one another compared to
        # points in other clusters.
        # The vertical red dashed line indicates the average silhouette score for all samples.
        # You want the majority of silhouette values to be above the average line, indicating that most points
        # are well-clustered.

        # 以下代码不用再跑一次了
        # n_clusters = (
        #     np.argmax(silhouette_avg_scores) + 2
        # )  # Optimal clusters based on max silhouette score
        # kmeans = KMeans(n_clusters=n_clusters, random_state=1)
        # cluster_labels = kmeans.fit_predict(X)
        silhouette_vals = silhouette_samples(X, cluster_labels)

        plt.figure(figsize=inch2cm(10, 6))
        y_lower = 10
        for i in range(n_clusters):
            # Aggregate the silhouette scores for samples belonging to cluster i
            ith_cluster_silhouette_values = silhouette_vals[cluster_labels == i]

            # Sort the values
            ith_cluster_silhouette_values.sort()

            size_cluster_i = ith_cluster_silhouette_values.shape[0]
            y_upper = y_lower + size_cluster_i

            # Create a horizontal bar plot for the silhouette scores
            plt.barh(range(y_lower, y_upper), ith_cluster_silhouette_values, height=0.5)

            # Label the silhouette scores
            plt.text(-0.05, (y_lower + y_upper) / 2, str(i + 2))
            y_lower = y_upper + 10  # 10 for the 0 samples

        plt.title("Silhouette Plot for the Various Clusters")
        plt.xlabel("Silhouette Coefficient Values")
        plt.ylabel("Cluster Label")
        plt.axvline(x=np.mean(silhouette_vals), color="red", linestyle="--")

        df_clusters = pd.DataFrame(
            X, columns=[f"Feature {i+1}" for i in range(X.shape[1])]
        )
        df_clusters["Cluster"] = cluster_labels
        # ! pairplot of the clusters
        # Overview of the Pairplot
        # Axes and Grid:
        # The pairplot creates a grid of scatter plots for each pair of features in your dataset.
        # Each point in the scatter plots represents a sample from your dataset, colored according to its cluster assignment.

        # Diagonal Elements:
        # The diagonal plots usually show the distribution of each feature. In this case, since X.shape[1] <= 4,
        # there will be a maximum of four features plotted against each other. The diagonal could display histograms or
        # kernel density estimates (KDE) for each feature.

        # Interpretation of the Pairplot

        # Feature Relationships:
        # Look at each scatter plot in the off-diagonal plots. Each plot shows the relationship between two features. Points that
        # are close together in the scatter plot suggest similar values for those features.
        # Cluster Separation: You want to see clusters of different colors (representing different clusters) that are visually distinct.
        # Good separation indicates that the clustering algorithm effectively identified different groups within your data.
        # Overlapping Points: If points from different clusters overlap significantly in any scatter plot, it indicates that those clusters
        # might not be distinct in terms of the two features being compared.
        # Cluster Characteristics:
        # Shape and Distribution: Observe the shape of the clusters. Are they spherical, elongated, or irregular? This can give insights
        # into how well the K-Means (or other clustering methods) has performed:
        # Spherical Clusters: Indicates that clusters are well defined and separated.
        # Elongated Clusters: May suggest that the algorithm is capturing variations along specific axes but could benefit from adjustments
        # in clustering parameters or methods.
        # Feature Influence: Identify which features contribute most to cluster separation. For instance, if you see that one feature
        # consistently separates two clusters, it may be a key factor for clustering.
        # Diagonal Histograms/KDE:
        # The diagonal plots show the distribution of individual features across all samples. Look for:
        # Distribution Shape: Is the distribution unimodal, bimodal, skewed, or uniform?
        # Concentration: Areas with a high density of points may indicate that certain values are more common among samples.
        # Differences Among Clusters: If you see distinct peaks in the histograms for different clusters, it suggests that those clusters are
        # characterized by specific ranges of feature values.
        # Example Observations
        # Feature 1 vs. Feature 2: If there are clear, well-separated clusters in this scatter plot, it suggests that these two features
        # effectively distinguish between the clusters.
        # Feature 3 vs. Feature 4: If you observe significant overlap between clusters in this plot, it may indicate that these features do not
        # provide a strong basis for clustering.
        # Diagonal Plots: If you notice that one cluster has a higher density of points at lower values for a specific feature, while another
        # cluster is concentrated at higher values, this suggests that this feature is critical for differentiating those clusters.

        # Pairplot of the clusters
        # * 为什么要限制到4个features?
        # 2 features=1 scatter plot            # 3 features=3 scatter plots
        # 4 features=6 scatter plots           # 5 features=10 scatter plots
        # 6 features=15 scatter plots          # 10 features=45 scatter plots
        # Pairplot works well with low-dimensional data, 如果维度比较高的话, 子图也很多,失去了它的意义
        if X.shape[1] <= 6:
            plt.figure(figsize=(8, 4))
            sns.pairplot(df_clusters, hue="Cluster", palette="tab10")
            plt.suptitle("Pairplot of Clusters", y=1.02)

    # Add cluster labels to the DataFrame or modify in-place
    if inplace:  # replace the oringinal data
        data["Cluster"] = cluster_labels
        return None, n_clusters, kmeans, ax  # Return None when inplace is True
    else:
        data_copy = data.copy()
        data_copy["Cluster"] = cluster_labels
        return data_copy, n_clusters, kmeans, ax


# example:
# clustering_features = [marker + "_log" for marker in markers]
# df_cluster(data, columns=clustering_features, n_clusters=3,range_n_clusters=np.arange(3, 7))

"""
# You're on the right track, but let's clarify how PCA and clustering (like KMeans) work, especially 
# in the context of your dataset with 7 columns and 23,121 rows.

# Principal Component Analysis (PCA)
# Purpose of PCA:
# PCA is a dimensionality reduction technique. It transforms your dataset from a high-dimensional space 
# (in your case, 7 dimensions corresponding to your 7 columns) to a lower-dimensional space while 
# retaining as much variance (information) as possible.
# How PCA Works:
# PCA computes new features called "principal components" that are linear combinations of the original 
# features.
# The first principal component captures the most variance, the second captures the next most variance 
# (orthogonal to the first), and so on.
# If you set n_components=2, for example, PCA will reduce your dataset from 7 columns to 2 columns. 
# This helps in visualizing and analyzing the data with fewer dimensions.
# Result of PCA:
# After applying PCA, your original dataset with 7 columns will be transformed into a new dataset with 
# the specified number of components (e.g., 2 or 3).
# The transformed dataset will have fewer columns but should capture most of the important information 
# from the original dataset.

# Clustering (KMeans)
# Purpose of Clustering:
# Clustering is used to group data points based on their similarities. KMeans, specifically, partitions 
# your data into a specified number of clusters (groups).
# How KMeans Works:
# KMeans assigns each data point to one of the k clusters based on the feature space (original or 
# PCA-transformed).
# It aims to minimize the variance within each cluster while maximizing the variance between clusters.
# It does not classify the data in each column independently; instead, it considers the overall similarity 
# between data points based on their features.
# Result of KMeans:
# The output will be cluster labels for each data point (e.g., which cluster a particular observation 
# belongs to).
# You can visualize how many groups were formed and analyze the characteristics of each cluster.

# Summary
# PCA reduces the number of features (columns) in your dataset, transforming it into a lower-dimensional
# space.
# KMeans then classifies data points based on the features of the transformed dataset (or the original 
# if you choose) into different subgroups (clusters).
# By combining these techniques, you can simplify the complexity of your data and uncover patterns that 
# might not be visible in the original high-dimensional space. Let me know if you have further questions!
"""


def df_reducer(
    data: pd.DataFrame,
    columns: Optional[List[str]] = None,
    method: str = "umap",  # 'pca', 'umap'
    n_components: int = 2,  # Default for umap, but 50 for PCA
    umap_neighbors: int = 15,  # UMAP-specific
    umap_min_dist: float = 0.1,  # UMAP-specific
    tsne_perplexity: int = 30,  # t-SNE-specific
    hue: str = None,  # lda-specific
    scale: bool = True,
    fill_missing: bool = True,
    size=2,  # for plot marker size
    markerscale=4,  # for plot, legend marker size scale
    edgecolor="none",  # for plot,
    legend_loc="best",  # for plot,
    bbox_to_anchor=None,
    ncols=1,
    debug: bool = False,
    inplace: bool = True,  # replace the oringinal data
    plot_: bool = False,  # plot scatterplot, but no 'hue',so it is meaningless
    random_state=1,
    ax=None,
    figsize=None,
    verbose=True,
    **kwargs,
) -> pd.DataFrame:
    dict_methods = {
        #!Linear Dimensionality Reduction: For simplifying data with techniques that assume linearity.
        "pca": "pca(Principal Component Analysis): \n\tUseful for reducing dimensionality of continuous data while retaining variance. Advantage: Simplifies data, speeds up computation, reduces noise. Limitation: Assumes linear relationships, may lose interpretability in transformed dimensions.",
        "lda": "lda(Linear Discriminant Analysis):\n\tUseful for supervised dimensionality reduction when class separability is important. Advantage: Enhances separability between classes, can improve classification performance. Limitation: Assumes normal distribution and equal class covariances, linear boundaries only.",
        "factor": "factor(Factor Analysis):\n\tSuitable for datasets with observed and underlying latent variables. Advantage: Reveals hidden structure in correlated data, dimensionality reduction with interpretable factors. Limitation: Assumes factors are linear combinations, less effective for nonlinear data.",
        "svd": "svd(Singular Value Decomposition):\n\tSuitable for matrix decomposition, dimensionality reduction in tasks like topic modeling or image compression. Advantage: Efficient, preserves variance, useful in linear transformations. Limitation: Assumes linear relationships, sensitive to noise, may not capture non-linear structure.",
        #! Non-linear Dimensionality Reduction (Manifold Learning)
        "umap": "umap(Uniform Manifold Approximation and Projection):\n\tBest for high-dimensional data visualization (e.g., embeddings). Advantage: Captures complex structure while preserving both local and global data topology. Limitation: Non-deterministic results can vary, sensitive to parameter tuning.",
        "tsne": "tsne(t-Distributed Stochastic Neighbor Embedding):\n\tt-SNE excels at preserving local structure (i.e., clusters), but it often loses global. relationships, causing clusters to appear in arbitrary proximities to each other.  Ideal for clustering and visualizing high-dimensional data, especially for clear cluster separation. Advantage: Captures local relationships effectively. Limitation: Computationally intensive, does not preserve global structure well, requires parameter tuning.",
        "mds": "mds(Multidimensional Scaling):\n\tAppropriate for visualizing pairwise similarity or distance in data. Advantage: Maintains the perceived similarity or dissimilarity between points. Limitation: Computationally expensive for large datasets, less effective for complex, high-dimensional structures.",
        "lle": "lle(Locally Linear Embedding):\n\tUseful for non-linear dimensionality reduction when local relationships are important (e.g., manifold learning). Advantage: Preserves local data structure, good for manifold-type data. Limitation: Sensitive to noise and number of neighbors, not effective for global structure.",
        "kpca": "kpca(Kernel Principal Component Analysis):\n\tGood for non-linear data with complex structure, enhancing separability. Advantage: Extends PCA to capture non-linear relationships. Limitation: Computationally expensive, sensitive to kernel and parameter choice, less interpretable.",
        "ica": "ica(Independent Component Analysis):\n\tEffective for blind source separation (e.g., EEG, audio signal processing).is generally categorized under Non-linear Dimensionality Reduction, but it also serves a distinct role in Blind Source Separation. While ICA is commonly used for dimensionality reduction, particularly in contexts where data sources need to be disentangled (e.g., separating mixed signals like EEG or audio data), it focuses on finding statistically independent components rather than maximizing variance (like PCA) or preserving distances (like MDS or UMAP). Advantage: Extracts independent signals/components, useful in mixed signal scenarios. Limitation: Assumes statistical independence, sensitive to noise and algorithm choice.",
        #! Anomaly Detection: Specialized for detecting outliers or unusual patterns
        "isolation_forest": "Isolation Forest:\n\tDesigned for anomaly detection, especially in high-dimensional data. Advantage: Effective in detecting outliers, efficient for large datasets. Limitation: Sensitive to contamination ratio parameter, not ideal for highly structured or non-anomalous data.",
        #! more methods
        "truncated_svd": "Truncated Singular Value Decomposition (SVD):\n\tEfficient for large sparse datasets, useful for feature reduction in natural language processing (e.g., Latent Semantic Analysis). Advantage: Efficient in memory usage for large datasets. Limitation: Limited in non-linear transformation.",
        "spectral_embedding": "Spectral Embedding:\n\tBased on graph theory, it can be useful for clustering and visualization, especially for data with connected structures. Advantage: Preserves global structure, good for graph-type data. Limitation: Sensitive to parameter choice, not ideal for arbitrary non-connected data.",
        "autoencoder": "Autoencoder:\n\tA neural network-based approach for complex feature learning and non-linear dimensionality reduction. Advantage: Can capture very complex relationships. Limitation: Computationally expensive, requires neural network expertise for effective tuning.",
        "nmf": "Non-negative Matrix Factorization:\n\tEffective for parts-based decomposition, commonly used for sparse and non-negative data, e.g., text data or images. Advantage: Interpretability with non-negativity, efficient with sparse data. Limitation: Less effective for negative or zero-centered data.",
        "umap_hdbscan": "UMAP + HDBSCAN:\n\tCombination of UMAP for dimensionality reduction and HDBSCAN for density-based clustering, suitable for cluster discovery in high-dimensional data. Advantage: Effective in discovering clusters in embeddings. Limitation: Requires careful tuning of both UMAP and HDBSCAN parameters.",
        "manifold_learning": "Manifold Learning (Isomap, Hessian LLE, etc.):\n\tMethods designed to capture intrinsic geometrical structure. Advantage: Preserves non-linear relationships in low dimensions. Limitation: Computationally expensive and sensitive to noise.",
    }

    from sklearn.preprocessing import StandardScaler
    from sklearn.impute import SimpleImputer

    if plot_:
        import matplotlib.pyplot as plt
        import seaborn as sns
    # Check valid method input
    methods = [
        "pca",
        "umap",
        "umap_hdbscan",
        "tsne",
        "factor",
        "isolation_forest",
        "manifold_learning",
        "lda",
        "kpca",
        "ica",
        "mds",
        "lle",
        "svd",
        "truncated_svd",
        "spectral_embedding",
        # "autoencoder","nmf",
    ]
    method = strcmp(method, methods)[0]
    if run_once_within(reverse=True):
        print(f"support methods:{methods}")

    if verbose:
        print(f"\nprocessing with using {dict_methods[method]}:")
    xlabel, ylabel = None, None
    if columns is None:
        columns = data.select_dtypes(include="number").columns.tolist()
    if hue is None:
        hue = data.select_dtypes(exclude="number").columns.tolist()
        print(f"auto select the non-number as 'hue':{hue}")
    if isinstance(hue, list):
        print("Warning: hue is a list, only select the 1st one")
        hue = hue[0]
    if not any(hue):
        # Select columns if specified, else use all columns
        X = data[columns].values if columns else data.values
    else:
        # Select columns to reduce and hue for LDA
        try:
            X = data[columns].values if columns else data.drop(columns=[hue]).values
            y = data[hue].values
        except:
            pass
    print(X.shape)
    # Handle missing values
    if fill_missing:
        imputer = SimpleImputer(strategy="mean")
        X = imputer.fit_transform(X)

    # Optionally scale the data
    if scale:
        scaler = StandardScaler()
        X = scaler.fit_transform(X)

    # Apply PCA if selected
    if method == "pca":
        from sklearn.decomposition import PCA

        pca = PCA(n_components=n_components)
        X_reduced = pca.fit_transform(X)

        # Additional PCA information
        explained_variance = pca.explained_variance_ratio_
        singular_values = pca.singular_values_
        loadings = pca.components_.T * np.sqrt(pca.explained_variance_)

        if debug:
            print(f"PCA completed: Reduced to {n_components} components.")
            print(f"Explained Variance: {explained_variance}")
            print(f"Singular Values: {singular_values}")

        # Plot explained variance if debug=True
        if debug:
            # Plot explained variance
            cumulative_variance = np.cumsum(explained_variance)
            plt.figure(figsize=(8, 5))
            plt.plot(
                range(1, len(cumulative_variance) + 1), cumulative_variance, marker="o"
            )
            plt.title("Cumulative Explained Variance by Principal Components")
            plt.xlabel("Number of Principal Components")
            plt.ylabel("Cumulative Explained Variance")
            plt.axhline(y=0.95, color="r", linestyle="--", label="Threshold (95%)")
            plt.axvline(
                x=n_components,
                color="g",
                linestyle="--",
                label=f"n_components = {n_components}",
            )
            plt.legend()
            plt.grid()
            plt.show()

        # Prepare reduced DataFrame with additional PCA info
        pca_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"PC_{i+1}" for i in range(n_components)],
        )
        # pca_df["Explained Variance"] = np.tile(explained_variance[:n_components], (pca_df.shape[0], 1))
        # pca_df["Singular Values"] = np.tile(singular_values[:n_components], (pca_df.shape[0], 1))
        # Expand explained variance to multiple columns if needed
        for i in range(n_components):
            pca_df[f"Explained Variance PC_{i+1}"] = np.tile(
                format(explained_variance[i] * 100, ".3f") + "%", (pca_df.shape[0], 1)
            )
        for i in range(n_components):
            pca_df[f"Singular Values PC_{i+1}"] = np.tile(
                singular_values[i], (pca_df.shape[0], 1)
            )
        if hue:
            pca_df[hue] = y
    elif method == "lda":
        from sklearn.discriminant_analysis import LinearDiscriminantAnalysis

        if "hue" not in locals() or hue is None:
            raise ValueError(
                "LDA requires a 'hue' col parameter to specify class labels."
            )

        lda_reducer = LinearDiscriminantAnalysis(n_components=n_components)
        X_reduced = lda_reducer.fit_transform(X, y)

        # Prepare reduced DataFrame with additional LDA info
        lda_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"LDA_{i+1}" for i in range(n_components)],
        )
        if debug:
            print(f"LDA completed: Reduced to {n_components} components.")
            print("Class separability achieved by LDA.")
        if hue:
            lda_df[hue] = y
    # Apply UMAP if selected
    elif method == "umap":
        import umap

        umap_reducer = umap.UMAP(
            n_neighbors=umap_neighbors,
            min_dist=umap_min_dist,
            n_components=n_components,
        )
        X_reduced = umap_reducer.fit_transform(X)

        # Additional UMAP information
        embedding = umap_reducer.embedding_
        trustworthiness = umap_reducer._raw_data[:, :n_components]

        if debug:
            print(f"UMAP completed: Reduced to {n_components} components.")
            print(f"Embedding Shape: {embedding.shape}")
            print(f"Trustworthiness: {trustworthiness}")

        # Prepare reduced DataFrame with additional UMAP info
        umap_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"UMAP_{i+1}" for i in range(n_components)],
        )
        umap_df["Embedding"] = embedding[:, 0]  # Example of embedding data
        umap_df["Trustworthiness"] = trustworthiness[:, 0]  # Trustworthiness metric
        if hue:
            umap_df[hue] = y
    elif method == "tsne":
        from sklearn.manifold import TSNE

        tsne = TSNE(
            n_components=n_components,
            perplexity=tsne_perplexity,
            random_state=random_state,
        )
        X_reduced = tsne.fit_transform(X)
        tsne_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"tSNE_{i+1}" for i in range(n_components)],
        )
        tsne_df["Perplexity"] = np.tile(
            f"Perplexity: {tsne_perplexity}", (tsne_df.shape[0], 1)
        )
        if hue:
            tsne_df[hue] = y
    # Apply Factor Analysis if selected
    elif method == "factor":
        from sklearn.decomposition import FactorAnalysis

        factor = FactorAnalysis(n_components=n_components, random_state=random_state)
        X_reduced = factor.fit_transform(X)
        # Factor Analysis does not directly provide explained variance, but we can approximate it
        fa_variance = factor.noise_variance_
        # Prepare reduced DataFrame with additional Factor Analysis info
        factor_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"Factor_{i+1}" for i in range(n_components)],
        )
        factor_df["Noise Variance"] = np.tile(
            format(np.mean(fa_variance) * 100, ".3f") + "%", (factor_df.shape[0], 1)
        )
        if hue:
            factor_df[hue] = y
    # Apply Isolation Forest for outlier detection if selected
    elif method == "isolation_forest":
        from sklearn.decomposition import PCA
        from sklearn.ensemble import IsolationForest

        # Step 1: Apply PCA for dimensionality reduction to 2 components
        pca = PCA(n_components=n_components)
        X_pca = pca.fit_transform(X)

        explained_variance = pca.explained_variance_ratio_
        singular_values = pca.singular_values_

        # Prepare reduced DataFrame with additional PCA info
        iso_forest_df = pd.DataFrame(
            X_pca, index=data.index, columns=[f"PC_{i+1}" for i in range(n_components)]
        )

        isolation_forest = IsolationForest(
            n_estimators=100, contamination="auto", random_state=1
        )
        isolation_forest.fit(X)
        anomaly_scores = isolation_forest.decision_function(
            X
        )  # Anomaly score: larger is less anomalous
        # Predict labels: 1 (normal), -1 (anomaly)
        anomaly_labels = isolation_forest.fit_predict(X)
        # Add anomaly scores and labels to the DataFrame
        iso_forest_df["Anomaly Score"] = anomaly_scores
        iso_forest_df["Anomaly Label"] = anomaly_labels
        # add info from pca
        for i in range(n_components):
            iso_forest_df[f"Explained Variance PC_{i+1}"] = np.tile(
                format(explained_variance[i] * 100, ".3f") + "%",
                (iso_forest_df.shape[0], 1),
            )
        for i in range(n_components):
            iso_forest_df[f"Singular Values PC_{i+1}"] = np.tile(
                singular_values[i], (iso_forest_df.shape[0], 1)
            )
        if hue:
            iso_forest_df[hue] = y
    # * Apply Kernel PCA if selected
    elif method == "kpca":
        from sklearn.decomposition import KernelPCA

        kpca = KernelPCA(
            n_components=n_components, kernel="rbf", random_state=random_state
        )
        X_reduced = kpca.fit_transform(X)

        # Prepare reduced DataFrame with KPCA info
        kpca_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"KPCA_{i+1}" for i in range(n_components)],
        )
        if debug:
            print("Kernel PCA completed with RBF kernel.")
        if hue:
            kpca_df[hue] = y
    # * Apply ICA if selected
    elif method == "ica":
        from sklearn.decomposition import FastICA

        ica = FastICA(n_components=n_components, random_state=random_state)
        X_reduced = ica.fit_transform(X)

        # Prepare reduced DataFrame with ICA info
        ica_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"ICA_{i+1}" for i in range(n_components)],
        )
        if debug:
            print("Independent Component Analysis (ICA) completed.")
        if hue:
            ica_df[hue] = y
    # * Apply MDS if selected
    elif method == "mds":
        from sklearn.manifold import MDS

        mds = MDS(n_components=n_components, random_state=random_state)
        X_reduced = mds.fit_transform(X)

        # Prepare reduced DataFrame with MDS info
        mds_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"MDS_{i+1}" for i in range(n_components)],
        )
        if debug:
            print("Multidimensional Scaling (MDS) completed.")
        if hue:
            mds_df[hue] = y
    # * Apply Locally Linear Embedding (LLE) if selected
    elif method == "lle":
        from sklearn.manifold import LocallyLinearEmbedding

        lle = LocallyLinearEmbedding(
            n_components=n_components,
            n_neighbors=umap_neighbors,
            random_state=random_state,
        )
        X_reduced = lle.fit_transform(X)

        # Prepare reduced DataFrame with LLE info
        lle_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"LLE_{i+1}" for i in range(n_components)],
        )
        if debug:
            print("Locally Linear Embedding (LLE) completed.")
        if hue:
            lle_df[hue] = y
    # * Apply Singular Value Decomposition (SVD) if selected
    elif method == "svd":
        # Using NumPy's SVD for dimensionality reduction
        U, s, Vt = np.linalg.svd(X, full_matrices=False)
        X_reduced = U[:, :n_components] * s[:n_components]

        # Prepare reduced DataFrame with SVD info
        svd_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"SVD_{i+1}" for i in range(n_components)],
        )
        colname_met = "SVD_"
        if hue:
            svd_df[hue] = y
        if debug:
            print("Singular Value Decomposition (SVD) completed.")
    elif method == "truncated_svd":
        from sklearn.decomposition import TruncatedSVD

        svd = TruncatedSVD(n_components=n_components, random_state=random_state)
        X_reduced = svd.fit_transform(X)
        reduced_df = pd.DataFrame(
            X_reduced,
            columns=[f"SVD Component {i+1}" for i in range(n_components)],
            index=data.index,
        )
        colname_met = "SVD Component "

        if debug:
            print("Truncated SVD completed.")
            print("Explained Variance Ratio:", svd.explained_variance_ratio_)
        if hue:
            reduced_df[hue] = y

    elif method == "spectral_embedding":
        from sklearn.manifold import SpectralEmbedding

        spectral = SpectralEmbedding(
            n_components=n_components, random_state=random_state
        )
        X_reduced = spectral.fit_transform(X)
        reduced_df = pd.DataFrame(
            X_reduced,
            columns=[f"Dimension_{i+1}" for i in range(n_components)],
            index=data.index,
        )
        colname_met = "Dimension_"

        if debug:
            print("Spectral Embedding completed.")
        if hue:
            reduced_df[hue] = y

    elif method == "autoencoder":
        from tensorflow.keras.models import Model
        from tensorflow.keras.layers import Input, Dense

        input_dim = X.shape[1]
        input_layer = Input(shape=(input_dim,))
        encoded = Dense(n_components * 2, activation="relu")(input_layer)
        encoded = Dense(n_components, activation="relu")(encoded)
        autoencoder = Model(input_layer, encoded)
        autoencoder.compile(optimizer="adam", loss="mean_squared_error")
        autoencoder.fit(X, X, epochs=50, batch_size=256, shuffle=True, verbose=0)

        X_reduced = autoencoder.predict(X)
        reduced_df = pd.DataFrame(
            X_reduced,
            columns=[f"Score_{i+1}" for i in range(n_components)],
            index=data.index,
        )
        colname_met = "Score_"

        if debug:
            print("Autoencoder reduction completed.")
        if hue:
            reduced_df[hue] = y

    elif method == "nmf":
        from sklearn.decomposition import NMF

        nmf = NMF(n_components=n_components, random_state=random_state)
        X_reduced = nmf.fit_transform(X)
        reduced_df = pd.DataFrame(
            X_reduced,
            columns=[f"NMF_{i+1}" for i in range(n_components)],
            index=data.index,
        )
        colname_met = "NMF_"

        if debug:
            print("Non-negative Matrix Factorization completed.")
        if hue:
            reduced_df[hue] = y

    elif method == "umap_hdbscan":
        import umap
        import hdbscan

        umap_model = umap.UMAP(
            n_neighbors=umap_neighbors,
            min_dist=umap_min_dist,
            n_components=n_components,
        )
        X_umap = umap_model.fit_transform(X)

        clusterer = hdbscan.HDBSCAN()
        clusters = clusterer.fit_predict(X_umap)

        reduced_df = pd.DataFrame(
            X_umap,
            columns=[f"UMAP_{i+1}" for i in range(n_components)],
            index=data.index,
        )
        reduced_df["Cluster"] = clusters
        colname_met = "UMAP_"
        if debug:
            print("UMAP + HDBSCAN reduction and clustering completed.")
        if hue:
            reduced_df[hue] = y

    elif method == "manifold_learning":
        from sklearn.manifold import Isomap

        isomap = Isomap(n_components=n_components)
        X_reduced = isomap.fit_transform(X)
        reduced_df = pd.DataFrame(
            X_reduced,
            columns=[f"Manifold_{i+1}" for i in range(n_components)],
            index=data.index,
        )
        colname_met = "Manifold_"

        if debug:
            print("Manifold Learning (Isomap) completed.")
        if hue:
            reduced_df[hue] = y

    #! Return reduced data and info as a new DataFrame with the same index
    if method == "pca":
        reduced_df = pca_df
        colname_met = "PC_"
        xlabel = f"PC_1 ({pca_df['Explained Variance PC_1'].tolist()[0]})"
        ylabel = f"PC_2 ({pca_df['Explained Variance PC_2'].tolist()[0]})"
    elif method == "umap":
        reduced_df = umap_df
        colname_met = "UMAP_"
    elif method == "tsne":
        reduced_df = tsne_df
        colname_met = "tSNE_"
    elif method == "factor":
        reduced_df = factor_df
        colname_met = "Factor_"
    elif method == "isolation_forest":
        reduced_df = iso_forest_df  # Already a DataFrame for outliers
        colname_met = "PC_"
        if plot_:
            ax = sns.scatterplot(
                data=iso_forest_df[iso_forest_df["Anomaly Label"] == 1],
                x="PC_1",
                y="PC_2",
                label="normal",
                c="b",
            )
            ax = sns.scatterplot(
                ax=ax,
                data=iso_forest_df[iso_forest_df["Anomaly Label"] == -1],
                x="PC_1",
                y="PC_2",
                c="r",
                label="outlier",
                marker="+",
                s=30,
            )
    elif method == "lda":
        reduced_df = lda_df
        colname_met = "LDA_"
    elif method == "kpca":
        reduced_df = kpca_df
        colname_met = "KPCA_"
    elif method == "ica":
        reduced_df = ica_df
        colname_met = "ICA_"
    elif method == "mds":
        reduced_df = mds_df
        colname_met = "MDS_"
    elif method == "lle":
        reduced_df = lle_df
        colname_met = "LLE_"
    elif method == "svd":
        reduced_df = svd_df
        colname_met = "SVD_"
    # Quick plots
    if plot_ and (not method in ["isolation_forest"]):
        from .plot import plotxy, figsets, get_color

        # if ax is None:
        #     if figsize is None:
        #         _, ax = plt.subplots(figsize=cm2inch(8, 8))
        #     else:
        #         _, ax = plt.subplots(figsize=figsize)
        # else:
        #     ax = ax.cla()
        xlabel = f"{colname_met}1" if xlabel is None else xlabel
        ylabel = f"{colname_met}2" if ylabel is None else ylabel
        palette = get_color(len(flatten(data[hue], verbose=0)))

        reduced_df = reduced_df.sort_values(by=hue)
        print(flatten(reduced_df[hue]))
        ax = plotxy(
            data=reduced_df,
            x=colname_met + "1",
            y=colname_met + "2",
            hue=hue,
            palette=palette,
            # size=size,
            edgecolor=edgecolor,
            kind_=[
                "joint",
                #    "kde",
                "ell",
            ],
            kws_kde=dict(
                hue=hue,
                levels=2,
                common_norm=False,
                fill=True,
                alpha=0.05,
            ),
            kws_joint=dict(kind="scatter", joint_kws=dict(s=size)),
            kws_ellipse=dict(alpha=0.1, lw=1, label=None),
            verbose=False,
            **kwargs,
        )
        figsets(
            legend=dict(
                loc=legend_loc,
                markerscale=markerscale,
                bbox_to_anchor=bbox_to_anchor,
                ncols=ncols,
                fontsize=8,
            ),
            xlabel=xlabel if xlabel else None,
            ylabel=ylabel if ylabel else None,
        )

    if inplace:
        # If inplace=True, add components back into the original data
        for col_idx in range(n_components):
            data.loc[:, f"{colname_met}{col_idx+1}"] = reduced_df.iloc[:, col_idx]
        # Add extra info for PCA/UMAP
        if method == "pca":
            for i in range(n_components):
                data.loc[:, f"Explained Variance PC_{i+1}"] = reduced_df.loc[
                    :, f"Explained Variance PC_{i+1}"
                ]
            for i in range(n_components):
                data.loc[:, f"Singular Values PC_{i+1}"] = reduced_df.loc[
                    :, f"Singular Values PC_{i+1}"
                ]
        elif method == "umap":
            for i in range(n_components):
                data.loc[:, f"UMAP_{i+1}"] = reduced_df.loc[:, f"UMAP_{i+1}"]
            data.loc[:, "Embedding"] = reduced_df.loc[:, "Embedding"]
            data.loc[:, "Trustworthiness"] = reduced_df.loc[:, "Trustworthiness"]

        return None  # No return when inplace=True

    return reduced_df


# example:
# df_reducer(data=data_log, columns=markers, n_components=2)


def get_df_format(data, threshold_unique=0.5, verbose=False):
    """
    检测表格: long, wide or uncertain.

    Parameters:
    - data (pd.DataFrame): DataFrame to check.
    - threshold_unique (float): Proportion threshold for detecting categorical columns.

    Returns:
    - "long" if detected as long format,
    - "wide" if detected as wide format
    - "uncertain" if ambiguous.
    """
    from scipy.stats import entropy
    from sklearn.cluster import AgglomerativeClustering
    from sklearn.preprocessing import StandardScaler

    long_score,wide_score,fs = 0,0,500
    n_rows, n_cols = data.shape
    # -----to reduce memory, only check 500 rows/columns----
    if n_rows > fs:
        if verbose:
            print(f"Sampling {fs} rows from {n_rows} rows.")
        data = data.sample(n=fs, random_state=1)
    if n_cols > fs:
        if verbose:
            print(f"Using first {fs} columns out of {n_cols} columns.")
        data = data.iloc[:, :fs]
    n_rows, n_cols = data.shape
    
    # Step 1: Row-Column Ratio Heuristic
    if n_rows > 3 * n_cols:
        long_score += 2
        if verbose:
            print(
                "Row-Column Ratio suggests long format (many rows relative to columns)."
            )
    elif n_cols > 3 * n_rows:
        wide_score += 2
        if verbose:
            print(
                "Row-Column Ratio suggests wide format (many columns relative to rows)."
            )

    # Step 2: Unique-to-duplicate ratio and entropy for categorical variables
    unique_counts = data.apply(lambda x: x.nunique())
    duplicate_ratio = 1 - unique_counts / n_rows
    if (duplicate_ratio > 0.2).sum() > 0.5 * n_cols:
        wide_score += 2
        if verbose:
            print("High duplicate values in columns suggest wide format.")
    else:
        long_score += 1
        if verbose:
            print(
                "Lower duplicate ratio suggests long format (higher row variability)."
            )

    # Calculate entropy for categorical columns
    categorical_cols = data.select_dtypes(include=["object", "category"]).columns
    if len(categorical_cols) > 0:
        for col in categorical_cols:
            counts = data[col].value_counts(normalize=True)
            col_entropy = entropy(counts)
            if col_entropy < 1.5:
                long_score += 1
                if verbose:
                    print(
                        f"Column '{col}' entropy suggests categorical, supporting long format."
                    )
            else:
                wide_score += 1
                if verbose:
                    print(f"Column '{col}' entropy is higher, supporting wide format.")

    # Step 3: Column grouping analysis for patterns in suffixes/prefixes
    col_names = data.columns.astype(str)
    suffix_count = sum("_" in col or col[-1].isdigit() for col in col_names)
    if suffix_count > 0.3 * n_cols:
        wide_score += 2
        if verbose:
            print(
                "Detected suffix/prefix patterns in column names, suggesting wide format."
            )

    # Step 4: Entity identifier detection for long format with categorical columns
    if len(categorical_cols) > 0 and n_rows > n_cols:
        entity_identifier_count = sum(
            data.duplicated(subset=categorical_cols, keep=False)
        )
        if entity_identifier_count > 0.2 * n_rows:
            long_score += 2
            if verbose:
                print(
                    "Significant duplicate rows based on categorical columns, suggesting long format."
                )

    # Step 5: Clustering analysis on numerical columns for correlation in wide format
    numeric_cols = data.select_dtypes(include="number").columns
    if len(numeric_cols) > 1:
        try:
            scaled_data = StandardScaler().fit_transform(data[numeric_cols].dropna())
            clustering = AgglomerativeClustering(n_clusters=2).fit(scaled_data.T)
            cluster_labels = pd.Series(clustering.labels_)
            if cluster_labels.nunique() < len(numeric_cols) * 0.5:
                wide_score += 2
                if verbose:
                    print(
                        "Clustering on columns shows grouping, suggesting wide format."
                    )
        except Exception as e:
            print(e) if verbose else None

    # Step 6: Inter-column correlation analysis
    if len(numeric_cols) > 1:
        corr_matrix = data[numeric_cols].corr().abs()
        avg_corr = (
            corr_matrix.where(~np.eye(len(corr_matrix), dtype=bool)).mean().mean()
        )
        if avg_corr > 0.6:
            wide_score += 2
            if verbose:
                print("High inter-column correlation suggests wide format.")

    # Step 7: Missing value pattern analysis
    missing_patterns = data.isna().sum(axis=1)
    if missing_patterns.std() < 2:
        wide_score += 1
        if verbose:
            print(
                "Low variation in missing patterns across rows, supporting wide format."
            )
    elif missing_patterns.mean() < 1:
        long_score += 1
        if verbose:
            print("Lower missing pattern suggests long format (less structured).")

    # Step 8: Multi-level clustering on rows to detect block structure for wide format
    if len(numeric_cols) > 1 and n_rows > 5:
        try:
            clustering_rows = AgglomerativeClustering(n_clusters=2).fit(scaled_data)
            if pd.Series(clustering_rows.labels_).nunique() < 2:
                wide_score += 2
                if verbose:
                    print("Row clustering reveals homogeneity, suggesting wide format.")
        except Exception as e:
            print(e) if verbose else None

    # Step 9: Sequential name detection for time-series pattern in wide format
    if any(col.isdigit() or col.startswith("T") for col in col_names):
        wide_score += 1
        if verbose:
            print("Detected time-like sequential column names, supporting wide format.")

    # Step 10: Entropy of numeric columns
    try:
        numeric_entropy = data[numeric_cols].apply(
            lambda x: entropy(pd.cut(x, bins=10).value_counts(normalize=True))
        )
        if numeric_entropy.mean() < 2:
            wide_score += 2
            if verbose:
                print(
                    "Low entropy in numeric columns indicates stability across columns, supporting wide format."
                )
    except Exception as e:
        print(e) if verbose else None

    # Step 11: Tie-breaking strategy if scores are equal
    if wide_score == long_score:
        if n_cols > n_rows:
            wide_score += 1
            if verbose:
                print(
                    "Tie-breaking based on column-major structure, favoring wide format."
                )
        elif n_rows > n_cols:
            long_score += 1
            if verbose:
                print(
                    "Tie-breaking based on row-major structure, favoring long format."
                )
        else:
            if verbose:
                print("Tie-breaking inconclusive; returning 'uncertain'.")
            return "uncertain"

    # Final decision
    if wide_score > long_score:
        if verbose:
            print("Final decision: Wide format.")
        return "wide"
    elif long_score > wide_score:
        if verbose:
            print("Final decision: Long format.")
        return "long"
    else:
        if verbose:
            print("Final decision: Uncertain format.")
        return "uncertain"


def plot_cluster(
    data: pd.DataFrame,
    labels: np.ndarray,
    metrics: dict = None,
    cmap="tab20",
    true_labels: Optional[np.ndarray] = None,
) -> None:
    """
    Visualize clustering results with various plots.

    Parameters:
    -----------
    data : pd.DataFrame
        The input data used for clustering.
    labels : np.ndarray
        Cluster labels assigned to each point.
    metrics : dict
        Dictionary containing evaluation metrics from evaluate_cluster function.
    true_labels : Optional[np.ndarray], default=None
        Ground truth labels, if available.
    """
    import seaborn as sns
    from sklearn.metrics import silhouette_samples
    import matplotlib.pyplot as plt

    if metrics is None:
        metrics = evaluate_cluster(data=data, labels=labels, true_labels=true_labels)

    # 1. Scatter Plot of Clusters
    plt.figure(figsize=(15, 6))
    plt.subplot(1, 3, 1)
    plt.scatter(data.iloc[:, 0], data.iloc[:, 1], c=labels, cmap=cmap, s=20)
    plt.title("Cluster Scatter Plot")
    plt.xlabel("Component 1")
    plt.ylabel("Component 2")
    plt.colorbar(label="Cluster Label")
    plt.grid()

    # 2. Silhouette Plot
    if "Silhouette Score" in metrics:
        silhouette_vals = silhouette_samples(data, labels)
        plt.subplot(1, 3, 2)
        y_lower = 10
        for i in range(len(set(labels))):
            # Aggregate the silhouette scores for samples belonging to the current cluster
            cluster_silhouette_vals = silhouette_vals[labels == i]
            cluster_silhouette_vals.sort()
            size_cluster_i = cluster_silhouette_vals.shape[0]
            y_upper = y_lower + size_cluster_i

            plt.fill_betweenx(np.arange(y_lower, y_upper), 0, cluster_silhouette_vals)
            plt.text(-0.05, y_lower + 0.5 * size_cluster_i, str(i))
            y_lower = y_upper + 10  # 10 for the 0 samples

        plt.title("Silhouette Plot")
        plt.xlabel("Silhouette Coefficient Values")
        plt.ylabel("Cluster Label")
        plt.axvline(x=metrics["Silhouette Score"], color="red", linestyle="--")
        plt.grid()

    # 3. Metrics Plot
    plt.subplot(1, 3, 3)
    metric_names = ["Davies-Bouldin Index", "Calinski-Harabasz Index"]
    metric_values = [
        metrics["Davies-Bouldin Index"],
        metrics["Calinski-Harabasz Index"],
    ]

    if true_labels is not None:
        metric_names += ["Homogeneity Score", "Completeness Score", "V-Measure"]
        metric_values += [
            metrics["Homogeneity Score"],
            metrics["Completeness Score"],
            metrics["V-Measure"],
        ]

    plt.barh(metric_names, metric_values, color="lightblue")
    plt.title("Clustering Metrics")
    plt.xlabel("Score")
    plt.axvline(x=0, color="gray", linestyle="--")
    plt.grid()
    plt.tight_layout()


def evaluate_cluster(
    data: pd.DataFrame, labels: np.ndarray, true_labels: Optional[np.ndarray] = None
) -> dict:
    """
    Evaluate clustering performance using various metrics.

    Parameters:
    -----------
    data : pd.DataFrame
        The input data used for clustering.
    labels : np.ndarray
        Cluster labels assigned to each point.
    true_labels : Optional[np.ndarray], default=None
        Ground truth labels, if available.

    Returns:
    --------
    metrics : dict
        Dictionary containing evaluation metrics.

    1. Silhouette Score:
    The silhouette score measures how similar an object is to its own cluster (cohesion) compared to
    how similar it is to other clusters (separation). The score ranges from -1 to +1:
    +1: Indicates that the data point is very far from the neighboring clusters and well clustered.
    0: Indicates that the data point is on or very close to the decision boundary between two neighboring
    clusters.
    -1: Indicates that the data point might have been assigned to the wrong cluster.

    Interpretation:
    A higher average silhouette score indicates better-defined clusters.
    If the score is consistently high (above 0.5), it suggests that the clusters are well separated.
    A score near 0 may indicate overlapping clusters, while negative scores suggest points may have
    been misclassified.

    2. Davies-Bouldin Index:
    The Davies-Bouldin Index (DBI) measures the average similarity ratio of each cluster with its
    most similar cluster. The index values range from 0 to ∞, with lower values indicating better clustering.
    It is defined as the ratio of within-cluster distances to between-cluster distances.

    Interpretation:
    A lower DBI value indicates that the clusters are compact and well-separated.
    Ideally, you want to minimize the Davies-Bouldin Index. If your DBI value is above 1, this indicates
    that your clusters might not be well-separated.

    3. Adjusted Rand Index (ARI):
    The Adjusted Rand Index (ARI) is a measure of the similarity between two data clusterings. The ARI
    score ranges from -1 to +1:
    1: Indicates perfect agreement between the two clusterings.
    0: Indicates that the clusterings are no better than random.
    Negative values: Indicate less agreement than expected by chance.

    Interpretation:
    A higher ARI score indicates better clustering, particularly if it's close to 1.
    An ARI score of 0 or lower suggests that the clustering results do not represent the true labels
    well, indicating a poor clustering performance.

    4. Calinski-Harabasz Index:
    The Calinski-Harabasz Index (also known as the Variance Ratio Criterion) evaluates the ratio of
    the sum of between-cluster dispersion to within-cluster dispersion. Higher values indicate better clustering.

    Interpretation:
    A higher Calinski-Harabasz Index suggests that clusters are dense and well-separated. It is typically
    used to validate the number of clusters, with higher values favoring more distinct clusters.

    5. Homogeneity Score:
    The homogeneity score measures how much a cluster contains only members of a single class (if true labels are provided).
    A score of 1 indicates perfect homogeneity, where all points in a cluster belong to the same class.

    Interpretation:
    A higher homogeneity score indicates that the clustering result is pure, meaning the clusters are composed
    of similar members. Lower values indicate mixed clusters, suggesting poor clustering performance.

    6. Completeness Score:
    The completeness score evaluates how well all members of a given class are assigned to the same cluster.
    A score of 1 indicates perfect completeness, meaning all points in a true class are assigned to a single cluster.

    Interpretation:
    A higher completeness score indicates that the clustering effectively groups all instances of a class together.
    Lower values suggest that some instances of a class are dispersed among multiple clusters.

    7. V-Measure:
    The V-measure is the harmonic mean of homogeneity and completeness, giving a balanced measure of clustering performance.

    Interpretation:
    A higher V-measure score indicates that the clusters are both homogenous (pure) and complete (cover all members of a class).
    Scores closer to 1 indicate better clustering quality.
    """
    from sklearn.metrics import (
        silhouette_score,
        davies_bouldin_score,
        adjusted_rand_score,
        calinski_harabasz_score,
        homogeneity_score,
        completeness_score,
        v_measure_score,
    )

    metrics = {}
    unique_labels = set(labels)
    if len(unique_labels) > 1 and len(unique_labels) < len(data):
        # Calculate Silhouette Score
        try:
            metrics["Silhouette Score"] = silhouette_score(data, labels)
        except Exception as e:
            metrics["Silhouette Score"] = np.nan
            print(f"Silhouette Score calculation failed: {e}")

        # Calculate Davies-Bouldin Index
        try:
            metrics["Davies-Bouldin Index"] = davies_bouldin_score(data, labels)
        except Exception as e:
            metrics["Davies-Bouldin Index"] = np.nan
            print(f"Davies-Bouldin Index calculation failed: {e}")

        # Calculate Calinski-Harabasz Index
        try:
            metrics["Calinski-Harabasz Index"] = calinski_harabasz_score(data, labels)
        except Exception as e:
            metrics["Calinski-Harabasz Index"] = np.nan
            print(f"Calinski-Harabasz Index calculation failed: {e}")

        # Calculate Adjusted Rand Index if true labels are provided
        if true_labels is not None:
            try:
                metrics["Adjusted Rand Index"] = adjusted_rand_score(
                    true_labels, labels
                )
            except Exception as e:
                metrics["Adjusted Rand Index"] = np.nan
                print(f"Adjusted Rand Index calculation failed: {e}")

            # Calculate Homogeneity Score
            try:
                metrics["Homogeneity Score"] = homogeneity_score(true_labels, labels)
            except Exception as e:
                metrics["Homogeneity Score"] = np.nan
                print(f"Homogeneity Score calculation failed: {e}")

            # Calculate Completeness Score
            try:
                metrics["Completeness Score"] = completeness_score(true_labels, labels)
            except Exception as e:
                metrics["Completeness Score"] = np.nan
                print(f"Completeness Score calculation failed: {e}")

            # Calculate V-Measure
            try:
                metrics["V-Measure"] = v_measure_score(true_labels, labels)
            except Exception as e:
                metrics["V-Measure"] = np.nan
                print(f"V-Measure calculation failed: {e}")
    else:
        # Metrics cannot be computed with 1 cluster or all points as noise
        metrics["Silhouette Score"] = np.nan
        metrics["Davies-Bouldin Index"] = np.nan
        metrics["Calinski-Harabasz Index"] = np.nan
        if true_labels is not None:
            metrics["Adjusted Rand Index"] = np.nan
            metrics["Homogeneity Score"] = np.nan
            metrics["Completeness Score"] = np.nan
            metrics["V-Measure"] = np.nan

    return metrics


def df_qc(
    data: pd.DataFrame,
    columns=None,
    skim=False,
    plot_=True,
    max_cols=20,  # only for plots
    hue=None,
    output=False,
    verbose=True,
    dir_save=None,
):
    """
    Usage example:
    df = pd.DataFrame(...)  # Your DataFrameres_qc = df_qc(df)
    """
    from statsmodels.stats.outliers_influence import variance_inflation_factor
    from scipy.stats import skew, kurtosis, entropy

    pd.options.display.max_seq_items = 10
    #! display(data.select_dtypes(include=[np.number]).describe())
    #!skim
    if columns is not None:
        if isinstance(columns, (list, pd.core.indexes.base.Index)):
            data = data[columns]
    if skim:
        try:
            import skimpy

            skimpy.skim(data)
        except:
            numerical_data = data.select_dtypes(include=[np.number])
            skimpy.skim(numerical_data)
    # Fill completely NaN columns with a default value (e.g., 0)
    data = data.copy()
    data.loc[:, data.isna().all()] = 0
    res_qc = {}
    print(f"⤵ data.shape:{data.shape}\n⤵ data.sample(10):")
    display(data.sample(10).style.background_gradient(cmap="coolwarm", axis=1))

    # Missing values
    res_qc["missing_values"] = data.isnull().sum()
    res_qc["missing_percentage"] = round(
        (res_qc["missing_values"] / len(data)) * 100, 2
    )
    res_qc["rows_with_missing"] = data.isnull().any(axis=1).sum()

    # Data types and unique values
    res_qc["data_types"] = data.dtypes
    res_qc["unique_counts"] = (
        data.select_dtypes(exclude=np.number).nunique().sort_values()
    )
    res_qc["unique_values"] = data.select_dtypes(exclude=np.number).apply(
        lambda x: x.unique()
    )
    res_qc["constant_columns"] = [
        col for col in data.columns if data[col].nunique() <= 1
    ]

    # Duplicate rows and columns
    res_qc["duplicate_rows"] = data.duplicated().sum()
    res_qc["duplicate_columns"] = data.columns[data.columns.duplicated()].tolist()

    # Empty columns
    res_qc["empty_columns"] = [col for col in data.columns if data[col].isnull().all()]

    # outliers
    data_outliers = df_outlier(data)
    outlier_num = data_outliers.isna().sum() - data.isnull().sum()
    res_qc["outlier_num"] = outlier_num[outlier_num > 0]
    outlier_percentage = round((outlier_num / len(data_outliers)) * 100, 2)
    res_qc["outlier_percentage"] = outlier_percentage[outlier_percentage > 0]
    try:
        # Correlation and multicollinearity (VIF)
        if any(data.dtypes.apply(pd.api.types.is_numeric_dtype)):
            numeric_df = data.select_dtypes(include=[np.number]).dropna()
            corr_matrix = numeric_df.corr()
            high_corr_pairs = [
                (col1, col2)
                for col1 in corr_matrix.columns
                for col2 in corr_matrix.columns
                if col1 != col2 and abs(corr_matrix[col1][col2]) > 0.9
            ]
            res_qc["high_correlations"] = high_corr_pairs

            # VIF for multicollinearity check
            numeric_df = data.select_dtypes(include=[np.number]).dropna()
            if isinstance(numeric_df.columns, pd.MultiIndex):
                numeric_df.columns = [
                    "_".join(col).strip() if isinstance(col, tuple) else col
                    for col in numeric_df.columns
                ]

            vif_data = pd.DataFrame()
            res_qc["vif"] = vif_data
            if numeric_df.shape[1] > 1 and not numeric_df.empty:
                vif_data["feature"] = numeric_df.columns.tolist()
                vif_data["VIF"] = [
                    round(variance_inflation_factor(numeric_df.values, i), 2)
                    for i in range(numeric_df.shape[1])
                ]
                res_qc["vif"] = vif_data[
                    vif_data["VIF"] > 5
                ]  # Typically VIF > 5 indicates multicollinearity
    except Exception as e:
        print(e)
    # Skewness and Kurtosis
    skewness = data.skew(numeric_only=True)
    kurtosis_vals = data.kurt(numeric_only=True)
    res_qc["skewness"] = skewness[abs(skewness) > 1]
    res_qc["kurtosis"] = kurtosis_vals[abs(kurtosis_vals) > 3]

    # Entropy for categorical columns (higher entropy suggests more disorder)
    categorical_cols = data.select_dtypes(include=["object", "category"]).columns
    res_qc["entropy_categoricals"] = {
        col: entropy(data[col].value_counts(normalize=True), base=2)
        for col in categorical_cols
    }

    # dtypes counts
    res_qc["dtype_counts"] = data.dtypes.value_counts()

    # Distribution Analysis (mean, median, mode, std dev, IQR for numeric columns)
    
    distribution_stats = data.select_dtypes(include=[np.number]).describe().T
    iqr = data.select_dtypes(include=[np.number]).apply(
        lambda x: x.quantile(0.75) - x.quantile(0.25)
    )
    distribution_stats["IQR"] = iqr
    res_qc["distribution_analysis"] = distribution_stats

    # Variance Check: Identify low-variance columns
    variance_threshold = 0.01
    low_variance_cols = [
        col
        for col in data.select_dtypes(include=[np.number]).columns
        if data[col].var() < variance_threshold
    ]
    res_qc["low_variance_features"] = low_variance_cols

    # Categorical columns and cardinality
    categorical_cols = data.select_dtypes(include=["object", "category"]).columns
    high_cardinality = {
        col: data[col].nunique() for col in categorical_cols if data[col].nunique() > 50
    }
    res_qc["high_cardinality_categoricals"] = high_cardinality

    # Feature-type inconsistency (mixed types in columns)
    inconsistent_types = {}
    for col in data.columns:
        unique_types = set(type(val) for val in data[col].dropna())
        if len(unique_types) > 1:
            inconsistent_types[col] = unique_types
    res_qc["inconsistent_types"] = inconsistent_types

    # Text length analysis for text fields
    text_lengths = {}
    for col in categorical_cols:
        text_lengths[col] = {
            "avg_length": data[col].dropna().apply(len).mean(),
            "length_variance": data[col].dropna().apply(len).var(),
        }
    res_qc["text_length_analysis"] = text_lengths

    # Summary statistics
    res_qc["summary_statistics"] = data.describe().T.style.background_gradient(
        cmap="coolwarm", axis=0
    )

    # Automated warnings
    warnings = []
    if res_qc["duplicate_rows"] > 0:
        warnings.append("Warning: Duplicate rows detected.")
    if len(res_qc["empty_columns"]) > 0:
        warnings.append("Warning: Columns with only NaN values detected.")
    if len(res_qc["constant_columns"]) > 0:
        warnings.append("Warning: Columns with a single constant value detected.")
    if len(high_corr_pairs) > 0:
        warnings.append("Warning: Highly correlated columns detected.")
    if len(res_qc["vif"]) > 0:
        warnings.append("Warning: Multicollinearity detected in features.")
    if len(high_cardinality) > 0:
        warnings.append("Warning: High cardinality in categorical columns.")
    if len(inconsistent_types) > 0:
        warnings.append("Warning: Columns with mixed data types detected.")
    res_qc["warnings"] = warnings

    # Report generation
    if verbose:
        print("\n⤵  Summary Statistics:")
        display(res_qc["summary_statistics"])
        print("\n⤵  Data Types:")
        display(res_qc["data_types"])
        if any(res_qc["missing_values"][res_qc["missing_values"] > 0]):
            print(" ⤵  Missing Values Counts:")
            display(
                pd.DataFrame(
                    {
                        "missing_values": res_qc["missing_values"][
                            res_qc["missing_values"] > 0
                        ],
                        "missing_percent(%)": res_qc["missing_percentage"][
                            res_qc["missing_percentage"] > 0
                        ],
                    }
                ).style.background_gradient(cmap="coolwarm", axis=0)
            )
            # print(res_qc["missing_percentage"][res_qc["missing_percentage"] > 0])
            print("\n⤵  Rows with Missing Values:", res_qc["rows_with_missing"])

        (
            print("\n⤵  Constant Columns:", res_qc["constant_columns"])
            if any(res_qc["constant_columns"])
            else None
        )
        (
            print("⤵  Duplicate Rows:", res_qc["duplicate_rows"])
            if res_qc["duplicate_rows"]
            else None
        )
        (
            print("⤵  Duplicate Columns:", res_qc["duplicate_columns"])
            if any(res_qc["duplicate_columns"])
            else None
        )

        if any(res_qc["outlier_num"]):
            print("\n⤵  Outlier Report:")
            display(
                pd.DataFrame(
                    {
                        "outlier_num": res_qc["outlier_num"][res_qc["outlier_num"] > 0],
                        "outlier_percentage(%)": res_qc["outlier_percentage"][
                            res_qc["outlier_percentage"] > 0
                        ],
                    }
                ).style.background_gradient(cmap="coolwarm", axis=0)
            )

        if any(res_qc["unique_counts"]):
            print("\n⤵  Unique Values per Column:")
            display(
                pd.DataFrame(
                    {
                        "unique_counts": res_qc["unique_counts"],
                        "unique_values": res_qc["unique_values"],
                    }
                ).style.background_gradient(cmap="coolwarm", axis=0)
            )

        if res_qc["empty_columns"]:
            print("\n⤵  Empty Columns:", res_qc["empty_columns"])

        if any(res_qc["high_correlations"]):
            print("\n⤵  High Correlations (>|0.9|):")
            for col1, col2 in res_qc["high_correlations"]:
                print(f"  {col1} and {col2}")

        if "vif" in res_qc:
            print("\n⤵  Features with High VIF (>|5|):")
            display(res_qc["vif"].style.background_gradient(cmap="coolwarm", axis=0))

        if any(res_qc["high_cardinality_categoricals"]):
            print("\n⤵  High Cardinality Categorical Columns (>|50 unique|):")
            print(res_qc["high_cardinality_categoricals"])
        if any(res_qc["inconsistent_types"]):
            print("\n⤵  Inconsistent Data Types:")
            display(res_qc["inconsistent_types"])
        if any(res_qc["text_length_analysis"]):
            print("\n⤵  Text Length Analysis:")
            for col, stats in res_qc["text_length_analysis"].items():
                print(
                    f"{col}: Avg Length={round(stats['avg_length'],1)}, Length Variance={round(stats['length_variance'],1)}"
                )

        if res_qc["warnings"]:
            print("\nWarnings:")
            for warning in res_qc["warnings"]:
                print("  -", warning)

    pd.reset_option("display.max_seq_items")
    if plot_:
        df_qc_plots(
            data=data, res_qc=res_qc, max_cols=max_cols, hue=hue, dir_save=dir_save
        )
    if output or not plot_:
        return res_qc
    return None


def df_qc_plots(
    data: pd.DataFrame,
    columns=None,
    res_qc: dict = None,
    max_cols=20,
    hue=None,
    dir_save=None,
):
    import matplotlib.pyplot as plt
    import seaborn as sns
    from .plot import subplot, figsets, get_color
    from datetime import datetime

    now_ = datetime.now().strftime("%y%m%d_%H%M%S")

    if columns is not None:
        if isinstance(columns, (list, pd.core.indexes.base.Index)):
            data = data[columns]
    len_total = len(res_qc)
    n_row, n_col = int((len_total + 10)), 3
    nexttile = subplot(n_row, n_col, figsize=[5 * n_col, 5 * n_row], verbose=False)

    missing_data = res_qc["missing_values"][res_qc["missing_values"] > 0].sort_values(
        ascending=False
    )
    if len(missing_data) > max_cols:
        missing_data = missing_data[:max_cols]
    ax_missing_data = sns.barplot(
        y=missing_data.index,
        x=missing_data.values,
        hue=missing_data.index,
        palette=get_color(len(missing_data), cmap="coolwarm")[::-1],
        ax=nexttile(),
    )
    figsets(
        title="Missing (#)",
        xlabel="#",
        ax=ax_missing_data,
        ylabel=None,
        fontsize=8 if len(missing_data) <= 20 else 6,
    )

    outlier_num = res_qc["outlier_num"].sort_values(ascending=False)
    if len(outlier_num) > max_cols:
        outlier_num = outlier_num[:max_cols]
    ax_outlier_num = sns.barplot(
        y=outlier_num.index,
        x=outlier_num.values,
        hue=outlier_num.index,
        palette=get_color(len(outlier_num), cmap="coolwarm")[::-1],
        ax=nexttile(),
    )
    figsets(
        ax=ax_outlier_num,
        title="Outliers (#)",
        xlabel="#",
        ylabel=None,
        fontsize=8 if len(outlier_num) <= 20 else 6,
    )

    #!
    try:
        for col in data.select_dtypes(include="category").columns:
            sns.countplot(
                y=data[col],
                palette=get_color(
                    data.select_dtypes(include="category").shape[1], cmap="coolwarm"
                )[::-1],
                ax=nexttile(),
            )
            figsets(title=f"Count Plot: {col}", xlabel="Count", ylabel=col)
    except Exception as e:
        pass

    # Skewness and Kurtosis Plots
    skewness = res_qc["skewness"].sort_values(ascending=False)
    kurtosis = res_qc["kurtosis"].sort_values(ascending=False)
    if not skewness.empty:
        ax_skewness = sns.barplot(
            y=skewness.index,
            x=skewness.values,
            hue=skewness.index,
            palette=get_color(len(skewness), cmap="coolwarm")[::-1],
            ax=nexttile(),
        )
        figsets(
            title="Highly Skewed Numeric Columns (Skewness > 1)",
            xlabel="Skewness",
            ylabel=None,
            ax=ax_skewness,
            fontsize=8 if len(skewness) <= 20 else 6,
        )
    if not kurtosis.empty:
        ax_kurtosis = sns.barplot(
            y=kurtosis.index,
            x=kurtosis.values,
            hue=kurtosis.index,
            palette=get_color(len(kurtosis), cmap="coolwarm")[::-1],
            ax=nexttile(),
        )
        figsets(
            title="Highly Kurtotic Numeric Columns (Kurtosis > 3)",
            xlabel="Kurtosis",
            ylabel=None,
            ax=ax_kurtosis,
            fontsize=8 if len(kurtosis) <= 20 else 6,
        )

    # Entropy for Categorical Variables
    entropy_data = pd.Series(res_qc["entropy_categoricals"]).sort_values(
        ascending=False
    )
    ax_entropy_data = sns.barplot(
        y=entropy_data.index,
        x=entropy_data.values,
        hue=entropy_data.index,
        palette=get_color(len(entropy_data), cmap="coolwarm")[::-1],
        ax=nexttile(),
    )
    figsets(
        ylabel="Categorical Columns",
        title="Entropy of Categorical Variables",
        xlabel="Entropy (bits)",
        ax=ax_entropy_data,
        fontsize=8 if len(entropy_data) <= 20 else 6,
    )

    # unique counts
    unique_counts = res_qc["unique_counts"].sort_values(ascending=False)
    ax_unique_counts_ = sns.barplot(
        y=unique_counts.index,
        x=unique_counts.values,
        hue=unique_counts.index,
        palette=get_color(len(unique_counts), cmap="coolwarm")[::-1],
        ax=nexttile(),
    )
    figsets(
        title="Unique Counts",
        ylabel=None,
        xlabel="#",
        ax=ax_unique_counts_,
        fontsize=8 if len(unique_counts) <= 20 else 6,
    )
    # Binary Checking
    ax_unique_counts = sns.barplot(
        y=unique_counts[unique_counts < 8].index,
        x=unique_counts[unique_counts < 8].values,
        hue=unique_counts[unique_counts < 8].index,
        palette=get_color(len(unique_counts[unique_counts < 8].index), cmap="coolwarm")[
            ::-1
        ],
        ax=nexttile(),
    )
    plt.axvline(x=2, color="r", linestyle="--", lw=2)
    figsets(
        ylabel=None,
        title="Binary Checking",
        xlabel="#",
        ax=ax_unique_counts,
        fontsize=8 if len(unique_counts[unique_counts < 10].index) <= 20 else 6,
    )

    # dtypes counts
    dtype_counts = res_qc["dtype_counts"]
    txt = []
    for tp in dtype_counts.index:
        txt.append(list(data.select_dtypes(include=tp).columns))

    ax_dtype_counts = sns.barplot(
        x=dtype_counts.index,
        y=dtype_counts.values,
        color="#F3C8B2",
        ax=nexttile(),
    )
    max_columns_per_row = 1  # Maximum number of columns per row
    for i, tp in enumerate(dtype_counts.index):
        if i <= 20:
            column_names = txt[i]
            # Split the column names into multiple lines if too long
            column_name_str = ", ".join(column_names)
            if len(column_name_str) > 40:  # If column names are too long, split them
                column_name_str = "\n".join(
                    [
                        ", ".join(column_names[j : j + max_columns_per_row])
                        for j in range(0, len(column_names), max_columns_per_row)
                    ]
                )
            # Place text annotation with line breaks and rotate the text if needed
            ax_dtype_counts.text(
                i,
                dtype_counts.values[i],
                f"{column_name_str}",
                ha="center",
                va="top",
                c="k",
                fontsize=8 if len(dtype_counts.index) <= 20 else 6,
                rotation=0,
            )
    figsets(
        xlabel=None,
        title="Dtypes",
        ylabel="#",
        ax=ax_dtype_counts,
        fontsize=8 if len(dtype_counts.index) <= 20 else 6,
    )
    # from .plot import pie
    # pie()

    # High cardinality: Show top categorical columns by unique value count
    high_cardinality = res_qc["high_cardinality_categoricals"]
    if high_cardinality and len(high_cardinality) > max_cols:
        high_cardinality = dict(
            sorted(high_cardinality.items(), key=lambda x: x[1], reverse=True)[
                :max_cols
            ]
        )

    if high_cardinality:
        ax_high_cardinality = sns.barplot(
            y=list(high_cardinality.keys()),
            x=list(high_cardinality.values()),
            hue=list(high_cardinality.keys()),
            palette=get_color(len(list(high_cardinality.keys())), cmap="coolwarm")[
                ::-1
            ],
            ax=nexttile(),
        )
        figsets(
            title="High Cardinality Categorical Columns",
            xlabel="Unique Value Count",
            ax=ax_high_cardinality,
            fontsize=8 if len(list(high_cardinality.keys())) <= 20 else 6,
        )
    if res_qc["low_variance_features"]:
        low_variance_data = data[res_qc["low_variance_features"]].copy()
        for col in low_variance_data.columns:
            ax_low_variance_features = sns.histplot(
                low_variance_data[col], bins=20, kde=True, color="coral", ax=nexttile()
            )
            figsets(
                title=f"Low Variance Feature: {col}",
                ax=ax_low_variance_features,
                fontsize=8 if len(low_variance_data[col]) <= 20 else 6,
            )

    # VIF plot for multicollinearity detection
    if "vif" in res_qc and not res_qc["vif"].empty:
        vif_data = res_qc["vif"].sort_values(by="VIF", ascending=False)
        if len(vif_data) > max_cols:
            vif_data = vif_data[:max_cols]
        ax_vif = sns.barplot(
            data=vif_data,
            x="VIF",
            y="feature",
            hue="VIF",
            palette=get_color(len(vif_data), cmap="coolwarm")[::-1],
            ax=nexttile(),
        )
        figsets(
            title="Variance Inflation Factor(VIF)",
            xlabel="VIF",
            ylabel="Features",
            legend=None,
            ax=ax_vif,
            fontsize=8 if len(vif_data) <= 20 else 6,
        )

    # Correlation heatmap for numeric columns with high correlation pairs
    if any(data.dtypes.apply(pd.api.types.is_numeric_dtype)):
        corr = data.select_dtypes(include=[np.number]).corr()
        if corr.shape[1] <= 33:
            mask = np.triu(np.ones_like(corr, dtype=bool))
            num_columns = corr.shape[1]
            fontsize = max(
                6, min(12, 12 - (num_columns - 10) * 0.2)
            )  # Scale between 8 and 12

            ax_heatmap = sns.heatmap(
                corr,
                mask=mask,
                annot=True,
                cmap="coolwarm",
                center=0,
                fmt=".1f",
                linewidths=0.5,
                vmin=-1,
                vmax=1,
                ax=nexttile(2, 2),
                cbar_kws=dict(shrink=0.2, ticks=np.arange(-1, 2, 1)),
                annot_kws={"size": fontsize},
            )

            figsets(xangle=45, title="Correlation Heatmap", ax=ax_heatmap)
    # # save figure
    # if dir_save:
    #     figsave(dir_save,f"qc_plot_{now_}.pdf")

    if columns is not None:
        if isinstance(columns, (list, pd.core.indexes.base.Index)):
            data = data[columns]

    # len_total = len(res_qc)
    # n_row, n_col = int((len_total + 10) / 3), 3
    # nexttile = subplot(n_row, n_col, figsize=[5 * n_col, 5 * n_row],verbose=False)
    #! check distribution
    data_num = data.select_dtypes(include=np.number)
    if len(data_num) > max_cols:
        data_num = data_num.iloc[:, :max_cols]

    data_num = df_scaler(data=data_num, method="standard")

    import scipy.stats as stats

    for column in data_num.columns:
        # * Shapiro-Wilk test for normality
        stat, p_value = stats.shapiro(data_num[column])
        normality = "norm" if p_value > 0.05 else "not_norm"
        # * Plot histogram
        ax_hist = sns.histplot(data_num[column], kde=True, ax=nexttile())
        x_min, x_max = ax_hist.get_xlim()
        y_min, y_max = ax_hist.get_ylim()
        ax_hist.text(
            x_min + (x_max - x_min) * 0.5,
            y_min + (y_max - y_min) * 0.75,
            f"p(Shapiro-Wilk)={p_value:.3f}\n{normality}",
            ha="center",
            va="top",
        )
        figsets(title=column, ax=ax_hist)
        ax_twin = ax_hist.twinx()
        # * Q-Q plot
        stats.probplot(data_num[column], dist="norm", plot=ax_twin)
        figsets(ylabel=f"Q-Q Plot:{column}", title=None)
    # save figure
    if dir_save:
        figsave(dir_save, f"qc_plot_{now_}.pdf")


def df_corr(df: pd.DataFrame, method="pearson"):
    """
    Compute correlation coefficients and p-values for a DataFrame.

    Parameters:
    - df (pd.DataFrame): Input DataFrame with numeric data.
    - method (str): Correlation method ("pearson", "spearman", "kendall").

    Returns:
    - corr_matrix (pd.DataFrame): Correlation coefficient matrix.
    - pval_matrix (pd.DataFrame): P-value matrix.
    """
    from scipy.stats import pearsonr, spearmanr, kendalltau

    methods = ["pearson", "spearman", "kendall"]
    method = strcmp(method, methods)[0]
    methods_dict = {"pearson": pearsonr, "spearman": spearmanr, "kendall": kendalltau}

    cols = df.columns
    corr_matrix = pd.DataFrame(index=cols, columns=cols, dtype=float)
    pval_matrix = pd.DataFrame(index=cols, columns=cols, dtype=float)
    correlation_func = methods_dict[method]

    for col1 in cols:
        for col2 in cols:
            if col1 == col2:
                corr_matrix.loc[col1, col2] = 1.0
                pval_matrix.loc[col1, col2] = 0.0
            else:
                corr, pval = correlation_func(df[col1], df[col2])
                corr_matrix.loc[col1, col2] = corr
                pval_matrix.loc[col1, col2] = pval

    return corr_matrix, pval_matrix


def use_pd(
    func_name="excel",
    verbose=True,
    dir_json="/Users/macjianfeng/Dropbox/github/python/py2ls/py2ls/data/usages_pd.json",
):
    try:
        default_settings = fload(dir_json, output="json")
        valid_kinds = list(default_settings.keys())
        kind = strcmp(func_name, valid_kinds)[0]
        usage = default_settings[kind]
        if verbose:
            for i, i_ in enumerate(ssplit(usage, by=",")):
                i_ = i_.replace("=", "\t= ") + ","
                print(i_) if i == 0 else print("\t", i_)
        else:
            print(usage)
    except Exception as e:
        if verbose:
            print(e)


def get_phone(phone_number: str, region: str = None, verbose=True):
    """
    usage:
        info = get_phone(15237654321, "DE")
        preview(info)

    Extremely advanced phone number analysis function.

    Args:
        phone_number (str): The phone number to analyze.
        region (str): None (Default). Tries to work with international numbers including country codes; otherwise, uses the specified region.

    Returns:
        dict: Comprehensive information about the phone number.
    """
    import phonenumbers
    from phonenumbers import geocoder, carrier, timezone, number_type
    from datetime import datetime
    import pytz
    from tzlocal import get_localzone

    if not isinstance(phone_number, str):
        phone_number = str(phone_number)
    if isinstance(region, str):
        region = region.upper()

    try:
        # Parse the phone number
        parsed_number = phonenumbers.parse(phone_number, region)

        # Validate the phone number
        valid = phonenumbers.is_valid_number(parsed_number)
        possible = phonenumbers.is_possible_number(parsed_number)

        if not valid:
            suggested_fix = phonenumbers.example_number(region) if region else "Unknown"
            return {
                "valid": False,
                "error": "Invalid phone number",
                "suggested_fix": suggested_fix,
            }

        # Basic details
        formatted_international = phonenumbers.format_number(
            parsed_number, phonenumbers.PhoneNumberFormat.INTERNATIONAL
        )
        formatted_national = phonenumbers.format_number(
            parsed_number, phonenumbers.PhoneNumberFormat.NATIONAL
        )
        formatted_e164 = phonenumbers.format_number(
            parsed_number, phonenumbers.PhoneNumberFormat.E164
        )
        country_code = parsed_number.country_code
        region_code = geocoder.region_code_for_number(parsed_number)
        country_name = geocoder.country_name_for_number(parsed_number, "en")

        location = geocoder.description_for_number(parsed_number, "en")
        carrier_name = carrier.name_for_number(parsed_number, "en") or "Unknown Carrier"
        time_zones = timezone.time_zones_for_number(parsed_number)[0]
        current_times = datetime.now(pytz.timezone(time_zones)).strftime(
            "%Y-%m-%d %H:%M:%S %Z"
        )
        number_type_str = {
            phonenumbers.PhoneNumberType.FIXED_LINE: "Fixed Line",
            phonenumbers.PhoneNumberType.MOBILE: "Mobile",
            phonenumbers.PhoneNumberType.FIXED_LINE_OR_MOBILE: "Fixed Line or Mobile",
            phonenumbers.PhoneNumberType.TOLL_FREE: "Toll Free",
            phonenumbers.PhoneNumberType.PREMIUM_RATE: "Premium Rate",
            phonenumbers.PhoneNumberType.SHARED_COST: "Shared Cost",
            phonenumbers.PhoneNumberType.VOIP: "VOIP",
            phonenumbers.PhoneNumberType.PERSONAL_NUMBER: "Personal Number",
            phonenumbers.PhoneNumberType.PAGER: "Pager",
            phonenumbers.PhoneNumberType.UAN: "UAN",
            phonenumbers.PhoneNumberType.UNKNOWN: "Unknown",
        }.get(number_type(parsed_number), "Unknown")

        # Advanced Features
        is_toll_free = (
            number_type(parsed_number) == phonenumbers.PhoneNumberType.TOLL_FREE
        )
        is_premium_rate = (
            number_type(parsed_number) == phonenumbers.PhoneNumberType.PREMIUM_RATE
        )

        # Dialing Information
        dialing_instructions = f"Dial {formatted_national} within {country_name}. Dial {formatted_e164} from abroad."

        # Advanced Timezone Handling
        gmt_offsets = (
            pytz.timezone(time_zones).utcoffset(datetime.now()).total_seconds() / 3600
        )
        # Get the local timezone (current computer's time)
        local_timezone = get_localzone()
        # local_timezone = pytz.timezone(pytz.country_timezones[region_code][0])
        local_offset = local_timezone.utcoffset(datetime.now()).total_seconds() / 3600
        offset_diff = local_offset - gmt_offsets
        head_time = "earlier" if offset_diff < 0 else "later" if offset_diff > 0 else ""
        res = {
            "valid": True,
            "possible": possible,
            "formatted": {
                "international": formatted_international,
                "national": formatted_national,
                "e164": formatted_e164,
            },
            "country_code": country_code,
            "country_name": country_name,
            "region_code": region_code,
            "location": location if location else "Unknown",
            "carrier": carrier_name,
            "time_zone": time_zones,
            "current_times": current_times,
            "local_offset": f"{local_offset} utcoffset",
            "time_zone_diff": f"{head_time} {int(np.abs(offset_diff))} h",
            "number_type": number_type_str,
            "is_toll_free": is_toll_free,
            "is_premium_rate": is_premium_rate,
            "dialing_instructions": dialing_instructions,
            "suggested_fix": None,  # Use phonenumbers.example_number if invalid
            "logs": {
                "number_analysis_completed": datetime.now().strftime(
                    "%Y-%m-%d %H:%M:%S"
                ),
                "raw_input": phone_number,
                "parsed_number": str(parsed_number),
            },
        }

    except phonenumbers.NumberParseException as e:
        res = {"valid": False, "error": str(e)}
    if verbose:
        preview(res)
    return res


def decode_pluscode(
    pluscode: str, reference: tuple = (52.5200, 13.4050), return_bbox: bool = False
):
    """
    Decodes a Plus Code into latitude and longitude (and optionally returns a bounding box).

    Parameters:
        pluscode (str): The Plus Code to decode. Can be full or short.
        reference (tuple, optional): Reference latitude and longitude for decoding short Plus Codes.
                                     Default is None, required if Plus Code is short.
        return_bbox (bool): If True, returns the bounding box coordinates (latitude/longitude bounds).
                            Default is False.

    Returns:
        tuple: (latitude, longitude) if `return_bbox` is False.
               (latitude, longitude, bbox) if `return_bbox` is True.
               bbox = (latitudeLo, latitudeHi, longitudeLo, longitudeHi)
    Raises:
        ValueError: If the Plus Code is invalid or reference is missing for a short code.

    Usage:
    lat, lon = decode_pluscode("7FG6+89")
    print(f"Decoded Short Plus Code: Latitude: {lat}, Longitude: {lon}, Bounding Box: {bbox}")

    lat, lon = decode_pluscode("9F4M7FG6+89")
    print(f"Decoded Full Plus Code: Latitude: {lat}, Longitude: {lon}")
    """
    from openlocationcode import openlocationcode as olc

    # Validate Plus Code
    if not olc.isValid(pluscode):
        raise ValueError(f"Invalid Plus Code: {pluscode}")

    # Handle Short Plus Codes
    if olc.isShort(pluscode):
        if reference is None:
            raise ValueError(
                "Reference location (latitude, longitude) is required for decoding short Plus Codes."
            )
        # Recover the full Plus Code using the reference location
        pluscode = olc.recoverNearest(pluscode, reference[0], reference[1])

    # Decode the Plus Code
    decoded = olc.decode(pluscode)

    # Calculate the center point of the bounding box
    latitude = (decoded.latitudeLo + decoded.latitudeHi) / 2
    longitude = (decoded.longitudeLo + decoded.longitudeHi) / 2

    if return_bbox:
        bbox = (
            decoded.latitudeLo,
            decoded.latitudeHi,
            decoded.longitudeLo,
            decoded.longitudeHi,
        )
        return latitude, longitude, bbox

    return latitude, longitude


def get_loc(input_data, user_agent="0413@mygmail.com)", verbose=True):
    """
        Determine if the input is a city name, lat/lon, or DMS and perform geocoding or reverse geocoding.
    Usage:
        get_loc("Berlin, Germany")  # Example city
        # get_loc((48.8566, 2.3522))  # Example latitude and longitude
        # get_loc("48 51 24.3 N")  # Example DMS input
    """
    from geopy.geocoders import Nominatim
    import re

    def dms_to_decimal(dms):
        """
        Convert DMS (Degrees, Minutes, Seconds) to Decimal format.
        Input should be in the format of "DD MM SS" or "D M S".
        """
        # Regex pattern for DMS input
        pattern = r"(\d{1,3})[^\d]*?(\d{1,2})[^\d]*?(\d{1,2})"
        match = re.match(pattern, dms)

        if match:
            degrees, minutes, seconds = map(float, match.groups())
            decimal = degrees + (minutes / 60) + (seconds / 3600)
            return decimal
        else:
            raise ValueError("Invalid DMS format")

    geolocator = Nominatim(user_agent="0413@mygmail.com)")
    # Case 1: Input is a city name (string)
    if isinstance(input_data, str) and not re.match(r"^\d+(\.\d+)?$", input_data):
        location = geolocator.geocode(input_data)
        try:
            if verbose:
                print(
                    f"Latitude and Longitude for {input_data}: {location.latitude}, {location.longitude}"
                )
            else:
                print(f"Could not find {input_data}.")
            return location
        except Exception as e:
            print(f"Error: {e}")
            return

    # Case 2: Input is latitude and longitude (float or tuple)
    elif isinstance(input_data, (float, tuple)):
        if isinstance(input_data, tuple) and len(input_data) == 2:
            latitude, longitude = input_data
        elif isinstance(input_data, float):
            latitude = input_data
            longitude = None  # No longitude provided for a single float

        # Reverse geocoding
        location_reversed = geolocator.reverse(
            (latitude, longitude) if longitude else latitude
        )
        if verbose:
            print(
                f"Address from coordinates ({latitude}, {longitude if longitude else ''}): {location_reversed.address}"
            )
        else:
            print("Could not reverse geocode the coordinates.")
        return location_reversed

    # Case 3: Input is a DMS string
    elif isinstance(input_data, str):
        try:
            decimal_lat = dms_to_decimal(input_data)
            print(f"Converted DMS to decimal latitude: {decimal_lat}")

            location_reversed = geolocator.reverse(decimal_lat)
            if verbose:
                print(f"Address from coordinates: {location_reversed.address}")
            else:
                print("Could not reverse geocode the coordinates.")
            return location_reversed
        except ValueError:
            print(
                "Invalid input format. Please provide a city name, latitude/longitude, or DMS string."
            )


def enpass(code: str, method: str = "AES", key: str = None):
    """
    usage: enpass("admin")
    Master encryption function that supports multiple methods: AES, RSA, and SHA256.
    :param code: The input data to encrypt or hash.
    :param method: The encryption or hashing method ('AES', 'RSA', or 'SHA256').
    :param key: The key to use for encryption. For AES and RSA, it can be a password or key in PEM format.
    :return: The encrypted data or hashed value.
    """
    import hashlib

    # AES Encryption (Advanced)
    def aes_encrypt(data: str, key: str):
        """
        Encrypts data using AES algorithm in CBC mode.
        :param data: The data to encrypt.
        :param key: The key to use for AES encryption.
        :return: The encrypted data, base64 encoded.
        """
        from cryptography.hazmat.primitives.ciphers import Cipher, algorithms, modes
        from cryptography.hazmat.backends import default_backend
        from cryptography.hazmat.primitives import padding
        import base64
        import os

        # Generate a 256-bit key from the provided password
        key = hashlib.sha256(key.encode()).digest()

        # Generate a random initialization vector (IV)
        iv = os.urandom(16)  # 16 bytes for AES block size

        # Pad the data to be a multiple of 16 bytes using PKCS7
        padder = padding.PKCS7(128).padder()
        padded_data = padder.update(data.encode()) + padder.finalize()

        # Create AES cipher object using CBC mode
        cipher = Cipher(algorithms.AES(key), modes.CBC(iv), backend=default_backend())
        encryptor = cipher.encryptor()
        encrypted_data = encryptor.update(padded_data) + encryptor.finalize()

        # Return the base64 encoded result (IV + encrypted data)
        return base64.b64encode(iv + encrypted_data).decode()

    # RSA Encryption (Advanced)
    def rsa_encrypt(data: str, public_key: str):
        """
        Encrypts data using RSA encryption with OAEP padding.
        :param data: The data to encrypt.
        :param public_key: The public key in PEM format.
        :return: The encrypted data, base64 encoded.
        """
        import base64
        from Crypto.PublicKey import RSA
        from Crypto.Cipher import PKCS1_OAEP

        public_key_obj = RSA.import_key(public_key)
        cipher_rsa = PKCS1_OAEP.new(public_key_obj)
        encrypted_data = cipher_rsa.encrypt(data.encode())
        return base64.b64encode(encrypted_data).decode()

    # SHA256 Hashing (Non-reversible)
    def sha256_hash(data: str):
        """
        Generates a SHA256 hash of the data.
        :param data: The data to hash.
        :return: The hashed value (hex string).
        """
        return hashlib.sha256(data.encode()).hexdigest()

    if key is None:
        key = "worldpeace"
    method = strcmp(method, ["AES", "RSA", "SHA256"])[0]
    if method == "AES":
        return aes_encrypt(code, key)
    elif method == "RSA":
        return rsa_encrypt(code, key)
    elif method == "SHA256":
        return sha256_hash(code)
    else:
        raise ValueError("Unsupported encryption method")


# Master Decryption Function (Supports AES, RSA)
def depass(encrypted_code: str, method: str = "AES", key: str = None):
    """
    Master decryption function that supports multiple methods: AES and RSA.
    :param encrypted_code: The encrypted data to decrypt.
    :param method: The encryption method ('AES' or 'RSA').
    :param key: The key to use for decryption. For AES and RSA, it can be a password or key in PEM format.
    :return: The decrypted data.
    """
    import hashlib

    def aes_decrypt(encrypted_data: str, key: str):
        """
        Decrypts data encrypted using AES in CBC mode.
        :param encrypted_data: The encrypted data, base64 encoded.
        :param key: The key to use for AES decryption.
        :return: The decrypted data (string).
        """
        from cryptography.hazmat.primitives.ciphers import Cipher, algorithms, modes
        from cryptography.hazmat.backends import default_backend
        from cryptography.hazmat.primitives import padding
        import base64

        # Generate the same 256-bit key from the password
        key = hashlib.sha256(key.encode()).digest()

        # Decode the encrypted data from base64
        encrypted_data = base64.b64decode(encrypted_data)

        # Extract the IV and the actual encrypted data
        iv = encrypted_data[:16]  # First 16 bytes are the IV
        encrypted_data = encrypted_data[16:]  # Remaining data is the encrypted message

        # Create AES cipher object using CBC mode
        cipher = Cipher(algorithms.AES(key), modes.CBC(iv), backend=default_backend())
        decryptor = cipher.decryptor()
        decrypted_data = decryptor.update(encrypted_data) + decryptor.finalize()

        # Unpad the decrypted data using PKCS7
        unpadder = padding.PKCS7(128).unpadder()
        unpadded_data = unpadder.update(decrypted_data) + unpadder.finalize()

        return unpadded_data.decode()

    def rsa_decrypt(encrypted_data: str, private_key: str):
        """
        Decrypts RSA-encrypted data using the private key.
        :param encrypted_data: The encrypted data, base64 encoded.
        :param private_key: The private key in PEM format.
        :return: The decrypted data (string).
        """
        from Crypto.PublicKey import RSA
        from Crypto.Cipher import PKCS1_OAEP
        import base64

        encrypted_data = base64.b64decode(encrypted_data)
        private_key_obj = RSA.import_key(private_key)
        cipher_rsa = PKCS1_OAEP.new(private_key_obj)
        decrypted_data = cipher_rsa.decrypt(encrypted_data)
        return decrypted_data.decode()

    if key is None:
        key = "worldpeace"
    method = strcmp(method, ["AES", "RSA", "SHA256"])[0]
    if method == "AES":
        return aes_decrypt(encrypted_code, key)
    elif method == "RSA":
        return rsa_decrypt(encrypted_code, key)
    elif method == "SHA256":
        raise ValueError("SHA256 is a hash function and cannot be decrypted.")
    else:
        raise ValueError("Unsupported decryption method")


def get_clip(dir_save=None):
    """
    Master function to extract content from the clipboard (text, URL, or image).

    Parameters:
        dir_save (str, optional): If an image is found, save it to this path.

    Returns:
        dict: A dictionary with extracted content:
              {
                  "type": "text" | "url" | "image" | "none",
                  "content": <str|Image|None>,
                  "saved_to": <str|None>  # Path if an image is saved
              }
    """
    result = {"type": "none", "content": None, "saved_to": None}

    try:
        import pyperclip
        from PIL import ImageGrab, Image
        import validators

        # 1. Check for text in the clipboard
        clipboard_content = pyperclip.paste()
        if clipboard_content:
            if validators.url(clipboard_content.strip()):
                result["type"] = "url"
                result["content"] = clipboard_content.strip()

            else:
                result["type"] = "text"
                result["content"] = clipboard_content.strip()
            return clipboard_content.strip()

        # 2. Check for image in the clipboard
        image = ImageGrab.grabclipboard()
        if isinstance(image, Image.Image):
            result["type"] = "image"
            result["content"] = image
            if dir_save:
                image.save(dir_save)
                result["saved_to"] = dir_save
                print(f"Image saved to {dir_save}.")
            else:
                print("Image detected in clipboard but not saved.")
            return image
        print("No valid text, URL, or image found in clipboard.")
        return result

    except Exception as e:
        print(f"An error occurred: {e}")
        return result


def keyboard(*args, action="press", n_click=1, interval=0, verbose=False, **kwargs):
    """
    Simulates keyboard input using pyautogui.

    Parameters:
        input_key (str): The key to simulate. Check the list of supported keys with verbose=True.
        action (str): The action to perform. Options are 'press', 'keyDown', or 'keyUp'.
        n_click (int): Number of times to press the key (only for 'press' action).
        interval (float): Time interval between key presses for 'press' action.
        verbose (bool): Print detailed output, including supported keys and debug info.
        kwargs: Additional arguments (reserved for future extensions).

    keyboard("command", "d", action="shorcut")
    """
    import pyautogui

    input_key = args

    actions = ["press", "keyDown", "keyUp", "hold", "release", "hotkey", "shortcut"]
    action = strcmp(action, actions)[0]
    keyboard_keys_ = [
        "\t",
        "\n",
        "\r",
        " ",
        "!",
        '"',
        "#",
        "$",
        "%",
        "&",
        "'",
        "(",
        ")",
        "*",
        "+",
        ",",
        "-",
        ".",
        "/",
        "0",
        "1",
        "2",
        "3",
        "4",
        "5",
        "6",
        "7",
        "8",
        "9",
        ":",
        ";",
        "<",
        "=",
        ">",
        "?",
        "@",
        "[",
        "\\",
        "]",
        "^",
        "_",
        "`",
        "a",
        "b",
        "c",
        "d",
        "e",
        "f",
        "g",
        "h",
        "i",
        "j",
        "k",
        "l",
        "m",
        "n",
        "o",
        "p",
        "q",
        "r",
        "s",
        "t",
        "u",
        "v",
        "w",
        "x",
        "y",
        "z",
        "{",
        "|",
        "}",
        "~",
        "accept",
        "add",
        "alt",
        "altleft",
        "altright",
        "apps",
        "backspace",
        "browserback",
        "browserfavorites",
        "browserforward",
        "browserhome",
        "browserrefresh",
        "browsersearch",
        "browserstop",
        "capslock",
        "clear",
        "convert",
        "ctrl",
        "ctrlleft",
        "ctrlright",
        "decimal",
        "del",
        "delete",
        "divide",
        "down",
        "end",
        "enter",
        "esc",
        "escape",
        "execute",
        "f1",
        "f10",
        "f11",
        "f12",
        "f13",
        "f14",
        "f15",
        "f16",
        "f17",
        "f18",
        "f19",
        "f2",
        "f20",
        "f21",
        "f22",
        "f23",
        "f24",
        "f3",
        "f4",
        "f5",
        "f6",
        "f7",
        "f8",
        "f9",
        "final",
        "fn",
        "hanguel",
        "hangul",
        "hanja",
        "help",
        "home",
        "insert",
        "junja",
        "kana",
        "kanji",
        "launchapp1",
        "launchapp2",
        "launchmail",
        "launchmediaselect",
        "left",
        "modechange",
        "multiply",
        "nexttrack",
        "nonconvert",
        "num0",
        "num1",
        "num2",
        "num3",
        "num4",
        "num5",
        "num6",
        "num7",
        "num8",
        "num9",
        "numlock",
        "pagedown",
        "pageup",
        "pause",
        "pgdn",
        "pgup",
        "playpause",
        "prevtrack",
        "print",
        "printscreen",
        "prntscrn",
        "prtsc",
        "prtscr",
        "return",
        "right",
        "scrolllock",
        "select",
        "separator",
        "shift",
        "shiftleft",
        "shiftright",
        "sleep",
        "space",
        "stop",
        "subtract",
        "tab",
        "up",
        "volumedown",
        "volumemute",
        "volumeup",
        "win",
        "winleft",
        "winright",
        "yen",
        "command",
        "option",
        "optionleft",
        "optionright",
    ]
    if verbose:
        print(f"supported keys: {keyboard_keys_}")

    if action not in ["hotkey", "shortcut"]:
        if not isinstance(input_key, list):
            input_key = list(input_key)
        input_key = [strcmp(i, keyboard_keys_)[0] for i in input_key]

    # correct action
    cmd_keys = [
        "command",
        "option",
        "optionleft",
        "optionright",
        "win",
        "winleft",
        "winright",
        "ctrl",
        "ctrlleft",
        "ctrlright",
    ]
    try:
        if any([i in cmd_keys for i in input_key]):
            action = "hotkey"
    except:
        pass

    print(f"\n{action}: {input_key}")
    # keyboard
    if action in ["press"]:
        # pyautogui.press(input_key, presses=n_click,interval=interval)
        for _ in range(n_click):
            for key in input_key:
                pyautogui.press(key)
                pyautogui.sleep(interval)
    elif action in ["keyDown", "hold"]:
        # pyautogui.keyDown(input_key)
        for _ in range(n_click):
            for key in input_key:
                pyautogui.keyDown(key)
                pyautogui.sleep(interval)

    elif action in ["keyUp", "release"]:
        # pyautogui.keyUp(input_key)
        for _ in range(n_click):
            for key in input_key:
                pyautogui.keyUp(key)
                pyautogui.sleep(interval)

    elif action in ["hotkey", "shortcut"]:
        pyautogui.hotkey(input_key)


def mouse(
    *args,  # loc
    action: str = "move",
    duration: float = 0.5,
    loc_type: str = "absolute",  # 'absolute', 'relative'
    region: tuple = None,  # (tuple, optional): A region (x, y, width, height) to search for the image.
    image_path: str = None,
    wait: float = 0,
    text: str = None,
    confidence: float = 0.8,
    button: str = "left",
    n_click: int = 1,  # number of clicks
    interval: float = 0.25,  # time between clicks
    scroll_amount: int = -500,
    fail_safe: bool = True,
    grayscale: bool = False,
    n_try: int = 10,
    verbose: bool = True,
    **kwargs,
):
    """
    Master function to handle pyautogui actions.

    Parameters:
        action (str): The action to perform ('click', 'double_click', 'type', 'drag', 'scroll', 'move', 'locate', etc.).
        image_path (str, optional): Path to the image for 'locate' or 'click' actions.
        text (str, optional): Text to type for 'type' action.
        confidence (float, optional): Confidence level for image recognition (default 0.8).
        duration (float, optional): Duration for smooth movements in seconds (default 0.5).
        region (tuple, optional): A region (x, y, width, height) to search for the image.
        button (str, optional): Mouse button to use ('left', 'right', 'middle').
        n_click (int, optional): Number of times to click for 'click' actions.
        interval (float, optional): Interval between clicks for 'click' actions.
        offset (tuple, optional): Horizontal offset from the located image. y_offset (int, optional): Vertical offset from the located image.
        scroll_amount (int, optional): Amount to scroll (positive for up, negative for down).
        fail_safe (bool, optional): Enable/disable pyautogui's fail-safe feature.
        grayscale (bool, optional): Search for the image in grayscale mode.

    Returns:
        tuple or None: Returns coordinates for 'locate' actions, otherwise None.
    """
    import pyautogui
    import time

    # import logging
    # logging.basicConfig(level=logging.DEBUG, filename="debug.log")

    pyautogui.FAILSAFE = fail_safe  # Enable/disable fail-safe
    loc_type = "absolute" if "abs" in loc_type else "relative"
    if len(args) == 1:
        if isinstance(args[0], str):
            image_path = args[0]
            x_offset, y_offset = None, None
        else:
            x_offset, y_offset = args
    elif len(args) == 2:
        x_offset, y_offset = args
    elif len(args) == 3:
        x_offset, y_offset, action = args
    elif len(args) == 4:
        x_offset, y_offset, action, duration = args
    else:
        x_offset, y_offset = None, None

    what_action = [
        "locate",
        "click",
        "double_click",
        "triple_click",
        "input",
        "write",
        "type",
        "drag",
        "move",
        "scroll",
        "down",
        "up",
        "hold",
        "press",
        "release",
    ]
    action = strcmp(action, what_action)[0]
    # get the locations
    location = None
    if any([x_offset is None, y_offset is None]):
        if region is None:
            w, h = pyautogui.size()
            region = (0, 0, w, h)
        retries = 0
        while location is None and retries <= n_try:
            try:
                confidence_ = round(float(confidence - 0.05 * retries), 2)
                location = pyautogui.locateOnScreen(
                    image_path,
                    confidence=confidence_,
                    region=region,
                    grayscale=grayscale,
                )
            except Exception as e:
                if verbose:
                    print(f"confidence={confidence_},{e}")
                location = None
            retries += 1

    # try:
    if location:
        x, y = pyautogui.center(location)
        x += x_offset if x_offset else 0
        if x_offset is not None:
            x += x_offset
        if y_offset is not None:
            y += y_offset
        x_offset, y_offset = x, y
    print(action) if verbose else None
    if action in ["locate"]:
        x, y = pyautogui.position()
    elif action in ["click", "double_click", "triple_click"]:
        if action == "click":
            pyautogui.moveTo(x_offset, y_offset, duration=duration)
            time.sleep(wait)
            pyautogui.click(
                x=x_offset, y=y_offset, clicks=n_click, interval=interval, button=button
            )
        elif action == "double_click":
            pyautogui.moveTo(x_offset, y_offset, duration=duration)
            time.sleep(wait)
            pyautogui.doubleClick(
                x=x_offset, y=y_offset, interval=interval, button=button
            )
        elif action == "triple_click":
            pyautogui.moveTo(x_offset, y_offset, duration=duration)
            time.sleep(wait)
            pyautogui.tripleClick(
                x=x_offset, y=y_offset, interval=interval, button=button
            )

    elif action in ["type", "write", "input"]:
        pyautogui.moveTo(x_offset, y_offset, duration=duration)
        time.sleep(wait)
        if text is not None:
            pyautogui.typewrite(text, interval=interval)
        else:
            print("Text must be provided for the 'type' action.") if verbose else None

    elif action == "drag":
        if loc_type == "absolute":
            pyautogui.dragTo(x_offset, y_offset, duration=duration, button=button)
        else:
            pyautogui.dragRel(x_offset, y_offset, duration=duration, button=button)

    elif action in ["move"]:
        if loc_type == "absolute":
            pyautogui.moveTo(x_offset, y_offset, duration=duration)
        else:
            pyautogui.moveRel(x_offset, y_offset, duration=duration)

    elif action == "scroll":
        pyautogui.moveTo(x_offset, y_offset, duration=duration)
        time.sleep(wait)
        pyautogui.scroll(scroll_amount)

    elif action in ["down", "hold", "press"]:
        pyautogui.moveTo(x_offset, y_offset, duration=duration)
        time.sleep(wait)
        pyautogui.mouseDown(x_offset, y_offset, button=button, duration=duration)

    elif action in ["up", "release"]:
        pyautogui.moveTo(x_offset, y_offset, duration=duration)
        time.sleep(wait)
        pyautogui.mouseUp(x_offset, y_offset, button=button, duration=duration)

    else:
        raise ValueError(f"Unsupported action: {action}")


def py2installer(
    script_path: str = None,
    flatform: str = "mingw64",
    output_dir: str = "dist",
    icon_path: str = None,
    include_data: list = None,
    include_import: list = None,
    exclude_import: list = None,
    plugins: list = None,
    use_nuitka: bool = True,
    console: bool = True,
    clean_build: bool = False,
    additional_args: list = None,
    verbose: bool = True,
    standalone: bool = True,
    onefile: bool = False,
    use_docker: bool = False,
    docker_image: str = "python:3.12-slim",
):
    """
    to package Python scripts into standalone application.

    script_path (str): Path to the Python script to package.
    output_dir (str): Directory where the executable will be stored.
    icon_path (str): Path to the .ico file for the executable icon.
    include_data (list): List of additional data files or directories in "source:dest" format.
    exclude_import (list): List of hidden imports to include.
    plugins (list): List of plugins imports to include.e.g., 'tk-inter'
    use_nuitka (bool): Whether to use Nuitka instead of PyInstaller.
    console (bool): If False, hides the console window (GUI mode).
    clean_build (bool): If True, cleans previous build and dist directories.
    additional_args (list): Additional arguments for PyInstaller/Nuitka.
    verbose (bool): If True, provides detailed logs.
    use_docker (bool): If True, uses Docker to package the script.
    docker_image (str): Docker image to use for packaging.

    """

    import os
    import sys
    import shutil
    import subprocess
    import sys
    import glob
    from pathlib import Path

    if run_once_within():
        usage_str = """
            # build locally
            py2installer(
                script_path="update_tab.py",
                output_dir="dist",
                icon_path="icon4app.ico",
                include_data=["dat/*.xlsx:dat"],
                exclude_import=["msoffcrypto", "tkinter", "pandas", "numpy"],
                onefile=True,
                console=False,
                clean_build=True,
                verbose=True,
            )
            # build via docker
            py2installer(
                "my_script.py",
                output_dir="dist",
                onefile=True,
                clean_build=True,
                use_docker=True,
                docker_image="python:3.12-slim"
            )
            # 尽量不要使用--include-package,这可能导致冲突
            py2installer(
                script_path="update_tab.py",
                # flatform=None,
                output_dir="dist_simp_subprocess",
                icon_path="icon4app.ico",
                standalone=True,
                onefile=False,
                include_data=["dat/*.xlsx=dat"],
                plugins=[
                    "tk-inter",
                ],
                use_nuitka=True,
                console=True,
                clean_build=False,
                verbose=0,
            )
            # 最终文件大小对比
            900 MB: nuitka --mingw64 --standalone --windows-console-mode=attach --show-progress --output-dir=dist --macos-create-app-bundle --macos-app-icon=icon4app.ico --nofollow-import-to=timm,paddle,torch,torchmetrics,torchvision,tensorflow,tensorboard,tensorboardx,tensorboard-data-server,textblob,PIL,sklearn,scienceplots,scikit-image,scikit-learn,scikit-surprise,scipy,spikeinterface,spike-sort-lfpy,stanza,statsmodels,streamlit,streamlit-autorefresh,streamlit-folium,pkg2ls,plotly --include-package=msoffcrypto,tkinter,datetime,pandas,numpy --enable-plugin=tk-inter update_tab.py;
            470 MB: nuitka --mingw64 --standalone --windows-console-mode=attach --show-progress --output-dir=dist_direct_nuitka --macos-create-app-bundle --macos-app-icon=icon4app.ico --enable-plugin=tk-inter update_taby ;
        """
        print(usage_str)
        if verbose:
            return
        else:
            pass
    # Check if the script path exists
    script_path = Path(script_path)
    if not script_path.exists():
        raise FileNotFoundError(f"Script '{script_path}' not found.")

    # Clean build and dist directories if requested
    if clean_build:
        for folder in ["build", "dist"]:
            folder_path = Path(folder)
            if folder_path.exists():
                shutil.rmtree(folder_path, ignore_errors=True)
        # Recreate the folders
        for folder in ["build", "dist"]:
            folder_path = Path(folder)
            folder_path.mkdir(parents=True, exist_ok=True)

    if use_docker:
        # Ensure Docker is installed
        try:
            subprocess.run(
                ["docker", "--version"], check=True, capture_output=True, text=True
            )
        except FileNotFoundError:
            raise EnvironmentError("Docker is not installed or not in the PATH.")

        # Prepare Docker volume mappings
        script_dir = script_path.parent.resolve()
        dist_path = Path(output_dir).resolve()
        volumes = [
            f"{script_dir}:/app:rw",
            f"{dist_path}:/output:rw",
        ]
        docker_cmd = [
            "docker",
            "run",
            "--rm",
            "-v",
            volumes[0],
            "-v",
            volumes[1],
            docker_image,
            "bash",
            "-c",
        ]

        # Build the packaging command inside the container
        cmd = ["nuitka"] if use_nuitka else ["pyinstaller"]
        if onefile:
            cmd.append("--onefile")
        if not console:
            cmd.append("--windowed")
        cmd.extend(["--distpath", "/output"])
        if icon_path:
            cmd.extend(["--icon", f"/app/{Path(icon_path).name}"])
        if include_data:
            for data in include_data:
                cmd.extend(["--add-data", f"/app/{data}"])
        if exclude_import:
            for hidden in exclude_import:
                cmd.extend(["--hidden-import", hidden])
        if additional_args:
            cmd.extend(additional_args)
        cmd.append(f"/app/{script_path.name}")

        # Full command to execute inside the container
        docker_cmd.append(" ".join(cmd))

        if verbose:
            print(f"Running Docker command: {' '.join(docker_cmd)}")

        # Run Docker command
        try:
            subprocess.run(
                docker_cmd,
                capture_output=not verbose,
                text=True,
                check=True,
            )
        except subprocess.CalledProcessError as e:
            print(f"Error during Docker packaging:\n{e.stderr}", file=sys.stderr)
            raise
    else:
        # Handle local packaging (native build)
        cmd = ["nuitka"] if use_nuitka else ["pyinstaller"]
        if "min" in flatform.lower() and use_nuitka:
            cmd.append("--mingw64")
        cmd.append("--standalone") if use_nuitka and standalone else None
        cmd.append("--onefile") if onefile else None
        if not console:
            cmd.append("--windows-console-mode=disable")
        else:
            cmd.append("--windows-console-mode=attach")

        cmd.extend(["--show-progress", f"--output-dir={output_dir}"])
        if icon_path:
            icon_path = Path(icon_path)
            if not icon_path.exists():
                raise FileNotFoundError(f"Icon file '{icon_path}' not found.")
            if sys.platform == "darwin":  # macOS platform
                cmd.extend(
                    ["--macos-create-app-bundle", f"--macos-app-icon={icon_path}"]
                )
            elif sys.platform == "win32":  # Windows platform
                cmd.extend([f"--windows-icon-from-ico={icon_path}"])
            elif sys.platform == "linux":  # Linux platform
                cmd.append("--linux-onefile")

        if include_data:
            for data in include_data:
                if "*" in data:
                    matches = glob.glob(data.split(":")[0])
                    for match in matches:
                        dest = data.split(":")[1]
                        cmd.extend(
                            [
                                "--include-data-file=" if use_nuitka else "--add-data",
                                f"{match}:{dest}",
                            ]
                        )
                else:
                    cmd.extend(
                        ["--include-data-file=" if use_nuitka else "--add-data", data]
                    )
        if exclude_import is not None:
            if any(exclude_import):
                cmd.extend([f"--nofollow-import-to={','.join(exclude_import)}"])
        if include_import is not None:
            if any(
                include_import
            ):  # are included in the final build. Some packages may require manual inclusion.
                cmd.extend([f"--include-package={','.join(include_import)}"])
        if plugins:
            for plugin in plugins:
                # Adds support for tkinter, ensuring it works correctly in the standalone build.
                cmd.extend([f"--enable-plugin={plugin}"])

        if additional_args:
            cmd.extend(additional_args)

        # # clean
        # cmd.extend(
        #     [   "--noinclude-numba-mode=nofollow", #Prevents the inclusion of the numba library and its dependencies, reducing the executable size.
        #         "--noinclude-dask-mode=nofollow",#Excludes the dask library
        #         "--noinclude-IPython-mode=nofollow",#Excludes the IPython library and its dependencies.
        #         "--noinclude-unittest-mode=nofollow",#Excludes the unittest module (used for testing) from the build
        #         "--noinclude-pytest-mode=nofollow",#Excludes the pytest library (used for testing) from the build.
        #         "--noinclude-setuptools-mode=nofollow",#Excludes setuptools, which is not needed for the standalone executable.
        #         "--lto=no",#Disables Link-Time Optimization (LTO), which reduces the compilation time but may slightly increase the size of the output.
        #     ]
        # )

        if clean_build:
            cmd.append(
                "--remove-output"
            )  # Removes intermediate files created during the build process, keeping only the final executable.
        # Add the script path (final positional argument)
        cmd.append(str(script_path))
        # Ensure Windows shell compatibility
        shell_flag = sys.platform.startswith("win")
        print(f"Running command: ⤵ \n{' '.join(cmd)}\n")
        try:
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                shell=shell_flag,
                check=True,
            )
            if verbose:
                print(result.stdout)
        except subprocess.CalledProcessError as e:
            print(f"Error during packaging:\n{e.stderr}", file=sys.stderr)
            print(" ".join(cmd))
            raise

    print("\nPackaging complete. Check the output directory for the executable.")


def set_theme(
    context="paper",
    style="whitegrid",
    palette="deep",
    font="sans-serif",
    font_scale=1.0,
    color_codes=True,
    grid_alpha=0.5,
    grid_linewidth=0.8,
    grid_linestyle="--",
    tick_direction="out",
    # tick_length=4,
    spine_visibility=False,
    # figsize=(8, 6),
    # linewidth=2,
    dpi=100,
    rc=None,
):
    """
    to configure Seaborn theme with maximum flexibility.

    # Example Usage
    set_sns_theme(font_scale=1.2, grid_alpha=0.8, tick_direction="in", dpi=150)

        Parameters:
        - context: Plotting context ('notebook', 'paper', 'talk', 'poster')
        - style: Style of the plot ('darkgrid', 'whitegrid', 'dark', 'white', 'ticks')
        - palette: Color palette (string or list of colors)
        - font: Font family ('sans-serif', 'serif', etc.)
        - font_scale: Scaling factor for fonts
        - color_codes: Boolean, whether to use seaborn color codes
        - grid_alpha: Opacity of the grid lines
        - grid_linewidth: Thickness of grid lines
        - grid_linestyle: Style of grid lines ('-', '--', '-.', ':')
        - tick_direction: Direction of ticks ('in', 'out', 'inout')
        - tick_length: Length of ticks
        - spine_visibility: Whether to show plot spines (True/False)
        - figsize: Default figure size as tuple (width, height)
        - linewidth: Default line width for plots
        - dpi: Resolution of the figure
        - rc: Dictionary of additional rc settings
    """
    import seaborn as sns
    import matplotlib.pyplot as plt

    # Define additional rc parameters for fine-tuning
    rc_params = {
        # "axes.grid": True,
        "grid.alpha": grid_alpha,
        "grid.linewidth": grid_linewidth,
        "grid.linestyle": grid_linestyle,
        "xtick.direction": tick_direction,
        "ytick.direction": tick_direction,
        # "xtick.major.size": tick_length,
        # "ytick.major.size": tick_length,
        # "axes.linewidth": linewidth,
        # "figure.figsize": figsize,
        "figure.dpi": dpi,
        "axes.spines.top": spine_visibility,
        "axes.spines.right": spine_visibility,
        "axes.spines.bottom": spine_visibility,
        "axes.spines.left": spine_visibility,
    }

    # Merge user-provided rc settings
    if rc:
        rc_params.update(rc)

    # Apply the theme settings
    sns.set_theme(
        context=context,
        style=style,
        palette=palette,
        font=font,
        font_scale=font_scale,
        color_codes=color_codes,
        rc=rc_params,
    )